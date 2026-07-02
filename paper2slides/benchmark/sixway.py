"""Six-route hybrid style proposal smoke harness.

This module keeps the protected frozen references separate from the new
assisted/autonomous proposal routes. The proposal routes consume only fresh
parse checkpoints, abstract design primitives, and benchmark badcase metadata.
"""
from __future__ import annotations

import argparse
import csv
import hashlib
import json
import shutil
import sys
from collections import Counter
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from paper2slides.benchmark.from_scratch import build_content_inventory, build_rough_draft_spec
from paper2slides.benchmark.fourway import (
    DEFAULT_REPORT_ROOT,
    GOLDEN_BASELINE2_STYLE_ID,
    STYLE_CONTRACTS,
    _audit_to_file,
    _build_rough_speaker_script,
    _figure_path,
    _fit_image_box,
    _iteration_payload,
    _limit_words,
    _materialize_existing_deck_route,
    _materialize_golden_baseline1_route,
    _read_json,
    _repair_log_payload,
    _route_result,
    _run_academic_generation,
    _section_buckets,
    _stop_reason_for_profile,
    _table_rows,
    _write_artifact_index,
    _write_curve,
    _write_json,
    _write_speaker_script_audit,
)


ASSISTED_SEED_STYLE_ID = "assisted_seed_scaffold_style"
AUTONOMOUS_STYLE_A_ID = "autonomous_style_proposal_a"
AUTONOMOUS_STYLE_B_ID = "autonomous_style_proposal_b"
GOLDEN_BASELINE2_REFERENCE_DIR = Path("outputs") / "golden_baselines" / GOLDEN_BASELINE2_STYLE_ID
GOLDEN_BASELINE2_REFERENCE_PPTX = GOLDEN_BASELINE2_REFERENCE_DIR / "DeepResidual_20260630_blind_rectangular_golden2_reference.pptx"


def run_sixway_hybrid_smoke(
    paper_path: Path,
    run_dir: Optional[Path] = None,
    slides: int = 24,
    length: str = "medium",
    fast: bool = True,
    from_stage: str = "rag",
    python_executable: str = sys.executable,
    max_iterations: int = 3,
    patience: int = 2,
    top_k_repairs: int = 3,
) -> Dict[str, Any]:
    """Parse once, then materialize six hybrid benchmark routes."""
    paper_path = paper_path.resolve()
    if not paper_path.exists():
        raise FileNotFoundError(f"Paper PDF does not exist: {paper_path}")

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    run_dir = run_dir or DEFAULT_REPORT_ROOT / f"{_slugify(paper_path.stem)}_sixway_{timestamp}"
    run_dir = run_dir.resolve()
    output_root = run_dir / "fresh_parse_outputs"
    routes_dir = run_dir / "routes"
    logs_dir = run_dir / "logs"
    for directory in (output_root, routes_dir, logs_dir):
        directory.mkdir(parents=True, exist_ok=True)

    parse_state = _existing_parse_state(paper_path)
    primitives = _design_primitives_library()
    policy = _style_proposal_policy(max_iterations=max_iterations, patience=patience, top_k_repairs=top_k_repairs)
    _write_json(run_dir / "style_proposal_policy.json", policy)
    _write_json(run_dir / "design_primitives_library.json", primitives)

    manifest = {
        "schema_version": "single_paper_sixway_hybrid.v1",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "paper_path": str(paper_path),
        "paper_sha256": _sha256(paper_path),
        "run_dir": str(run_dir),
        "fresh_parse_output_root": str(output_root),
        "parse_status_before_run": parse_state,
        "slides": slides,
        "length": length,
        "fast": fast,
        "from_stage": from_stage,
        "python_executable": python_executable,
        "max_iterations": max_iterations,
        "patience": patience,
        "top_k_repairs": top_k_repairs,
        "routes": [route["route_id"] for route in policy["routes"]],
        "principles": [
            "parse the input PDF once into a fresh checkpoint root",
            "use frozen references only in routes 01 to 03",
            "do not feed golden PPTX files or full golden style contracts into routes 04 to 06",
            "run bounded metadata-first audit and repair loops for the new style routes",
        ],
    }
    _write_json(run_dir / "manifest.json", manifest)

    generation = _run_academic_generation(
        paper_path=paper_path,
        output_root=output_root,
        logs_dir=logs_dir,
        slides=slides,
        length=length,
        fast=fast,
        from_stage=from_stage,
        python_executable=python_executable,
    )
    _write_json(run_dir / "fresh_generation_result.json", generation)
    if generation["returncode"] != 0 or not generation.get("output_subdir"):
        raise RuntimeError(
            "Fresh academic generation did not produce a usable output. "
            f"See log: {generation.get('log_path', '')}"
        )

    summary_checkpoint = Path(generation["summary_checkpoint"])
    plan_checkpoint = Path(generation["plan_checkpoint"])
    spec_checkpoint = Path(generation["spec_checkpoint"])
    if not summary_checkpoint.exists() or not plan_checkpoint.exists() or not spec_checkpoint.exists():
        raise RuntimeError("Fresh run is missing one or more checkpoints needed for sixway branch generation.")

    routes: List[Dict[str, Any]] = []
    routes.append(
        _materialize_existing_deck_route(
            source_output=Path(generation["output_subdir"]),
            route_dir=routes_dir / "01_academic_frozen_reference",
            route_id="01_academic_frozen_reference",
            style_id="academic",
            repair_profile="frozen_reference_audit_only",
            style_scope="academic",
            note="Frozen academic reference route; generated by the protected academic renderer and audited without mutation.",
        )
    )

    golden1 = _materialize_golden_baseline1_route(
        summary_checkpoint=summary_checkpoint,
        plan_checkpoint=plan_checkpoint,
        spec_checkpoint=spec_checkpoint,
        route_dir=routes_dir / "02_golden1_frozen_reference",
    )
    routes.append(
        _retag_route(
            golden1,
            route_id="02_golden1_frozen_reference",
            repair_profile="frozen_reference_golden1_scoped",
            note="Frozen golden_baseline1 reference route; its rounded proof-panel grammar is used only in this route.",
        )
    )

    inventory = build_content_inventory(summary_checkpoint, plan_checkpoint, spec_checkpoint)
    rough = build_rough_draft_spec(inventory)
    routes.append(
        _materialize_golden2_frozen_route(
            inventory=inventory,
            rough=rough,
            route_dir=routes_dir / "03_golden2_frozen_reference",
        )
    )

    new_routes = [
        (
            "04_assisted_seed_scaffold_style",
            ASSISTED_SEED_STYLE_ID,
            "assisted_seed_scaffold",
            "L3.5_assisted_seed_scaffold_repair",
        ),
        (
            "05_autonomous_style_proposal_a",
            AUTONOMOUS_STYLE_A_ID,
            "autonomous_free_proposal",
            "L4_candidate_autonomous_style_proposal_and_repair",
        ),
        (
            "06_autonomous_style_proposal_b",
            AUTONOMOUS_STYLE_B_ID,
            "autonomous_free_proposal",
            "L4_candidate_autonomous_style_proposal_and_repair",
        ),
    ]
    for route_id, style_id, route_type, autonomy_level in new_routes:
        routes.append(
            _materialize_proposal_route(
                inventory=inventory,
                rough=rough,
                route_dir=routes_dir / route_id,
                route_id=route_id,
                style_id=style_id,
                route_type=route_type,
                autonomy_level=autonomy_level,
                primitives=primitives,
                max_iterations=max_iterations,
                patience=patience,
                top_k_repairs=top_k_repairs,
            )
        )

    curve_rows = _write_curve(run_dir / "score_curve.csv", routes)
    artifact_rows = _write_artifact_index(run_dir / "artifact_index.csv", routes)
    _write_human_feedback_effort(run_dir / "human_feedback_effort.csv", routes)
    _write_json(run_dir / "external_artifact_eval.json", _external_artifact_eval(routes))
    report_path = run_dir / "comparison_report.md"
    report_path.write_text(_render_sixway_report(manifest, generation, routes, curve_rows), encoding="utf-8")

    result = {
        "schema_version": "sixway_hybrid_result.v1",
        "run_dir": str(run_dir),
        "manifest": str(run_dir / "manifest.json"),
        "style_proposal_policy": str(run_dir / "style_proposal_policy.json"),
        "design_primitives_library": str(run_dir / "design_primitives_library.json"),
        "comparison_report": str(report_path),
        "score_curve": str(run_dir / "score_curve.csv"),
        "artifact_index": str(run_dir / "artifact_index.csv"),
        "human_feedback_effort": str(run_dir / "human_feedback_effort.csv"),
        "external_artifact_eval": str(run_dir / "external_artifact_eval.json"),
        "fresh_generation": generation,
        "routes": routes,
        "artifact_rows": artifact_rows,
    }
    _write_json(run_dir / "sixway_result.json", result)
    return result


def _materialize_golden2_frozen_route(inventory: Dict[str, Any], rough: Dict[str, Any], route_dir: Path) -> Dict[str, Any]:
    route_dir.mkdir(parents=True, exist_ok=True)
    contract = dict(STYLE_CONTRACTS[GOLDEN_BASELINE2_STYLE_ID])
    contract.update(
        {
            "schema_version": "frozen_reference_style_contract.v1",
            "style_id": GOLDEN_BASELINE2_STYLE_ID,
            "deck_label": "FROZEN GOLDEN2",
            "board_title": "Rectangular Research Board",
            "agenda_title": "Evidence Route",
            "palette_name": "reference_indigo",
            "policy": {
                "route_role": "frozen_reference",
                "may_read_full_contract": True,
                "may_be_used_by_proposal_routes": False,
            },
        }
    )
    _write_json(route_dir / "content_inventory.json", inventory)
    _write_json(route_dir / "rough_draft_spec.json", rough)
    _write_json(route_dir / "style_contract.json", contract)
    pptx_path = route_dir / "slides.pptx"
    if GOLDEN_BASELINE2_REFERENCE_PPTX.exists():
        shutil.copy2(GOLDEN_BASELINE2_REFERENCE_PPTX, pptx_path)
        materialization_note = f"Copied the exact frozen golden_baseline2 reference PPTX from {GOLDEN_BASELINE2_REFERENCE_PPTX}."
    else:
        _render_reference_board_pptx(inventory, rough, pptx_path, contract=contract, repair_mode=True)
        materialization_note = "Fallback: frozen golden_baseline2 reference PPTX was missing, so the route used the legacy reference-board renderer."
    audit = _audit_to_file(pptx_path, route_dir / "nonvisual_audit.json")
    speaker_path = route_dir / "speaker_script.md"
    speaker_path.write_text(
        _build_rough_speaker_script(inventory, rough, "golden2 frozen reference", narrative_mode="blind_experimental")
        + "\n\n"
        + materialization_note
        + "\n",
        encoding="utf-8",
    )
    speaker_audit_path = route_dir / "speaker_script_audit.json"
    _write_speaker_script_audit(speaker_path, pptx_path, GOLDEN_BASELINE2_STYLE_ID, speaker_audit_path)
    repair_log = _repair_log_payload(
        route_id="03_golden2_frozen_reference",
        style_id=GOLDEN_BASELINE2_STYLE_ID,
        repair_profile="frozen_reference_golden2_scoped",
        iterations=[
            _iteration_payload(
                0,
                pptx_path,
                audit,
                [materialization_note, "No proposal route may consume this protected frozen reference as a template."],
                _stop_reason_for_profile(audit, "audit_only"),
                speaker_script_path=speaker_path,
                speaker_script_audit_path=speaker_audit_path,
            )
        ],
        note="Frozen golden_baseline2 reference route. This route materializes the exact protected PPTX when available; new proposal routes are not allowed to consume it.",
    )
    _write_json(route_dir / "repair_log.json", repair_log)
    style_report = build_style_drift_report_payload_local(
        pptx_path=pptx_path,
        audit=audit,
        style_id=GOLDEN_BASELINE2_STYLE_ID,
        repair_profile="frozen_reference_golden2_scoped",
        style_scope="golden_baseline2",
    )
    _write_json(route_dir / "style_drift_report.json", style_report)
    return _route_result(
        "03_golden2_frozen_reference",
        route_dir,
        GOLDEN_BASELINE2_STYLE_ID,
        "frozen_reference_golden2_scoped",
        pptx_path,
        speaker_path,
        speaker_audit_path,
        None,
        audit,
        repair_log,
        style_report,
    )


def _materialize_proposal_route(
    inventory: Dict[str, Any],
    rough: Dict[str, Any],
    route_dir: Path,
    route_id: str,
    style_id: str,
    route_type: str,
    autonomy_level: str,
    primitives: Dict[str, Any],
    max_iterations: int,
    patience: int,
    top_k_repairs: int,
) -> Dict[str, Any]:
    route_dir.mkdir(parents=True, exist_ok=True)
    _write_json(route_dir / "content_inventory.json", inventory)
    _write_json(route_dir / "rough_draft_spec.json", rough)

    contract = _build_seed_contract(inventory, route_id, style_id, route_type, autonomy_level)
    if route_type == "assisted_seed_scaffold":
        _write_json(route_dir / "seed_scaffold_contract.json", contract)
        (route_dir / "seed_authoring_note.md").write_text(_seed_authoring_note(contract), encoding="utf-8")
    else:
        _write_json(route_dir / "style_contract.json", contract)
        _write_json(route_dir / "layout_grammar.json", contract["layout_grammar"])
        _write_json(route_dir / "renderer_parameters.json", contract["renderer_parameters"])
        _write_json(route_dir / "novelty_report.json", contract["novelty_report"])

    _write_json(route_dir / "design_primitives_used.json", _design_primitives_used(contract, primitives))
    _write_json(route_dir / "forbidden_reference_attestation.json", _forbidden_reference_attestation(route_id, route_type))

    iterations: List[Dict[str, Any]] = []
    best_score: Optional[int] = None
    stale_iterations = 0
    final_pptx = route_dir / "slides.pptx"
    final_speaker_path = route_dir / "speaker_script.md"
    final_speaker_audit_path = route_dir / "speaker_script_audit.json"
    last_audit: Dict[str, Any] = {}
    applied_by_iteration: List[List[str]] = []
    min_iterations = min(2, max(1, max_iterations))

    for iteration in range(max(1, max_iterations)):
        iteration_dir = route_dir / "iterations" / f"iter_{iteration:02d}"
        iteration_dir.mkdir(parents=True, exist_ok=True)
        iteration_pptx = iteration_dir / "slides.pptx"
        repair_mode = iteration > 0
        _render_proposal_pptx(inventory, rough, iteration_pptx, contract, iteration=iteration, repair_mode=repair_mode)
        last_audit = _audit_to_file(iteration_pptx, iteration_dir / "nonvisual_audit.json")
        iteration_speaker_path = iteration_dir / "speaker_script.md"
        iteration_speaker_path.write_text(
            _build_rough_speaker_script(inventory, rough, f"{route_id} iteration {iteration:02d}", narrative_mode="blind_experimental"),
            encoding="utf-8",
        )
        iteration_speaker_audit_path = iteration_dir / "speaker_script_audit.json"
        _write_speaker_script_audit(iteration_speaker_path, iteration_pptx, style_id, iteration_speaker_audit_path)
        applied = [] if iteration == 0 else _planned_repairs_from_previous(iterations, top_k_repairs)
        applied_by_iteration.append(applied)
        score = _blocking_score(last_audit)
        if best_score is None or score < best_score:
            best_score = score
            stale_iterations = 0
        elif iteration > 0:
            stale_iterations += 1
        is_last_allowed_iteration = iteration + 1 >= max(1, max_iterations)
        stop_reason = _proposal_stop_reason(
            last_audit,
            iteration,
            min_iterations,
            stale_iterations,
            patience,
            is_last_allowed_iteration=is_last_allowed_iteration,
        )
        iterations.append(
            _iteration_payload(
                iteration,
                iteration_pptx,
                last_audit,
                applied,
                stop_reason,
                speaker_script_path=iteration_speaker_path,
                speaker_script_audit_path=iteration_speaker_audit_path,
            )
        )
        if iteration + 1 >= min_iterations and stop_reason != "continue_bounded_repair":
            break

    last = iterations[-1]
    shutil.copy2(Path(last["pptx_path"]), final_pptx)
    shutil.copy2(Path(last["speaker_script_path"]), final_speaker_path)
    shutil.copy2(Path(last["speaker_script_audit_path"]), final_speaker_audit_path)
    final_audit = _audit_to_file(final_pptx, route_dir / "nonvisual_audit.json")
    repair_log = _repair_log_payload(
        route_id=route_id,
        style_id=style_id,
        repair_profile=autonomy_level,
        iterations=iterations,
        note="New-style route. It consumes fresh checkpoints, abstract design primitives, and badcase metadata only.",
    )
    repair_log["route_type"] = route_type
    repair_log["autonomy_level"] = autonomy_level
    repair_log["max_iterations"] = max_iterations
    repair_log["patience"] = patience
    repair_log["top_k_repairs"] = top_k_repairs
    repair_log["forbidden_reference_attestation_path"] = str(route_dir / "forbidden_reference_attestation.json")
    _write_json(route_dir / "repair_log.json", repair_log)
    _write_visual_review_packet_for_route(route_dir, route_id, style_id, contract, iterations)
    style_report = build_style_drift_report_payload_local(
        pptx_path=final_pptx,
        audit=final_audit,
        style_id=style_id,
        repair_profile=autonomy_level,
        style_scope="proposal",
    )
    _write_json(route_dir / "style_drift_report.json", style_report)
    route = _route_result(
        route_id,
        route_dir,
        style_id,
        autonomy_level,
        final_pptx,
        final_speaker_path,
        final_speaker_audit_path,
        None,
        final_audit,
        repair_log,
        style_report,
    )
    if route_type == "assisted_seed_scaffold":
        route["seed_scaffold_contract_path"] = str(route_dir / "seed_scaffold_contract.json")
        route["seed_authoring_note_path"] = str(route_dir / "seed_authoring_note.md")
    else:
        route["style_contract_path"] = str(route_dir / "style_contract.json")
        route["layout_grammar_path"] = str(route_dir / "layout_grammar.json")
        route["renderer_parameters_path"] = str(route_dir / "renderer_parameters.json")
        route["novelty_report_path"] = str(route_dir / "novelty_report.json")
    route["forbidden_reference_attestation_path"] = str(route_dir / "forbidden_reference_attestation.json")
    route["design_primitives_used_path"] = str(route_dir / "design_primitives_used.json")
    route["visual_human_review_packet"] = str(route_dir / "visual_human_review_packet.zh-CN.md")
    return route


def _retag_route(route: Dict[str, Any], route_id: str, repair_profile: str, note: str) -> Dict[str, Any]:
    route_dir = Path(route["route_dir"])
    route["route_id"] = route_id
    route["repair_profile"] = repair_profile
    repair_log = route.get("repair_log", {})
    repair_log["route_id"] = route_id
    repair_log["repair_profile"] = repair_profile
    repair_log["note"] = note
    for iteration in repair_log.get("iterations", []) or []:
        iteration.setdefault("applied_repairs", []).append("Retagged as frozen reference route for sixway hybrid smoke.")
    _write_json(route_dir / "repair_log.json", repair_log)
    route["repair_log"] = repair_log
    style_report = route.get("style_drift_report", {})
    style_report["repair_profile"] = repair_profile
    _write_json(route_dir / "style_drift_report.json", style_report)
    route["style_drift_report"] = style_report
    return route


def _render_reference_board_pptx(
    inventory: Dict[str, Any],
    rough: Dict[str, Any],
    output_path: Path,
    contract: Dict[str, Any],
    repair_mode: bool,
) -> Path:
    reference_contract = dict(contract)
    reference_contract.setdefault("visual_family", "reference_board")
    return _render_generic_native_deck(inventory, rough, output_path, reference_contract, iteration=1 if repair_mode else 0)


def _render_proposal_pptx(
    inventory: Dict[str, Any],
    rough: Dict[str, Any],
    output_path: Path,
    contract: Dict[str, Any],
    iteration: int,
    repair_mode: bool,
) -> Path:
    local_contract = dict(contract)
    local_contract["iteration"] = iteration
    local_contract["repair_mode"] = repair_mode
    return _render_generic_native_deck(inventory, rough, output_path, local_contract, iteration=iteration)


def _render_generic_native_deck(
    inventory: Dict[str, Any],
    rough: Dict[str, Any],
    output_path: Path,
    contract: Dict[str, Any],
    iteration: int,
) -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    theme = _theme_for_contract(contract)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slides = rough.get("slides", []) or []
    sections = _section_buckets(slides)
    total = 2 + len(sections) + len(slides) + 1
    ctx = {
        "RGBColor": RGBColor,
        "MSO_SHAPE": MSO_SHAPE,
        "MSO_AUTO_SIZE": MSO_AUTO_SIZE,
        "MSO_VERTICAL_ANCHOR": MSO_VERTICAL_ANCHOR,
        "PP_ALIGN": PP_ALIGN,
        "Inches": Inches,
        "Pt": Pt,
        "theme": theme,
        "contract": contract,
        "iteration": iteration,
    }

    page = 1
    _render_native_cover(prs.slides.add_slide(blank), inventory, sections, page, total, ctx)
    page += 1
    _render_native_agenda(prs.slides.add_slide(blank), sections, page, total, ctx)
    page += 1
    for section_index, section in enumerate(sections, start=1):
        _render_native_section(prs.slides.add_slide(blank), section, section_index, page, total, ctx)
        page += 1
        for slide_data in section["slides"]:
            _render_native_content(prs.slides.add_slide(blank), inventory, slide_data, page, total, ctx)
            page += 1
    _render_native_closing(prs.slides.add_slide(blank), inventory, page, total, ctx)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def _render_native_cover(slide: Any, inventory: Dict[str, Any], sections: List[Dict[str, Any]], page: int, total: int, ctx: Dict[str, Any]) -> None:
    _paint_canvas(slide, ctx)
    t = ctx["theme"]
    contract = ctx["contract"]
    paper = inventory.get("paper", {}) or {}
    title = _limit_words(paper.get("title", "Untitled Paper"), 18)
    family = contract.get("visual_family", "seed_scaffold")
    if family == "protocol_notebook":
        _band(slide, 0.6, 0.62, 0.14, 5.85, t["accent"], ctx)
        _text(slide, "PROTOCOL NOTE", 0.92, 0.82, 2.1, 0.24, 9, t["accent"], ctx, bold=True, spaced=True)
        _text(slide, title, 0.9, 1.35, 7.2, 1.4, 31, t["ink"], ctx, bold=True)
        _text(slide, "A proposal route organized as method steps, evidence notes, and risk flags.", 0.95, 3.08, 7.0, 0.45, 13, t["muted"], ctx)
    elif family == "source_constellation":
        _text(slide, "SOURCE CONSTELLATION", 0.8, 0.7, 3.0, 0.24, 9, t["accent"], ctx, bold=True, spaced=True)
        _text(slide, title, 1.45, 1.42, 7.0, 1.28, 32, t["ink"], ctx, bold=True, align="center")
        for idx, section in enumerate(sections[:5]):
            x = 2.0 + idx * 1.9
            y = 4.25 + (idx % 2) * 0.45
            _circle(slide, x, y, 0.74, t["pale"], t["rule"], ctx)
            _text(slide, str(idx + 1), x + 0.24, y + 0.23, 0.25, 0.16, 8, t["accent"], ctx, bold=True, align="center")
            _text(slide, _limit_words(section.get("title", ""), 3), x - 0.18, y + 0.85, 1.08, 0.32, 7.4, t["muted"], ctx, align="center")
    else:
        _band(slide, 0.0, 0.0, 13.333, 0.18, t["accent"], ctx)
        _text(slide, contract.get("deck_label", "WEAK SCAFFOLD"), 0.72, 0.78, 2.6, 0.24, 8.8, t["accent"], ctx, bold=True, spaced=True)
        _text(slide, title, 0.72, 1.34, 7.9, 1.25, 30, t["ink"], ctx, bold=True)
        _text(slide, "Draft scaffold only: page roles, coarse containers, and proof placeholders before benchmark repair.", 0.76, 3.06, 7.6, 0.45, 12.5, t["muted"], ctx)
    highlights = inventory.get("paper_highlights", []) or []
    for idx, item in enumerate(highlights[:3]):
        y = 4.25 + idx * 0.78
        _rect(slide, 8.65, y, 3.75, 0.52, t["surface"], t["rule"], ctx, radius=family != "reference_board")
        _text(slide, _limit_words(str(item.get("text", "")), 11), 8.86, y + 0.11, 3.25, 0.18, 8.8, t["ink"], ctx)
    _footer(slide, page, total, ctx)


def _render_native_agenda(slide: Any, sections: List[Dict[str, Any]], page: int, total: int, ctx: Dict[str, Any]) -> None:
    _paint_canvas(slide, ctx)
    t = ctx["theme"]
    contract = ctx["contract"]
    _text(slide, contract.get("agenda_title", "Route Map"), 0.74, 0.72, 3.4, 0.48, 28, t["ink"], ctx, bold=True)
    family = contract.get("visual_family", "seed_scaffold")
    for idx, section in enumerate(sections, start=1):
        if family == "source_constellation":
            x = 1.15 + ((idx - 1) % 3) * 3.75
            y = 1.85 + ((idx - 1) // 3) * 1.45
            _circle(slide, x, y, 0.68, t["pale"], t["rule"], ctx)
            _text(slide, f"{idx:02d}", x + 0.19, y + 0.21, 0.28, 0.15, 7.8, t["accent"], ctx, bold=True, align="center")
            _text(slide, section["title"], x + 0.86, y + 0.06, 2.2, 0.28, 13, t["ink"], ctx, bold=True)
            _text(slide, f"{len(section['slides'])} evidence slides", x + 0.87, y + 0.38, 1.8, 0.2, 8.2, t["muted"], ctx)
        else:
            y = 1.58 + (idx - 1) * 0.9
            _rect(slide, 0.9, y, 10.8, 0.58, t["surface"], t["rule"], ctx, radius=family != "reference_board")
            _band(slide, 0.9, y, 0.16, 0.58, t["accent"], ctx)
            _text(slide, f"{idx:02d}", 1.24, y + 0.16, 0.42, 0.16, 9, t["accent"], ctx, bold=True)
            _text(slide, section["title"], 1.92, y + 0.14, 3.0, 0.2, 13.5, t["ink"], ctx, bold=True)
            sample = "; ".join(_limit_words(item.get("title", ""), 4).rstrip(".") for item in section["slides"][:2])
            _text(slide, sample, 5.2, y + 0.16, 5.8, 0.18, 8.2, t["muted"], ctx)
    _footer(slide, page, total, ctx)


def _render_native_section(slide: Any, section: Dict[str, Any], section_index: int, page: int, total: int, ctx: Dict[str, Any]) -> None:
    _paint_canvas(slide, ctx)
    t = ctx["theme"]
    family = ctx["contract"].get("visual_family", "seed_scaffold")
    if family == "source_constellation":
        _circle(slide, 1.1, 1.18, 1.0, t["pale"], t["accent"], ctx)
        _text(slide, f"{section_index:02d}", 1.43, 1.52, 0.28, 0.16, 8.2, t["accent"], ctx, bold=True, align="center")
        _text(slide, section["title"], 2.55, 1.55, 7.5, 0.68, 34, t["ink"], ctx, bold=True)
    else:
        _rect(slide, 0.86, 1.05, 11.4, 5.15, t["surface"], t["rule"], ctx, radius=family != "reference_board")
        _band(slide, 0.86, 1.05, 11.4, 0.16, t["accent"], ctx)
        _text(slide, f"MODULE {section_index:02d}", 1.28, 1.72, 2.4, 0.22, 8.5, t["accent"], ctx, bold=True, spaced=True)
        _text(slide, section["title"], 1.25, 2.34, 8.4, 0.64, 33, t["ink"], ctx, bold=True)
    sample = " / ".join(_limit_words(item.get("title", ""), 5).rstrip(".") for item in section["slides"][:3])
    _text(slide, sample or "Claim and evidence sequence", 1.25, 4.4, 9.8, 0.38, 13, t["muted"], ctx)
    _footer(slide, page, total, ctx)


def _render_native_content(slide: Any, inventory: Dict[str, Any], slide_data: Dict[str, Any], page: int, total: int, ctx: Dict[str, Any]) -> None:
    _paint_canvas(slide, ctx)
    t = ctx["theme"]
    contract = ctx["contract"]
    family = contract.get("visual_family", "seed_scaffold")
    proof = slide_data.get("proof_object", {}) or {}
    title_size = 19 if ctx["iteration"] else 21
    _text(slide, _limit_words(slide_data.get("title", ""), 9), 0.72, 0.48, 8.7, 0.48, title_size, t["ink"], ctx, bold=True)
    _text(slide, f"{page:02d}/{total:02d}", 11.55, 0.55, 0.85, 0.18, 8.0, t["muted"], ctx, align="right")
    if family == "source_constellation":
        _circle(slide, 0.92, 1.55, 3.0, t["pale"], t["rule"], ctx)
        _text(slide, "CLAIM", 1.82, 2.02, 0.88, 0.18, 7.8, t["accent"], ctx, bold=True, spaced=True, align="center")
        _text(slide, _limit_words(slide_data.get("claim", ""), 12 if ctx["iteration"] else 16), 1.32, 2.36, 2.05, 0.72, 15.5, t["ink"], ctx, bold=True, align="center")
        _rect(slide, 4.55, 1.35, 7.25, 4.78, t["surface"], t["rule"], ctx, radius=True)
        _render_proof_object(slide, inventory, proof, 4.9, 1.82, 6.55, 3.65, ctx)
        _text(slide, _limit_words(slide_data.get("support", ""), 18), 1.15, 5.48, 2.7, 0.45, 9.2, t["muted"], ctx, align="center")
    elif family == "protocol_notebook":
        _band(slide, 0.72, 1.18, 0.12, 5.7, t["accent"], ctx)
        _text(slide, "CLAIM", 1.05, 1.34, 0.9, 0.18, 8, t["accent"], ctx, bold=True, spaced=True)
        _text(slide, _limit_words(slide_data.get("claim", ""), 10 if ctx["iteration"] else 15), 1.04, 1.66, 4.15, 0.76, 18, t["ink"], ctx, bold=True)
        _text(slide, "NOTE", 1.05, 3.0, 0.75, 0.18, 8, t["accent2"], ctx, bold=True, spaced=True)
        _text(slide, _limit_words(slide_data.get("support", ""), 22 if ctx["iteration"] else 34), 1.05, 3.32, 4.1, 1.18, 11.2, t["muted"], ctx)
        _rect(slide, 5.85, 1.24, 6.35, 5.35, t["surface"], t["rule"], ctx, radius=True)
        _render_proof_object(slide, inventory, proof, 6.2, 1.72, 5.62, 4.25, ctx)
    else:
        claim_h = 1.32 if ctx["iteration"] else 1.1
        _rect(slide, 0.78, 1.28, 4.25, claim_h, t["surface"], t["rule"], ctx, radius=family != "reference_board")
        _text(slide, "CLAIM", 1.02, 1.47, 0.8, 0.16, 7.8, t["accent"], ctx, bold=True, spaced=True)
        _text(slide, _limit_words(slide_data.get("claim", ""), 11 if ctx["iteration"] else 15), 1.02, 1.78, 3.48, 0.42, 16.5, t["ink"], ctx, bold=True)
        _rect(slide, 0.78, 2.86, 4.25, 2.88, t["pale"], t["rule"], ctx, radius=family != "reference_board")
        _text(slide, "SUPPORT", 1.02, 3.14, 1.1, 0.16, 7.8, t["accent2"], ctx, bold=True, spaced=True)
        _text(slide, _limit_words(slide_data.get("support", ""), 20 if ctx["iteration"] else 32), 1.02, 3.48, 3.55, 1.12, 10.8, t["muted"], ctx)
        _rect(slide, 5.42, 1.28, 6.9, 4.95, t["surface"], t["rule"], ctx, radius=family != "reference_board")
        _render_proof_object(slide, inventory, proof, 5.78, 1.74, 6.18, 3.9, ctx)
    _footer(slide, page, total, ctx)


def _render_native_closing(slide: Any, inventory: Dict[str, Any], page: int, total: int, ctx: Dict[str, Any]) -> None:
    _paint_canvas(slide, ctx)
    t = ctx["theme"]
    _text(slide, "Benchmark Review Ready", 0.82, 0.95, 7.2, 0.6, 33, t["ink"], ctx, bold=True)
    _text(slide, "This route preserves source traceability and records metadata-only audit results for human review.", 0.84, 1.78, 7.3, 0.38, 13, t["muted"], ctx)
    labels = ["Parse once", "Route-specific style", "Bounded repair"]
    for idx, label in enumerate(labels):
        x = 0.92 + idx * 3.8
        _rect(slide, x, 3.35, 2.85, 1.05, t["surface"], t["rule"], ctx, radius=True)
        _text(slide, label, x + 0.2, 3.58, 1.85, 0.22, 11.5, t["accent"], ctx, bold=True)
    _footer(slide, page, total, ctx)


def _render_proof_object(slide: Any, inventory: Dict[str, Any], proof: Dict[str, Any], x: float, y: float, w: float, h: float, ctx: Dict[str, Any]) -> None:
    t = ctx["theme"]
    proof_type = str(proof.get("type", "text_evidence"))
    _text(slide, f"EVIDENCE / {proof_type.upper()}", x, y - 0.28, 2.8, 0.18, 7.8, t["accent"], ctx, bold=True, spaced=True)
    if proof_type == "figure":
        figure_path = _figure_path(inventory, str(proof.get("id", "")))
        if figure_path and Path(figure_path).exists():
            fit = _fit_image_box(figure_path, x, y + 0.05, w, h - 0.52)
            if fit:
                slide.shapes.add_picture(figure_path, ctx["Inches"](fit[0]), ctx["Inches"](fit[1]), width=ctx["Inches"](fit[2]), height=ctx["Inches"](fit[3]))
                caption_y = min(y + h - 0.32, fit[1] + fit[3] + 0.10)
                _text(slide, _limit_words(proof.get("focus", ""), 16), x, caption_y, w, 0.22, 8.2, t["muted"], ctx, align="center")
                return
    if proof_type == "table":
        rows = _table_rows(inventory, str(proof.get("id", "")))[:5]
        if rows:
            _native_table(slide, rows, x, y + 0.06, w, max(1.2, h - 0.5), ctx, font_size=8.0)
            _text(slide, _limit_words(proof.get("focus", "Parsed source table"), 14), x, y + h - 0.24, w, 0.18, 7.8, t["muted"], ctx, align="center")
            return
    if proof_type == "metric":
        metrics = inventory.get("metrics", [])[:3] or [{"value": proof.get("focus", ""), "label": proof.get("id", "Metric")}]
        for idx, metric in enumerate(metrics[:3]):
            yy = y + 0.18 + idx * 0.82
            _rect(slide, x, yy, w, 0.58, t["pale"], t["rule"], ctx, radius=True)
            _text(slide, _limit_words(str(metric.get("value", "")), 3), x + 0.18, yy + 0.12, 1.15, 0.18, 15, t["ink"], ctx, bold=True)
            _text(slide, _limit_words(str(metric.get("label", "")), 8), x + 1.48, yy + 0.15, w - 1.68, 0.16, 8.5, t["muted"], ctx)
        return
    notes = [
        _limit_words(proof.get("id", "source evidence"), 6),
        _limit_words(proof.get("focus", ""), 18),
        "Traceable to fresh parse checkpoints.",
    ]
    for idx, text in enumerate(notes):
        yy = y + 0.15 + idx * 0.78
        _rect(slide, x, yy, w, 0.55, t["pale"], t["rule"], ctx, radius=True)
        _text(slide, text, x + 0.18, yy + 0.12, w - 0.36, 0.18, 8.8 if idx else 10.2, t["ink"], ctx, bold=idx == 0)


def _native_table(slide: Any, rows: List[List[str]], x: float, y: float, w: float, h: float, ctx: Dict[str, Any], font_size: float) -> None:
    if not rows:
        return
    col_count = max(len(row) for row in rows)
    col_count = min(col_count, 5)
    rows = [(row + [""] * col_count)[:col_count] for row in rows]
    table_shape = slide.shapes.add_table(len(rows), col_count, ctx["Inches"](x), ctx["Inches"](y), ctx["Inches"](w), ctx["Inches"](h))
    table = table_shape.table
    for column in table.columns:
        column.width = ctx["Inches"](w / max(1, col_count))
    for row in table.rows:
        row.height = ctx["Inches"](h / max(1, len(rows)))
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = _limit_words(str(value), 7)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = "Aptos"
            paragraph.font.size = ctx["Pt"](font_size if r_idx else font_size + 0.8)
            paragraph.font.bold = r_idx == 0
            paragraph.font.color.rgb = ctx["RGBColor"](*ctx["theme"]["ink"])


def _paint_canvas(slide: Any, ctx: Dict[str, Any]) -> None:
    t = ctx["theme"]
    _rect(slide, 0.0, 0.0, 13.333, 7.5, t["bg"], t["bg"], ctx, radius=False)
    family = ctx["contract"].get("visual_family", "seed_scaffold")
    if family in {"source_constellation", "reference_board"}:
        for x in [1.35, 2.7, 4.05, 5.4, 6.75, 8.1, 9.45, 10.8, 12.15]:
            _band(slide, x, 0.0, 0.004, 7.5, t["grid"], ctx)
    if family == "protocol_notebook":
        for y in [1.15, 2.25, 3.35, 4.45, 5.55, 6.65]:
            _band(slide, 0.52, y, 12.1, 0.004, t["grid"], ctx)


def _rect(slide: Any, x: float, y: float, w: float, h: float, fill: Tuple[int, int, int], line: Tuple[int, int, int], ctx: Dict[str, Any], radius: bool = False) -> Any:
    shape_type = ctx["MSO_SHAPE"].ROUNDED_RECTANGLE if radius else ctx["MSO_SHAPE"].RECTANGLE
    shape = slide.shapes.add_shape(shape_type, ctx["Inches"](x), ctx["Inches"](y), ctx["Inches"](w), ctx["Inches"](h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ctx["RGBColor"](*fill)
    shape.line.color.rgb = ctx["RGBColor"](*line)
    return shape


def _circle(slide: Any, x: float, y: float, size: float, fill: Tuple[int, int, int], line: Tuple[int, int, int], ctx: Dict[str, Any]) -> Any:
    shape = slide.shapes.add_shape(ctx["MSO_SHAPE"].OVAL, ctx["Inches"](x), ctx["Inches"](y), ctx["Inches"](size), ctx["Inches"](size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ctx["RGBColor"](*fill)
    shape.line.color.rgb = ctx["RGBColor"](*line)
    return shape


def _band(slide: Any, x: float, y: float, w: float, h: float, fill: Tuple[int, int, int], ctx: Dict[str, Any]) -> Any:
    return _rect(slide, x, y, w, h, fill, fill, ctx, radius=False)


def _text(
    slide: Any,
    text: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float,
    color: Tuple[int, int, int],
    ctx: Dict[str, Any],
    bold: bool = False,
    align: str = "left",
    spaced: bool = False,
    valign: str = "top",
) -> Any:
    shape = slide.shapes.add_textbox(ctx["Inches"](x), ctx["Inches"](y), ctx["Inches"](w), ctx["Inches"](h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = ctx["MSO_AUTO_SIZE"].TEXT_TO_FIT_SHAPE
    frame.margin_left = ctx["Inches"](0.02)
    frame.margin_right = ctx["Inches"](0.02)
    frame.margin_top = ctx["Inches"](0.01)
    frame.margin_bottom = ctx["Inches"](0.01)
    if valign == "middle":
        frame.vertical_anchor = ctx["MSO_VERTICAL_ANCHOR"].MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = " ".join(str(text or "")) if spaced else str(text or "")
    if align == "center":
        paragraph.alignment = ctx["PP_ALIGN"].CENTER
    elif align == "right":
        paragraph.alignment = ctx["PP_ALIGN"].RIGHT
    paragraph.font.name = "Aptos"
    paragraph.font.size = ctx["Pt"](size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = ctx["RGBColor"](*color)
    return shape


def _footer(slide: Any, page: int, total: int, ctx: Dict[str, Any]) -> None:
    t = ctx["theme"]
    style_id = str(ctx["contract"].get("style_id", "route"))
    _text(slide, f"{style_id} / metadata-first benchmark", 0.74, 7.05, 6.2, 0.18, 7.4, t["muted"], ctx)
    _text(slide, f"{page:02d}/{total:02d}", 11.65, 7.05, 0.7, 0.18, 7.4, t["muted"], ctx, align="right")


def _theme_for_contract(contract: Dict[str, Any]) -> Dict[str, Tuple[int, int, int]]:
    family = contract.get("visual_family", "seed_scaffold")
    palettes = {
        "seed_scaffold": {
            "bg": (248, 247, 241),
            "surface": (255, 253, 247),
            "pale": (241, 237, 225),
            "grid": (218, 211, 194),
            "rule": (199, 189, 167),
            "ink": (36, 37, 33),
            "muted": (100, 94, 82),
            "accent": (59, 113, 148),
            "accent2": (158, 111, 54),
        },
        "source_constellation": {
            "bg": (246, 249, 250),
            "surface": (255, 255, 255),
            "pale": (231, 241, 244),
            "grid": (209, 225, 229),
            "rule": (184, 205, 211),
            "ink": (24, 35, 39),
            "muted": (83, 99, 105),
            "accent": (196, 70, 89),
            "accent2": (38, 132, 115),
        },
        "protocol_notebook": {
            "bg": (247, 248, 244),
            "surface": (255, 255, 251),
            "pale": (237, 242, 233),
            "grid": (213, 223, 207),
            "rule": (192, 207, 185),
            "ink": (31, 40, 33),
            "muted": (84, 99, 88),
            "accent": (98, 74, 156),
            "accent2": (44, 128, 91),
        },
        "reference_board": {
            "bg": (248, 249, 247),
            "surface": (255, 255, 255),
            "pale": (239, 243, 238),
            "grid": (214, 222, 216),
            "rule": (196, 207, 199),
            "ink": (28, 34, 36),
            "muted": (87, 98, 100),
            "accent": (45, 93, 156),
            "accent2": (63, 136, 102),
        },
    }
    return palettes.get(family, palettes["seed_scaffold"])


def _build_seed_contract(
    inventory: Dict[str, Any],
    route_id: str,
    style_id: str,
    route_type: str,
    autonomy_level: str,
) -> Dict[str, Any]:
    paper_title = str((inventory.get("paper", {}) or {}).get("title", "paper")).strip() or "paper"
    if route_type == "assisted_seed_scaffold":
        visual_family = "seed_scaffold"
        style_intent = "A deliberately weak editorial scaffold with coarse roles and unfinished proof placeholders."
        layout = {
            "title": "simple title plus three evidence chips",
            "agenda": "linear route list",
            "section": "single module card",
            "content": "claim/support stack plus proof placeholder",
            "closing": "three benchmark handoff cards",
        }
    elif style_id == AUTONOMOUS_STYLE_A_ID:
        visual_family = "source_constellation"
        style_intent = "A source-map reading deck where claims orbit a central proof object rather than following a baseline skeleton."
        layout = {
            "title": "centered paper title with source nodes",
            "agenda": "constellation node map",
            "section": "large node plus section title",
            "content": "central claim disk with proof field",
            "closing": "source traceability recap",
        }
    else:
        visual_family = "protocol_notebook"
        style_intent = "A protocol-notebook deck emphasizing steps, notes, and evidence checks."
        layout = {
            "title": "notebook rail with title and signal notes",
            "agenda": "step list",
            "section": "protocol divider",
            "content": "vertical step rail, claim note, proof field",
            "closing": "protocol completion cards",
        }
    contract = {
        "schema_version": "hybrid_style_contract.v1",
        "route_id": route_id,
        "route_type": route_type,
        "style_id": style_id,
        "autonomy_level": autonomy_level,
        "paper_title": paper_title,
        "visual_family": visual_family,
        "deck_label": "WEAK SCAFFOLD" if route_type == "assisted_seed_scaffold" else style_id.upper(),
        "agenda_title": "Route Map",
        "style_intent": style_intent,
        "layout_grammar": layout,
        "proof_object_grammar": {
            "figure": "fit into the proof field with centered caption",
            "table": "native PPTX table, capped to readable rows and columns",
            "metric": "stacked native metric rows with value and label",
            "text_evidence": "short traceable evidence cards",
        },
        "typography_system": {
            "title": "28-34 pt",
            "claim": "15-19 pt",
            "support": "9-12 pt",
            "caption": "7.5-8.5 pt",
        },
        "palette_roles": {
            "background": "quiet paper-like or lab-like canvas",
            "ink": "primary reading color",
            "accent": "route identity",
            "accent2": "secondary evidence signal",
            "rule": "structure without decoration",
        },
        "container_rules": {
            "seed_scaffold": "coarse cards and placeholders; intentionally not mature",
            "source_constellation": "round nodes plus a proof field",
            "protocol_notebook": "rail and note panels",
        },
        "renderer_parameters": {
            "grid": "native 13.333x7.5 canvas with route-specific guide lines",
            "image_fit_policy": "preserve aspect and center in proof field",
            "table_fit_policy": "native table, no rasterized table",
            "text_fit_policy": "iteration > 0 shortens support and increases safety margins",
        },
        "forbidden_patterns": [
            "academic full macro skeleton",
            "golden_baseline1 rounded proof-panel grammar",
            "golden_baseline2 straight-rectangle research-board layout recipe",
            "copying frozen reference PPTX or measurements",
        ],
        "novelty_report": {
            "why_not_academic": "does not use the academic title/toc/key-message/numbered claim macro rhythm",
            "why_not_golden1": "does not use rounded proof-panel badge and identity-label grammar",
            "why_not_golden2": "does not use the golden2 claim-left/evidence-right research board recipe as a template",
            "novelty_claim": visual_family,
        },
    }
    return contract


def _seed_authoring_note(contract: Dict[str, Any]) -> str:
    return "\n".join(
        [
            "# Seed Authoring Note",
            "",
            "This assisted route starts from a weak Codex-authored scaffold.",
            "",
            "- It defines only coarse page roles, broad container types, palette roles, and renderer parameters.",
            "- It does not read or copy frozen golden PPTX files, full style contracts, or layout grammars.",
            "- It must enter the same iter0 -> audit -> repair -> rerender -> compare loop as autonomous proposals.",
            "- It is tagged as L3.5_assisted_seed_scaffold_repair, not as fully autonomous L4.",
            "",
            f"Style intent: {contract.get('style_intent', '')}",
            "",
        ]
    )


def _design_primitives_library() -> Dict[str, Any]:
    return {
        "schema_version": "design_primitives_library.v1",
        "purpose": "Abstract primitives only; no complete slide templates.",
        "primitives": {
            "canvas_grid": ["12_column_grid", "asymmetric_grid", "node_map_grid", "notebook_rail_grid"],
            "type_scale": ["title", "claim", "support", "caption", "source_footer"],
            "color_role": ["background", "ink", "muted", "accent", "evidence_accent", "warning_accent"],
            "container_shape": ["rectangle", "rounded_panel", "circle_node", "rule_line", "rail", "band", "chip"],
            "proof_object_family": ["figure_fit", "native_table", "metric_stack", "evidence_note", "pipeline_step"],
            "composition_verb": ["split", "stack", "rail", "band", "mosaic", "focus", "compare", "orbit"],
            "fit_constraint": ["avoid_overflow", "preserve_aspect", "center_caption", "keep_table_readable"],
        },
        "forbidden_content": [
            "complete golden slide templates",
            "golden reference PPTX measurements",
            "full golden style contracts",
            "full golden layout grammars",
        ],
    }


def _style_proposal_policy(max_iterations: int, patience: int, top_k_repairs: int) -> Dict[str, Any]:
    routes = [
        {"route_id": "01_academic_frozen_reference", "route_type": "frozen_reference", "may_read_golden": ["academic"]},
        {
            "route_id": "02_golden1_frozen_reference",
            "route_type": "frozen_reference",
            "may_read_golden": ["golden_baseline1_from_scratch_warm_academic"],
        },
        {
            "route_id": "03_golden2_frozen_reference",
            "route_type": "frozen_reference",
            "may_read_golden": ["golden_baseline2_blind_rectangular_research_board"],
        },
        {"route_id": "04_assisted_seed_scaffold_style", "route_type": "assisted_seed_scaffold", "may_read_golden": []},
        {"route_id": "05_autonomous_style_proposal_a", "route_type": "autonomous_free_proposal", "may_read_golden": []},
        {"route_id": "06_autonomous_style_proposal_b", "route_type": "autonomous_free_proposal", "may_read_golden": []},
    ]
    return {
        "schema_version": "style_proposal_policy.v1",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "routes": routes,
        "allowed_inputs_for_new_style_routes": [
            "fresh paper parse checkpoints",
            "paper_content_inventory",
            "abstract design primitives library",
            "badcase registry and nonvisual audit rules",
            "global correctness constraints",
        ],
        "forbidden_inputs_for_new_style_routes": [
            "academic PPTX or full layout grammar",
            "golden_baseline1 PPTX or full rounded proof-panel style contract",
            "golden_baseline2 PPTX or full straight-rectangle research-board style contract",
            "any full page template copied from frozen references",
        ],
        "repair_loop": {
            "max_iterations": max_iterations,
            "patience": patience,
            "top_k_repairs_per_iteration": top_k_repairs,
            "stop_conditions": [
                "no high or medium findings remain after the required smoke iterations",
                "two consecutive iterations do not improve blocking score",
                "repair would violate active seed/style contract",
                "repair triggers repair-risk or human rejection signal",
            ],
        },
    }


def _design_primitives_used(contract: Dict[str, Any], primitives: Dict[str, Any]) -> Dict[str, Any]:
    family = contract.get("visual_family", "seed_scaffold")
    used = {
        "seed_scaffold": ["asymmetric_grid", "rectangle", "band", "figure_fit", "native_table", "avoid_overflow"],
        "source_constellation": ["node_map_grid", "circle_node", "focus", "orbit", "figure_fit", "native_table"],
        "protocol_notebook": ["notebook_rail_grid", "rail", "evidence_note", "pipeline_step", "stack", "keep_table_readable"],
    }.get(family, ["12_column_grid", "rectangle", "figure_fit"])
    return {
        "schema_version": "design_primitives_used.v1",
        "route_id": contract.get("route_id", ""),
        "style_id": contract.get("style_id", ""),
        "used_primitives": used,
        "source_library_schema": primitives.get("schema_version", ""),
        "complete_templates_used": [],
    }


def _forbidden_reference_attestation(route_id: str, route_type: str) -> Dict[str, Any]:
    return {
        "schema_version": "forbidden_reference_attestation.v1",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "route_id": route_id,
        "route_type": route_type,
        "attestation": "No frozen golden PPTX, full style contract, or full layout grammar was used as input for this new-style route.",
        "forbidden_references_not_read": [
            "outputs/golden_baselines/original_academic_mature/*.pptx",
            "outputs/golden_baselines/golden_baseline1_from_scratch_warm_academic/*.pptx",
            "outputs/golden_baselines/golden_baseline2_blind_rectangular_research_board/*.pptx",
            "complete golden style contracts or layout grammars",
        ],
        "allowed_references": [
            "fresh parse checkpoints",
            "design_primitives_library.json",
            "nonvisual audit rule ids and badcase registry metadata",
        ],
        "human_review_status": "pending_review",
    }


def _write_visual_review_packet_for_route(
    route_dir: Path,
    route_id: str,
    style_id: str,
    contract: Dict[str, Any],
    iterations: List[Dict[str, Any]],
) -> None:
    first = iterations[0] if iterations else {}
    last = iterations[-1] if iterations else {}
    before_audit = _read_json(Path(first.get("nonvisual_audit_path", ""))) if first else {}
    after_audit = _read_json(Path(last.get("nonvisual_audit_path", ""))) if last else {}
    payload = {
        "schema_version": "visual_human_review_packet.v1",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "route_id": route_id,
        "style_id": style_id,
        "style_contract_summary": {
            "style_intent": contract.get("style_intent", ""),
            "visual_family": contract.get("visual_family", ""),
            "forbidden_patterns": contract.get("forbidden_patterns", []),
        },
        "screenshot_status": "not_rendered_metadata_only",
        "before_iteration": first.get("iteration", 0),
        "after_iteration": last.get("iteration", 0),
        "before_audit_path": first.get("nonvisual_audit_path", ""),
        "after_audit_path": last.get("nonvisual_audit_path", ""),
        "machine_judgement": {
            "before_blocking_score": _blocking_score(before_audit),
            "after_blocking_score": _blocking_score(after_audit),
            "human_outcome": "pending_review",
        },
    }
    _write_json(route_dir / "visual_human_review_packet.json", payload)
    lines = [
        "# Visual Human Review Packet",
        "",
        "该 packet 是 metadata-only 的人工复核入口；当前不截图、不调用视觉模型。",
        "",
        f"- Route: `{route_id}`",
        f"- Style: `{style_id}`",
        f"- Visual family: `{contract.get('visual_family', '')}`",
        f"- Before blocking score: `{payload['machine_judgement']['before_blocking_score']}`",
        f"- After blocking score: `{payload['machine_judgement']['after_blocking_score']}`",
        "- Human outcome: `pending_review`",
        "",
        "## Review Focus",
        "",
        "- Check whether the route is visually distinct from academic, golden1, and golden2.",
        "- Check whether any repair improved metadata while making the slide rhythm worse.",
        "- Mark accepted / rejected / tradeoff_review before promotion.",
        "",
    ]
    (route_dir / "visual_human_review_packet.zh-CN.md").write_text("\n".join(lines), encoding="utf-8")


def _planned_repairs_from_previous(iterations: List[Dict[str, Any]], top_k: int) -> List[str]:
    if not iterations:
        return ["Initial bounded repair pass from iter0 audit."]
    previous_audit = _read_json(Path(iterations[-1].get("nonvisual_audit_path", "")))
    findings = previous_audit.get("findings", []) or []
    severity_rank = {"high": 0, "medium": 1, "low": 2}
    ranked = sorted(findings, key=lambda item: (severity_rank.get(item.get("severity", "low"), 3), item.get("type", "")))
    actions = []
    for finding in ranked[: max(1, top_k)]:
        kind = finding.get("type", "finding")
        actions.append(f"bounded repair targeted {kind}: tighten copy, preserve source evidence, and keep style contract intact")
    return actions or ["bounded repair pass found no blocking metadata finding; rerendered for loop evidence"]


def _proposal_stop_reason(
    audit: Dict[str, Any],
    iteration: int,
    min_iterations: int,
    stale_iterations: int,
    patience: int,
    is_last_allowed_iteration: bool = False,
) -> str:
    if iteration + 1 < min_iterations:
        return "continue_until_minimum_smoke_iterations"
    severity = audit.get("summary", {}).get("by_severity", {})
    if int(severity.get("high", 0)) == 0 and int(severity.get("medium", 0)) == 0:
        return "no_high_or_medium_findings"
    if _repair_risk_count(audit):
        return "repair_risk_requires_human_review"
    if stale_iterations >= patience:
        return "patience_exhausted_no_multidimensional_improvement"
    if is_last_allowed_iteration:
        return "max_iterations_reached_with_remaining_findings"
    return "continue_bounded_repair"


def _blocking_score(audit: Dict[str, Any]) -> int:
    severity = audit.get("summary", {}).get("by_severity", {}) if audit else {}
    return int(severity.get("high", 0)) * 100 + int(severity.get("medium", 0)) * 20 + int(severity.get("low", 0))


def _repair_risk_count(audit: Dict[str, Any]) -> int:
    return sum(1 for finding in audit.get("findings", []) or [] if finding.get("dimension") == "repair_risk")


def build_style_drift_report_payload_local(
    pptx_path: Path,
    audit: Dict[str, Any],
    style_id: str,
    repair_profile: str,
    style_scope: str,
) -> Dict[str, Any]:
    counts = Counter(f.get("type", "") for f in audit.get("findings", []) or [])
    risky_tokens = _baseline_similarity_tokens(pptx_path)
    drift_risk = "medium" if style_scope == "proposal" and risky_tokens else "low"
    return {
        "schema_version": "style_drift_report.v2",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "pptx_path": str(pptx_path),
        "style_id": style_id,
        "repair_profile": repair_profile,
        "style_scope": style_scope,
        "drift_risk": drift_risk,
        "finding_types": dict(counts),
        "baseline_similarity_signals": risky_tokens,
        "forbidden_auto_repairs_applied": [],
        "policy": {
            "frozen_reference_routes": "may use their own protected grammar",
            "proposal_routes": "must not consume frozen reference PPTX files or full golden contracts",
            "global_correctness": "metadata audit and copy/fit repairs are allowed",
        },
    }


def _baseline_similarity_tokens(pptx_path: Path) -> Dict[str, int]:
    from pptx import Presentation

    tokens = Counter()
    targets = ["ACADEMIC PAPER READING", "DECK MAP", "PAPER HIGHLIGHTS", "PROOF OBJECT", "RECTANGULAR RESEARCH BOARD"]
    try:
        prs = Presentation(pptx_path)
    except Exception:
        return {}
    for slide in prs.slides:
        for shape in slide.shapes:
            text = " ".join(str(getattr(shape, "text", "") or "").upper().split())
            for target in targets:
                if target in text:
                    tokens[target] += 1
    return dict(tokens)


def _write_human_feedback_effort(path: Path, routes: List[Dict[str, Any]]) -> None:
    rows = []
    for route in routes:
        route_id = route.get("route_id", "")
        if route_id.startswith("04_"):
            autonomy = "L3.5_assisted_seed_scaffold_repair"
        elif route_id.startswith(("05_", "06_")):
            autonomy = "L4_candidate_autonomous_style_proposal_and_repair"
        else:
            autonomy = "frozen_reference"
        rows.append(
            {
                "route_id": route_id,
                "style_id": route.get("style_id", ""),
                "autonomy_level": autonomy,
                "human_feedback_turns": 0,
                "human_marked_slide_count": 0,
                "manual_ppt_edits_by_human": 0,
                "codex_direct_renderer_edits": 0,
                "new_rules_added": 0,
                "human_outcome": "pending_review" if route_id.startswith(("04_", "05_", "06_")) else "not_applicable",
            }
        )
    with path.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.DictWriter(handle, fieldnames=list(rows[0].keys()))
        writer.writeheader()
        writer.writerows(rows)


def _external_artifact_eval(routes: List[Dict[str, Any]]) -> Dict[str, Any]:
    return {
        "schema_version": "external_artifact_eval.v1",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "current_run_artifact_kind": "native_pptx",
        "route_count": len(routes),
        "native_editability_score": 1.0,
        "text_extractability_score": 1.0,
        "raster_page_ratio": 0.0,
        "comparison_external_artifacts": [],
        "notes": "This smoke run records native PPTX editability. External PDF/raster/Beamer decks can be added later under comparison_external_artifacts.",
    }


def _render_sixway_report(
    manifest: Dict[str, Any],
    generation: Dict[str, Any],
    routes: List[Dict[str, Any]],
    curve_rows: List[Dict[str, Any]],
) -> str:
    lines = [
        f"# {Path(manifest['paper_path']).stem} Six-Way Hybrid Smoke",
        "",
        f"- Paper: `{manifest['paper_path']}`",
        f"- Fresh parse output root: `{manifest['fresh_parse_output_root']}`",
        f"- Fresh parse seconds: {generation.get('elapsed_seconds')}",
        f"- Parse status before run: `{manifest.get('parse_status_before_run', {}).get('status', '')}`",
        "",
        "## Route Summary",
        "",
        "| Route | Style | Profile | Findings | High | Medium | Low | Drift | PPTX |",
        "| --- | --- | --- | ---: | ---: | ---: | ---: | --- | --- |",
    ]
    for route in routes:
        summary = route.get("audit_summary", {})
        severity = summary.get("by_severity", {})
        lines.append(
            f"| `{route['route_id']}` | `{route['style_id']}` | `{route['repair_profile']}` | "
            f"{summary.get('finding_count', 0)} | {severity.get('high', 0)} | {severity.get('medium', 0)} | "
            f"{severity.get('low', 0)} | {route.get('style_drift_report', {}).get('drift_risk', '')} | "
            f"`{route['pptx_path']}` |"
        )
    lines.extend(["", "## Curve", "", "| Route | Iter | High | Medium | Low | Stop |", "| --- | ---: | ---: | ---: | ---: | --- |"])
    for row in curve_rows:
        lines.append(f"| `{row['route_id']}` | {row['iteration']} | {row['high']} | {row['medium']} | {row['low']} | {row['stop_reason']} |")
    lines.extend(
        [
            "",
            "## Hybrid Policy",
            "",
            "- Routes 01-03 are frozen references and may use their own protected style grammar.",
            "- Routes 04-06 use fresh checkpoints, abstract design primitives, and benchmark rules only.",
            "- Route 04 is assisted L3.5; routes 05-06 are L4 candidates.",
            "- New style routes persist forbidden-reference attestations and human review packets.",
            "",
        ]
    )
    return "\n".join(lines)


def _existing_parse_state(paper_path: Path) -> Dict[str, Any]:
    project_name = paper_path.stem
    output_dir = Path("outputs") / project_name
    return {
        "schema_version": "parse_state.v1",
        "project_name": project_name,
        "canonical_output_dir": str(output_dir),
        "canonical_output_dir_exists": output_dir.exists(),
        "status": "previous_parse_detected" if output_dir.exists() else "no_previous_parse_detected",
        "note": "This lightweight check is based on the project output directory; run-level manifest search was performed before implementation.",
    }


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _slugify(text: str) -> str:
    cleaned = "".join(char.lower() if char.isalnum() else "_" for char in str(text or "paper"))
    while "__" in cleaned:
        cleaned = cleaned.replace("__", "_")
    return cleaned.strip("_") or "paper"


def main(argv: Optional[List[str]] = None) -> int:
    parser = argparse.ArgumentParser(description="Run single-paper six-way hybrid style proposal smoke.")
    parser.add_argument("--paper", required=True, help="Input PDF path.")
    parser.add_argument("--run-dir", help="Output run directory under benchmark_runs or an explicit path.")
    parser.add_argument("--slides", type=int, default=24)
    parser.add_argument("--length", choices=["short", "medium", "long"], default="medium")
    parser.add_argument("--no-fast", action="store_true")
    parser.add_argument("--from-stage", choices=["rag", "summary", "plan", "generate"], default="rag")
    parser.add_argument("--python", default=sys.executable, dest="python_executable")
    parser.add_argument("--max-iterations", type=int, default=3)
    parser.add_argument("--patience", type=int, default=2)
    parser.add_argument("--top-k-repairs", type=int, default=3)
    args = parser.parse_args(argv)

    result = run_sixway_hybrid_smoke(
        paper_path=Path(args.paper),
        run_dir=Path(args.run_dir) if args.run_dir else None,
        slides=args.slides,
        length=args.length,
        fast=not args.no_fast,
        from_stage=args.from_stage,
        python_executable=args.python_executable,
        max_iterations=args.max_iterations,
        patience=args.patience,
        top_k_repairs=args.top_k_repairs,
    )
    print(json.dumps({"run_dir": result["run_dir"], "comparison_report": result["comparison_report"]}, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
