"""Single-paper four-way validation harness for from-scratch benchmark work."""
from __future__ import annotations

import argparse
import csv
import hashlib
import json
import math
import os
import re
import shutil
import subprocess
import sys
import time
from collections import Counter
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple

from paper2slides.core.paths import get_base_dir, get_config_dir, get_plan_checkpoint, get_summary_checkpoint
from paper2slides.generator.pptx_qa import evaluate_presentation_spec, inspect_pptx_layout
from paper2slides.generator.pptx_renderer import PptxRenderer
from paper2slides.generator.slide_schema import PresentationSpec
from paper2slides.generator.text_pptx_workflow import _build_speaker_script
from paper2slides.utils.path_utils import get_project_name, parse_style

from .from_scratch import build_content_inventory, build_rough_draft_spec, write_from_scratch_artifacts
from .nonvisual_audit import inspect_pptx_nonvisual


DEFAULT_REPORT_ROOT = Path("benchmark_runs")

GLOBAL_CORRECTNESS_RULES = {
    "estimated_text_overflow",
    "text_exceeds_container_bounds",
    "shape_overlap_risk",
    "weak_table_grammar",
    "dense_table_readability_risk",
    "table_exceeds_container_bounds",
    "table_container_height_mismatch",
    "table_sparse_columns_rendered",
    "table_cell_text_wrapping_risk",
    "figure_picture_aspect_distortion",
    "low_font_size",
}

GOLDEN_BASELINE1_SCOPED_RULES = {
    "agenda_read_path_header_too_close",
    "card_internal_spacing_not_scaled_to_frame",
    "figure_badge_identity_label_conflation",
    "figure_panel_aspect_mismatch",
    "figure_image_off_center_in_panel",
    "figure_label_anchor_drift",
    "figure_label_text_alignment_off_center",
    "panel_identity_label_anchor_drift",
    "panel_identity_label_text_alignment_off_center",
    "stacked_figure_identity_label_overcorrection",
    "table_support_band_off_balance",
}

STYLE_AWARE_REVIEW_RULES = {
    "figure_panel_aspect_mismatch",
    "image_underutilized_in_wide_panel",
    "figure_caption_not_centered_in_wide_panel",
    "table_container_height_mismatch",
    "table_view_label_missing",
    "table_caption_missing_or_not_centered",
    "table_underutilized_in_evidence_panel",
    "text_card_vertical_alignment_top_heavy",
}

REPAIR_RISK_RULES = {
    "metric_improved_visual_regressed",
    "likely_overcorrection",
    "style_scope_mismatch",
    "repair_introduced_new_findings",
    "image_legibility_regression",
    "layout_rhythm_regression",
}

EXPERIMENTAL_STYLE_PREFIX = "blind_experimental"
BLIND_RECTANGULAR_STYLE_ID = "blind_rectangular_research_board"
GOLDEN_BASELINE2_STYLE_ID = "golden_baseline2_blind_rectangular_research_board"

STYLE_CONTRACTS = {
    "academic": {
        "schema_version": "style_contract.v1",
        "style_id": "academic",
        "promotion_status": "protected_mature_baseline",
        "container_shape": "mixed_academic_baseline",
        "container_padding_model": {"horizontal_in": 0.28, "vertical_in": 0.20, "space_efficiency": "medium"},
        "preferred_layouts": {
            "table": "bottom_table_with_right_evidence_region",
            "figure": "academic_visual_region",
            "metric": "numbered_claim_or_metric_card",
        },
        "protected_rhythm": [
            "title, contents, section divider, content, and closing rhythm",
            "key message card plus numbered claim/support blocks",
            "do not inherit rounded proof-panel polish unless the style contract explicitly changes",
        ],
    },
    "golden_baseline1_from_scratch_warm_academic": {
        "schema_version": "style_contract.v1",
        "style_id": "golden_baseline1_from_scratch_warm_academic",
        "promotion_status": "promoted_golden_reference",
        "container_shape": "rounded_proof_panel",
        "container_padding_model": {"horizontal_in": 0.30, "vertical_in": 0.24, "space_efficiency": "medium"},
        "preferred_layouts": {
            "moderately_wide_figure": "aspect_aware_proof_panel",
            "extremely_wide_figure": "bottom_wide_panel_if_readable",
            "tall_figure": "vertical_proof_panel",
            "table": "native_table_inside_rounded_panel",
        },
        "protected_rhythm": [
            "warm paper background with restrained teal, gold, and clay accents",
            "one claim plus one proof object per content slide",
            "green proof type badge stays at the panel corner while black identity label anchors to fitted content",
        ],
    },
    BLIND_RECTANGULAR_STYLE_ID: {
        "schema_version": "style_contract.v1",
        "style_id": BLIND_RECTANGULAR_STYLE_ID,
        "promotion_status": "promoted_to_golden_baseline2_human_tuned_reference",
        "promoted_style_id": GOLDEN_BASELINE2_STYLE_ID,
        "container_shape": "straight_rectangle",
        "container_padding_model": {"horizontal_in": 0.30, "vertical_in": 0.22, "space_efficiency": "high"},
        "preferred_layouts": {
            "moderately_wide_figure": "right_panel_large",
            "extremely_wide_figure": "bottom_band_if_readable",
            "tall_figure": "asymmetric_left_wide_or_right_narrow",
            "table": "container_fit_with_min_padding",
        },
        "protected_rhythm": [
            "straight rectangular left claim panel pairs cleanly with straight rectangular right evidence panel",
            "right-side large figure is preferred when it preserves readability",
        ],
        "autonomous_proposal_input_policy": "forbidden_as_full_template_or_layout_recipe",
    },
    GOLDEN_BASELINE2_STYLE_ID: {
        "schema_version": "style_contract.v1",
        "style_id": GOLDEN_BASELINE2_STYLE_ID,
        "style_family": BLIND_RECTANGULAR_STYLE_ID,
        "promotion_status": "human_tuned_frozen_reference",
        "container_shape": "straight_rectangle",
        "container_padding_model": {"horizontal_in": 0.24, "vertical_in": 0.18, "space_efficiency": "high"},
        "preferred_layouts": {
            "text_evidence": "right_panel_stacked_cards_with_middle_body_anchor",
            "moderately_wide_figure": "right_panel_large_or_bottom_panel_if_readability_improves",
            "extremely_wide_figure": "bottom_wide_panel_with_large_fitted_image",
            "tall_figure": "large_asymmetric_panel_preserving_source_aspect",
            "dense_table": "focused_table_view_inside_large_evidence_panel",
            "sparse_table": "native_table_fit_with_caption",
        },
        "protected_rhythm": [
            "white research-board canvas with faint vertical grid lines and a restrained top rail",
            "straight rectangular claim and evidence containers with crisp academic spacing",
            "large evidence panels should feel used, not decorative or empty",
        ],
        "autonomous_proposal_input_policy": "forbidden_as_full_template_or_layout_recipe",
    },
}


def run_fourway_validation(
    paper_path: Path,
    run_dir: Optional[Path] = None,
    slides: int = 24,
    length: str = "medium",
    fast: bool = True,
    from_stage: str = "rag",
    python_executable: str = sys.executable,
    max_iterations: int = 2,
) -> Dict[str, Any]:
    """Parse/generate once, then materialize the four validation branches."""
    paper_path = paper_path.resolve()
    if not paper_path.exists():
        raise FileNotFoundError(f"Paper PDF does not exist: {paper_path}")

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    run_dir = run_dir or DEFAULT_REPORT_ROOT / f"agi_wordle_fourway_{timestamp}"
    run_dir = run_dir.resolve()
    output_root = run_dir / "fresh_parse_outputs"
    routes_dir = run_dir / "routes"
    logs_dir = run_dir / "logs"
    for directory in (output_root, routes_dir, logs_dir):
        directory.mkdir(parents=True, exist_ok=True)

    manifest = {
        "schema_version": "single_paper_fourway.v1",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "paper_path": str(paper_path),
        "run_dir": str(run_dir),
        "fresh_parse_output_root": str(output_root),
        "slides": slides,
        "length": length,
        "fast": fast,
        "from_stage": from_stage,
        "python_executable": python_executable,
        "principles": [
            "parse the PDF into a fresh output root",
            "reuse only this run's checkpoints for style branches",
            "auto-repair global correctness only for academic",
            "auto-repair rounded proof-panel polish only for golden_baseline1",
            "keep blind experimental style independent from academic and golden_baseline1 visual grammar",
        ],
        "style_contracts": STYLE_CONTRACTS,
    }
    _write_json(run_dir / "manifest.json", manifest)

    generation = _run_academic_generation(
        paper_path=paper_path,
        output_root=output_root,
        logs_dir=logs_dir,
        slides=slides,
        length=length,
        fast=fast,
        from_stage=from_stage,
        python_executable=python_executable,
    )
    _write_json(run_dir / "fresh_generation_result.json", generation)
    if generation["returncode"] != 0 or not generation.get("output_subdir"):
        raise RuntimeError(
            "Fresh academic generation did not produce a usable output. "
            f"See log: {generation.get('log_path', '')}"
        )

    summary_checkpoint = Path(generation["summary_checkpoint"])
    plan_checkpoint = Path(generation["plan_checkpoint"])
    spec_checkpoint = Path(generation["spec_checkpoint"])
    if not summary_checkpoint.exists() or not plan_checkpoint.exists() or not spec_checkpoint.exists():
        raise RuntimeError("Fresh run is missing one or more checkpoints needed for branch generation.")

    routes: List[Dict[str, Any]] = []
    routes.append(
        _materialize_existing_deck_route(
            source_output=Path(generation["output_subdir"]),
            route_dir=routes_dir / "01_academic_audit_only",
            route_id="academic_audit_only",
            style_id="academic",
            repair_profile="audit_only",
            style_scope="academic",
            note="Original golden baseline route. Nonvisual audit only; no golden_baseline1 polish is auto-applied.",
        )
    )
    routes.append(
        _materialize_golden_baseline1_route(
            summary_checkpoint=summary_checkpoint,
            plan_checkpoint=plan_checkpoint,
            spec_checkpoint=spec_checkpoint,
            route_dir=routes_dir / "02_golden_baseline1_scoped",
        )
    )
    routes.append(
        _materialize_academic_global_repair_route(
            source_output=Path(generation["output_subdir"]),
            spec_checkpoint=spec_checkpoint,
            route_dir=routes_dir / "03_academic_global_repair",
        )
    )
    routes.append(
        _materialize_blind_experimental_route(
            summary_checkpoint=summary_checkpoint,
            plan_checkpoint=plan_checkpoint,
            spec_checkpoint=spec_checkpoint,
            route_dir=routes_dir / "04_blind_experimental_loop",
            max_iterations=max_iterations,
        )
    )

    curve_rows = _write_curve(run_dir / "score_curve.csv", routes)
    artifact_rows = _write_artifact_index(run_dir / "artifact_index.csv", routes)
    report_path = run_dir / "comparison_report.md"
    report_path.write_text(_render_comparison_report(manifest, generation, routes, curve_rows), encoding="utf-8")
    result = {
        "run_dir": str(run_dir),
        "manifest": str(run_dir / "manifest.json"),
        "comparison_report": str(report_path),
        "score_curve": str(run_dir / "score_curve.csv"),
        "artifact_index": str(run_dir / "artifact_index.csv"),
        "fresh_generation": generation,
        "routes": routes,
        "artifact_rows": artifact_rows,
    }
    _write_json(run_dir / "fourway_result.json", result)
    return result


def _run_academic_generation(
    paper_path: Path,
    output_root: Path,
    logs_dir: Path,
    slides: int,
    length: str,
    fast: bool,
    from_stage: str,
    python_executable: str,
) -> Dict[str, Any]:
    staged_paper_path = _stage_paper_with_short_name(paper_path, output_root.parent)
    config = _config_for_style("academic", length, slides, fast)
    project_name = get_project_name(str(staged_paper_path))
    base_dir = get_base_dir(str(output_root), project_name, "paper")
    config_dir = get_config_dir(base_dir, config)
    before = {path.resolve() for path in config_dir.iterdir() if path.is_dir()} if config_dir.exists() else set()
    command = [
        python_executable,
        "-m",
        "paper2slides",
        "--input",
        str(staged_paper_path),
        "--output",
        "slides",
        "--style",
        "academic",
        "--length",
        length,
        "--slides",
        str(slides),
        "--output-dir",
        str(output_root),
    ]
    if fast:
        command.append("--fast")
    if from_stage:
        command.extend(["--from-stage", from_stage])

    log_path = logs_dir / "fresh_academic_generation.log"
    env = os.environ.copy()
    env["PYTHONDONTWRITEBYTECODE"] = "1"
    start = time.perf_counter()
    with log_path.open("w", encoding="utf-8", errors="replace") as log:
        log.write("Command: " + " ".join(command) + "\n\n")
        log.flush()
        process = subprocess.run(command, cwd=Path.cwd(), env=env, stdout=log, stderr=subprocess.STDOUT, text=True)
    elapsed = time.perf_counter() - start
    output_subdir = _find_new_or_latest_output(config_dir, before)
    return {
        "command": command,
        "original_paper_path": str(paper_path),
        "staged_paper_path": str(staged_paper_path),
        "returncode": process.returncode,
        "elapsed_seconds": round(elapsed, 2),
        "log_path": str(log_path),
        "base_dir": str(base_dir),
        "config_dir": str(config_dir),
        "output_subdir": str(output_subdir) if output_subdir else "",
        "pptx_path": str(output_subdir / "slides.pptx") if output_subdir else "",
        "speaker_script_path": str(output_subdir / "speaker_script.md") if output_subdir else "",
        "layout_qa_path": str(output_subdir / "layout_qa.json") if output_subdir else "",
        "summary_checkpoint": str(get_summary_checkpoint(base_dir, config)),
        "plan_checkpoint": str(get_plan_checkpoint(config_dir)),
        "spec_checkpoint": str(config_dir / "checkpoint_slide_spec.json"),
        "error_tail": _tail_text(log_path, 36) if process.returncode != 0 or not output_subdir else "",
    }


def _stage_paper_with_short_name(paper_path: Path, run_dir: Path) -> Path:
    staged_path = run_dir / f"source{paper_path.suffix.lower() or '.pdf'}"
    staged_path.parent.mkdir(parents=True, exist_ok=True)
    if not staged_path.exists() or staged_path.stat().st_size != paper_path.stat().st_size:
        shutil.copy2(paper_path, staged_path)
    return staged_path


def _materialize_existing_deck_route(
    source_output: Path,
    route_dir: Path,
    route_id: str,
    style_id: str,
    repair_profile: str,
    style_scope: str,
    note: str,
) -> Dict[str, Any]:
    route_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = _copy_if_exists(source_output / "slides.pptx", route_dir / "slides.pptx")
    speaker_path = _copy_if_exists(source_output / "speaker_script.md", route_dir / "speaker_script.md")
    speaker_audit_path = route_dir / "speaker_script_audit.json"
    speaker_audit = _write_speaker_script_audit(speaker_path, pptx_path, style_id, speaker_audit_path) if speaker_path else {}
    qa_path = _copy_if_exists(source_output / "layout_qa.json", route_dir / "layout_qa.json")
    audit = _audit_to_file(pptx_path, route_dir / "nonvisual_audit.json")
    repair_log = _repair_log_payload(
        route_id=route_id,
        style_id=style_id,
        repair_profile=repair_profile,
        iterations=[
            _iteration_payload(
                iteration=0,
                pptx_path=pptx_path,
                audit=audit,
                applied_repairs=[],
                stop_reason=_stop_reason_for_profile(audit, repair_profile),
                speaker_script_path=speaker_path,
                speaker_script_audit_path=speaker_audit_path if speaker_audit else None,
            )
        ],
        note=note,
    )
    _write_json(route_dir / "repair_log.json", repair_log)
    style_report = build_style_drift_report_payload(
        pptx_path=pptx_path,
        audit=audit,
        style_id=style_id,
        repair_profile=repair_profile,
        style_scope=style_scope,
    )
    _write_json(route_dir / "style_drift_report.json", style_report)
    return _route_result(route_id, route_dir, style_id, repair_profile, pptx_path, speaker_path, speaker_audit_path if speaker_audit else None, qa_path, audit, repair_log, style_report)


def _materialize_academic_global_repair_route(source_output: Path, spec_checkpoint: Path, route_dir: Path) -> Dict[str, Any]:
    route_dir.mkdir(parents=True, exist_ok=True)
    original_iter = route_dir / "iterations" / "iter_00"
    repaired_iter = route_dir / "iterations" / "iter_01"
    original_iter.mkdir(parents=True, exist_ok=True)
    repaired_iter.mkdir(parents=True, exist_ok=True)

    original_pptx = _copy_if_exists(source_output / "slides.pptx", original_iter / "slides.pptx")
    if original_pptx is None:
        raise RuntimeError(f"Missing source PPTX for academic repair route: {source_output}")
    original_audit = _audit_to_file(original_pptx, original_iter / "nonvisual_audit.json")

    spec = PresentationSpec.from_dict(_read_json(spec_checkpoint))
    repaired_spec, applied = _compact_academic_spec_for_global_repair(spec)
    repaired_pptx = repaired_iter / "slides.pptx"
    PptxRenderer(style="academic").render(repaired_spec, repaired_pptx)
    repaired_audit = _audit_to_file(repaired_pptx, repaired_iter / "nonvisual_audit.json")

    final_pptx = route_dir / "slides.pptx"
    shutil.copy2(repaired_pptx, final_pptx)
    final_audit = _audit_to_file(final_pptx, route_dir / "nonvisual_audit.json")
    speaker_path = route_dir / "speaker_script.md"
    speaker_path.write_text(_build_speaker_script(repaired_spec, applied), encoding="utf-8")
    speaker_audit_path = route_dir / "speaker_script_audit.json"
    _write_speaker_script_audit(speaker_path, final_pptx, "academic", speaker_audit_path)
    layout_result = inspect_pptx_layout(final_pptx)
    qa_report = evaluate_presentation_spec(repaired_spec, layout_result)
    _write_json(route_dir / "layout_qa.json", qa_report)
    _write_json(route_dir / "checkpoint_slide_spec_repaired.json", repaired_spec.to_dict())

    repair_log = _repair_log_payload(
        route_id="academic_global_repair",
        style_id="academic",
        repair_profile="global_correctness_repair",
        iterations=[
            _iteration_payload(0, original_pptx, original_audit, [], _stop_reason_for_profile(original_audit, "global_correctness_repair")),
            _iteration_payload(
                1,
                repaired_pptx,
                repaired_audit,
                applied,
                _stop_reason_for_profile(repaired_audit, "global_correctness_repair"),
                speaker_script_path=speaker_path,
                speaker_script_audit_path=speaker_audit_path,
            ),
        ],
        note="Academic route repaired only through global spec-level correctness/copy fitting; golden_baseline1 polish stayed detect-only.",
    )
    _write_json(route_dir / "repair_log.json", repair_log)
    style_report = build_style_drift_report_payload(
        pptx_path=final_pptx,
        audit=final_audit,
        style_id="academic",
        repair_profile="global_correctness_repair",
        style_scope="academic",
    )
    _write_json(route_dir / "style_drift_report.json", style_report)
    return _route_result(
        "academic_global_repair",
        route_dir,
        "academic",
        "global_correctness_repair",
        final_pptx,
        speaker_path,
        speaker_audit_path,
        route_dir / "layout_qa.json",
        final_audit,
        repair_log,
        style_report,
    )


def _materialize_golden_baseline1_route(
    summary_checkpoint: Path,
    plan_checkpoint: Path,
    spec_checkpoint: Path,
    route_dir: Path,
) -> Dict[str, Any]:
    route_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = route_dir / "slides.pptx"
    paths = write_from_scratch_artifacts(
        summary_checkpoint=summary_checkpoint,
        plan_checkpoint=plan_checkpoint,
        spec_checkpoint=spec_checkpoint,
        output_dir=route_dir,
        pptx_output=pptx_path,
        render_review_dir=None,
    )
    inventory = _read_json(route_dir / "content_inventory.json")
    rough = _read_json(route_dir / "rough_draft_spec.json")
    speaker_path = route_dir / "speaker_script.md"
    speaker_path.write_text(
        _build_rough_speaker_script(
            inventory,
            rough,
            "golden_baseline1 scoped deck",
            narrative_mode="golden_baseline1",
        ),
        encoding="utf-8",
    )
    speaker_audit_path = route_dir / "speaker_script_audit.json"
    _write_speaker_script_audit(speaker_path, pptx_path, "golden_baseline1_from_scratch_warm_academic", speaker_audit_path)
    audit = _read_json(route_dir / "nonvisual_audit.json")
    if not audit:
        audit = _audit_to_file(pptx_path, route_dir / "nonvisual_audit.json")
    repair_log = _repair_log_payload(
        route_id="golden_baseline1_scoped",
        style_id="golden_baseline1_from_scratch_warm_academic",
        repair_profile="golden_baseline1_repair",
        iterations=[
            _iteration_payload(
                iteration=0,
                pptx_path=pptx_path,
                audit=audit,
                applied_repairs=[
                    "Rendered with the current golden_baseline1 rounded proof-panel renderer and its already-scoped v25 badcase fixes."
                ],
                stop_reason=_stop_reason_for_profile(audit, "golden_baseline1_repair"),
                speaker_script_path=speaker_path,
                speaker_script_audit_path=speaker_audit_path,
            )
        ],
        note="Golden baseline1 route generated from this run's fresh checkpoints; rounded proof-panel polish is scoped to this style.",
    )
    repair_log["source_artifacts"] = paths
    _write_json(route_dir / "repair_log.json", repair_log)
    style_report = build_style_drift_report_payload(
        pptx_path=pptx_path,
        audit=audit,
        style_id="golden_baseline1_from_scratch_warm_academic",
        repair_profile="golden_baseline1_repair",
        style_scope="golden_baseline1",
    )
    _write_json(route_dir / "style_drift_report.json", style_report)
    return _route_result(
        "golden_baseline1_scoped",
        route_dir,
        "golden_baseline1_from_scratch_warm_academic",
        "golden_baseline1_repair",
        pptx_path,
        speaker_path,
        speaker_audit_path,
        None,
        audit,
        repair_log,
        style_report,
    )


def _materialize_blind_experimental_route(
    summary_checkpoint: Path,
    plan_checkpoint: Path,
    spec_checkpoint: Path,
    route_dir: Path,
    max_iterations: int,
) -> Dict[str, Any]:
    route_dir.mkdir(parents=True, exist_ok=True)
    inventory = build_content_inventory(summary_checkpoint, plan_checkpoint, spec_checkpoint)
    rough = build_rough_draft_spec(inventory)
    style_contract = _build_blind_style_contract(inventory, route_dir)
    _write_json(route_dir / "content_inventory.json", inventory)
    _write_json(route_dir / "rough_draft_spec.json", rough)
    _write_json(route_dir / "style_contract.json", style_contract)

    iterations: List[Dict[str, Any]] = []
    current_audit: Dict[str, Any] = {}
    final_pptx = route_dir / "slides.pptx"
    final_speaker_path = route_dir / "speaker_script.md"
    final_speaker_audit_path = route_dir / "speaker_script_audit.json"
    for iteration in range(max(1, max_iterations)):
        iteration_dir = route_dir / "iterations" / f"iter_{iteration:02d}"
        iteration_dir.mkdir(parents=True, exist_ok=True)
        iteration_pptx = iteration_dir / "slides.pptx"
        repair_mode = iteration > 0
        render_blind_blueprint_pptx(inventory, rough, iteration_pptx, repair_mode=repair_mode, style_contract=style_contract)
        current_audit = _audit_to_file(iteration_pptx, iteration_dir / "nonvisual_audit.json")
        iteration_speaker_path = iteration_dir / "speaker_script.md"
        iteration_speaker_path.write_text(
            _build_rough_speaker_script(
                inventory,
                rough,
                f"{style_contract['style_id']} iteration {iteration:02d}",
                narrative_mode="blind_experimental",
            ),
            encoding="utf-8",
        )
        iteration_speaker_audit_path = iteration_dir / "speaker_script_audit.json"
        _write_speaker_script_audit(
            iteration_speaker_path,
            iteration_pptx,
            style_contract["style_id"],
            iteration_speaker_audit_path,
        )
        applied = []
        if repair_mode:
            applied.append("benchmark-guided copy fitting: shortened prose, widened gutters, and lifted small evidence text")
        stop_reason = _stop_reason_for_profile(current_audit, "experimental_from_scratch_loop")
        iterations.append(
            _iteration_payload(
                iteration,
                iteration_pptx,
                current_audit,
                applied,
                stop_reason,
                speaker_script_path=iteration_speaker_path,
                speaker_script_audit_path=iteration_speaker_audit_path,
            )
        )
        if iteration == 0 and max_iterations > 1 and _needs_experimental_repair(current_audit):
            continue
        if iteration == 0 and max_iterations > 1:
            # Run one bounded polish pass anyway so the report demonstrates detect -> repair -> rescore.
            continue
        shutil.copy2(iteration_pptx, final_pptx)
        shutil.copy2(iteration_speaker_path, final_speaker_path)
        shutil.copy2(iteration_speaker_audit_path, final_speaker_audit_path)
        break

    if not final_pptx.exists():
        last_pptx = Path(iterations[-1]["pptx_path"])
        shutil.copy2(last_pptx, final_pptx)
        last_speaker = Path(iterations[-1]["speaker_script_path"])
        last_speaker_audit = Path(iterations[-1]["speaker_script_audit_path"])
        shutil.copy2(last_speaker, final_speaker_path)
        shutil.copy2(last_speaker_audit, final_speaker_audit_path)
    final_audit = _audit_to_file(final_pptx, route_dir / "nonvisual_audit.json")
    review_packet = _write_visual_human_review_packet(route_dir, style_contract, iterations)
    repair_log = _repair_log_payload(
        route_id="blind_experimental_loop",
        style_id=style_contract["style_id"],
        repair_profile="experimental_from_scratch_loop",
        iterations=iterations,
        note="Blind route does not use the academic renderer or golden_baseline1 rounded proof-panel grammar; this run gets its own experimental style contract.",
    )
    repair_log["style_contract_path"] = str(route_dir / "style_contract.json")
    repair_log["visual_human_review_packet"] = review_packet
    _write_json(route_dir / "repair_log.json", repair_log)
    style_report = build_style_drift_report_payload(
        pptx_path=final_pptx,
        audit=final_audit,
        style_id=style_contract["style_id"],
        repair_profile="experimental_from_scratch_loop",
        style_scope="experimental",
    )
    _write_json(route_dir / "style_drift_report.json", style_report)
    route_result = _route_result(
        "blind_experimental_loop",
        route_dir,
        style_contract["style_id"],
        "experimental_from_scratch_loop",
        final_pptx,
        final_speaker_path,
        final_speaker_audit_path,
        None,
        final_audit,
        repair_log,
        style_report,
    )
    route_result["visual_human_review_packet"] = review_packet
    route_result["visual_compare_dir"] = str(route_dir / "visual_compare")
    return route_result


def _write_visual_human_review_packet(
    route_dir: Path,
    style_contract: Dict[str, Any],
    iterations: List[Dict[str, Any]],
) -> Dict[str, Any]:
    compare_dir = route_dir / "visual_compare"
    compare_dir.mkdir(parents=True, exist_ok=True)
    before_iter = iterations[0] if iterations else {}
    after_iter = iterations[-1] if iterations else {}
    before_audit = _read_json(Path(before_iter.get("nonvisual_audit_path", ""))) if before_iter else {}
    after_audit = _read_json(Path(after_iter.get("nonvisual_audit_path", ""))) if after_iter else {}
    page_deltas = _page_level_audit_delta(before_audit, after_audit)
    repair_risks = _repair_risk_findings_from_delta(page_deltas, style_contract)
    payload = {
        "schema_version": "visual_human_review_packet.v1",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "route_id": "blind_experimental_loop",
        "style_id": style_contract.get("style_id", ""),
        "style_instance_id": style_contract.get("style_instance_id", ""),
        "style_contract": style_contract,
        "before_iteration": before_iter.get("iteration", 0),
        "after_iteration": after_iter.get("iteration", 0),
        "before_audit_path": before_iter.get("nonvisual_audit_path", ""),
        "after_audit_path": after_iter.get("nonvisual_audit_path", ""),
        "screenshot_status": "not_rendered_metadata_only",
        "expected_visual_artifacts": [
            "visual_compare/iter_00_slide_XX.png",
            "visual_compare/iter_01_slide_XX.png",
            "visual_compare/slide_XX_before_after_human_review.png",
        ],
        "page_level_delta_path": str(compare_dir / "page_level_audit_delta.json"),
        "page_level_delta_csv": str(compare_dir / "page_level_audit_delta.csv"),
        "packet_markdown": str(route_dir / "visual_human_review_packet.md"),
        "packet_markdown_zh": str(route_dir / "visual_human_review_packet.zh-CN.md"),
        "human_outcome": "pending_review",
        "machine_judgement": _machine_judgement(before_audit, after_audit, repair_risks),
        "page_deltas": page_deltas,
        "repair_risk_findings": repair_risks,
    }
    _write_json(compare_dir / "page_level_audit_delta.json", payload)
    _write_delta_csv(compare_dir / "page_level_audit_delta.csv", page_deltas)
    (route_dir / "visual_human_review_packet.md").write_text(_render_visual_packet_md(payload, zh=False), encoding="utf-8")
    (route_dir / "visual_human_review_packet.zh-CN.md").write_text(_render_visual_packet_md(payload, zh=True), encoding="utf-8")
    return {
        "schema_version": payload["schema_version"],
        "visual_compare_dir": str(compare_dir),
        "page_level_delta_path": str(compare_dir / "page_level_audit_delta.json"),
        "page_level_delta_csv": str(compare_dir / "page_level_audit_delta.csv"),
        "packet_markdown": str(route_dir / "visual_human_review_packet.md"),
        "packet_markdown_zh": str(route_dir / "visual_human_review_packet.zh-CN.md"),
        "human_outcome": payload["human_outcome"],
        "machine_judgement": payload["machine_judgement"],
        "repair_risk_count": len(repair_risks),
    }


def _page_level_audit_delta(before_audit: Dict[str, Any], after_audit: Dict[str, Any]) -> List[Dict[str, Any]]:
    before_pages = _findings_by_page(before_audit)
    after_pages = _findings_by_page(after_audit)
    pages = sorted(set(before_pages) | set(after_pages))
    rows = []
    for page in pages:
        before = before_pages.get(page, [])
        after = after_pages.get(page, [])
        before_types = Counter(f.get("type", "") for f in before)
        after_types = Counter(f.get("type", "") for f in after)
        removed = sorted((before_types - after_types).elements())
        added = sorted((after_types - before_types).elements())
        before_severity = Counter(f.get("severity", "low") for f in before)
        after_severity = Counter(f.get("severity", "low") for f in after)
        rows.append(
            {
                "page": page,
                "before_count": len(before),
                "after_count": len(after),
                "before_by_severity": dict(before_severity),
                "after_by_severity": dict(after_severity),
                "removed_types": removed,
                "added_types": added,
                "machine_delta": len(after) - len(before),
                "human_outcome": "pending_review",
            }
        )
    return rows


def _findings_by_page(audit: Dict[str, Any]) -> Dict[int, List[Dict[str, Any]]]:
    pages: Dict[int, List[Dict[str, Any]]] = {}
    for finding in audit.get("findings", []) or []:
        page = int(finding.get("slide_page") or 0)
        if page <= 0:
            continue
        pages.setdefault(page, []).append(finding)
    return pages


def _repair_risk_findings_from_delta(page_deltas: List[Dict[str, Any]], style_contract: Dict[str, Any]) -> List[Dict[str, Any]]:
    risks = []
    style_id = style_contract.get("style_id", "")
    for row in page_deltas:
        removed = set(row.get("removed_types", []) or [])
        added = set(row.get("added_types", []) or [])
        machine_improved = int(row.get("after_count", 0)) < int(row.get("before_count", 0))
        if machine_improved and (
            "image_underutilized_in_wide_panel" in added
            or "figure_panel_aspect_mismatch" in removed
            or added.intersection(STYLE_AWARE_REVIEW_RULES)
        ):
            risks.append(
                _repair_risk_entry(
                    "metric_improved_visual_regressed",
                    row,
                    style_id,
                    "Machine finding count improved, but the changed slide touches a style-aware figure/layout rule and needs human visual review.",
                    "tradeoff_review",
                )
            )
        if int(row.get("after_count", 0)) > int(row.get("before_count", 0)):
            risks.append(
                _repair_risk_entry(
                    "repair_introduced_new_findings",
                    row,
                    style_id,
                    "Repair increased page-level metadata findings.",
                    "pending_review",
                )
            )
        if style_id == BLIND_RECTANGULAR_STYLE_ID and "figure_panel_aspect_mismatch" in removed:
            risks.append(
                _repair_risk_entry(
                    "style_scope_mismatch",
                    row,
                    style_id,
                    "A figure-panel aspect repair may have applied a bottom-band preference too aggressively for the straight-rectangle style.",
                    "pending_review",
                )
            )
    return risks


def _repair_risk_entry(kind: str, row: Dict[str, Any], style_id: str, message: str, human_outcome: str) -> Dict[str, Any]:
    return {
        "type": kind,
        "dimension": "repair_risk",
        "scope": "human_feedback" if kind != "style_scope_mismatch" else "style_aware",
        "style_scope": [style_id] if style_id else [],
        "repair_mode": "detect_only",
        "confidence": 0.64,
        "human_outcome": human_outcome,
        "slide_page": row.get("page"),
        "severity": "medium",
        "message": message,
        "evidence": {
            "before_count": row.get("before_count", 0),
            "after_count": row.get("after_count", 0),
            "removed_types": row.get("removed_types", []),
            "added_types": row.get("added_types", []),
        },
        "repair_strategy": "Do not count this repair as accepted until the human review packet is approved.",
    }


def _machine_judgement(before_audit: Dict[str, Any], after_audit: Dict[str, Any], repair_risks: List[Dict[str, Any]]) -> Dict[str, Any]:
    before_count = int(before_audit.get("summary", {}).get("finding_count", 0) or 0)
    after_count = int(after_audit.get("summary", {}).get("finding_count", 0) or 0)
    return {
        "before_finding_count": before_count,
        "after_finding_count": after_count,
        "finding_delta": after_count - before_count,
        "metric_improved": after_count < before_count,
        "requires_human_review": bool(repair_risks),
        "reason": "metric improved but style-aware repair risk exists" if repair_risks else "metadata-only comparison has no repair-risk trigger",
    }


def _write_delta_csv(path: Path, rows: List[Dict[str, Any]]) -> None:
    fieldnames = [
        "page",
        "before_count",
        "after_count",
        "machine_delta",
        "removed_types",
        "added_types",
        "human_outcome",
    ]
    with path.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.DictWriter(handle, fieldnames=fieldnames)
        writer.writeheader()
        for row in rows:
            writer.writerow(
                {
                    "page": row.get("page", ""),
                    "before_count": row.get("before_count", 0),
                    "after_count": row.get("after_count", 0),
                    "machine_delta": row.get("machine_delta", 0),
                    "removed_types": ";".join(row.get("removed_types", []) or []),
                    "added_types": ";".join(row.get("added_types", []) or []),
                    "human_outcome": row.get("human_outcome", "pending_review"),
                }
            )


def _render_visual_packet_md(payload: Dict[str, Any], zh: bool) -> str:
    if zh:
        title = "Visual Human Review Packet"
        intro = "该 packet 是 04 blind experimental route 的标准输出。当前 harness 不渲染截图，因此这里记录 metadata before/after、规则 delta 和需要人工复核的页面。"
        columns = "| 页 | Before | After | Removed | Added | 人工裁决 |\n| ---: | ---: | ---: | --- | --- | --- |"
        risk_title = "## 误修风险"
        empty_risk = "暂无自动触发的误修风险，但仍建议人工抽查高风险页面。"
    else:
        title = "Visual Human Review Packet"
        intro = "This packet is the standard output for route 04. The current harness is metadata-only, so it records before/after audit deltas and pages that need human visual review without pretending screenshots exist."
        columns = "| Page | Before | After | Removed | Added | Human outcome |\n| ---: | ---: | ---: | --- | --- | --- |"
        risk_title = "## Repair Risks"
        empty_risk = "No automatic repair-risk trigger fired; still review representative high-risk pages."
    lines = [
        f"# {title}",
        "",
        intro,
        "",
        f"- Style: `{payload.get('style_id', '')}`",
        f"- Before iteration: `{payload.get('before_iteration', 0)}`",
        f"- After iteration: `{payload.get('after_iteration', 0)}`",
        f"- Screenshot status: `{payload.get('screenshot_status', '')}`",
        f"- Machine judgement: `{payload.get('machine_judgement', {}).get('reason', '')}`",
        "",
        "## Page Delta",
        "",
        columns,
    ]
    for row in payload.get("page_deltas", []) or []:
        if not row.get("removed_types") and not row.get("added_types") and row.get("before_count") == row.get("after_count"):
            continue
        lines.append(
            f"| {row.get('page')} | {row.get('before_count')} | {row.get('after_count')} | "
            f"{', '.join(row.get('removed_types', []) or []) or '-'} | "
            f"{', '.join(row.get('added_types', []) or []) or '-'} | "
            f"{row.get('human_outcome', 'pending_review')} |"
        )
    lines.extend(["", risk_title, ""])
    risks = payload.get("repair_risk_findings", []) or []
    if not risks:
        lines.append(empty_risk)
    for risk in risks:
        lines.append(f"- Slide {risk.get('slide_page')}: `{risk.get('type')}` - {risk.get('message')}")
    lines.extend(
        [
            "",
            "## Review Fields",
            "",
            "- human_outcome: pending_review | accepted | rejected | tradeoff_review | likely_overcorrection",
            "- visual note: fill this after opening the saved PPTX or generated screenshots.",
            "",
        ]
    )
    return "\n".join(lines)


def _compact_academic_spec_for_global_repair(spec: PresentationSpec) -> Tuple[PresentationSpec, List[str]]:
    """Apply conservative global correctness repairs without changing style grammar."""
    repaired = PresentationSpec.from_dict(spec.to_dict())
    changed = 0
    for slide in repaired.slides:
        original = slide.to_dict()
        if slide.section_type == "opening" and repaired.title:
            slide.title = repaired.title
        else:
            slide.title = _limit_words(slide.title, 10)
        slide.takeaway = _limit_words(slide.takeaway, 16 if slide.section_type == "opening" else 13)
        for block in slide.text_blocks:
            block.text = _limit_words(block.text, 24)
            block.claim = _limit_words(block.claim, 12)
            block.detail = _limit_words(block.detail, 18)
            block.evidence = _limit_words(block.evidence, 10)
        for metric in slide.metric_blocks:
            metric.label = _limit_words(metric.label, 7)
            metric.value = _limit_words(metric.value, 4)
            metric.note = _limit_words(metric.note, 10)
        if slide.to_dict() != original:
            changed += 1
    return repaired, [
        f"global correctness copy fitting: compacted title/takeaway/claim/detail/evidence on {changed} slide(s)",
        "preserved academic renderer and did not apply rounded proof-panel polish",
    ]


def render_blind_blueprint_pptx(
    inventory: Dict[str, Any],
    rough: Dict[str, Any],
    output_path: Path,
    repair_mode: bool = False,
    style_contract: Optional[Dict[str, Any]] = None,
) -> Path:
    """Render a deliberately non-baseline experimental deck."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    style_contract = style_contract or _build_blind_style_contract(inventory, output_path.parent)
    theme = _blind_theme(style_contract)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slides = rough.get("slides", []) or []
    sections = _section_buckets(slides)
    total = 2 + len(sections) + len(slides) + 1
    ctx = {
        "RGBColor": RGBColor,
        "MSO_SHAPE": MSO_SHAPE,
        "MSO_AUTO_SIZE": MSO_AUTO_SIZE,
        "MSO_VERTICAL_ANCHOR": MSO_VERTICAL_ANCHOR,
        "PP_ALIGN": PP_ALIGN,
        "Inches": Inches,
        "Pt": Pt,
        "theme": theme,
        "repair_mode": repair_mode,
        "style_contract": style_contract,
    }

    page = 1
    _blind_cover(prs.slides.add_slide(blank), inventory, page, total, ctx)
    page += 1
    _blind_agenda(prs.slides.add_slide(blank), sections, page, total, ctx)
    page += 1
    for section_index, section in enumerate(sections, start=1):
        _blind_section_slide(prs.slides.add_slide(blank), section, section_index, page, total, ctx)
        page += 1
        for slide_data in section["slides"]:
            _blind_content_slide(prs.slides.add_slide(blank), inventory, slide_data, page, total, ctx)
            page += 1
    _blind_closing(prs.slides.add_slide(blank), inventory, page, total, ctx)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def build_style_drift_report_payload(
    pptx_path: Path,
    audit: Dict[str, Any],
    style_id: str,
    repair_profile: str,
    style_scope: str,
) -> Dict[str, Any]:
    """Build a style-scoped drift report from audit and deck text signals."""
    counts = Counter(f.get("type", "") for f in audit.get("findings", []) or [])
    scoped_hits = {key: counts[key] for key in sorted(GOLDEN_BASELINE1_SCOPED_RULES) if counts.get(key)}
    global_hits = {key: counts[key] for key in sorted(GLOBAL_CORRECTNESS_RULES) if counts.get(key)}
    style_aware_hits = {key: counts[key] for key in sorted(STYLE_AWARE_REVIEW_RULES) if counts.get(key)}
    text_signals = _pptx_text_signals(pptx_path)
    golden_tokens = {
        key: text_signals.get(key, 0)
        for key in ["ACADEMIC PAPER READING", "DECK MAP", "PAPER HIGHLIGHTS", "PROOF OBJECT"]
        if text_signals.get(key, 0)
    }
    forbidden_auto_repairs: List[str] = []
    if style_scope == "academic":
        forbidden_auto_repairs = sorted(scoped_hits)
    if style_scope == "experimental":
        forbidden_auto_repairs = sorted(golden_tokens)
    drift_risk = "low"
    if style_scope == "academic" and forbidden_auto_repairs:
        drift_risk = "medium"
    if style_scope == "experimental" and golden_tokens:
        drift_risk = "medium"
    return {
        "schema_version": "style_drift_report.v1",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "pptx_path": str(pptx_path),
        "style_id": style_id,
        "repair_profile": repair_profile,
        "style_scope": style_scope,
        "style_contract": STYLE_CONTRACTS.get(style_id, STYLE_CONTRACTS.get(style_scope, {})),
        "drift_risk": drift_risk,
        "global_correctness_findings": global_hits,
        "golden_baseline1_scoped_findings": scoped_hits,
        "style_aware_review_findings": style_aware_hits,
        "baseline_similarity_signals": golden_tokens,
        "policy": {
            "global_correctness": "auto_repair_allowed",
            "golden_baseline1_polish": "auto_repair_only_when_style_scope_is_golden_baseline1",
            "academic_handling": "detect_only_for_rounded_proof_panel_polish",
            "experimental_handling": "avoid academic header/key-message and golden_baseline1 rounded proof-panel grammar",
        },
        "forbidden_auto_repairs_detected": forbidden_auto_repairs,
        "forbidden_auto_repairs_applied": [],
    }


def _build_blind_style_contract(inventory: Dict[str, Any], route_dir: Path) -> Dict[str, Any]:
    paper_title = str((inventory.get("paper", {}) or {}).get("title", "paper")).strip() or "paper"
    paper_slug = _slugify(paper_title)
    seed = hashlib.sha1(f"{paper_title}|{route_dir}".encode("utf-8")).hexdigest()[:8]
    palette_name, deck_label, board_title, agenda_title = [
        ("ledger_blue", "PAPER TRACE", "Evidence Control Board", "Route Map"),
        ("lab_green", "LAB NOTE", "Claim Verification Board", "Experiment Map"),
        ("archive_red", "CLAIM LEDGER", "Evidence Review Board", "Audit Route"),
    ][int(seed, 16) % 3]
    contract = dict(STYLE_CONTRACTS[BLIND_RECTANGULAR_STYLE_ID])
    style_instance_id = f"{EXPERIMENTAL_STYLE_PREFIX}_{paper_slug}_{seed}"
    contract.update({
        "schema_version": "blind_style_contract.v2",
        "style_id": style_instance_id,
        "style_family": BLIND_RECTANGULAR_STYLE_ID,
        "style_instance_id": style_instance_id,
        "style_scope": "experimental",
        "paper_title": paper_title,
        "paper_slug": paper_slug,
        "seed": seed,
        "palette_name": palette_name,
        "deck_label": deck_label,
        "board_title": board_title,
        "agenda_title": agenda_title,
        "forgets_styles": [
            "academic",
            "golden_baseline1_from_scratch_warm_academic",
            GOLDEN_BASELINE2_STYLE_ID,
            "previous_blind_experimental_candidates",
        ],
        "policy": {
            "reuse_visual_grammar": False,
            "allowed_reuse": ["fresh parse checkpoints", "benchmark badcases", "bounded repair loop"],
            "forbidden_text_signals": ["ACADEMIC PAPER READING", "DECK MAP", "PAPER HIGHLIGHTS", "PROOF OBJECT"],
        },
    })
    return contract


def _blind_cover(slide: Any, inventory: Dict[str, Any], page: int, total: int, ctx: Dict[str, Any]) -> None:
    _blind_bg(slide, ctx)
    t = ctx["theme"]
    contract = ctx["style_contract"]
    paper = inventory.get("paper", {})
    title = _limit_words(paper.get("title", "Untitled Paper"), 18)
    highlights = inventory.get("paper_highlights", []) or []
    _rect(slide, 0.0, 0.0, 3.0, 7.5, t["ink"], t["ink"], ctx)
    _text(slide, contract["deck_label"], 0.42, 0.55, 1.8, 0.25, 9, t["accent"], ctx, bold=True, spaced=True)
    _text(slide, title, 0.42, 1.2, 2.25, 4.7, 28 if ctx["repair_mode"] else 30, t["white"], ctx, bold=True)
    _text(slide, "Fresh parse -> four route validation", 0.45, 6.45, 2.1, 0.5, 10, t["pale"], ctx)
    _text(slide, contract["board_title"], 4.0, 0.9, 7.8, 0.55, 24, t["ink"], ctx, bold=True)
    _text(slide, "This route uses a blueprint grammar: rectangular ledgers, axis lines, and evidence boards.", 4.02, 1.52, 7.4, 0.5, 14 if ctx["repair_mode"] else 13, t["muted"], ctx)
    for idx, item in enumerate(highlights[:3] or [{"label": "Core claim", "text": title}]):
        y = 2.35 + idx * 1.25
        _rect(slide, 4.05, y, 7.5, 0.9, t["white"], t["grid"], ctx)
        _rect(slide, 4.05, y, 0.1, 0.9, [t["blue"], t["green"], t["coral"]][idx % 3], [t["blue"], t["green"], t["coral"]][idx % 3], ctx)
        _text(slide, f"SIGNAL {idx + 1}", 4.32, y + 0.12, 1.7, 0.2, 8.5, t["muted"], ctx, bold=True, spaced=True)
        _text(slide, _limit_words(str(item.get("text", "")), 16 if ctx["repair_mode"] else 20), 4.32, y + 0.34, 6.65, 0.45, 12.5 if ctx["repair_mode"] else 13, t["ink"], ctx, bold=True)
    _footer(slide, page, total, ctx)


def _blind_agenda(slide: Any, sections: List[Dict[str, Any]], page: int, total: int, ctx: Dict[str, Any]) -> None:
    _blind_bg(slide, ctx)
    t = ctx["theme"]
    _text(slide, ctx["style_contract"]["agenda_title"], 0.72, 0.62, 3.2, 0.45, 30, t["ink"], ctx, bold=True)
    _text(slide, "The deck is organized as a sequence of claim/evidence ledgers.", 0.75, 1.16, 5.8, 0.35, 13, t["muted"], ctx)
    for idx, section in enumerate(sections, start=1):
        y = 1.75 + (idx - 1) * 1.12
        color = [t["blue"], t["green"], t["coral"], t["yellow"]][(idx - 1) % 4]
        _rect(slide, 0.82, y, 10.8, 0.78, t["white"], t["grid"], ctx)
        _rect(slide, 0.82, y, 0.2, 0.78, color, color, ctx)
        _text(slide, f"{idx:02d}", 1.18, y + 0.17, 0.55, 0.28, 15, color, ctx, bold=True)
        _text(slide, section["title"], 1.95, y + 0.14, 3.2, 0.3, 16, t["ink"], ctx, bold=True)
        sample = "; ".join(_limit_words(item.get("title", ""), 5).rstrip(".") for item in section["slides"][:2])
        _text(slide, sample, 5.45, y + 0.17, 4.75, 0.3, 11, t["muted"], ctx)
        _text(slide, f"{len(section['slides'])} slides", 11.0, y + 0.18, 0.75, 0.25, 9, color, ctx, bold=True)
    _footer(slide, page, total, ctx)


def _blind_section_slide(slide: Any, section: Dict[str, Any], section_index: int, page: int, total: int, ctx: Dict[str, Any]) -> None:
    _blind_bg(slide, ctx)
    t = ctx["theme"]
    color = [t["blue"], t["green"], t["coral"], t["yellow"]][(section_index - 1) % 4]
    _rect(slide, 0.8, 0.8, 11.6, 5.8, t["white"], t["grid"], ctx)
    _rect(slide, 0.8, 0.8, 11.6, 0.18, color, color, ctx)
    _text(slide, f"MODULE {section_index:02d}", 1.15, 1.45, 2.2, 0.3, 10, color, ctx, bold=True, spaced=True)
    _text(slide, section["title"], 1.12, 2.05, 8.7, 0.65, 34, t["ink"], ctx, bold=True)
    sample = " / ".join(_limit_words(item.get("title", ""), 6).rstrip(".") for item in section["slides"][:3])
    _text(slide, sample or "Claim and evidence ledger", 1.18, 4.0, 9.6, 0.45, 15, t["muted"], ctx)
    _footer(slide, page, total, ctx)


def _blind_content_slide(slide: Any, inventory: Dict[str, Any], slide_data: Dict[str, Any], page: int, total: int, ctx: Dict[str, Any]) -> None:
    _blind_bg(slide, ctx)
    t = ctx["theme"]
    repair = bool(ctx["repair_mode"])
    section = _section_title_for_slide(slide_data)
    proof = slide_data.get("proof_object", {}) or {}
    proof_type = str(proof.get("type", "text_evidence"))
    figure_aspect = _figure_aspect(inventory, str(proof.get("id", ""))) if proof_type == "figure" else None
    table_wide_layout = repair and proof_type == "table" and _proof_table_needs_wide_layout(inventory, proof)
    _text(slide, section.upper(), 0.72, 0.38, 2.6, 0.24, 8.5, t["blue"], ctx, bold=True, spaced=True)
    _text(slide, _limit_words(slide_data.get("title", ""), 8 if repair else 13), 0.72, 0.62, 8.4, 0.55, 20 if repair else 22, t["ink"], ctx, bold=True)
    _text(slide, f"{page:02d}/{total:02d}", 11.35, 0.47, 1.0, 0.22, 8.5, t["muted"], ctx, align="right")

    if repair and ((figure_aspect and figure_aspect >= 1.9) or table_wide_layout):
        _rect(slide, 0.72, 1.24, 11.9, 1.55, t["white"], t["grid"], ctx)
        _rect(slide, 0.72, 1.24, 0.12, 1.55, t["blue"], t["blue"], ctx)
        _text(slide, "CLAIM", 1.02, 1.43, 0.9, 0.2, 8, t["blue"], ctx, bold=True, spaced=True)
        _text(slide, _limit_words(slide_data.get("claim", ""), 9), 1.02, 1.68, 4.25, 0.42, 18, t["ink"], ctx, bold=True)
        _text(slide, _limit_words(slide_data.get("support", ""), 16), 5.6, 1.66, 6.2, 0.42, 11.5, t["muted"], ctx)
        _rect(slide, 0.72, 3.05, 11.9, 3.58, t["white"], t["grid"], ctx)
        _text(slide, f"EVIDENCE / {proof_type.upper()}", 1.02, 3.27, 2.6, 0.2, 8, t["coral"], ctx, bold=True, spaced=True)
        _render_blind_proof(slide, inventory, proof, 1.0, 3.42, 11.25, 3.0, ctx)
        _text(slide, _evidence_footer(slide_data), 0.75, 6.9, 8.6, 0.25, 7.5, t["muted"], ctx)
        return

    if repair and figure_aspect and figure_aspect <= 0.9:
        _rect(slide, 0.72, 1.28, 7.0, 5.35, t["white"], t["grid"], ctx)
        _rect(slide, 0.72, 1.28, 0.12, 5.35, t["blue"], t["blue"], ctx)
        _text(slide, "CLAIM", 1.05, 1.55, 1.2, 0.22, 8.5, t["blue"], ctx, bold=True, spaced=True)
        _text(slide, _limit_words(slide_data.get("claim", ""), 10), 1.05, 1.88, 5.8, 0.82, 19, t["ink"], ctx, bold=True)
        _text(slide, "READING NOTE", 1.05, 3.0, 1.8, 0.22, 8.5, t["green"], ctx, bold=True, spaced=True)
        _text(slide, _limit_words(slide_data.get("support", ""), 18), 1.05, 3.34, 5.9, 1.25, 12.5, t["muted"], ctx)
        _rect(slide, 8.2, 1.28, 4.25, 5.35, t["white"], t["grid"], ctx)
        _text(slide, f"EVIDENCE / {proof_type.upper()}", 8.48, 1.55, 2.6, 0.22, 8.5, t["coral"], ctx, bold=True, spaced=True)
        _render_blind_proof(slide, inventory, proof, 8.48, 1.95, 3.6, 4.05, ctx)
        _text(slide, _evidence_footer(slide_data), 0.75, 6.9, 8.6, 0.25, 7.5, t["muted"], ctx)
        return

    _rect(slide, 0.72, 1.28, 4.2, 5.35, t["white"], t["grid"], ctx)
    _rect(slide, 0.72, 1.28, 0.12, 5.35, t["blue"], t["blue"], ctx)
    _text(slide, "CLAIM", 1.05, 1.55, 1.2, 0.22, 8.5, t["blue"], ctx, bold=True, spaced=True)
    _text(slide, _limit_words(slide_data.get("claim", ""), 10 if repair else 15), 1.05, 1.88, 3.25, 1.28, 19 if repair else 21, t["ink"], ctx, bold=True)
    _text(slide, "READING NOTE", 1.05, 3.42, 1.8, 0.22, 8.5, t["green"], ctx, bold=True, spaced=True)
    _text(slide, _limit_words(slide_data.get("support", ""), 12 if repair else 38), 1.05, 3.76, 3.35, 1.95, 12.5 if repair else 13, t["muted"], ctx)

    _rect(slide, 5.25, 1.28, 7.38, 5.35, t["white"], t["grid"], ctx)
    _text(slide, f"EVIDENCE / {str(proof.get('type', 'text')).upper()}", 5.55, 1.55, 3.3, 0.22, 8.5, t["coral"], ctx, bold=True, spaced=True)
    _render_blind_proof(slide, inventory, proof, 5.55, 1.95, 6.65, 4.05, ctx)
    _text(slide, _evidence_footer(slide_data), 0.75, 6.9, 8.6, 0.25, 7.5, t["muted"], ctx)


def _blind_closing(slide: Any, inventory: Dict[str, Any], page: int, total: int, ctx: Dict[str, Any]) -> None:
    _blind_bg(slide, ctx)
    t = ctx["theme"]
    highlights = inventory.get("paper_highlights", []) or []
    _text(slide, "Inspection Ready", 0.72, 0.8, 6.5, 0.6, 34, t["ink"], ctx, bold=True)
    _text(slide, "The final branch keeps content traceability while avoiding both protected visual grammars.", 0.75, 1.5, 7.5, 0.4, 14, t["muted"], ctx)
    for idx, label in enumerate(["Parse once", "Benchmark", "Bounded repair"]):
        x = 0.85 + idx * 3.95
        _rect(slide, x, 3.0, 3.15, 1.2, t["white"], t["grid"], ctx)
        _text(slide, label, x + 0.24, 3.22, 1.6, 0.25, 13, [t["blue"], t["green"], t["coral"]][idx], ctx, bold=True)
        text = _limit_words(str((highlights[idx] if idx < len(highlights) else {}).get("text", "Evidence remains linked to parsed checkpoints.")), 10)
        _text(slide, text, x + 0.24, 3.58, 2.5, 0.45, 10.5, t["muted"], ctx)
    _footer(slide, page, total, ctx)


def _render_blind_proof(slide: Any, inventory: Dict[str, Any], proof: Dict[str, Any], x: float, y: float, w: float, h: float, ctx: Dict[str, Any]) -> None:
    proof_type = str(proof.get("type", "text_evidence"))
    t = ctx["theme"]
    if proof_type == "figure":
        figure_path = _figure_path(inventory, str(proof.get("id", "")))
        if figure_path and Path(figure_path).exists():
            top_pad = 0.12 if ctx["repair_mode"] else 0.55
            side_pad = 0.24 if ctx["repair_mode"] else 0.25
            caption_h = 0.34
            fit = _fit_image_box(figure_path, x + side_pad, y + top_pad, w - side_pad * 2, h - top_pad - caption_h - 0.18)
            if fit:
                slide.shapes.add_picture(figure_path, ctx["Inches"](fit[0]), ctx["Inches"](fit[1]), width=ctx["Inches"](fit[2]), height=ctx["Inches"](fit[3]))
                caption_y = min(y + h - caption_h, fit[1] + fit[3] + 0.12)
            else:
                caption_y = y + h - caption_h
            _text(slide, _limit_words(proof.get("focus", ""), 18), x + 0.35, caption_y, w - 0.7, 0.28, 8.5, t["muted"], ctx, align="center")
            return
    if proof_type == "table":
        rows = _table_rows(inventory, str(proof.get("id", "")))[:6]
        if rows:
            caption = _table_caption(inventory, proof)
            if _table_needs_summary_fallback(rows, w - 0.4):
                _table_summary_cards(slide, rows, x + 0.18, y + 0.12, w - 0.36, h - 0.28, ctx, caption)
            else:
                _table_with_context(slide, rows, x + 0.18, y + 0.12, w - 0.36, h - 0.28, ctx, caption)
            return
    if proof_type == "metric":
        metrics = inventory.get("metrics", [])[:3] or [{"label": proof.get("id", "Metric"), "value": proof.get("focus", "")}]
        for idx, metric in enumerate(metrics[:3]):
            yy = y + 0.55 + idx * 1.05
            _rect(slide, x + 0.22, yy, w - 0.44, 0.72, t["pale"], t["grid"], ctx)
            _text(slide, _limit_words(str(metric.get("value", "") or proof.get("focus", "")), 4), x + 0.42, yy + 0.14, 1.7, 0.25, 18, t["ink"], ctx, bold=True)
            _text(slide, _limit_words(str(metric.get("label", "") or proof.get("id", "")), 10), x + 2.35, yy + 0.18, w - 3.0, 0.25, 11, t["muted"], ctx)
        return
    notes = [
        _limit_words(str(proof.get("id", "source evidence")), 7),
        _limit_words(str(proof.get("focus", "")), 22),
        "Evidence is preserved from this run's fresh parse checkpoints.",
    ]
    card_specs = _text_evidence_card_specs(notes, w - 0.44, max(1.0, h - 0.72))
    yy = y + 0.58
    for idx, (text, card_h, font_size) in enumerate(card_specs):
        color = [t["blue"], t["green"], t["coral"]][idx % 3]
        _rect(slide, x + 0.22, yy, w - 0.44, card_h, t["pale"], t["grid"], ctx)
        _rect(slide, x + 0.22, yy, 0.08, card_h, color, color, ctx)
        _text(slide, text, x + 0.48, yy + 0.10, w - 0.85, max(0.32, card_h - 0.20), font_size, t["ink"], ctx, bold=idx == 0, valign="middle")
        yy += card_h + 0.25


def _text_evidence_card_specs(notes: List[str], card_w: float, available_h: float) -> List[Tuple[str, float, float]]:
    specs: List[Tuple[str, float, float]] = []
    for idx, text in enumerate(notes):
        font_size = 10.8 if idx else 12.2
        estimated_h = _estimated_text_box_height(text, max(1.0, card_w - 0.45), font_size) + 0.28
        min_card_h = 0.72
        if idx == 1:
            min_card_h = 0.96 if _effective_text_units(text) >= 72.0 else 0.84
        elif idx >= 2:
            min_card_h = 0.82
        card_h = min(1.36, max(min_card_h, estimated_h))
        specs.append((text, card_h, font_size))
    total_h = sum(item[1] for item in specs) + max(0, len(specs) - 1) * 0.25
    if total_h <= available_h:
        return specs
    scale = max(0.72, (available_h - max(0, len(specs) - 1) * 0.22) / max(1, len(specs)))
    compact: List[Tuple[str, float, float]] = []
    for idx, (text, _card_h, font_size) in enumerate(specs):
        compact_text = _limit_words(text, 14 if idx else 7)
        compact.append((compact_text, min(_card_h, scale), max(9.2, font_size - 1.4)))
    return compact


def _estimated_text_box_height(text: Any, width_in: float, font_pt: float) -> float:
    units = _effective_text_units(str(text or ""))
    chars_per_line = max(1.0, width_in * 72.0 / max(1.0, font_pt * 0.50))
    lines = max(1.0, math.ceil(units / chars_per_line))
    return lines * font_pt * 1.18 / 72.0


def _effective_text_units(text: str) -> float:
    total = 0.0
    for char in str(text or ""):
        if char.isspace():
            total += 0.35
        elif ord(char) > 127:
            total += 1.75
        elif char in ".,;:!|":
            total += 0.35
        else:
            total += 1.0
    return total


def _table_needs_summary_fallback(rows: List[List[str]], width_in: float) -> bool:
    if not rows:
        return False
    col_count = max(len(row) for row in rows)
    if col_count <= 6:
        return False
    total_cells = max(1, len(rows) * col_count)
    empty_cells = 0
    long_cells = 0
    for row in rows:
        padded = row + [""] * (col_count - len(row))
        for cell in padded:
            text = str(cell or "").strip()
            if not text:
                empty_cells += 1
            elif _effective_text_units(text) >= 12.0:
                long_cells += 1
    col_width = width_in / max(1, col_count)
    return col_width < 0.62 or empty_cells / total_cells >= 0.32 or long_cells >= 3


def _proof_table_needs_wide_layout(inventory: Dict[str, Any], proof: Dict[str, Any]) -> bool:
    rows = _table_rows(inventory, str(proof.get("id", "")))[:6]
    if not rows:
        return False
    col_count = max(len(row) for row in rows)
    return col_count >= 7 or _table_needs_summary_fallback(rows, 6.25)


def _table_summary_cards(
    slide: Any,
    rows: List[List[str]],
    x: float,
    y: float,
    w: float,
    h: float,
    ctx: Dict[str, Any],
    caption: str = "",
) -> None:
    t = ctx["theme"]
    summary_rows = _summarize_table_rows(rows)[:5]
    if not summary_rows:
        _text(slide, "Table evidence is available in the parsed source, but the native grid is too dense for this panel.", x + 0.2, y + 0.2, w - 0.4, 0.6, 10.5, t["muted"], ctx)
        return
    _text(slide, "Focused table view", x + 0.12, y + 0.03, w - 0.24, 0.22, 9, t["blue"], ctx, bold=True)
    table_h = min(max(2.25, len(summary_rows) * 0.54), max(1.2, h - 0.78))
    _native_table(slide, [["Stage", "Readable evidence"]] + summary_rows, x + 0.12, y + 0.34, w - 0.24, table_h, ctx, font_size=8.7)
    note = caption or "Dense source table was collapsed to avoid unreadable 13-column wrapping."
    _text(slide, _limit_words(note, 22), x + 0.12, y + 0.42 + table_h, w - 0.24, 0.32, 8.2, t["muted"], ctx, align="center")


def _table_with_context(
    slide: Any,
    rows: List[List[str]],
    x: float,
    y: float,
    w: float,
    h: float,
    ctx: Dict[str, Any],
    caption: str = "",
) -> None:
    t = ctx["theme"]
    _text(slide, "Focused table view", x + 0.12, y + 0.03, w - 0.24, 0.22, 9, t["blue"], ctx, bold=True)
    table_y = y + 0.36
    table_h = max(0.85, h - 0.86)
    _native_table(slide, rows, x + 0.12, table_y, w - 0.24, table_h, ctx, font_size=8.5)
    note = caption or "Table evidence summarized from parsed source."
    _text(slide, _limit_words(note, 22), x + 0.12, table_y + table_h + 0.09, w - 0.24, 0.32, 8.2, t["muted"], ctx, align="center")


def _summarize_table_rows(rows: List[List[str]]) -> List[List[str]]:
    summary: List[List[str]] = []
    for row in rows:
        cells = [str(cell or "").strip() for cell in row if str(cell or "").strip()]
        if not cells:
            continue
        if len(cells) == 1:
            label = cells[0]
            detail = "Source table row retained."
        else:
            label = " / ".join(cells[:2])
            detail = "; ".join(cells[2:]) if len(cells) > 2 else cells[-1]
        summary.append([_limit_words(label, 5), _limit_words(detail, 14)])
    return summary


def _table_caption(inventory: Dict[str, Any], proof: Dict[str, Any]) -> str:
    proof_id = str(proof.get("id", "") or "").strip()
    focus = str(proof.get("focus", "") or "").strip()
    if focus:
        return _limit_words(_clean_inline_text(focus), 22)
    table = _matching_table(inventory, proof_id)
    for key in ("caption", "title"):
        value = str((table or {}).get(key, "") or "").strip()
        if value and value != proof_id:
            return _limit_words(_clean_inline_text(value), 22)
    if proof_id:
        return _limit_words(_clean_inline_text(proof_id), 12)
    return "Table evidence summarized from parsed source."


def _blind_theme(style_contract: Dict[str, Any]) -> Dict[str, Tuple[int, int, int]]:
    palette_name = style_contract.get("palette_name", "ledger_blue")
    palettes = {
        "ledger_blue": {
            "bg": (247, 249, 245),
            "white": (255, 255, 255),
            "pale": (239, 244, 238),
            "grid": (213, 222, 214),
            "ink": (23, 30, 33),
            "muted": (82, 96, 98),
            "blue": (42, 103, 190),
            "green": (65, 136, 105),
            "coral": (190, 74, 84),
            "yellow": (188, 159, 58),
            "accent": (117, 207, 174),
        },
        "lab_green": {
            "bg": (245, 248, 244),
            "white": (255, 255, 255),
            "pale": (235, 243, 236),
            "grid": (207, 220, 210),
            "ink": (24, 34, 29),
            "muted": (84, 99, 91),
            "blue": (48, 108, 141),
            "green": (52, 137, 88),
            "coral": (176, 98, 62),
            "yellow": (173, 153, 72),
            "accent": (118, 193, 156),
        },
        "archive_red": {
            "bg": (249, 246, 244),
            "white": (255, 255, 255),
            "pale": (244, 237, 235),
            "grid": (223, 211, 206),
            "ink": (35, 28, 29),
            "muted": (99, 85, 86),
            "blue": (66, 97, 147),
            "green": (87, 127, 93),
            "coral": (170, 68, 76),
            "yellow": (176, 143, 67),
            "accent": (205, 143, 121),
        },
    }
    return palettes.get(palette_name, palettes["ledger_blue"])


def _blind_bg(slide: Any, ctx: Dict[str, Any]) -> None:
    t = ctx["theme"]
    _rect(slide, 0.0, 0.0, 13.333, 7.5, t["bg"], t["bg"], ctx)
    for x in [1.5, 3.0, 4.5, 6.0, 7.5, 9.0, 10.5, 12.0]:
        _rect(slide, x, 0.0, 0.006, 7.5, t["grid"], t["grid"], ctx)
    _rect(slide, 0.0, 0.0, 13.333, 0.05, t["blue"], t["blue"], ctx)


def _rect(slide: Any, x: float, y: float, w: float, h: float, fill: Tuple[int, int, int], line: Tuple[int, int, int], ctx: Dict[str, Any]) -> Any:
    shape = slide.shapes.add_shape(ctx["MSO_SHAPE"].RECTANGLE, ctx["Inches"](x), ctx["Inches"](y), ctx["Inches"](w), ctx["Inches"](h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ctx["RGBColor"](*fill)
    shape.line.color.rgb = ctx["RGBColor"](*line)
    return shape


def _text(
    slide: Any,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float,
    color: Tuple[int, int, int],
    ctx: Dict[str, Any],
    bold: bool = False,
    align: str = "left",
    spaced: bool = False,
    valign: str = "top",
) -> Any:
    shape = slide.shapes.add_textbox(ctx["Inches"](x), ctx["Inches"](y), ctx["Inches"](w), ctx["Inches"](h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = ctx["MSO_AUTO_SIZE"].TEXT_TO_FIT_SHAPE
    frame.margin_left = ctx["Inches"](0.02)
    frame.margin_right = ctx["Inches"](0.02)
    frame.margin_top = ctx["Inches"](0.01)
    frame.margin_bottom = ctx["Inches"](0.01)
    if valign == "middle":
        frame.vertical_anchor = ctx["MSO_VERTICAL_ANCHOR"].MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = " ".join(str(text or "")) if spaced else str(text or "")
    if align == "center":
        paragraph.alignment = ctx["PP_ALIGN"].CENTER
    elif align == "right":
        paragraph.alignment = ctx["PP_ALIGN"].RIGHT
    paragraph.font.name = "Aptos"
    paragraph.font.size = ctx["Pt"](size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = ctx["RGBColor"](*color)
    return shape


def _native_table(
    slide: Any,
    rows: List[List[str]],
    x: float,
    y: float,
    w: float,
    h: float,
    ctx: Dict[str, Any],
    font_size: float = 8.5,
) -> None:
    if not rows:
        return
    col_count = max(len(row) for row in rows)
    rows = [row + [""] * (col_count - len(row)) for row in rows]
    table_shape = slide.shapes.add_table(len(rows), col_count, ctx["Inches"](x), ctx["Inches"](y), ctx["Inches"](w), ctx["Inches"](h))
    table = table_shape.table
    t = ctx["theme"]
    if col_count == 2:
        first_width = min(max(w * 0.32, 1.25), w * 0.42)
        table.columns[0].width = ctx["Inches"](first_width)
        table.columns[1].width = ctx["Inches"](max(0.5, w - first_width))
    else:
        for column in table.columns:
            column.width = ctx["Inches"](w / max(1, col_count))
    for row in table.rows:
        row.height = ctx["Inches"](h / max(1, len(rows)))
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = _limit_words(value, 8)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = "Aptos"
            paragraph.font.size = ctx["Pt"](font_size)
            paragraph.font.bold = r_idx == 0
            paragraph.font.color.rgb = ctx["RGBColor"](*(t["ink"] if r_idx == 0 else t["muted"]))


def _footer(slide: Any, page: int, total: int, ctx: Dict[str, Any]) -> None:
    t = ctx["theme"]
    label = str(ctx["style_contract"].get("style_id", EXPERIMENTAL_STYLE_PREFIX)).replace("_", " ")
    _text(slide, f"{label} / no protected baseline grammar", 0.74, 7.05, 6.3, 0.2, 7.5, t["muted"], ctx)
    _text(slide, f"{page:02d}/{total:02d}", 11.62, 7.05, 0.7, 0.2, 7.5, t["muted"], ctx, align="right")


def _config_for_style(style: str, length: str, slides: int, fast: bool) -> Dict[str, Any]:
    style_type, custom_style = parse_style(style)
    return {
        "output_type": "slides",
        "style": style_type,
        "custom_style": custom_style,
        "slides_length": length,
        "target_slides": slides,
        "fast_mode": fast,
    }


def _section_buckets(slides: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    buckets = [
        {"id": "premise", "title": "Premise & Stakes", "slides": []},
        {"id": "mechanism", "title": "Mechanism & Method", "slides": []},
        {"id": "evidence", "title": "Evidence & Results", "slides": []},
        {"id": "implications", "title": "Implications", "slides": []},
    ]
    by_id = {bucket["id"]: bucket for bucket in buckets}
    for slide in slides:
        by_id[_section_id(slide)]["slides"].append(slide)
    return [bucket for bucket in buckets if bucket["slides"]]


def _section_id(slide: Dict[str, Any]) -> str:
    text = f"{slide.get('slide_role', '')} {slide.get('title', '')}".lower()
    if any(token in text for token in ["motivation", "problem", "thesis", "gap", "challenge"]):
        return "premise"
    if any(token in text for token in ["method", "mechanism", "architecture", "system", "training"]):
        return "mechanism"
    if any(token in text for token in ["result", "experiment", "metric", "table", "figure", "evaluation"]):
        return "evidence"
    return "implications"


def _section_title_for_slide(slide: Dict[str, Any]) -> str:
    return {
        "premise": "Premise",
        "mechanism": "Mechanism",
        "evidence": "Evidence",
        "implications": "Implications",
    }.get(_section_id(slide), "Evidence")


def _figure_path(inventory: Dict[str, Any], proof_id: str) -> str:
    for figure in inventory.get("assets", {}).get("figures", []) or []:
        if proof_id and proof_id in {str(figure.get("id", "")), str(figure.get("caption", ""))}:
            return str(figure.get("path", ""))
    figures = inventory.get("assets", {}).get("figures", []) or []
    return str((figures[0] if figures else {}).get("path", ""))


def _figure_aspect(inventory: Dict[str, Any], proof_id: str) -> Optional[float]:
    path = _figure_path(inventory, proof_id)
    if not path or not Path(path).exists():
        return None
    try:
        from PIL import Image

        with Image.open(path) as image:
            return image.width / image.height if image.height else None
    except Exception:
        return None


def _table_rows(inventory: Dict[str, Any], proof_id: str) -> List[List[str]]:
    table = _matching_table(inventory, proof_id)
    if table:
        rows = _normalize_rows(table.get("rows", []))
        if rows:
            return rows
    for table in _table_candidates(inventory):
        rows = _normalize_rows(table.get("rows", []))
        if rows:
            return rows
    return []


def _matching_table(inventory: Dict[str, Any], proof_id: str) -> Optional[Dict[str, Any]]:
    proof_id = str(proof_id or "").strip()
    if not proof_id:
        return None
    proof_norm = _clean_inline_text(proof_id).lower()
    for table in _table_candidates(inventory):
        keys = [
            str(table.get("id", "") or ""),
            str(table.get("caption", "") or ""),
            str(table.get("title", "") or ""),
        ]
        normalized = [_clean_inline_text(key).lower() for key in keys if key]
        if proof_norm in normalized:
            return table
        if any(proof_norm and proof_norm in key for key in normalized):
            return table
    return None


def _table_candidates(inventory: Dict[str, Any]) -> List[Dict[str, Any]]:
    candidates = []
    for table in inventory.get("assets", {}).get("tables", []) or []:
        candidates.append(table)
    for slide in inventory.get("curated_slides", []) or []:
        for table in slide.get("tables", []) or []:
            candidates.append(table)
    return candidates


def _clean_inline_text(value: Any) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip()
    text = text.replace("$", "")
    text = re.sub(r"\\[a-zA-Z]+", "", text)
    text = re.sub(r"\s+([,.;:)])", r"\1", text)
    text = re.sub(r"([(])\s+", r"\1", text)
    return text


def _normalize_rows(rows: Any) -> List[List[str]]:
    result = []
    if not isinstance(rows, list):
        return result
    for row in rows:
        if isinstance(row, list):
            cells = [str(cell or "").strip() for cell in row]
        elif isinstance(row, dict):
            cells = [str(value or "").strip() for value in row.values()]
        else:
            cells = [str(row or "").strip()]
        if any(cells):
            result.append(cells)
    return result


def _fit_image_box(path: str, x: float, y: float, w: float, h: float) -> Optional[Tuple[float, float, float, float]]:
    try:
        from PIL import Image

        with Image.open(path) as image:
            aspect = image.width / image.height if image.height else 1.0
    except Exception:
        return None
    target_w = w
    target_h = target_w / aspect
    if target_h > h:
        target_h = h
        target_w = target_h * aspect
    return (x + (w - target_w) / 2, y + (h - target_h) / 2, max(0.1, target_w), max(0.1, target_h))


def _build_rough_speaker_script(
    inventory: Dict[str, Any],
    rough: Dict[str, Any],
    label: str,
    narrative_mode: str = "rough_draft",
) -> str:
    title = inventory.get("paper", {}).get("title", rough.get("source_inventory_title", "Paper"))
    opener = {
        "rough_draft": "Speaker script generated for a from-scratch deck.",
        "golden_baseline1": "Speaker script generated for the warm academic proof-panel route.",
        "blind_experimental": "Speaker script generated for the blind experimental route; cue the claim/evidence ledger rather than academic baseline grammar.",
    }.get(narrative_mode, f"Speaker script generated for {label}.")
    lines = [f"# {title}", "", f"> {opener}", f"> Route label: {label}.", ""]
    slide_number = 1
    sections = _section_buckets(rough.get("slides", []) or [])
    lines.extend(
        [
            f"## Slide {slide_number}: {title}",
            "",
            "Suggested narration: Open with the paper title, why this route exists, and what the audience should look for in the deck.",
            "",
        ]
    )
    slide_number += 1
    agenda_title = "Route Map" if narrative_mode == "blind_experimental" else "Deck Map"
    section_names = ", ".join(section["title"] for section in sections[:6]) or "the paper story"
    lines.extend(
        [
            f"## Slide {slide_number}: {agenda_title}",
            "",
            f"Suggested narration: Preview the route through {section_names}.",
            "",
        ]
    )
    slide_number += 1
    for section in sections:
        lines.extend(
            [
                f"## Slide {slide_number}: {section['title']}",
                "",
                f"Suggested narration: Mark the transition into the {section['title']} module and state what evidence will matter next.",
                "",
            ]
        )
        slide_number += 1
        for slide in section.get("slides", []) or []:
            proof = slide.get("proof_object", {}) or {}
            if narrative_mode == "blind_experimental":
                narration = f"{slide.get('claim', '')} Then use the evidence board to support it: {slide.get('support', '')}".strip()
            else:
                narration = f"{slide.get('claim', '')} {slide.get('support', '')}".strip()
            lines.extend(
                [
                    f"## Slide {slide_number}: {slide.get('title', f'Slide {slide_number}')}",
                    "",
                    f"Suggested narration: {narration}".strip(),
                    "",
                    f"Evidence cue: {proof.get('type', 'text')} / {proof.get('id', '')} {proof.get('focus', '')}".strip(),
                    "",
                ]
            )
            slide_number += 1
    lines.extend(
        [
            f"## Slide {slide_number}: Closing",
            "",
            "Suggested narration: Re-state the main paper takeaway, mention the strongest evidence one more time, and close the route cleanly.",
            "",
        ]
    )
    return "\n".join(lines)


def _write_speaker_script_audit(
    speaker_path: Path,
    pptx_path: Path,
    style_id: str,
    output_path: Path,
) -> Dict[str, Any]:
    script_text = speaker_path.read_text(encoding="utf-8") if speaker_path.exists() else ""
    audit = _audit_speaker_script(script_text, pptx_path, style_id, speaker_path)
    _write_json(output_path, audit)
    return audit


def _audit_speaker_script(script_text: str, pptx_path: Path, style_id: str, speaker_path: Optional[Path] = None) -> Dict[str, Any]:
    slide_titles = _pptx_slide_titles(pptx_path)
    headings = re.findall(r"^##\s+Slide\s+\d+:\s+(.+)$", script_text, flags=re.MULTILINE)
    placeholder_phrases = [
        "The paper addresses the problem of",
        "Its goal is to make it",
        "In short",
        "Taken together",
        "Briefly state the slide message and move on.",
    ]
    findings: List[Dict[str, Any]] = []
    if len(headings) != len(slide_titles):
        findings.append({
            "type": "script_slide_count_mismatch",
            "severity": "medium",
            "message": f"Speaker script has {len(headings)} slide sections but PPTX has {len(slide_titles)} slides.",
        })
    missing_titles = [title for title in slide_titles if title not in headings]
    if missing_titles:
        findings.append({
            "type": "script_missing_slide_titles",
            "severity": "medium",
            "message": "Speaker script is missing one or more PPTX slide titles.",
            "evidence": {"missing_titles": missing_titles[:8]},
        })
    generic_hits = [phrase for phrase in placeholder_phrases if phrase.lower() in script_text.lower()]
    if generic_hits:
        findings.append({
            "type": "script_generic_placeholder_copy",
            "severity": "medium",
            "message": "Speaker script still contains generic or fragmented narration phrases.",
            "evidence": {"matches": generic_hits},
        })
    if _is_experimental_style(style_id) and "Evidence cue:" not in script_text:
        findings.append({
            "type": "script_missing_evidence_cue",
            "severity": "low",
            "message": "Blind experimental script should point the presenter to the evidence board on each slide.",
        })
    summary = {
        "finding_count": len(findings),
        "by_severity": dict(Counter(f.get("severity", "low") for f in findings)),
    }
    return {
        "schema_version": "speaker_script_audit.v1",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "speaker_script_path": str(speaker_path) if speaker_path else "",
        "pptx_path": str(pptx_path),
        "style_id": style_id,
        "slide_titles": slide_titles,
        "script_titles": headings,
        "summary": summary,
        "findings": findings,
    }


def _pptx_slide_titles(pptx_path: Path) -> List[str]:
    from pptx import Presentation

    titles: List[str] = []
    try:
        prs = Presentation(pptx_path)
    except Exception:
        return titles
    for slide in prs.slides:
        title = ""
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = " ".join(str(getattr(shape, "text", "") or "").split())
            if not text:
                continue
            if len(text.split()) <= 16 and len(text) <= 120:
                title = text
                break
        titles.append(title or f"Slide {len(titles) + 1}")
    return titles


def _write_curve(path: Path, routes: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    rows: List[Dict[str, Any]] = []
    for route in routes:
        for iteration in route.get("repair_log", {}).get("iterations", []):
            summary = iteration.get("audit_summary", {})
            severity = summary.get("by_severity", {})
            dimension = summary.get("by_dimension", {})
            dimension_scores = summary.get("dimension_scores", {})
            rows.append(
                {
                    "route_id": route["route_id"],
                    "style_id": route["style_id"],
                    "repair_profile": route["repair_profile"],
                    "iteration": iteration["iteration"],
                    "finding_count": summary.get("finding_count", 0),
                    "high": severity.get("high", 0),
                    "medium": severity.get("medium", 0),
                    "low": severity.get("low", 0),
                    "stop_reason": iteration.get("stop_reason", ""),
                    "content_findings": dimension.get("content", 0),
                    "evidence_findings": dimension.get("evidence", 0),
                    "layout_findings": dimension.get("layout", 0),
                    "typography_findings": dimension.get("typography", 0),
                    "component_fit_findings": dimension.get("component_fit", 0),
                    "style_findings": dimension.get("style", 0),
                    "repair_risk_findings": dimension.get("repair_risk", 0),
                    "dimension_scores_json": json.dumps(dimension_scores, ensure_ascii=False, sort_keys=True),
                    "pptx_path": iteration.get("pptx_path", ""),
                    "speaker_script_path": iteration.get("speaker_script_path", ""),
                    "speaker_script_audit_path": iteration.get("speaker_script_audit_path", ""),
                }
            )
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.DictWriter(handle, fieldnames=list(rows[0].keys()) if rows else ["route_id"])
        writer.writeheader()
        writer.writerows(rows)
    return rows


def _write_artifact_index(path: Path, routes: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    rows: List[Dict[str, Any]] = []
    for route in routes:
        for iteration in route.get("repair_log", {}).get("iterations", []):
            summary = iteration.get("audit_summary", {})
            rows.append(
                {
                    "route_id": route["route_id"],
                    "style_id": route["style_id"],
                    "repair_profile": route["repair_profile"],
                    "iteration": iteration["iteration"],
                    "pptx_path": iteration.get("pptx_path", ""),
                    "speaker_script_path": iteration.get("speaker_script_path", route.get("speaker_script_path", "")),
                    "speaker_script_audit_path": iteration.get("speaker_script_audit_path", route.get("speaker_script_audit_path", "")),
                    "nonvisual_audit_path": iteration.get("nonvisual_audit_path", ""),
                    "style_contract_path": route.get("style_contract_path", ""),
                    "finding_count": summary.get("finding_count", 0),
                    "high": summary.get("by_severity", {}).get("high", 0),
                    "medium": summary.get("by_severity", {}).get("medium", 0),
                    "low": summary.get("by_severity", {}).get("low", 0),
                    "stop_reason": iteration.get("stop_reason", ""),
                }
            )
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.DictWriter(handle, fieldnames=list(rows[0].keys()) if rows else ["route_id"])
        writer.writeheader()
        writer.writerows(rows)
    return rows


def _render_comparison_report(
    manifest: Dict[str, Any],
    generation: Dict[str, Any],
    routes: List[Dict[str, Any]],
    curve_rows: List[Dict[str, Any]],
) -> str:
    paper_title = Path(manifest["paper_path"]).stem
    lines = [
        f"# {paper_title} Four-Way Benchmark Validation",
        "",
        f"- Paper: `{manifest['paper_path']}`",
        f"- Fresh parse output root: `{manifest['fresh_parse_output_root']}`",
        f"- Fresh academic generation seconds: {generation.get('elapsed_seconds')}",
        "",
        "## Route Summary",
        "",
        "| Route | Style | Repair profile | Findings | High | Medium | Low | Drift risk | PPTX |",
        "| --- | --- | --- | ---: | ---: | ---: | ---: | --- | --- |",
    ]
    for route in routes:
        summary = route.get("audit_summary", {})
        severity = summary.get("by_severity", {})
        lines.append(
            f"| `{route['route_id']}` | `{route['style_id']}` | `{route['repair_profile']}` | "
            f"{summary.get('finding_count', 0)} | {severity.get('high', 0)} | {severity.get('medium', 0)} | "
            f"{severity.get('low', 0)} | {route.get('style_drift_report', {}).get('drift_risk', '')} | "
            f"`{route['pptx_path']}` |"
        )
    lines.extend(["", "## Benchmark Curve", "", "| Route | Iteration | High | Medium | Low | Stop reason |", "| --- | ---: | ---: | ---: | ---: | --- |"])
    for row in curve_rows:
        lines.append(
            f"| `{row['route_id']}` | {row['iteration']} | {row['high']} | {row['medium']} | {row['low']} | {row['stop_reason']} |"
        )
    lines.extend(["", "## Dimension Signals", "", "| Route | Iteration | Content | Layout | Typography | Component fit | Style | Repair risk |", "| --- | ---: | ---: | ---: | ---: | ---: | ---: | ---: |"])
    for row in curve_rows:
        lines.append(
            f"| `{row['route_id']}` | {row['iteration']} | {row.get('content_findings', 0)} | {row.get('layout_findings', 0)} | "
            f"{row.get('typography_findings', 0)} | {row.get('component_fit_findings', 0)} | "
            f"{row.get('style_findings', 0)} | {row.get('repair_risk_findings', 0)} |"
        )
    lines.extend(
        [
            "",
            "## Artifact Tracking",
            "",
            "- `score_curve.csv` stores the benchmark score curve by route and iteration.",
            "- `score_curve.csv` now also stores dimension-level finding counts and `dimension_scores_json`.",
            "- `artifact_index.csv` stores the PPTX, speaker script, speaker-script audit, and nonvisual audit path for every saved iteration.",
            "- Blind experimental runs also persist a per-run `style_contract.json` so the route does not silently reuse a prior promoted style.",
            "- Blind experimental runs emit `visual_human_review_packet.md` and `visual_compare/page_level_audit_delta.*` for human review of metric-improved but visually risky repairs.",
            "",
        ]
    )
    lines.extend(["", "## Representative Badcases", ""])
    for route in routes:
        by_type = route.get("audit_summary", {}).get("by_type", {})
        top = ", ".join(f"`{kind}` x{count}" for kind, count in Counter(by_type).most_common(6)) or "none"
        lines.append(f"- `{route['route_id']}`: {top}")
    lines.extend(
        [
            "",
            "## Human Quick Check",
            "",
            "- Open each `slides.pptx` and check the cover, agenda, one figure/table page, one text-evidence page, and closing page.",
            "- Confirm the academic route did not inherit rounded proof-panel polish.",
            "- Confirm golden_baseline1 uses rounded proof-panel grammar only inside its own scoped branch.",
            "- Confirm blind experimental branch does not use the protected baseline text signals or proof-panel labels.",
            "- Use `score_curve.csv` as the before/after evidence for benchmark-driven repair.",
            "",
        ]
    )
    return "\n".join(lines)


def _route_result(
    route_id: str,
    route_dir: Path,
    style_id: str,
    repair_profile: str,
    pptx_path: Path,
    speaker_path: Optional[Path],
    speaker_audit_path: Optional[Path],
    qa_path: Optional[Path],
    audit: Dict[str, Any],
    repair_log: Dict[str, Any],
    style_report: Dict[str, Any],
) -> Dict[str, Any]:
    return {
        "route_id": route_id,
        "route_dir": str(route_dir),
        "style_id": style_id,
        "repair_profile": repair_profile,
        "pptx_path": str(pptx_path),
        "speaker_script_path": str(speaker_path) if speaker_path else "",
        "speaker_script_audit_path": str(speaker_audit_path) if speaker_audit_path else "",
        "layout_qa_path": str(qa_path) if qa_path else "",
        "nonvisual_audit_path": str(route_dir / "nonvisual_audit.json"),
        "repair_log_path": str(route_dir / "repair_log.json"),
        "style_drift_report_path": str(route_dir / "style_drift_report.json"),
        "style_contract_path": str(route_dir / "style_contract.json") if (route_dir / "style_contract.json").exists() else "",
        "audit_summary": audit.get("summary", {}),
        "repair_log": repair_log,
        "style_drift_report": style_report,
    }


def _iteration_payload(
    iteration: int,
    pptx_path: Path,
    audit: Dict[str, Any],
    applied_repairs: List[str],
    stop_reason: str,
    speaker_script_path: Optional[Path] = None,
    speaker_script_audit_path: Optional[Path] = None,
) -> Dict[str, Any]:
    return {
        "iteration": iteration,
        "pptx_path": str(pptx_path),
        "nonvisual_audit_path": str(Path(pptx_path).with_name("nonvisual_audit.json")),
        "speaker_script_path": str(speaker_script_path) if speaker_script_path else "",
        "speaker_script_audit_path": str(speaker_script_audit_path) if speaker_script_audit_path else "",
        "audit_summary": audit.get("summary", {}),
        "applied_repairs": applied_repairs,
        "stop_reason": stop_reason,
    }


def _repair_log_payload(route_id: str, style_id: str, repair_profile: str, iterations: List[Dict[str, Any]], note: str) -> Dict[str, Any]:
    return {
        "schema_version": "benchmark_repair_log.v1",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "route_id": route_id,
        "style_id": style_id,
        "repair_profile": repair_profile,
        "note": note,
        "iterations": iterations,
    }


def _stop_reason_for_profile(audit: Dict[str, Any], repair_profile: str) -> str:
    severity = audit.get("summary", {}).get("by_severity", {})
    high = int(severity.get("high", 0))
    medium = int(severity.get("medium", 0))
    if repair_profile == "audit_only":
        return "audit_only_no_mutation"
    if high == 0 and medium == 0:
        return "no_high_or_medium_findings"
    if repair_profile == "global_correctness_repair":
        return "global_findings_need_spec_level_repair" if _global_findings(audit) else "style_specific_findings_detect_only"
    return "remaining_findings_reported_for_next_bounded_repair"


def _needs_experimental_repair(audit: Dict[str, Any]) -> bool:
    severity = audit.get("summary", {}).get("by_severity", {})
    return bool(severity.get("high") or severity.get("medium") or severity.get("low"))


def _global_findings(audit: Dict[str, Any]) -> List[Dict[str, Any]]:
    return [finding for finding in audit.get("findings", []) or [] if finding.get("type") in GLOBAL_CORRECTNESS_RULES]


def _audit_to_file(pptx_path: Path, output_path: Path) -> Dict[str, Any]:
    audit = inspect_pptx_nonvisual(pptx_path)
    _write_json(output_path, audit)
    return audit


def _copy_if_exists(source: Path, target: Path) -> Optional[Path]:
    if not source.exists():
        return None
    target.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source, target)
    return target


def _find_new_or_latest_output(config_dir: Path, before: set[Path]) -> Optional[Path]:
    if not config_dir.exists():
        return None
    dirs = [path for path in config_dir.iterdir() if path.is_dir()]
    new_dirs = [path for path in dirs if path.resolve() not in before]
    candidates = new_dirs or dirs
    return max(candidates, key=lambda item: item.stat().st_mtime) if candidates else None


def _pptx_text_signals(pptx_path: Path) -> Dict[str, int]:
    from pptx import Presentation

    signals = Counter()
    try:
        prs = Presentation(pptx_path)
    except Exception:
        return {}
    targets = ["ACADEMIC PAPER READING", "DECK MAP", "PAPER HIGHLIGHTS", "PROOF OBJECT"]
    for slide in prs.slides:
        for shape in slide.shapes:
            text = " ".join(str(getattr(shape, "text", "") or "").upper().split())
            for target in targets:
                if text == target or text.startswith(target + " /") or text.startswith(target + ":"):
                    signals[target] += 1
    return dict(signals)


def _evidence_footer(slide_data: Dict[str, Any]) -> str:
    refs = []
    for item in slide_data.get("source_evidence", []) or []:
        if isinstance(item, dict):
            refs.append(f"{item.get('source', '')}:{item.get('id', '')}".strip(":"))
    return "Sources: " + "; ".join(refs[:3]) if refs else "Sources: fresh parse checkpoints"


def _slugify(text: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9]+", "_", str(text or "").strip())
    cleaned = re.sub(r"_+", "_", cleaned).strip("_")
    return cleaned.lower() or "paper"


def _is_experimental_style(style_id: str) -> bool:
    return str(style_id or "").startswith(EXPERIMENTAL_STYLE_PREFIX) or style_id == BLIND_RECTANGULAR_STYLE_ID


def _limit_words(text: Any, max_words: int) -> str:
    words = str(text or "").replace("\n", " ").split()
    if len(words) <= max_words:
        return " ".join(words)
    return " ".join(words[:max_words]).rstrip(" ,;:") + "."


def _read_json(path: Path) -> Dict[str, Any]:
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return {}


def _write_json(path: Path, data: Dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def _tail_text(path: Path, lines: int) -> str:
    try:
        return "\n".join(path.read_text(encoding="utf-8", errors="replace").splitlines()[-lines:])
    except Exception:
        return ""


def main(argv: Optional[List[str]] = None) -> int:
    parser = argparse.ArgumentParser(description="Run single-paper four-way benchmark validation.")
    parser.add_argument("--paper", required=True, help="Input PDF path.")
    parser.add_argument("--run-dir", help="Output run directory under benchmark_runs or an explicit path.")
    parser.add_argument("--slides", type=int, default=24)
    parser.add_argument("--length", choices=["short", "medium", "long"], default="medium")
    parser.add_argument("--no-fast", action="store_true")
    parser.add_argument("--from-stage", choices=["rag", "summary", "plan", "generate"], default="rag")
    parser.add_argument("--python", default=sys.executable, dest="python_executable")
    parser.add_argument("--max-iterations", type=int, default=2)
    args = parser.parse_args(argv)

    result = run_fourway_validation(
        paper_path=Path(args.paper),
        run_dir=Path(args.run_dir) if args.run_dir else None,
        slides=args.slides,
        length=args.length,
        fast=not args.no_fast,
        from_stage=args.from_stage,
        python_executable=args.python_executable,
        max_iterations=args.max_iterations,
    )
    print(json.dumps({"run_dir": result["run_dir"], "comparison_report": result["comparison_report"]}, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
