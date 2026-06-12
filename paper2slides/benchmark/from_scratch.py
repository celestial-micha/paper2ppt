"""Prepare from-scratch template experiment inputs from existing checkpoints.

This module deliberately does not parse PDFs or call any LLM. It converts the
already-good paper understanding checkpoints into reusable content artifacts
that a new visual system can consume without inheriting the academic renderer's
page skeleton.
"""
from __future__ import annotations

import argparse
import html
import json
import re
import shutil
import subprocess
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional


SUMMARY_SECTIONS = [
    ("paper_info", "paper_metadata"),
    ("motivation", "motivation"),
    ("solution", "method"),
    ("results", "results"),
    ("contributions", "contribution"),
]


THEME_HEX = {
    "paper": "F7F2E8",
    "paper_alt": "FBF8F0",
    "ink": "1C2936",
    "muted_ink": "4F5A61",
    "soft_text": "74706A",
    "line": "D8D1C2",
    "panel": "EEE9DE",
    "panel_light": "F4EFE6",
    "teal": "4B7F78",
    "teal_deep": "2F5F5A",
    "teal_light": "DBE9E3",
    "gold": "C99A45",
    "gold_light": "F1E2BC",
    "clay": "B76B57",
    "clay_light": "ECD4CA",
    "sage": "DDE6D8",
    "white": "FFFDF8",
}


def build_content_inventory(
    summary_checkpoint: Path,
    plan_checkpoint: Optional[Path] = None,
    spec_checkpoint: Optional[Path] = None,
) -> Dict[str, Any]:
    """Build a content inventory from existing paper2ppt checkpoints."""
    summary_data = _read_json(summary_checkpoint)
    plan_data = _read_json(plan_checkpoint) if plan_checkpoint else {}
    spec_data = _read_json(spec_checkpoint) if spec_checkpoint else {}

    content = summary_data.get("content", {}) if isinstance(summary_data, dict) else {}
    origin = summary_data.get("origin", {}) if isinstance(summary_data, dict) else {}
    plan = plan_data.get("plan", {}) if isinstance(plan_data, dict) else {}

    summary_items = _summary_items(content)
    plan_slides = _plan_slides(plan)
    curated_slides = _curated_slides(spec_data)
    assets = _assets(origin)
    metrics = _metrics_from_curated_slides(curated_slides)
    paper_highlights = _paper_highlights(summary_items, plan_slides, curated_slides, metrics)

    title = _extract_title(content.get("paper_info", "")) or _first_non_empty(
        [spec_data.get("title", ""), plan_slides[0]["title"] if plan_slides else ""]
    )

    return {
        "schema_version": "content_inventory.v1",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "purpose": "from_scratch_template_experiment",
        "source_checkpoints": {
            "summary": str(summary_checkpoint),
            "plan": str(plan_checkpoint) if plan_checkpoint else "",
            "slide_spec": str(spec_checkpoint) if spec_checkpoint else "",
        },
        "paper": {
            "title": title or "Untitled Paper",
            "metadata_text": _clean_text(content.get("paper_info", "")),
            "content_type": summary_data.get("content_type", "paper"),
        },
        "summary_items": summary_items,
        "plan_slides": plan_slides,
        "curated_slides": curated_slides,
        "assets": assets,
        "metrics": metrics,
        "paper_highlights": paper_highlights,
        "coverage": _coverage(summary_items, plan_slides, curated_slides, assets, metrics),
        "design_constraints": {
            "reuse_allowed": [
                "paper understanding",
                "claims and evidence",
                "figure/table/metric references",
                "QA lessons and badcase rules",
            ],
            "reuse_forbidden": [
                "academic header skeleton",
                "academic key-message block as the primary page structure",
                "numbered-point rhythm as the default layout grammar",
                "academic title/toc/section/content macro page skeleton",
            ],
        },
    }


def build_rough_draft_spec(inventory: Dict[str, Any]) -> Dict[str, Any]:
    """Build a content-complete rough draft spec from the inventory."""
    curated_by_id = {
        slide.get("slide_id"): slide
        for slide in inventory.get("curated_slides", [])
        if slide.get("slide_id")
    }

    slides: List[Dict[str, Any]] = []
    for index, plan_slide in enumerate(inventory.get("plan_slides", []), start=1):
        slide_id = plan_slide.get("slide_id") or f"slide_{index:02d}"
        curated = curated_by_id.get(slide_id, {})
        proof = _proof_object(plan_slide, curated)
        claim = _first_non_empty(
            [
                curated.get("takeaway", ""),
                (curated.get("points") or [{}])[0].get("claim", "") if curated.get("points") else "",
                plan_slide.get("title", ""),
            ]
        )
        support = _first_non_empty(
            [
                " ".join(point.get("detail", "") for point in curated.get("points", [])[:3]),
                plan_slide.get("content", ""),
            ]
        )
        slides.append(
            {
                "slide_id": slide_id,
                "source_plan_id": plan_slide.get("source_id", slide_id),
                "title": plan_slide.get("title", "") or curated.get("title", "") or f"Slide {index}",
                "slide_role": _slide_role(plan_slide, curated, index, len(inventory.get("plan_slides", []))),
                "claim": _limit_words(claim, 18),
                "support": _limit_words(support, 80),
                "proof_object": proof,
                "source_evidence": _source_evidence(plan_slide, curated, proof),
                "content_priority": "must_keep" if index <= 2 or proof["type"] != "text_evidence" else "normal",
                "visual_instruction": "rough draft only; preserve content before aesthetic decisions",
            }
        )

    if not slides:
        slides = _fallback_slides_from_summary(inventory)

    return {
        "schema_version": "rough_draft_spec.v1",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "purpose": "content_complete_draft_before_visual_design",
        "source_inventory_title": inventory.get("paper", {}).get("title", "Untitled Paper"),
        "rules": {
            "aesthetic_goal": "none",
            "one_claim_per_slide": True,
            "proof_object_required": True,
            "baseline_skeleton_forbidden": True,
        },
        "slides": slides,
    }


def write_from_scratch_artifacts(
    summary_checkpoint: Path,
    plan_checkpoint: Optional[Path],
    spec_checkpoint: Optional[Path],
    output_dir: Path,
    pptx_output: Optional[Path] = None,
    render_review_dir: Optional[Path] = None,
) -> Dict[str, str]:
    """Write content inventory and rough draft JSON artifacts."""
    inventory = build_content_inventory(summary_checkpoint, plan_checkpoint, spec_checkpoint)
    rough = build_rough_draft_spec(inventory)
    output_dir.mkdir(parents=True, exist_ok=True)
    inventory_path = output_dir / "content_inventory.json"
    rough_path = output_dir / "rough_draft_spec.json"
    _write_json(inventory_path, inventory)
    _write_json(rough_path, rough)
    paths = {
        "content_inventory": str(inventory_path),
        "rough_draft_spec": str(rough_path),
    }
    if pptx_output:
        render_rough_draft_pptx(inventory, rough, pptx_output)
        from .nonvisual_audit import inspect_pptx_nonvisual

        nonvisual_audit = inspect_pptx_nonvisual(pptx_output)
        nonvisual_audit_path = output_dir / "nonvisual_audit.json"
        _write_json(nonvisual_audit_path, nonvisual_audit)
        audit = audit_rough_draft(inventory, rough, pptx_output)
        audit_path = output_dir / "visual_audit.json"
        _write_json(audit_path, audit)
        paths["rough_draft_pptx"] = str(pptx_output)
        paths["nonvisual_audit"] = str(nonvisual_audit_path)
        paths["visual_audit"] = str(audit_path)
        if render_review_dir:
            render_status = export_visual_review_pages(pptx_output, audit, render_review_dir)
            render_status_path = output_dir / "visual_render_status.json"
            _write_json(render_status_path, render_status)
            paths["visual_render_status"] = str(render_status_path)
    return paths


def export_visual_review_pages(pptx_path: Path, audit: Dict[str, Any], output_dir: Path) -> Dict[str, Any]:
    """Export selected PPTX pages to PNG when a local renderer is available."""
    output_dir.mkdir(parents=True, exist_ok=True)
    requests = audit.get("visual_review_manifest", {}).get("render_requests", []) or []
    pages = [int(item.get("page")) for item in requests if item.get("page")]
    pages = sorted(set(page for page in pages if page > 0))
    status: Dict[str, Any] = {
        "schema_version": "visual_render_status.v1",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "pptx_path": str(pptx_path),
        "output_dir": str(output_dir),
        "requested_pages": pages,
        "rendered_files": [],
        "text_snapshot_files": [],
        "status": "not_run",
        "renderer": "",
        "message": "",
    }
    if not pptx_path.exists():
        status.update({"status": "failed", "message": "PPTX path does not exist."})
        return status
    if not pages:
        status.update({"status": "skipped", "message": "No render requests were present in visual audit."})
        return status

    text_snapshots = _export_slide_text_snapshots(pptx_path, output_dir, pages)
    if text_snapshots:
        status["text_snapshot_files"] = text_snapshots

    for renderer in (_render_pages_with_powerpoint, _render_pages_with_powershell_powerpoint, _render_pages_with_libreoffice):
        rendered = renderer(pptx_path, output_dir, pages)
        if rendered.get("status") == "rendered":
            rendered.setdefault("text_snapshot_files", text_snapshots)
            status.update(rendered)
            return status
        status.setdefault("attempts", []).append(rendered)

    status.update(
        {
            "status": "renderer_unavailable",
            "message": "No supported PPTX-to-PNG renderer is available. Text snapshots were exported; install PowerPoint, pywin32, or LibreOffice/soffice for PNG review.",
        }
    )
    return status


def _export_slide_text_snapshots(pptx_path: Path, output_dir: Path, pages: List[int]) -> List[str]:
    try:
        from pptx import Presentation
    except Exception:
        return []

    try:
        prs = Presentation(pptx_path)
    except Exception:
        return []

    rendered: List[str] = []
    for page in pages:
        if page < 1 or page > len(prs.slides):
            continue
        lines = [f"slide: {page:02d}"]
        slide = prs.slides[page - 1]
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = _clean_text(shape.text)
                if text:
                    lines.append(text)
            if getattr(shape, "has_table", False):
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [_clean_text(cell.text) for cell in row.cells]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    lines.append("[table]")
                    lines.extend(rows)
        snapshot_path = output_dir / f"slide_{page:02d}.txt"
        snapshot_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
        rendered.append(str(snapshot_path))
    return rendered


def _render_pages_with_powerpoint(pptx_path: Path, output_dir: Path, pages: List[int]) -> Dict[str, Any]:
    try:
        import win32com.client  # type: ignore
    except Exception as exc:
        return {"renderer": "powerpoint_com", "status": "unavailable", "message": str(exc)}

    app = None
    presentation = None
    rendered: List[str] = []
    try:
        app = win32com.client.Dispatch("PowerPoint.Application")
        presentation = app.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
        slide_count = int(presentation.Slides.Count)
        for page in pages:
            if page > slide_count:
                continue
            output_path = output_dir / f"slide_{page:02d}.png"
            presentation.Slides(page).Export(str(output_path.resolve()), "PNG", 1920, 1080)
            if output_path.exists():
                rendered.append(str(output_path))
    except Exception as exc:
        return {"renderer": "powerpoint_com", "status": "failed", "message": str(exc), "rendered_files": rendered}
    finally:
        try:
            if presentation is not None:
                presentation.Close()
        except Exception:
            pass
        try:
            if app is not None:
                app.Quit()
        except Exception:
            pass

    if rendered:
        return {"renderer": "powerpoint_com", "status": "rendered", "message": "Rendered requested pages.", "rendered_files": rendered}
    return {"renderer": "powerpoint_com", "status": "failed", "message": "PowerPoint opened, but no requested pages were rendered.", "rendered_files": rendered}


def _render_pages_with_powershell_powerpoint(pptx_path: Path, output_dir: Path, pages: List[int]) -> Dict[str, Any]:
    command = shutil.which("powershell") or shutil.which("pwsh")
    if not command:
        return {"renderer": "powershell_powerpoint_com", "status": "unavailable", "message": "PowerShell command not found."}

    script_path = output_dir / "_render_powerpoint_pages.ps1"
    script = r'''
param(
  [Parameter(Mandatory=$true)][string]$PptxPath,
  [Parameter(Mandatory=$true)][string]$OutputDir,
  [Parameter(Mandatory=$true)][string]$PagesCsv
)
$ErrorActionPreference = "Stop"
[Console]::OutputEncoding = [System.Text.Encoding]::UTF8
$app = $null
$presentation = $null
$rendered = @()
try {
  $pages = $PagesCsv -split "," | Where-Object { $_ } | ForEach-Object { [int]$_ }
  $app = New-Object -ComObject PowerPoint.Application
  $presentation = $app.Presentations.Open($PptxPath, $true, $false, $false)
  $slideCount = [int]$presentation.Slides.Count
  foreach ($page in $pages) {
    if ($page -gt $slideCount) { continue }
    $outPath = Join-Path $OutputDir ("slide_{0:D2}.png" -f $page)
    $presentation.Slides.Item($page).Export($outPath, "PNG", 1920, 1080)
    if (Test-Path -LiteralPath $outPath) {
      $rendered += $outPath
    }
  }
  @{ status = "rendered"; message = "Rendered requested pages via PowerShell PowerPoint COM."; rendered_files = $rendered } | ConvertTo-Json -Compress
}
catch {
  @{ status = "failed"; message = $_.Exception.Message; rendered_files = $rendered } | ConvertTo-Json -Compress
  exit 2
}
finally {
  if ($presentation -ne $null) {
    try { $presentation.Close() } catch {}
    try { [System.Runtime.InteropServices.Marshal]::ReleaseComObject($presentation) | Out-Null } catch {}
  }
  if ($app -ne $null) {
    try { $app.Quit() } catch {}
    try { [System.Runtime.InteropServices.Marshal]::ReleaseComObject($app) | Out-Null } catch {}
  }
}
'''
    try:
        script_path.write_text(script.strip() + "\n", encoding="utf-8")
        result = subprocess.run(
            [
                command,
                "-NoProfile",
                "-ExecutionPolicy",
                "Bypass",
                "-File",
                str(script_path),
                str(pptx_path.resolve()),
                str(output_dir.resolve()),
                ",".join(str(page) for page in pages),
            ],
            capture_output=True,
            text=True,
            timeout=180,
            check=False,
        )
    except Exception as exc:
        return {"renderer": "powershell_powerpoint_com", "status": "failed", "message": str(exc), "rendered_files": []}
    finally:
        try:
            script_path.unlink()
        except Exception:
            pass

    payload = _last_json_line(result.stdout)
    if not payload:
        return {
            "renderer": "powershell_powerpoint_com",
            "status": "failed" if result.returncode else "unavailable",
            "message": (result.stderr or result.stdout or "PowerShell did not return renderer status.")[-1000:],
            "rendered_files": [],
        }
    status = str(payload.get("status", "failed"))
    rendered_files = [str(path) for path in payload.get("rendered_files", []) if Path(str(path)).exists()]
    if status == "rendered" and rendered_files:
        return {
            "renderer": "powershell_powerpoint_com",
            "status": "rendered",
            "message": payload.get("message", "Rendered requested pages via PowerShell PowerPoint COM."),
            "rendered_files": rendered_files,
        }
    return {
        "renderer": "powershell_powerpoint_com",
        "status": "failed",
        "message": payload.get("message", "PowerShell PowerPoint COM did not render any requested pages."),
        "rendered_files": rendered_files,
    }


def _render_pages_with_libreoffice(pptx_path: Path, output_dir: Path, pages: List[int]) -> Dict[str, Any]:
    command = shutil.which("soffice") or shutil.which("libreoffice")
    if not command:
        return {"renderer": "libreoffice", "status": "unavailable", "message": "soffice/libreoffice command not found."}

    before = {path.name for path in output_dir.glob("*.png")}
    try:
        result = subprocess.run(
            [command, "--headless", "--convert-to", "png", "--outdir", str(output_dir), str(pptx_path.resolve())],
            capture_output=True,
            text=True,
            timeout=120,
            check=False,
        )
    except Exception as exc:
        return {"renderer": "libreoffice", "status": "failed", "message": str(exc), "rendered_files": []}

    after_files = [path for path in output_dir.glob("*.png") if path.name not in before]
    rendered: List[str] = []
    for path in sorted(after_files):
        match = re.search(r"(\d+)", path.stem)
        page = int(match.group(1)) if match else None
        if page in pages:
            target = output_dir / f"slide_{page:02d}.png"
            if path != target:
                path.replace(target)
            rendered.append(str(target))

    if rendered:
        return {
            "renderer": "libreoffice",
            "status": "rendered",
            "message": "Rendered requested pages via LibreOffice.",
            "rendered_files": rendered,
            "stdout": result.stdout[-1000:],
            "stderr": result.stderr[-1000:],
        }
    return {
        "renderer": "libreoffice",
        "status": "failed",
        "message": "LibreOffice did not produce matching PNG pages.",
        "rendered_files": [],
        "stdout": result.stdout[-1000:],
        "stderr": result.stderr[-1000:],
    }


def _last_json_line(text: str) -> Dict[str, Any]:
    for line in reversed((text or "").splitlines()):
        line = line.strip()
        if not line or not line.startswith("{"):
            continue
        try:
            payload = json.loads(line)
        except Exception:
            continue
        if isinstance(payload, dict):
            return payload
    return {}


def render_rough_draft_pptx(inventory: Dict[str, Any], rough: Dict[str, Any], output_path: Path) -> Path:
    """Render a plain, content-complete rough PPTX without using baseline styles."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = rough.get("slides", []) or []
    title = inventory.get("paper", {}).get("title", rough.get("source_inventory_title", "Rough Draft"))
    sections = _organize_deck_sections(slides)
    total_pages = 2 + len(sections) + len(slides) + 1
    figure_paths = _figure_path_index(inventory)
    table_index = _table_index(inventory)

    page = 1
    _add_academic_title_slide(prs, inventory, title, page, total_pages)
    page += 1
    _add_agenda_slide(prs, sections, page, total_pages)
    page += 1

    content_index = 0
    for section_index, section in enumerate(sections, start=1):
        _add_section_divider_slide(prs, section, section_index, page, total_pages)
        page += 1
        for slide_data in section["slides"]:
            content_index += 1
            layout = _layout_family(slide_data, content_index)
            _add_content_slide(
                prs=prs,
                slide_data=slide_data,
                page=page,
                total_pages=total_pages,
                layout=layout,
                figure_paths=figure_paths,
                table_index=table_index,
            )
            slide_data["layout_family"] = layout
            page += 1

    _add_closing_slide(prs, inventory, page, total_pages)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def audit_rough_draft(inventory: Dict[str, Any], rough: Dict[str, Any], pptx_path: Optional[Path] = None) -> Dict[str, Any]:
    """Rule-based visual audit for the generated rough deck."""
    slides = rough.get("slides", []) or []
    table_index = _table_index(inventory)
    figure_paths = _figure_path_index(inventory)
    layout_families = [_layout_family(slide, index) for index, slide in enumerate(slides, start=1)]
    sections = _organize_deck_sections(slides)
    page_by_slide_object = _content_page_map(sections)
    warnings = []

    dominant = _dominant_layout_family(layout_families)
    if dominant and dominant["ratio"] > 0.45:
        warnings.append(
            {
                "type": "layout_monotony",
                "severity": "medium",
                "message": f"{dominant['family']} is used on {dominant['count']} of {len(layout_families)} content slides.",
            }
        )
    repeated = _max_consecutive_layout(layout_families)
    if repeated["count"] >= 3:
        warnings.append(
            {
                "type": "consecutive_layout_repeat",
                "severity": "medium",
                "message": f"{repeated['family']} repeats for {repeated['count']} consecutive content slides.",
            }
        )

    review_targets = [
        {"page": 1, "reason": "title page: cover composition and color balance", "severity": "high"},
        {"page": 2, "reason": "agenda page: module components and count rail", "severity": "high"},
    ]
    for index, slide in enumerate(slides, start=1):
        proof = slide.get("proof_object", {}) or {}
        proof_type = proof.get("type", "")
        proof_id = proof.get("id", "")
        slide_page = page_by_slide_object.get(id(slide), 0)
        layout = layout_families[index - 1] if index - 1 < len(layout_families) else ""
        if proof_type == "table" and not table_index.get(proof_id, {}).get("rows"):
            warnings.append(
                {
                    "type": "table_proof_missing_rows",
                    "severity": "high",
                    "slide_id": slide.get("slide_id", ""),
                    "message": f"{proof_id} is referenced but has no parsed rows.",
                }
            )
            review_targets.append({"page": slide_page, "reason": f"table rows missing for {proof_id}", "severity": "high"})
        if proof_type == "table" and table_index.get(proof_id, {}).get("row_count", 0) > 18:
            review_targets.append(
                {
                    "page": slide_page,
                    "reason": f"dense parsed table {proof_id}: check row grammar, overlap, and caption placement",
                    "severity": "medium",
                }
            )
        if layout == "table_bottom":
            review_targets.append(
                {
                    "page": slide_page,
                    "reason": "table_bottom layout: verify the proof panel starts below the support text",
                    "severity": "high",
                }
            )
        if proof_type == "figure" and not Path(figure_paths.get(proof_id, "")).exists():
            warnings.append(
                {
                    "type": "figure_proof_missing_file",
                    "severity": "medium",
                    "slide_id": slide.get("slide_id", ""),
                    "message": f"{proof_id} is referenced but the image file is missing.",
                }
            )
            review_targets.append({"page": slide_page, "reason": f"figure file missing for {proof_id}", "severity": "medium"})
        if proof_type == "text_evidence" and len(_clean_text(proof.get("focus", ""))) < 80:
            warnings.append(
                {
                    "type": "thin_text_proof",
                    "severity": "low",
                    "slide_id": slide.get("slide_id", ""),
                    "message": "Text proof is short; check whether this page needs a stronger proof object.",
                }
            )
            review_targets.append(
                {
                    "page": slide_page,
                    "reason": "short text evidence: verify compact cards avoid a large empty proof panel",
                    "severity": "medium",
                }
            )
        if layout in {"metric_compact_band", "metric_left_alt"}:
            review_targets.append(
                {
                    "page": slide_page,
                    "reason": f"{layout}: verify metric card grammar and no orphan labels",
                    "severity": "medium",
                }
            )
    for family in sorted(set(layout_families)):
        try:
            first_index = layout_families.index(family) + 1
        except ValueError:
            continue
        first_slide = slides[first_index - 1]
        review_targets.append(
            {
                "page": page_by_slide_object.get(id(first_slide), 0),
                "reason": f"first {family} layout",
                "severity": "low",
            }
        )

    deduped_targets = _dedupe_review_targets(review_targets)
    return {
        "schema_version": "rough_visual_audit.v2",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "pptx_path": str(pptx_path) if pptx_path else "",
        "content_slide_count": len(slides),
        "layout_family_counts": _count_items(layout_families),
        "max_consecutive_layout": repeated,
        "warnings": warnings,
        "review_targets": deduped_targets,
        "non_visual_review_manifest": {
            "review_mode": "pptx_metadata_only",
            "requires_rendered_screenshots": False,
            "requires_vision_model": False,
            "audit_artifact": "nonvisual_audit.json",
            "checks": [
                "shape bounding box overlap",
                "text capacity estimate from box size and font size",
                "font-size floor by inferred role",
                "low text density without auto-shrinking components",
                "table row/column readability estimates",
                "slide occupancy and whitespace risk",
            ],
            "repair_policy": [
                "preserve accepted component composition before changing card/panel geometry",
                "increase typography or improve copy allocation before resizing components",
                "treat low-density warnings as review hints, not automatic shrink commands",
            ],
        },
        "visual_review_manifest": {
            "requires_rendered_screenshots": False,
            "renderer_contract": "Legacy optional path only. Current route prefers nonvisual_audit.json and avoids screenshot/vision review by default.",
            "render_requests": [
                {
                    "page": target["page"],
                    "reason": target.get("reason", ""),
                    "severity": target.get("severity", "medium"),
                    "expected_artifact": f"slide_{int(target['page']):02d}.png",
                }
                for target in deduped_targets
                if target.get("page")
            ],
            "badcase_rules": [
                "cover must not use a full-width black top bar",
                "agenda module numbers must not be black circular buttons",
                "section dividers must use warm academic color, not a black/white block page",
                "short text evidence must render as compact notes rather than a large empty proof panel",
                "metric pages must keep value, label, and context visible on every card",
                "dense table pages must preserve row/column grammar and avoid title/table overlap",
                "table proof panels must begin below claim/support text and leave a readable gutter",
                "body prose should not shrink below readable presentation size except in footers or dense tables",
                "no component may visually cover readable text",
            ],
        },
        "checks": [
            "title page exists",
            "agenda page exists",
            "section divider pages exist",
            "layout family repetition",
            "table proof has parsed rows",
            "figure proof image path exists",
            "thin text proof candidates",
            "machine-readable non-visual geometry/text-capacity QA",
            "legacy optional render requests only when explicitly enabled",
        ],
    }


def _organize_deck_sections(slides: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    buckets = [
        {"id": "motivation", "title": "Motivation & Research Gap", "slides": []},
        {"id": "method", "title": "Method & System Design", "slides": []},
        {"id": "results", "title": "Experiments & Results", "slides": []},
        {"id": "takeaways", "title": "Summary & Takeaways", "slides": []},
    ]
    by_id = {bucket["id"]: bucket for bucket in buckets}
    for index, slide in enumerate(slides, start=1):
        section_id = _section_id_for_slide(slide, index)
        by_id[section_id]["slides"].append(slide)
        slide["deck_section"] = by_id[section_id]["title"]
    return [bucket for bucket in buckets if bucket["slides"]]


def _section_id_for_slide(slide: Dict[str, Any], index: int) -> str:
    role = str(slide.get("slide_role", "")).lower()
    title = f"{slide.get('title', '')} {slide.get('claim', '')}".lower()
    if role in {"title", "thesis"} or index <= 5:
        return "motivation"
    if role in {"table_interpretation", "metric"} or any(word in title for word in ("result", "benchmark", "evaluation", "safety")):
        return "results"
    if role == "conclusion" or any(word in title for word in ("contribution", "takeaway", "finding")):
        return "takeaways"
    return "method"


def _content_page_map(sections: List[Dict[str, Any]]) -> Dict[int, int]:
    page = 3
    result: Dict[int, int] = {}
    for section in sections:
        page += 1
        for slide in section.get("slides", []):
            result[id(slide)] = page
            page += 1
    return result


def _layout_family(slide_data: Dict[str, Any], sequence: int) -> str:
    proof_type = (slide_data.get("proof_object", {}) or {}).get("type", "")
    role = str(slide_data.get("slide_role", "")).lower()
    if proof_type == "table":
        return "table_bottom" if sequence % 2 else "table_left"
    if proof_type == "figure":
        return "visual_right" if sequence % 2 else "visual_left"
    if proof_type == "metric" or role == "metric":
        variants = ["metric_left", "metric_compact_band", "metric_left_alt"]
        return variants[sequence % len(variants)]
    if role == "thesis":
        if _text_proof_is_short(slide_data):
            variants = ["argument_mosaic", "argument_bottom_notes", "argument_cards"]
            return variants[sequence % len(variants)]
        variants = ["argument_strip", "argument_cards", "evidence_cards", "argument_mosaic"]
        return variants[sequence % len(variants)]
    if role == "conclusion":
        return "closing_summary"
    variants = ["evidence_cards", "argument_cards", "argument_strip", "argument_bottom_notes"]
    return variants[sequence % len(variants)]


def _add_academic_title_slide(prs: Any, inventory: Dict[str, Any], title: str, page: int, total_pages: int) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, _theme("paper"))
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = _theme("teal")
    accent.line.fill.background()
    soft_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.75), Inches(0), Inches(4.58), Inches(7.5))
    soft_block.fill.solid()
    soft_block.fill.fore_color.rgb = _theme("teal_light")
    soft_block.line.fill.background()
    gold_rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.08), Inches(1.35), Inches(0.045))
    gold_rule.fill.solid()
    gold_rule.fill.fore_color.rgb = _theme("gold")
    gold_rule.line.fill.background()
    _add_textbox(slide, "ACADEMIC PAPER READING", 0.72, 0.72, 4.8, 0.3, 9, _theme("teal_deep"), bold=True, letter_spaced=True)
    _add_textbox(slide, title, 0.72, 1.52, 7.55, 1.35, 34, _theme("ink"), bold=True)
    authors = _extract_authors(inventory.get("paper", {}).get("metadata_text", ""))
    _add_textbox(slide, authors or "Kimi Team", 0.78, 3.18, 7.5, 0.35, 13, _theme("muted_ink"))
    _add_textbox(
        slide,
        "A content-first paper reading deck built from existing parsed checkpoints.",
        0.78,
        3.72,
        7.2,
        0.45,
        15,
        _theme("teal_deep"),
        bold=True,
    )
    _add_textbox(
        slide,
        "Focus: motivation, system design, training recipe, and evaluation evidence.",
        0.78,
        4.28,
        6.85,
        0.55,
        13,
        _theme("muted_ink"),
    )
    _add_cover_highlight_rail(slide, inventory, 9.25, 1.22, 3.1, 5.12)
    _add_textbox(slide, "Generated from reusable paper-understanding checkpoints; no PDF reparse.", 0.78, 6.66, 8.8, 0.25, 8, _theme("soft_text"))
    _add_page_marker(slide, page, total_pages, "TITLE")


def _add_agenda_slide(prs: Any, sections: List[Dict[str, Any]], page: int, total_pages: int) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, _theme("paper_alt"))
    _add_page_marker(slide, page, total_pages, "AGENDA")
    _add_textbox(slide, "Roadmap", 0.75, 0.95, 4.0, 0.45, 20, _theme("teal_deep"), bold=True)
    _add_textbox(slide, "How this paper reading deck is organized", 0.75, 1.48, 8.25, 0.75, 30, _theme("ink"), bold=True)
    rail = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.05), Inches(1.35), Inches(3.15), Inches(4.9))
    rail.fill.solid()
    rail.fill.fore_color.rgb = _theme("panel")
    rail.line.color.rgb = _theme("line")
    _add_textbox(slide, "DECK MAP", 9.3, 1.7, 2.4, 0.25, 8, _theme("teal_deep"), bold=True, letter_spaced=True)
    _add_textbox(slide, f"{sum(len(section['slides']) for section in sections)}", 9.3, 2.15, 1.3, 0.55, 28, _theme("ink"), bold=True)
    _add_textbox(slide, "content slides", 10.45, 2.35, 1.45, 0.25, 9, _theme("soft_text"))
    _add_textbox(slide, f"{len(sections)}", 9.3, 3.0, 1.0, 0.45, 22, _theme("gold"), bold=True)
    _add_textbox(slide, "modules", 10.1, 3.17, 1.5, 0.22, 9, _theme("soft_text"))
    _add_textbox(slide, "Read path", 9.3, 4.05, 2.2, 0.25, 11, _theme("ink"), bold=True)
    _add_read_path_flow(slide, 9.25, 4.42, 2.65, 1.0)
    for idx, section in enumerate(sections, start=1):
        y = 2.3 + (idx - 1) * 0.9
        marker = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(y + 0.07), Inches(0.52), Inches(0.28))
        marker.fill.solid()
        marker.fill.fore_color.rgb = _theme("teal") if idx % 2 else _theme("gold")
        marker.line.fill.background()
        _add_textbox(slide, f"{idx:02d}", 0.83, y + 0.095, 0.42, 0.16, 7, _theme("white"), bold=True)
        _add_textbox(slide, section["title"], 1.55, y, 5.8, 0.33, 16, _theme("ink"), bold=True)
        _add_textbox(slide, f"{len(section['slides'])} slides", 7.25, y + 0.04, 1.0, 0.22, 9, _theme("teal_deep"), bold=True)
        sample = "; ".join(_limit_words(item.get("title", ""), 5).rstrip(".") for item in section["slides"][:2])
        _add_textbox(slide, sample, 1.55, y + 0.38, 6.85, 0.25, 9, _theme("soft_text"))
    _add_rule(slide, 0.75, 6.35, 11.45, _theme("line"))
    _add_textbox(slide, "Each module begins with a divider; tables, figures, and metric cards are treated as evidence rather than filler panels.", 0.75, 6.55, 11.2, 0.38, 10, _theme("muted_ink"))


def _add_read_path_flow(slide: Any, x: float, y: float, w: float, h: float) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    steps = [
        ("P", "Problem", _theme("teal")),
        ("M", "Method", _theme("gold")),
        ("E", "Evidence", _theme("clay")),
        ("T", "Takeaways", _theme("teal_deep")),
    ]
    node_d = 0.34
    available = max(0.1, w - node_d)
    gap = available / max(1, len(steps) - 1)
    for idx, (letter, label, color) in enumerate(steps):
        cx = x + idx * gap
        if idx:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + (idx - 1) * gap + node_d + 0.04),
                Inches(y + 0.16),
                Inches(max(0.06, gap - node_d - 0.08)),
                Inches(0.025),
            )
            line.fill.solid()
            line.fill.fore_color.rgb = _theme("line")
            line.line.fill.background()
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(y), Inches(node_d), Inches(node_d))
        node.fill.solid()
        node.fill.fore_color.rgb = color
        node.line.fill.background()
        _set_shape_text(node, letter, 8, _theme("white"), bold=True)
        _add_textbox(slide, label, cx - 0.12, y + 0.45, 0.65, 0.18, 5, _theme("muted_ink"))


def _add_section_divider_slide(prs: Any, section: Dict[str, Any], section_index: int, page: int, total_pages: int) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, _theme("paper"))
    colors = [_theme("teal"), _theme("gold"), _theme("clay"), _theme("teal_deep")]
    accent_color = colors[(section_index - 1) % len(colors)]
    left_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.36), Inches(7.5))
    left_band.fill.solid()
    left_band.fill.fore_color.rgb = accent_color
    left_band.line.fill.background()
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.36), Inches(5.72), Inches(12.94), Inches(1.78))
    band.fill.solid()
    band.fill.fore_color.rgb = _theme("panel")
    band.line.fill.background()
    _add_textbox(slide, f"SECTION {section_index:02d}", 0.8, 0.75, 3.6, 0.3, 9, accent_color, bold=True, letter_spaced=True)
    _add_textbox(slide, section["title"], 0.78, 1.55, 8.5, 1.1, 34, _theme("ink"), bold=True)
    first_titles = [item.get("title", "") for item in section.get("slides", [])[:3]]
    _add_textbox(slide, " / ".join(_limit_words(title, 7).rstrip(".") for title in first_titles), 0.84, 3.15, 10.5, 0.55, 14, _theme("muted_ink"))
    _add_textbox(slide, f"{len(section.get('slides', []))} content slides", 0.84, 6.2, 3.2, 0.3, 13, _theme("ink"), bold=True)
    _add_textbox(slide, "Module checkpoint", 9.1, 5.98, 2.2, 0.25, 8, accent_color, bold=True, letter_spaced=True)
    _add_textbox(slide, "The next pages keep one claim, one support note, and one proof object visible.", 9.1, 6.32, 3.1, 0.45, 10, _theme("muted_ink"))
    _add_page_marker(slide, page, total_pages, "SECTION")


def _add_content_slide(
    prs: Any,
    slide_data: Dict[str, Any],
    page: int,
    total_pages: int,
    layout: str,
    figure_paths: Dict[str, str],
    table_index: Dict[str, Dict[str, Any]],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, _theme("paper_alt"))
    _add_page_marker(slide, page, total_pages, slide_data.get("slide_role", "draft").upper())
    proof = slide_data.get("proof_object", {}) or {}
    proof_id = proof.get("id", "")

    if layout == "visual_left":
        _add_proof_object(slide, proof, figure_paths, table_index, 0.65, 1.18, 5.45, 5.15, compact=False)
        _add_claim_and_support(slide, slide_data, 6.55, 1.05, 5.9, 4.8, claim_size=25)
    elif layout == "visual_right":
        _add_claim_and_support(slide, slide_data, 0.65, 1.05, 5.9, 4.8, claim_size=25)
        _add_proof_object(slide, proof, figure_paths, table_index, 7.0, 1.18, 5.45, 5.15, compact=False)
    elif layout == "table_bottom":
        _add_claim_and_support(
            slide,
            slide_data,
            0.65,
            1.0,
            11.8,
            2.3,
            claim_size=22,
            support_offset=1.72,
            support_font_size=13,
        )
        _add_proof_object(slide, proof, figure_paths, table_index, 0.75, 3.5, 11.85, 2.95, compact=True)
    elif layout == "table_left":
        _add_proof_object(slide, proof, figure_paths, table_index, 0.65, 1.15, 6.35, 5.25, compact=True)
        _add_claim_and_support(slide, slide_data, 7.35, 1.05, 5.0, 4.85, claim_size=22)
    elif layout == "metric_grid":
        _add_claim_and_support(slide, slide_data, 0.65, 1.05, 6.0, 3.15, claim_size=25)
        _add_metric_proof_grid(slide, proof, slide_data, 7.05, 1.25, 5.25, 4.65)
    elif layout == "metric_left":
        _add_metric_proof_grid(slide, proof, slide_data, 0.75, 1.35, 4.8, 4.55)
        _add_claim_and_support(slide, slide_data, 6.1, 1.05, 6.0, 4.5, claim_size=25)
    elif layout == "metric_left_alt":
        _add_claim_and_support(slide, slide_data, 0.75, 1.05, 5.85, 4.4, claim_size=25)
        _add_metric_side_cluster(slide, proof, slide_data, 7.05, 1.35, 4.95, 4.35)
    elif layout == "metric_compact_band":
        _add_claim_and_support(slide, slide_data, 0.7, 1.05, 10.7, 3.0, claim_size=25)
        _add_compact_metric_band(slide, proof, slide_data, 0.85, 4.9, 11.3, 1.35)
    elif layout == "metric_strip":
        _add_claim_and_support(slide, slide_data, 0.7, 1.05, 11.2, 3.0, claim_size=25)
        _add_metric_strip(slide, proof, slide_data, 0.85, 5.05, 11.3, 1.05)
    elif layout == "argument_strip":
        _add_claim_and_support(slide, slide_data, 0.7, 1.05, 10.9, 3.1, claim_size=27)
        _add_bottom_evidence_strip(slide, proof, 0.75, 5.1, 11.7, 1.08)
    elif layout == "argument_split":
        _add_claim_and_support(slide, slide_data, 0.7, 1.05, 6.35, 4.2, claim_size=25)
        _add_evidence_card_stack(slide, proof, slide_data, 7.15, 1.35, 4.95, 3.95)
    elif layout == "argument_cards":
        _add_claim_and_support(slide, slide_data, 0.7, 1.05, 6.25, 4.35, claim_size=25)
        _add_evidence_card_stack(slide, proof, slide_data, 7.15, 1.3, 4.9, 4.15)
    elif layout == "argument_cards_left":
        _add_evidence_card_stack(slide, proof, slide_data, 0.75, 1.32, 4.55, 4.0)
        _add_claim_and_support(slide, slide_data, 5.85, 1.05, 6.25, 4.45, claim_size=25)
    elif layout == "argument_mosaic":
        _add_claim_and_support(
            slide,
            slide_data,
            0.75,
            1.02,
            11.3,
            2.35,
            claim_size=25,
            support_offset=1.72,
            support_font_size=13,
        )
        _add_evidence_mosaic(slide, proof, slide_data, 0.85, 4.02, 11.15, 1.75)
    elif layout == "argument_bottom_notes":
        _add_claim_and_support(
            slide,
            slide_data,
            0.75,
            1.02,
            10.9,
            2.65,
            claim_size=26,
            support_offset=1.92,
            support_font_size=13,
        )
        _add_evidence_bottom_cards(slide, proof, slide_data, 0.85, 4.55, 11.2, 1.48)
    elif layout == "evidence_cards":
        _add_claim_and_support(slide, slide_data, 0.72, 1.05, 7.0, 3.8, claim_size=25)
        _add_evidence_card_stack(slide, proof, slide_data, 8.15, 1.25, 3.95, 4.15)
    elif layout == "closing_summary":
        _add_claim_and_support(slide, slide_data, 0.85, 1.35, 7.8, 3.2, claim_size=27)
        _add_metric_proof_grid(slide, proof, slide_data, 9.05, 1.55, 3.2, 3.5)
    else:
        _add_claim_and_support(slide, slide_data, 0.75, 1.05, 6.1, 4.4, claim_size=24)
        _add_bottom_evidence_strip(slide, proof, 7.2, 1.45, 4.9, 3.8)

    _add_textbox(slide, _evidence_footer(slide_data), 0.72, 6.87, 11.7, 0.22, 7, _theme("soft_text"))


def _add_closing_slide(prs: Any, inventory: Dict[str, Any], page: int, total_pages: int) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, _theme("paper"))
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.25), Inches(0), Inches(5.08), Inches(7.5))
    band.fill.solid()
    band.fill.fore_color.rgb = _theme("teal_deep")
    band.line.fill.background()
    gold = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.25), Inches(0), Inches(0.12), Inches(7.5))
    gold.fill.solid()
    gold.fill.fore_color.rgb = _theme("gold")
    gold.line.fill.background()
    _add_textbox(slide, "SUMMARY", 0.82, 0.85, 3.0, 0.3, 9, _theme("teal_deep"), bold=True, letter_spaced=True)
    _add_textbox(slide, "Thanks for watching", 0.8, 1.62, 6.5, 0.8, 34, _theme("ink"), bold=True)
    _add_textbox(slide, "Kimi K2 shows how open-weight agentic intelligence depends on architecture, training stability, data synthesis, RL, and broad evaluation evidence.", 0.84, 2.85, 6.7, 1.2, 18, _theme("muted_ink"))
    _add_textbox(slide, "Content lineage", 8.75, 4.92, 3.5, 0.25, 8, _theme("gold_light"), bold=True, letter_spaced=True)
    _add_textbox(slide, inventory.get("paper", {}).get("title", ""), 8.75, 5.32, 3.6, 0.75, 16, _theme("white"), bold=True)
    _add_textbox(slide, "Parsed checkpoints -> authored slide narrative -> visual QA manifest", 8.75, 6.22, 3.55, 0.55, 10, _theme("teal_light"))
    _add_page_marker(slide, page, total_pages, "CLOSING")


def _table_index(inventory: Dict[str, Any]) -> Dict[str, Dict[str, Any]]:
    return {
        str(item.get("id", "")): item
        for item in inventory.get("assets", {}).get("tables", []) or []
        if isinstance(item, dict)
    }


def _add_page_marker(slide: Any, page: int, total_pages: int, label: str) -> None:
    _add_textbox(slide, f"{page:02d} / {total_pages:02d}  {label}", 0.6, 0.35, 4.9, 0.34, 8, _theme("soft_text"), bold=True)
    _add_rule(slide, 0.6, 0.72, 12.1, _theme("line"))


def _add_claim_and_support(
    slide: Any,
    slide_data: Dict[str, Any],
    x: float,
    y: float,
    w: float,
    h: float,
    claim_size: int = 24,
    support_offset: float = 2.3,
    support_font_size: int = 13,
) -> None:
    _add_textbox(slide, slide_data.get("title", ""), x, y, min(w, 5.4), 0.34, 14, _theme("teal_deep"), bold=True)
    _add_textbox(slide, slide_data.get("claim", ""), x, y + 0.55, w, min(1.45, h * 0.38), claim_size, _theme("ink"), bold=True)
    support_height = max(0.72, h - support_offset - 0.15)
    _add_textbox(slide, slide_data.get("support", ""), x + 0.02, y + support_offset, w * 0.92, support_height, support_font_size, _theme("muted_ink"))


def _add_proof_object(
    slide: Any,
    proof: Dict[str, Any],
    figure_paths: Dict[str, str],
    table_index: Dict[str, Dict[str, Any]],
    x: float,
    y: float,
    w: float,
    h: float,
    compact: bool = False,
) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    proof_type = str(proof.get("type", "text_evidence"))
    proof_id = str(proof.get("id", ""))
    proof_focus = str(proof.get("focus", ""))
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = _theme("panel_light")
    panel.line.color.rgb = _theme("line")
    _add_textbox(slide, proof_type.upper(), x + 0.22, y + 0.18, w - 0.45, 0.24, 8, _theme("teal_deep"), bold=True, letter_spaced=True)
    _add_textbox(slide, proof_id or "source evidence", x + 0.22, y + 0.55, w - 0.45, 0.38, 16 if compact else 18, _theme("ink"), bold=True)

    if proof_type == "figure":
        _add_figure_proof(slide, figure_paths.get(proof_id, ""), proof_focus, x + 0.25, y + 1.05, w - 0.5, h - 1.35)
    elif proof_type == "table":
        table = table_index.get(proof_id, {})
        rows = table.get("rows", []) or []
        caption = table.get("caption", "") or proof_focus
        _add_table_proof(slide, rows, caption, x + 0.25, y + 1.02, w - 0.5, h - 1.25)
    elif proof_type == "metric":
        _add_metric_card(slide, proof_focus or proof_id, proof_id if proof_focus else "Key metric", x + 0.35, y + 1.2, w - 0.7, min(1.55, h - 1.5))
        if proof_focus and proof_id:
            _add_textbox(slide, f"{proof_id}: {proof_focus}", x + 0.35, y + 2.92, w - 0.7, 0.5, 12, _theme("muted_ink"))
    else:
        _add_textbox(slide, proof_focus or "Evidence preserved from parsed checkpoints.", x + 0.35, y + 1.25, w - 0.7, h - 1.55, 12, _theme("muted_ink"))


def _add_figure_proof(slide: Any, figure_path: str, caption: str, x: float, y: float, w: float, h: float) -> None:
    from pptx.util import Inches

    caption_h = 0.48 if caption else 0
    image_h = max(0.6, h - caption_h)
    if figure_path and Path(figure_path).exists():
        try:
            slide.shapes.add_picture(figure_path, Inches(x), Inches(y), width=Inches(w), height=Inches(image_h))
        except Exception:
            _add_textbox(slide, "Figure image could not be embedded.", x, y + 0.2, w, 0.5, 11, _theme("clay"))
    else:
        _add_textbox(slide, "Figure image missing; check extracted asset path.", x, y + 0.2, w, 0.5, 11, _theme("clay"))
    if caption:
        _add_textbox(slide, caption, x, y + image_h + 0.08, w, caption_h, 9, _theme("soft_text"))


def _add_table_proof(slide: Any, rows: List[List[str]], caption: str, x: float, y: float, w: float, h: float) -> None:
    from pptx.util import Inches, Pt

    if not rows:
        _add_textbox(slide, caption or "Table rows missing in content inventory.", x, y, w, h, 12, _theme("clay"))
        return

    max_cols = min(max((len(row) for row in rows), default=1), 6)
    max_rows = max(3, min(len(rows), int(h / 0.3)))
    visible_rows = [row[:max_cols] + [""] * (max_cols - len(row[:max_cols])) for row in rows[:max_rows]]
    table_shape = slide.shapes.add_table(len(visible_rows), max_cols, Inches(x), Inches(y), Inches(w), Inches(max(0.8, h - 0.42)))
    table = table_shape.table
    for col in range(max_cols):
        table.columns[col].width = Inches(w / max_cols)
    for row_idx, row in enumerate(visible_rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = _limit_chars(value, 54)
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(6.5 if max_cols >= 5 else 7.5)
            paragraph.font.bold = row_idx == 0
            paragraph.font.color.rgb = _theme("ink")
            cell.fill.solid()
            if row_idx == 0:
                cell.fill.fore_color.rgb = _theme("teal_light")
            elif row_idx % 2:
                cell.fill.fore_color.rgb = _theme("white")
            else:
                cell.fill.fore_color.rgb = _theme("paper")
    note = caption or ""
    if len(rows) > max_rows:
        note = f"{note} Showing {max_rows}/{len(rows)} parsed rows.".strip()
    _add_textbox(slide, note, x, y + h - 0.33, w, 0.28, 8, _theme("soft_text"))


def _add_metric_proof_grid(slide: Any, proof: Dict[str, Any], slide_data: Dict[str, Any], x: float, y: float, w: float, h: float) -> None:
    value = proof.get("focus", "") or _first_metric_value(slide_data.get("support", ""))
    label = proof.get("id", "") or "Key metric"
    _add_textbox(slide, "KEY EVIDENCE", x, y, w, 0.25, 8, _theme("teal_deep"), bold=True, letter_spaced=True)
    _add_metric_card(slide, value or "n/a", label, x, y + 0.45, w, 1.35)
    metrics = _extract_metric_pairs(slide_data.get("support", ""))
    for idx, (metric_value, metric_label) in enumerate(metrics[:3]):
        _add_metric_card(slide, metric_value, metric_label, x + (idx % 2) * (w / 2 + 0.1), y + 2.05 + (idx // 2) * 1.25, w / 2 - 0.12, 0.9, small=True)


def _add_metric_side_cluster(slide: Any, proof: Dict[str, Any], slide_data: Dict[str, Any], x: float, y: float, w: float, h: float) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = _theme("panel")
    panel.line.color.rgb = _theme("line")
    _add_textbox(slide, "KEY EVIDENCE", x + 0.25, y + 0.25, w - 0.5, 0.24, 8, _theme("teal_deep"), bold=True, letter_spaced=True)
    value = proof.get("focus", "") or _first_metric_value(slide_data.get("support", "")) or "n/a"
    label = proof.get("id", "") or "Key metric"
    _add_metric_card(slide, value, label, x + 0.3, y + 0.75, w - 0.6, 1.15)
    metrics = _extract_metric_pairs(slide_data.get("support", ""))[:3]
    for idx, (metric_value, metric_label) in enumerate(metrics):
        _add_metric_card(slide, metric_value, metric_label, x + 0.3, y + 2.1 + idx * 0.72, w - 0.6, 0.55, small=True)


def _add_compact_metric_band(slide: Any, proof: Dict[str, Any], slide_data: Dict[str, Any], x: float, y: float, w: float, h: float) -> None:
    metrics = []
    if proof.get("focus") or proof.get("id"):
        metrics.append((str(proof.get("focus", "") or "n/a"), str(proof.get("id", "") or "Key metric")))
    metrics.extend(_extract_metric_pairs(slide_data.get("support", ""))[:3])
    if not metrics:
        metrics.append(("n/a", "Metric"))
    slot_count = min(4, len(metrics))
    slot_w = w / slot_count
    for idx, (value, label) in enumerate(metrics[:slot_count]):
        is_primary = idx == 0
        _add_metric_card(slide, value, label, x + idx * slot_w, y, slot_w - 0.12, h, small=not is_primary)


def _add_metric_strip(slide: Any, proof: Dict[str, Any], slide_data: Dict[str, Any], x: float, y: float, w: float, h: float) -> None:
    metrics = []
    if proof.get("focus") or proof.get("id"):
        metrics.append((str(proof.get("focus", "") or "n/a"), str(proof.get("id", "") or "Key metric")))
    metrics.extend(_extract_metric_pairs(slide_data.get("support", ""))[:3])
    if not metrics:
        metrics.append(("n/a", "Metric"))
    slot_w = w / min(4, len(metrics))
    for idx, (value, label) in enumerate(metrics[:4]):
        _add_metric_card(slide, value, label, x + idx * slot_w, y, slot_w - 0.12, h, small=idx > 0)


def _add_metric_card(slide: Any, value: str, label: str, x: float, y: float, w: float, h: float, small: bool = False) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = _theme("teal_deep") if not small else _theme("panel_light")
    card.line.color.rgb = _theme("line") if small else _theme("teal_deep")
    if small:
        _add_textbox(slide, str(value or "n/a"), x + 0.18, y + 0.1, w - 0.36, 0.22, 14, _theme("ink"), bold=True)
        _add_textbox(slide, str(label or "Metric"), x + 0.18, y + max(0.36, h * 0.64), w - 0.36, 0.16, 9, _theme("muted_ink"))
    else:
        _add_textbox(slide, str(value or "n/a"), x + 0.18, y + 0.17, w - 0.36, max(0.22, h * 0.38), 22, _theme("white"), bold=True)
        _add_textbox(slide, str(label or "Metric"), x + 0.18, y + h * 0.6, w - 0.36, max(0.18, h * 0.25), 10, _theme("teal_light"))


def _add_metric_chip(slide: Any, x: float, y: float, value: str, label: str) -> None:
    _add_metric_card(slide, value, label, x, y, 1.35, 0.75, small=True)


def _add_bottom_evidence_strip(slide: Any, proof: Dict[str, Any], x: float, y: float, w: float, h: float) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    strip.fill.solid()
    strip.fill.fore_color.rgb = _theme("panel")
    strip.line.color.rgb = _theme("line")
    _add_textbox(slide, str(proof.get("type", "evidence")).upper(), x + 0.22, y + 0.18, 1.8, 0.24, 8, _theme("teal_deep"), bold=True, letter_spaced=True)
    _add_textbox(slide, str(proof.get("id", "source evidence")), x + 2.1, y + 0.15, 2.4, 0.3, 13, _theme("ink"), bold=True)
    _add_textbox(slide, str(proof.get("focus", "Evidence preserved from parsed checkpoints.")), x + 0.22, y + 0.52, w - 0.45, h - 0.6, 11, _theme("muted_ink"))


def _add_evidence_note_card(
    slide: Any,
    item: Dict[str, str],
    x: float,
    y: float,
    w: float,
    h: float,
    index: int,
    compact: bool = False,
) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = [_theme("teal_light"), _theme("gold_light"), _theme("clay_light")][index % 3]
    card.line.color.rgb = _theme("line")
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = [_theme("teal"), _theme("gold"), _theme("clay")][index % 3]
    stripe.line.fill.background()
    label_size = 9 if compact else 10
    body_size = 9 if compact else 10
    _add_textbox(slide, item["label"], x + 0.18, y + 0.13, w - 0.36, 0.22, label_size, _theme("ink"), bold=True)
    _add_textbox(slide, item["body"], x + 0.18, y + 0.43, w - 0.36, max(0.28, h - 0.5), body_size, _theme("muted_ink"))


def _add_evidence_mosaic(
    slide: Any,
    proof: Dict[str, Any],
    slide_data: Dict[str, Any],
    x: float,
    y: float,
    w: float,
    h: float,
) -> None:
    items = _evidence_card_items(proof, slide_data)
    _add_textbox(slide, "EVIDENCE MOSAIC", x, y - 0.22, w, 0.2, 8, _theme("teal_deep"), bold=True, letter_spaced=True)
    gap = 0.16
    top_h = min(0.86, h * 0.52)
    _add_evidence_note_card(slide, items[0], x, y, w, top_h, 0)
    side_items = (items[1:] or items[:1])[:2]
    bottom_y = y + top_h + gap
    bottom_h = max(0.5, h - top_h - gap)
    card_w = (w - gap) / 2
    for idx, item in enumerate(side_items, start=1):
        _add_evidence_note_card(slide, item, x + (idx - 1) * (card_w + gap), bottom_y, card_w, bottom_h, idx, compact=True)


def _add_evidence_bottom_cards(
    slide: Any,
    proof: Dict[str, Any],
    slide_data: Dict[str, Any],
    x: float,
    y: float,
    w: float,
    h: float,
) -> None:
    items = _evidence_card_items(proof, slide_data)
    _add_textbox(slide, "EVIDENCE NOTES", x, y - 0.22, w, 0.2, 8, _theme("teal_deep"), bold=True, letter_spaced=True)
    count = max(1, min(3, len(items)))
    gap = 0.16
    card_w = (w - gap * (count - 1)) / count
    for idx, item in enumerate(items[:count]):
        _add_evidence_note_card(slide, item, x + idx * (card_w + gap), y, card_w, h, idx, compact=True)


def _add_evidence_card_stack(
    slide: Any,
    proof: Dict[str, Any],
    slide_data: Dict[str, Any],
    x: float,
    y: float,
    w: float,
    h: float,
) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    _add_textbox(slide, "EVIDENCE NOTES", x, y - 0.05, w, 0.25, 8, _theme("teal_deep"), bold=True, letter_spaced=True)
    items = _evidence_card_items(proof, slide_data)
    card_h = min(1.05, (h - 0.38) / max(1, len(items)) - 0.08)
    for idx, item in enumerate(items):
        cy = y + 0.35 + idx * (card_h + 0.17)
        _add_evidence_note_card(slide, item, x, cy, w, card_h, idx)


def _add_cover_highlight_rail(slide: Any, inventory: Dict[str, Any], x: float, y: float, w: float, h: float) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = _theme("white")
    panel.line.color.rgb = _theme("line")
    _add_textbox(slide, "PAPER HIGHLIGHTS", x + 0.28, y + 0.32, w - 0.56, 0.25, 8, _theme("teal_deep"), bold=True, letter_spaced=True)
    highlights = inventory.get("paper_highlights", []) or []
    for idx, item in enumerate(highlights[:3]):
        yy = y + 0.88 + idx * 1.22
        color = [_theme("teal"), _theme("gold"), _theme("clay")][idx % 3]
        marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.28), Inches(yy + 0.05), Inches(0.08), Inches(0.52))
        marker.fill.solid()
        marker.fill.fore_color.rgb = color
        marker.line.fill.background()
        _add_textbox(slide, item.get("label", f"Highlight {idx + 1}"), x + 0.48, yy, w - 0.76, 0.22, 9, _theme("ink"), bold=True)
        _add_textbox(slide, item.get("body", ""), x + 0.48, yy + 0.28, w - 0.76, 0.54, 9, _theme("muted_ink"))


def _extract_authors(metadata_text: str) -> str:
    match = re.search(r"Authors?\s*:\s*(.+)$", _clean_text(metadata_text), flags=re.IGNORECASE)
    if not match:
        return ""
    return _limit_words(match.group(1), 18)


def _first_metric_value(text: str) -> str:
    match = re.search(r"(?<![A-Za-z0-9])(?:\d+(?:\.\d+)?%|\d+(?:\.\d+)?[TBMK]?|\d+/\d+)", text or "")
    return match.group(0) if match else ""


def _extract_metric_pairs(text: str) -> List[tuple[str, str]]:
    pairs = []
    for match in re.finditer(r"([A-Za-z][A-Za-z0-9 +&/-]{2,24})\s+(\d+(?:\.\d+)?%?|\d+/\d+)", text or ""):
        label = _clean_text(match.group(1)).strip(" ,.;:")
        value = match.group(2)
        if label and value:
            pairs.append((value, _limit_words(label, 4)))
    return pairs


def _text_proof_is_short(slide_data: Dict[str, Any]) -> bool:
    proof = slide_data.get("proof_object", {}) or {}
    if proof.get("type") != "text_evidence":
        return False
    combined = " ".join([str(proof.get("focus", "")), str(slide_data.get("support", ""))])
    return len(_clean_text(combined)) < 260


def _evidence_card_items(proof: Dict[str, Any], slide_data: Dict[str, Any]) -> List[Dict[str, str]]:
    proof_id = _clean_text(proof.get("id", "")) or "Parsed checkpoint"
    proof_focus = _clean_text(proof.get("focus", ""))
    support = _clean_text(slide_data.get("support", ""))
    items: List[Dict[str, str]] = []
    if proof_focus:
        items.append({"label": proof_id, "body": _limit_words(proof_focus, 16)})
    sentences = [
        part.strip()
        for part in re.split(r"(?<=[.!?])\s+", support)
        if len(part.strip()) > 12
    ]
    for idx, sentence in enumerate(sentences[:3], start=1):
        items.append({"label": f"Reading note {idx}", "body": _limit_words(sentence, 13)})
    if not items:
        items.append({"label": proof_id, "body": "Evidence preserved from parsed checkpoints."})
    return items[:3]


def _count_items(items: Iterable[str]) -> Dict[str, int]:
    counts: Dict[str, int] = {}
    for item in items:
        counts[item] = counts.get(item, 0) + 1
    return counts


def _dominant_layout_family(layouts: List[str]) -> Dict[str, Any]:
    if not layouts:
        return {}
    counts = _count_items(layouts)
    family, count = max(counts.items(), key=lambda item: item[1])
    return {"family": family, "count": count, "ratio": round(count / len(layouts), 4)}


def _max_consecutive_layout(layouts: List[str]) -> Dict[str, Any]:
    best_family = ""
    best_count = 0
    current_family = ""
    current_count = 0
    for family in layouts:
        if family == current_family:
            current_count += 1
        else:
            current_family = family
            current_count = 1
        if current_count > best_count:
            best_family = current_family
            best_count = current_count
    return {"family": best_family, "count": best_count}


def _dedupe_review_targets(targets: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    result = []
    seen = set()
    for item in targets:
        page = item.get("page")
        reason = item.get("reason", "")
        key = (page, reason)
        if key in seen or not page:
            continue
        seen.add(key)
        result.append(item)
    return result[:18]


def _summary_items(content: Dict[str, Any]) -> List[Dict[str, Any]]:
    items: List[Dict[str, Any]] = []
    for section_name, category in SUMMARY_SECTIONS:
        text = _clean_text(content.get(section_name, ""))
        if not text:
            continue
        chunks = _split_summary_section(text)
        for chunk_index, chunk in enumerate(chunks, start=1):
            items.append(
                {
                    "id": f"{category}_{chunk_index:02d}",
                    "category": category,
                    "title": _infer_item_title(chunk, category),
                    "text": chunk,
                    "source": {
                        "checkpoint": "summary",
                        "section": section_name,
                    },
                }
            )
    return items


def _plan_slides(plan: Dict[str, Any]) -> List[Dict[str, Any]]:
    slides = []
    for index, item in enumerate(plan.get("sections", []) or [], start=1):
        tables = [ref for ref in item.get("tables", []) if isinstance(ref, dict)]
        figures = [ref for ref in item.get("figures", []) if isinstance(ref, dict)]
        slide_id = item.get("id") or f"slide_{index:02d}"
        slides.append(
            {
                "slide_id": slide_id,
                "source_id": slide_id,
                "title": _clean_text(item.get("title", "")),
                "section": item.get("section", "") or item.get("chapter", ""),
                "section_type": item.get("type", "content"),
                "content": _clean_text(item.get("content", "")),
                "figures": figures,
                "tables": tables,
                "source": {
                    "checkpoint": "plan",
                    "index": index,
                },
            }
        )
    return slides


def _curated_slides(spec: Dict[str, Any]) -> List[Dict[str, Any]]:
    slides = []
    for index, item in enumerate(spec.get("slides", []) or [], start=1):
        points = []
        for block in item.get("text_blocks", []) or []:
            if not isinstance(block, dict):
                continue
            points.append(
                {
                    "text": _clean_text(block.get("text", "")),
                    "claim": _clean_text(block.get("claim", "")),
                    "detail": _clean_text(block.get("detail", "")),
                    "evidence": _clean_text(block.get("evidence", "")),
                }
            )
        slides.append(
            {
                "slide_id": item.get("slide_id") or f"slide_{index:02d}",
                "title": _clean_text(item.get("title", "")),
                "layout": item.get("layout", ""),
                "takeaway": _clean_text(item.get("takeaway", "")),
                "section_type": item.get("section_type", "content"),
                "section_label": item.get("section_label", ""),
                "points": points,
                "figures": item.get("image_blocks", []) or [],
                "tables": item.get("table_blocks", []) or [],
                "metrics": item.get("metric_blocks", []) or [],
                "notes": item.get("notes", []) or [],
                "source": {
                    "checkpoint": "slide_spec",
                    "index": index,
                },
            }
        )
    return slides


def _assets(origin: Dict[str, Any]) -> Dict[str, Any]:
    tables = []
    for item in origin.get("tables", []) or []:
        if not isinstance(item, dict):
            continue
        html_content = item.get("html", "")
        rows = _html_table_to_rows(html_content)
        tables.append(
            {
                "id": item.get("id", ""),
                "caption": _clean_text(item.get("caption", "")),
                "rows": rows,
                "row_count": len(rows),
                "column_count": max((len(row) for row in rows), default=0),
                "html_preview": _clean_text(html_content)[:1200],
                "asset_type": "table",
            }
        )
    return {
        "figures": [
            {
                "id": item.get("id", ""),
                "caption": _clean_text(item.get("caption", "")),
                "path": item.get("path", ""),
                "asset_type": "figure",
            }
            for item in origin.get("figures", []) or []
            if isinstance(item, dict)
        ],
        "tables": tables,
        "base_path": origin.get("base_path", ""),
    }


def _metrics_from_curated_slides(slides: Iterable[Dict[str, Any]]) -> List[Dict[str, Any]]:
    metrics = []
    seen = set()
    for slide in slides:
        for metric in slide.get("metrics", []) or []:
            if not isinstance(metric, dict):
                continue
            value = _clean_text(metric.get("value", ""))
            label = _clean_text(metric.get("label", ""))
            if not value:
                continue
            key = (label.lower(), value.lower())
            if key in seen:
                continue
            seen.add(key)
            metrics.append(
                {
                    "label": label,
                    "value": value,
                    "note": _clean_text(metric.get("note", "")),
                    "source_slide": slide.get("slide_id", ""),
                }
            )
    return metrics


def _paper_highlights(
    summary_items: List[Dict[str, Any]],
    plan_slides: List[Dict[str, Any]],
    curated_slides: List[Dict[str, Any]],
    metrics: List[Dict[str, Any]],
) -> List[Dict[str, str]]:
    """Derive cover-ready paper highlights without reparsing the PDF."""
    text_pool = _highlight_text_pool(summary_items, plan_slides, curated_slides)
    highlights: List[Dict[str, str]] = []

    result_text = _curated_takeaway_highlight(
        curated_slides,
        ("sota", "state-of-the-art", "surpass", "lead", "best", "competitive"),
    ) or _best_highlight_sentence(
        text_pool,
        ("sota", "state-of-the-art", "surpass", "lead", "best", "competitive", "benchmark", "agentic"),
    )
    if result_text:
        highlights.append({"label": "Core result", "body": _limit_words(result_text, 15)})

    scale_text = _scale_highlight(text_pool, metrics)
    if scale_text:
        highlights.append({"label": "Scale", "body": scale_text})

    method_text = _best_highlight_sentence(
        text_pool,
        ("qk-clip", "qk clip", "muon", "optimizer", "moe", "sparsity", "experts", "rl", "self-critique", "tool"),
    )
    if method_text:
        highlights.append({"label": "Design edge", "body": _limit_words(method_text, 15)})

    contribution_text = _best_highlight_sentence(
        text_pool,
        ("contribution", "pipeline", "data synthesis", "training recipe", "safety", "evaluation", "open"),
    )
    if contribution_text:
        highlights.append({"label": "Evidence scope", "body": _limit_words(contribution_text, 15)})

    for fallback in _fallback_highlight_sentences(text_pool):
        highlights.append({"label": "Takeaway", "body": _limit_words(fallback, 15)})
        if len(highlights) >= 4:
            break

    return _dedupe_highlights(highlights)[:3]


def _curated_takeaway_highlight(curated_slides: List[Dict[str, Any]], keywords: Iterable[str]) -> str:
    keyword_list = [keyword.lower() for keyword in keywords]
    for slide in curated_slides:
        text = _clean_text(slide.get("takeaway", ""))
        lowered = text.lower()
        if 20 <= len(text) <= 180 and any(keyword in lowered for keyword in keyword_list):
            return text
    return ""


def _highlight_text_pool(
    summary_items: List[Dict[str, Any]],
    plan_slides: List[Dict[str, Any]],
    curated_slides: List[Dict[str, Any]],
) -> List[Dict[str, str]]:
    pool: List[Dict[str, str]] = []
    for slide in curated_slides:
        for field in ("takeaway", "title"):
            text = _clean_text(slide.get(field, ""))
            if text:
                pool.append({"category": slide.get("section_label", "") or "slide", "source": "curated", "text": text})
        for point in slide.get("points", []) or []:
            if not isinstance(point, dict):
                continue
            text = _first_non_empty([point.get("claim", ""), point.get("detail", ""), point.get("text", "")])
            if text:
                pool.append({"category": slide.get("section_label", "") or "slide_point", "source": "curated", "text": text})
    for slide in plan_slides:
        text = _first_non_empty([slide.get("content", ""), slide.get("title", "")])
        if text:
            pool.append({"category": slide.get("section", "") or "plan", "source": "plan", "text": text})
    for item in summary_items:
        text = _clean_text(item.get("text", ""))
        if text:
            pool.append({"category": item.get("category", ""), "source": "summary", "text": text})
    return pool


def _best_highlight_sentence(pool: List[Dict[str, str]], keywords: Iterable[str]) -> str:
    keyword_list = [keyword.lower() for keyword in keywords]
    best = ("", -1)
    high_value = {"sota", "state-of-the-art", "surpass", "lead", "best", "qk-clip", "qk clip", "muon", "rl", "self-critique"}
    for entry in pool:
        for sentence in _highlight_sentences(entry.get("text", "")):
            if len(sentence) > 240:
                continue
            lowered = sentence.lower()
            score = sum(5 if keyword in high_value else 3 for keyword in keyword_list if keyword in lowered)
            if entry.get("source") == "curated":
                score += 3
            if entry.get("category", "").lower() in {"results", "contribution", "method"}:
                score += 1
            if 45 <= len(sentence) <= 180:
                score += 1
            if score > best[1]:
                best = (sentence, score)
    return best[0] if best[1] > 0 else ""


def _scale_highlight(pool: List[Dict[str, str]], metrics: List[Dict[str, Any]]) -> str:
    combined = " ".join(entry.get("text", "") for entry in pool)
    total_param = _first_regex_value(combined, r"\d+(?:\.\d+)?\s*T")
    active_param = ""
    for match in re.finditer(r"\d+(?:\.\d+)?\s*B", combined, flags=re.IGNORECASE):
        window = combined[max(0, match.start() - 48) : match.end() + 60].lower()
        if "activat" in window:
            active_param = _clean_text(match.group(0))
            break
    if not active_param:
        active_param = _first_regex_value(combined, r"\d+(?:\.\d+)?\s*B")

    metric_bits = []
    for metric in metrics:
        label = _clean_text(metric.get("label", ""))
        value = _clean_text(metric.get("value", ""))
        if not label or not value:
            continue
        if any(word in label.lower() for word in ("parameter", "expert", "token")):
            metric_bits.append(f"{value} {label}")
        if len(metric_bits) >= 2:
            break

    bits = []
    if total_param:
        bits.append(f"{total_param} total parameters")
    if active_param:
        bits.append(f"{active_param} activated per token")
    bits.extend(metric_bits)
    return "; ".join(bits[:2]) + "." if bits else ""


def _fallback_highlight_sentences(pool: List[Dict[str, str]]) -> List[str]:
    result = []
    for entry in pool:
        for sentence in _highlight_sentences(entry.get("text", "")):
            if len(sentence) >= 35:
                result.append(sentence)
                break
    return result


def _highlight_sentences(text: str) -> List[str]:
    clean = re.sub(r"\s+", " ", _clean_text(text))
    if not clean:
        return []
    parts = re.split(r"(?<=[.!?])\s+|\s+-\s+|\s*\n+\s*", clean)
    return [_clean_text(part).strip(" -") for part in parts if len(_clean_text(part)) > 20]


def _dedupe_highlights(highlights: List[Dict[str, str]]) -> List[Dict[str, str]]:
    result: List[Dict[str, str]] = []
    seen = set()
    for item in highlights:
        body = _normalize_highlight_text(item.get("body", ""))
        if not body:
            continue
        key = re.sub(r"[^a-z0-9]+", " ", body.lower()).strip()[:80]
        if key in seen:
            continue
        seen.add(key)
        result.append({"label": _clean_text(item.get("label", "Highlight")), "body": body})
    return result


def _normalize_highlight_text(text: str) -> str:
    clean = _clean_text(text)
    clean = clean.replace("\u00a6\u00d3", "tau")
    clean = clean.replace("\u03c4", "tau")
    clean = clean.replace("tau=", "tau = ")
    clean = re.sub(r"\s+", " ", clean)
    return clean.strip()


def _first_regex_value(text: str, pattern: str) -> str:
    match = re.search(pattern, text or "", flags=re.IGNORECASE)
    return _clean_text(match.group(0)) if match else ""


def _coverage(
    summary_items: List[Dict[str, Any]],
    plan_slides: List[Dict[str, Any]],
    curated_slides: List[Dict[str, Any]],
    assets: Dict[str, Any],
    metrics: List[Dict[str, Any]],
) -> Dict[str, Any]:
    categories = sorted({item["category"] for item in summary_items})
    return {
        "summary_item_count": len(summary_items),
        "summary_categories": categories,
        "plan_slide_count": len(plan_slides),
        "curated_slide_count": len(curated_slides),
        "figure_count": len(assets.get("figures", [])),
        "table_count": len(assets.get("tables", [])),
        "metric_count": len(metrics),
        "has_core_sections": {
            "motivation": "motivation" in categories,
            "method": "method" in categories,
            "results": "results" in categories,
            "contribution": "contribution" in categories,
        },
    }


def _proof_object(plan_slide: Dict[str, Any], curated_slide: Dict[str, Any]) -> Dict[str, Any]:
    slide_text = " ".join(
        [
            plan_slide.get("title", ""),
            plan_slide.get("section", ""),
            curated_slide.get("section_label", ""),
        ]
    ).lower()
    prefers_text_evidence = any(
        word in slide_text
        for word in ("limitation", "problem", "challenge", "motivation", "research gap")
    )
    if plan_slide.get("figures"):
        first = plan_slide["figures"][0]
        return {
            "type": "figure",
            "id": first.get("figure_id", ""),
            "focus": first.get("focus", ""),
        }
    if curated_slide.get("figures"):
        first = curated_slide["figures"][0]
        return {
            "type": "figure",
            "id": first.get("title", "") or first.get("path", ""),
            "focus": first.get("caption", "") or first.get("placeholder_text", ""),
        }
    if plan_slide.get("tables"):
        first = plan_slide["tables"][0]
        return {
            "type": "table",
            "id": first.get("table_id", ""),
            "focus": first.get("focus", "") or first.get("extract", ""),
        }
    if curated_slide.get("tables"):
        first = curated_slide["tables"][0]
        return {
            "type": "table",
            "id": first.get("title", ""),
            "focus": first.get("caption", ""),
        }
    if curated_slide.get("metrics") and not prefers_text_evidence:
        first = curated_slide["metrics"][0]
        return {
            "type": "metric",
            "id": _clean_text(first.get("label", "")),
            "focus": _clean_text(first.get("value", "")),
        }
    return {
        "type": "text_evidence",
        "id": plan_slide.get("section", "") or curated_slide.get("section_label", "") or "source_text",
        "focus": _limit_words(plan_slide.get("content", "") or curated_slide.get("takeaway", ""), 24),
    }


def _slide_role(plan_slide: Dict[str, Any], curated_slide: Dict[str, Any], index: int, total: int) -> str:
    section_type = (plan_slide.get("section_type") or curated_slide.get("section_type") or "").lower()
    title = f"{plan_slide.get('title', '')} {curated_slide.get('title', '')}".lower()
    proof_type = _proof_object(plan_slide, curated_slide)["type"]
    if section_type == "opening" or index == 1:
        return "title"
    if section_type == "ending" or index == total:
        return "conclusion"
    if any(word in title for word in ("problem", "motivation", "challenge", "limitation", "research gap")):
        return "thesis"
    if proof_type == "figure":
        return "figure_explainer"
    if proof_type == "table":
        return "table_interpretation"
    if proof_type == "metric" or any(word in title for word in ("result", "benchmark", "evaluation")):
        return "metric"
    if any(word in title for word in ("architecture", "pipeline", "training", "method", "system")):
        return "mechanism"
    return "evidence"


def _source_evidence(plan_slide: Dict[str, Any], curated_slide: Dict[str, Any], proof: Dict[str, Any]) -> List[Dict[str, str]]:
    evidence = [
        {"source": "plan", "id": plan_slide.get("source_id", plan_slide.get("slide_id", ""))},
    ]
    if curated_slide:
        evidence.append({"source": "slide_spec", "id": curated_slide.get("slide_id", "")})
    if proof.get("type") != "text_evidence":
        evidence.append({"source": proof["type"], "id": proof.get("id", "")})
    return evidence


def _fallback_slides_from_summary(inventory: Dict[str, Any]) -> List[Dict[str, Any]]:
    slides = []
    items = inventory.get("summary_items", [])[:12]
    for index, item in enumerate(items, start=1):
        slides.append(
            {
                "slide_id": f"slide_{index:02d}",
                "source_plan_id": "",
                "title": item.get("title", f"Slide {index}"),
                "slide_role": "title" if index == 1 else "evidence",
                "claim": _limit_words(item.get("title", ""), 18),
                "support": _limit_words(item.get("text", ""), 80),
                "proof_object": {
                    "type": "text_evidence",
                    "id": item.get("id", ""),
                    "focus": item.get("category", ""),
                },
                "source_evidence": [{"source": "summary", "id": item.get("id", "")}],
                "content_priority": "must_keep" if index == 1 else "normal",
                "visual_instruction": "rough draft only; preserve content before aesthetic decisions",
            }
        )
    return slides


def _figure_path_index(inventory: Dict[str, Any]) -> Dict[str, str]:
    figures = inventory.get("assets", {}).get("figures", []) or []
    result = {}
    for item in figures:
        if not isinstance(item, dict):
            continue
        figure_id = str(item.get("id", ""))
        path = str(item.get("path", ""))
        if figure_id and path:
            result[figure_id] = path
    return result


def _rgb(hex_color: str) -> Any:
    from pptx.dml.color import RGBColor

    value = hex_color.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _theme(name: str) -> Any:
    return _rgb(THEME_HEX[name])


def _set_slide_background(slide: Any, color: Any) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide: Any,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    font_size: int,
    color: Any,
    bold: bool = False,
    letter_spaced: bool = False,
) -> Any:
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.util import Inches, Pt

    raw_text = str(text or "")
    if h < 0.3 and len(raw_text) > 32:
        h = 0.34
    if h < 0.75 and len(raw_text.split()) > 22:
        h = 0.78

    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    paragraph = frame.paragraphs[0]
    paragraph.text = _letter_space(raw_text) if letter_spaced else raw_text
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return shape


def _set_shape_text(shape: Any, text: str, font_size: int, color: Any, bold: bool = False) -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    paragraph = frame.paragraphs[0]
    paragraph.text = str(text or "")
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def _add_rule(slide: Any, x: float, y: float, w: float, color: Any) -> Any:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.01))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_proof_panel(
    slide: Any,
    proof_type: str,
    proof_id: str,
    proof_focus: str,
    figure_path: str,
    table_caption: str,
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.25), Inches(1.15), Inches(5.35), Inches(5.35))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(239, 239, 235)
    panel.line.color.rgb = RGBColor(210, 210, 205)
    panel.line.width = Inches(0.006)

    _add_textbox(
        slide,
        text=f"PROOF OBJECT / {str(proof_type or 'text').upper()}",
        x=7.55,
        y=1.42,
        w=4.6,
        h=0.28,
        font_size=8,
        color=RGBColor(112, 112, 106),
        bold=True,
        letter_spaced=True,
    )
    _add_textbox(
        slide,
        text=str(proof_id or "source evidence"),
        x=7.55,
        y=1.85,
        w=4.6,
        h=0.5,
        font_size=18,
        color=RGBColor(33, 43, 56),
        bold=True,
    )

    if proof_type == "figure" and figure_path and Path(figure_path).exists():
        try:
            slide.shapes.add_picture(figure_path, Inches(7.55), Inches(2.45), width=Inches(4.75), height=Inches(2.75))
            y = 5.35
            h = 0.62
        except Exception:
            y = 2.55
            h = 2.8
    else:
        y = 2.55
        h = 2.8

    body = _first_non_empty([table_caption, proof_focus, "Evidence preserved from parsed checkpoints."])
    _add_textbox(
        slide,
        text=body,
        x=7.55,
        y=y,
        w=4.6,
        h=h,
        font_size=12,
        color=RGBColor(63, 68, 73),
    )


def _evidence_footer(slide_data: Dict[str, Any]) -> str:
    evidence = []
    for item in slide_data.get("source_evidence", []) or []:
        if not isinstance(item, dict):
            continue
        source = item.get("source", "")
        item_id = item.get("id", "")
        if source or item_id:
            evidence.append(f"{source}:{item_id}".strip(":"))
    return "Sources: " + "; ".join(evidence[:4]) if evidence else "Sources: parsed checkpoints"


def _letter_space(text: str) -> str:
    return " ".join(str(text or ""))


def _html_table_to_rows(table_html: str) -> List[List[str]]:
    """Best-effort conversion of extracted HTML tables into plain rows."""
    if not table_html:
        return []
    row_matches = re.findall(r"<tr[^>]*>(.*?)</tr>", table_html, flags=re.IGNORECASE | re.DOTALL)
    rows: List[List[str]] = []
    for row_html in row_matches:
        cell_matches = re.findall(r"<t[dh][^>]*>(.*?)</t[dh]>", row_html, flags=re.IGNORECASE | re.DOTALL)
        cells = [_clean_table_cell(cell) for cell in cell_matches]
        if cells:
            rows.append(cells)
    max_columns = max((len(row) for row in rows), default=0)
    if max_columns:
        rows = [row + [""] * (max_columns - len(row)) for row in rows]
    return rows


def _clean_table_cell(value: str) -> str:
    text = re.sub(r"<br\s*/?>", "\n", value or "", flags=re.IGNORECASE)
    text = re.sub(r"<[^>]+>", " ", text)
    text = html.unescape(text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _split_summary_section(text: str) -> List[str]:
    text = _clean_text(text)
    if not text:
        return []
    headings = list(re.finditer(r"(?m)^#{2,}\s+(.+?)\s*$", text))
    chunks: List[str] = []
    if headings:
        for index, heading in enumerate(headings):
            end = headings[index + 1].start() if index + 1 < len(headings) else len(text)
            chunk = _clean_text(text[heading.start():end])
            if chunk:
                chunks.append(chunk)
    else:
        paragraphs = [part.strip() for part in re.split(r"\n\s*\n|---+", text) if part.strip()]
        chunks = paragraphs or [text]
    return [_limit_chars(chunk, 2200) for chunk in chunks if len(chunk) > 20]


def _infer_item_title(text: str, fallback: str) -> str:
    heading = re.search(r"(?m)^#{1,6}\s+(.+?)\s*$", text)
    if heading:
        return _limit_words(heading.group(1), 10)
    label = re.match(r"^([A-Z][A-Z0-9 /&-]{4,80})", text)
    if label:
        return _limit_words(label.group(1).title(), 10)
    sentence = re.split(r"[.!?]\s+", text, maxsplit=1)[0]
    return _limit_words(sentence or fallback.replace("_", " ").title(), 10)


def _extract_title(paper_info: str) -> str:
    text = _clean_text(paper_info)
    patterns = [
        r"\*\*Title\*\*:\s*(.+?)(?:\s+\*\*Authors?\*\*:|$)",
        r"Title:\s*(.+?)(?:\s+Authors?:|$)",
    ]
    for pattern in patterns:
        match = re.search(pattern, text, flags=re.IGNORECASE)
        if match:
            return _clean_text(match.group(1))
    return ""


def _clean_text(text: Any) -> str:
    cleaned = str(text or "")
    cleaned = cleaned.replace("\r\n", "\n")
    cleaned = re.sub(r"<[^>]+>", " ", cleaned)
    cleaned = re.sub(r"\*\*(.*?)\*\*", r"\1", cleaned)
    cleaned = re.sub(r"[ \t]+", " ", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def _limit_words(text: str, max_words: int) -> str:
    words = _clean_text(text).split()
    if len(words) <= max_words:
        return " ".join(words)
    return " ".join(words[:max_words]).rstrip(" ,;:") + "."


def _limit_chars(text: str, max_chars: int) -> str:
    text = _clean_text(text)
    if len(text) <= max_chars:
        return text
    return text[:max_chars].rsplit(" ", 1)[0].rstrip(" ,;:") + "."


def _first_non_empty(candidates: Iterable[str]) -> str:
    for candidate in candidates:
        text = _clean_text(candidate)
        if text:
            return text
    return ""


def _read_json(path: Optional[Path]) -> Dict[str, Any]:
    if not path:
        return {}
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return {}


def _write_json(path: Path, data: Dict[str, Any]) -> None:
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def main(argv: Optional[List[str]] = None) -> int:
    parser = argparse.ArgumentParser(description="Prepare from-scratch template experiment artifacts.")
    parser.add_argument("--summary-checkpoint", required=True, help="Path to checkpoint_summary.json.")
    parser.add_argument("--plan-checkpoint", help="Optional path to checkpoint_plan.json.")
    parser.add_argument("--spec-checkpoint", help="Optional path to checkpoint_slide_spec.json.")
    parser.add_argument("--output-dir", required=True, help="Directory for content_inventory.json and rough_draft_spec.json.")
    parser.add_argument("--pptx-output", help="Optional path for a plain rough draft PPTX.")
    parser.add_argument("--render-review-dir", help="Optional directory for exported PNG review pages.")
    args = parser.parse_args(argv)

    paths = write_from_scratch_artifacts(
        summary_checkpoint=Path(args.summary_checkpoint),
        plan_checkpoint=Path(args.plan_checkpoint) if args.plan_checkpoint else None,
        spec_checkpoint=Path(args.spec_checkpoint) if args.spec_checkpoint else None,
        output_dir=Path(args.output_dir),
        pptx_output=Path(args.pptx_output) if args.pptx_output else None,
        render_review_dir=Path(args.render_review_dir) if args.render_review_dir else None,
    )
    print(json.dumps(paths, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
