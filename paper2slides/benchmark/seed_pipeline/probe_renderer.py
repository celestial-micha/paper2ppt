"""Render a small native PPTX visual probe from a seed spec."""
from __future__ import annotations

import argparse
import json
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .strategist import load_seed_template_package


def render_visual_probe_pptx(
    visual_probe_spec: Dict[str, Any],
    package: Dict[str, Any],
    output_pptx: Path,
    *,
    promotion_gate: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    """Render a spec-only visual probe as native editable PPTX."""
    _assert_render_allowed(visual_probe_spec, promotion_gate)
    spec_lock = package.get("spec_lock", {}) or {}
    brand = package.get("brand", {}) or {}
    palette = _palette(spec_lock, brand)
    fonts = brand.get("font_families", []) or ["Aptos", "Arial", "Calibri"]
    font_name = fonts[0]
    canvas = spec_lock.get("canvas", {}) or {}
    width = float(canvas.get("width_in", 13.333))
    height = float(canvas.get("height_in", 7.5))

    prs = Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)
    blank_layout = prs.slide_layouts[6]
    trace_slides = []
    for slide_spec in visual_probe_spec.get("slides", []) or []:
        slide = prs.slides.add_slide(blank_layout)
        _set_background(slide, palette["background"])
        role = slide_spec.get("role", "")
        if role == "cover":
            trace = _render_cover(slide, slide_spec, palette, font_name, width, height)
        elif role == "metric_ledger":
            trace = _render_metric_ledger(slide, slide_spec, palette, font_name, width, height)
        elif role == "evidence_wall":
            trace = _render_evidence_wall(slide, slide_spec, palette, font_name, width, height)
        elif role == "figure_or_table_focus":
            trace = _render_focus_slide(slide, slide_spec, palette, font_name, width, height)
        else:
            trace = _render_content_slide(slide, slide_spec, palette, font_name, width, height)
        trace_slides.append(trace)

    output_pptx = Path(output_pptx)
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    return {
        "schema_version": "visual_probe_render_trace.v0",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "output_pptx": str(output_pptx),
        "template_id": package.get("provenance", {}).get("template_id", ""),
        "paper_title": visual_probe_spec.get("paper_title", ""),
        "slide_count": len(trace_slides),
        "slides": trace_slides,
        "renderer_policy": {
            "scope": "visual_probe_only",
            "native_editable_required": True,
            "full_deck_generation_allowed": False,
            "main_renderer_modified": False,
        },
    }


def write_visual_probe_render_artifacts(
    spec_path: Path,
    package_dir: Path,
    output_pptx: Path,
    *,
    trace_path: Optional[Path] = None,
    promotion_gate_path: Optional[Path] = None,
) -> Dict[str, str]:
    """Render the PPTX and write a trace JSON next to it."""
    spec = _read_json(spec_path)
    package = load_seed_template_package(package_dir)
    gate = _read_json(promotion_gate_path) if promotion_gate_path else None
    trace = render_visual_probe_pptx(spec, package, output_pptx, promotion_gate=gate)
    trace_path = Path(trace_path) if trace_path else Path(output_pptx).with_suffix(".render_trace.json")
    trace_path.parent.mkdir(parents=True, exist_ok=True)
    _write_json(trace_path, trace)
    return {
        "pptx": str(output_pptx),
        "render_trace": str(trace_path),
    }


def _assert_render_allowed(spec: Dict[str, Any], gate: Optional[Dict[str, Any]]) -> None:
    slides = spec.get("slides", []) or []
    if not 1 <= len(slides) <= 8:
        raise ValueError("visual probe renderer only accepts 1-8 slides")
    if spec.get("rules", {}).get("full_deck_generation_allowed") is True:
        raise ValueError("visual probe renderer refuses specs that allow full deck generation")
    if gate and not gate.get("decision", {}).get("renderer_prototype_allowed", False):
        raise ValueError("promotion gate does not allow renderer prototype")


def _render_cover(slide: Any, spec: Dict[str, Any], palette: Dict[str, str], font_name: str, width: float, height: float) -> Dict[str, Any]:
    _accent_rule(slide, palette, 0.75, 0.7, 2.1)
    _text(slide, _cover_title(spec.get("title", "")), 0.75, 1.0, 7.95, 1.6, 36, palette["ink"], font_name, bold=True)
    _text(slide, spec.get("claim", ""), 0.82, 3.05, 7.2, 1.25, 19, palette["ink"], font_name)
    _panel(slide, 8.85, 1.0, 3.65, 4.85, palette, fill="surface")
    _text(slide, "Visual probe", 9.15, 1.3, 2.8, 0.35, 13, palette["accent"], font_name, bold=True)
    _text(slide, spec.get("support", ""), 9.15, 1.85, 2.85, 2.0, 14, palette["ink"], font_name)
    _source_chip(slide, spec, 9.15, 5.1, palette, font_name)
    _footer(slide, spec, palette, font_name, width, height)
    return _trace(spec, "cover_editorial", 7)


def _render_content_slide(slide: Any, spec: Dict[str, Any], palette: Dict[str, str], font_name: str, width: float, height: float) -> Dict[str, Any]:
    _slide_title(slide, spec, palette, font_name)
    _text(slide, spec.get("claim", ""), 0.75, 1.42, 5.85, 1.55, 24, palette["ink"], font_name, bold=True)
    _text(slide, spec.get("support", ""), 0.8, 3.18, 5.25, 1.87, 14, palette["ink"], font_name)
    _panel(slide, 6.85, 1.45, 5.65, 4.7, palette, fill="surface")
    _text(slide, "Evidence note", 7.15, 1.75, 2.4, 0.35, 12, palette["accent"], font_name, bold=True)
    _proof_note(slide, spec, 7.15, 2.2, 4.95, 2.75, palette, font_name)
    _source_chip(slide, spec, 7.15, 5.35, palette, font_name)
    _footer(slide, spec, palette, font_name, width, height)
    return _trace(spec, "claim_evidence_split", 8)


def _render_metric_ledger(slide: Any, spec: Dict[str, Any], palette: Dict[str, str], font_name: str, width: float, height: float) -> Dict[str, Any]:
    _slide_title(slide, spec, palette, font_name)
    _text(slide, spec.get("claim", ""), 0.75, 1.3, 8.2, 0.85, 24, palette["ink"], font_name, bold=True)
    items = spec.get("proof_object", {}).get("items", []) or []
    card_boxes = [(0.75, 2.25), (4.0, 2.25), (7.25, 2.25), (10.5, 2.25)]
    for index, (x, y) in enumerate(card_boxes):
        item = items[index] if index < len(items) else {}
        _metric_card(slide, item, x, y, 2.45, 2.25, palette, font_name, index + 1)
    _text(slide, spec.get("support", ""), 0.8, 5.15, 11.5, 0.75, 12.5, palette["muted"], font_name)
    _footer(slide, spec, palette, font_name, width, height)
    return _trace(spec, "metric_ledger_cards", 14)


def _render_evidence_wall(slide: Any, spec: Dict[str, Any], palette: Dict[str, str], font_name: str, width: float, height: float) -> Dict[str, Any]:
    _slide_title(slide, spec, palette, font_name)
    _text(slide, spec.get("claim", ""), 0.75, 1.28, 6.4, 0.95, 24, palette["ink"], font_name, bold=True)
    items = spec.get("proof_object", {}).get("items", []) or []
    boxes = [(0.75, 2.35), (4.92, 2.35), (9.08, 2.35), (0.75, 4.55), (4.92, 4.55), (9.08, 4.55)]
    for index, (x, y) in enumerate(boxes):
        item = items[index] if index < len(items) else {}
        _evidence_card(slide, item, x, y, 3.35, 1.55, palette, font_name, index + 1)
    _footer(slide, spec, palette, font_name, width, height)
    return _trace(spec, "evidence_wall_cards", 15)


def _render_focus_slide(slide: Any, spec: Dict[str, Any], palette: Dict[str, str], font_name: str, width: float, height: float) -> Dict[str, Any]:
    _slide_title(slide, spec, palette, font_name)
    _text(slide, spec.get("claim", ""), 0.75, 1.3, 5.35, 1.0, 24, palette["ink"], font_name, bold=True)
    _text(slide, spec.get("support", ""), 0.78, 2.45, 4.9, 2.3, 13, palette["ink"], font_name)
    _panel(slide, 6.15, 1.35, 6.35, 4.8, palette, fill="surface")
    _text(slide, "Source evidence focus", 6.48, 1.68, 3.3, 0.35, 12, palette["accent"], font_name, bold=True)
    _proof_note(slide, spec, 6.48, 2.18, 5.55, 2.95, palette, font_name)
    _source_chip(slide, spec, 6.48, 5.4, palette, font_name)
    _footer(slide, spec, palette, font_name, width, height)
    return _trace(spec, "single_proof_focus", 9)


def _slide_title(slide: Any, spec: Dict[str, Any], palette: Dict[str, str], font_name: str) -> None:
    _accent_rule(slide, palette, 0.75, 0.62, 1.45)
    _text(slide, spec.get("title", ""), 0.75, 0.77, 8.8, 0.42, 22, palette["ink"], font_name, bold=True)
    _text(slide, spec.get("role", "").replace("_", " "), 10.1, 0.78, 2.45, 0.35, 10, palette["muted"], font_name, align=PP_ALIGN.RIGHT)


def _metric_card(slide: Any, item: Dict[str, Any], x: float, y: float, w: float, h: float, palette: Dict[str, str], font_name: str, index: int) -> None:
    _panel(slide, x, y, w, h, palette, fill="surface")
    value = str(item.get("value", "") or f"{index:02d}")
    label = str(item.get("label", "") or "Metric")
    note = str(item.get("note", "") or item.get("source_slide", "") or "Source-traceable metric")
    _text(slide, value, x + 0.22, y + 0.25, w - 0.45, 0.55, 26, palette["accent"], font_name, bold=True)
    _text(slide, label, x + 0.24, y + 0.83, w - 0.48, 0.6, 12, palette["ink"], font_name, bold=True)
    _text(slide, note, x + 0.24, y + 1.48, w - 0.48, 0.48, 9.8, palette["muted"], font_name)


def _evidence_card(slide: Any, item: Dict[str, Any], x: float, y: float, w: float, h: float, palette: Dict[str, str], font_name: str, index: int) -> None:
    _panel(slide, x, y, w, h, palette, fill="surface")
    kind = str(item.get("kind", "") or item.get("type", "") or "evidence").title()
    text = _proof_item_text(item) or f"Evidence item {index}"
    _text(slide, kind, x + 0.22, y + 0.18, w - 0.44, 0.25, 10, palette["accent"], font_name, bold=True)
    _text(slide, text, x + 0.22, y + 0.5, w - 0.44, 0.78, 10, palette["ink"], font_name)


def _proof_note(slide: Any, spec: Dict[str, Any], x: float, y: float, w: float, h: float, palette: Dict[str, str], font_name: str) -> None:
    proof = spec.get("proof_object", {}) or {}
    items = proof.get("items", []) or []
    item = proof.get("item", {}) or (items[0] if items else {})
    text = _proof_item_text(item) or spec.get("support", "") or proof.get("type", "")
    _text(slide, text, x, y, w, h, 13, palette["ink"], font_name)


def _proof_item_text(item: Dict[str, Any]) -> str:
    if not isinstance(item, dict):
        return str(item or "")
    parts = [
        item.get("caption", ""),
        item.get("label", ""),
        item.get("value", ""),
        item.get("note", ""),
    ]
    text = " ".join(str(part) for part in parts if part)
    return _limit(text, 180)


def _cover_title(text: Any) -> str:
    value = str(text or "")
    if ": " in value:
        return value.replace(": ", ":\n", 1)
    return value


def _panel(slide: Any, x: float, y: float, w: float, h: float, palette: Dict[str, str], *, fill: str = "surface") -> Any:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(palette[fill])
    shape.line.color.rgb = _rgb(palette["line"])
    shape.line.width = Pt(0.75)
    return shape


def _accent_rule(slide: Any, palette: Dict[str, str], x: float, y: float, w: float) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.035))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(palette["accent"])
    shape.line.fill.background()


def _source_chip(slide: Any, spec: Dict[str, Any], x: float, y: float, palette: Dict[str, str], font_name: str) -> None:
    chip = f"{spec.get('slide_id', '')} / {spec.get('role', '')}"
    _text(slide, chip, x, y, 3.5, 0.28, 10, palette["muted"], font_name)


def _footer(slide: Any, spec: Dict[str, Any], palette: Dict[str, str], font_name: str, width: float, height: float) -> None:
    _text(slide, spec.get("slide_id", ""), 0.75, height - 0.48, 1.5, 0.2, 8, palette["muted"], font_name)
    _text(slide, "native editable visual probe", width - 3.1, height - 0.48, 2.35, 0.2, 8, palette["muted"], font_name, align=PP_ALIGN.RIGHT)


def _text(
    slide: Any,
    text: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float,
    color: str,
    font_name: str,
    *,
    bold: bool = False,
    align: Optional[int] = None,
) -> Any:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    para = frame.paragraphs[0]
    if align is not None:
        para.alignment = align
    para.space_after = Pt(0)
    run = para.add_run()
    run.text = _limit(str(text or ""), 520)
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return shape


def _set_background(slide: Any, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _palette(spec_lock: Dict[str, Any], brand: Dict[str, Any]) -> Dict[str, str]:
    source = spec_lock.get("palette", {}) or brand.get("palette", {}) or {}
    return {
        "background": source.get("background", "#F4F0E8"),
        "surface": "#FFFFFF",
        "ink": source.get("ink", "#171717"),
        "accent": source.get("accent", "#E24A2B"),
        "muted": source.get("muted", "#6B7280"),
        "line": "#D8D1C5",
    }


def _trace(spec: Dict[str, Any], layout: str, primitive_count: int) -> Dict[str, Any]:
    return {
        "slide_id": spec.get("slide_id", ""),
        "role": spec.get("role", ""),
        "title": spec.get("title", ""),
        "layout_used": layout,
        "layout_candidate": spec.get("layout_candidate", ""),
        "primitive_count_estimate": primitive_count,
        "native_only": True,
    }


def _rgb(hex_color: str) -> RGBColor:
    value = str(hex_color or "#000000").strip().lstrip("#")
    if len(value) != 6:
        value = "000000"
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _limit(text: str, chars: int) -> str:
    text = " ".join(str(text or "").split())
    if len(text) <= chars:
        return text
    return text[: max(0, chars - 1)].rstrip() + "..."


def _read_json(path: Optional[Path]) -> Dict[str, Any]:
    if path is None:
        return {}
    return json.loads(Path(path).read_text(encoding="utf-8"))


def _write_json(path: Path, payload: Dict[str, Any]) -> None:
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def main(argv: List[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Render a native editable PPTX from a visual probe spec.")
    parser.add_argument("--spec", required=True, help="Path to visual_probe_spec.json")
    parser.add_argument("--package-dir", required=True, help="Path to seed_template_package directory")
    parser.add_argument("--output-pptx", required=True, help="Output PPTX path")
    parser.add_argument("--trace", help="Optional render trace JSON path")
    parser.add_argument("--promotion-gate", help="Optional promotion_gate.json")
    args = parser.parse_args(argv)

    paths = write_visual_probe_render_artifacts(
        Path(args.spec),
        Path(args.package_dir),
        Path(args.output_pptx),
        trace_path=Path(args.trace) if args.trace else None,
        promotion_gate_path=Path(args.promotion_gate) if args.promotion_gate else None,
    )
    print(json.dumps(paths, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
