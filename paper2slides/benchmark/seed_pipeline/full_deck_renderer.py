"""Render a full-deck prototype from seed strategy artifacts."""
from __future__ import annotations

import argparse
import json
from collections import OrderedDict
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .probe_renderer import _accent_rule, _limit, _palette, _panel, _read_json, _rgb, _set_background, _text, _write_json
from .strategist import load_seed_template_package


DEFAULT_STYLE_VARIANT = "editorial_data_report"

STYLE_VARIANTS: Dict[str, Dict[str, Any]] = {
    "editorial_data_report": {
        "label": "Editorial Data Report",
        "visual_language": "warm editorial data-reporting with restrained source evidence blocks",
        "cover_label": "Seed full-deck prototype",
        "section_label": "MODULE",
        "closing_label": "Takeaway",
        "footer_label": "seed full-deck prototype",
        "chrome": "none",
        "palette": {},
    },
    "blueprint_system_map": {
        "label": "Blueprint System Map",
        "visual_language": "cool technical blueprint for systems, safeguards, and evaluation pathways",
        "cover_label": "Blueprint seed candidate",
        "section_label": "SYSTEM NODE",
        "closing_label": "System Takeaway",
        "footer_label": "blueprint system map",
        "chrome": "blueprint_grid",
        "palette": {
            "background": "#EEF4F8",
            "surface": "#FFFFFF",
            "ink": "#102A43",
            "accent": "#2563EB",
            "muted": "#52606D",
            "line": "#B6C7D6",
        },
    },
    "dark_evidence_console": {
        "label": "Dark Evidence Console",
        "visual_language": "dark command-center treatment for evidence, safety controls, and operating signals",
        "cover_label": "Evidence console candidate",
        "section_label": "CONSOLE",
        "closing_label": "Operating Takeaway",
        "footer_label": "dark evidence console",
        "chrome": "console_guides",
        "palette": {
            "background": "#111827",
            "surface": "#1F2937",
            "ink": "#F9FAFB",
            "accent": "#38BDF8",
            "muted": "#CBD5E1",
            "line": "#374151",
        },
    },
}


def build_seed_full_deck_spec(
    inventory: Dict[str, Any],
    contract: Dict[str, Any],
    package: Dict[str, Any],
    *,
    slide_budget: int = 24,
    style_variant: str = DEFAULT_STYLE_VARIANT,
) -> Dict[str, Any]:
    """Build a deterministic full-deck prototype spec from checkpoint inventory."""
    variant = _style_variant(style_variant)
    curated_by_id = {
        slide.get("slide_id"): slide
        for slide in inventory.get("curated_slides", []) or []
        if slide.get("slide_id")
    }
    plan_slides = list(inventory.get("plan_slides", []) or [])
    sections = _section_groups(plan_slides)

    slides: List[Dict[str, Any]] = [
        _cover_slide(inventory, contract),
        _agenda_slide(sections),
    ]
    for section_index, (section, section_items) in enumerate(sections.items(), start=1):
        slides.append(_section_slide(section_index, section, section_items))
        for index, plan_slide in enumerate(section_items, start=1):
            curated = curated_by_id.get(plan_slide.get("slide_id"), {})
            slides.append(_content_slide(plan_slide, curated, inventory, section, index))

    closing = _closing_slide(inventory, contract)
    while len(slides) < max(3, slide_budget - 1):
        extra_index = len(slides) + 1
        if not any(slide.get("role") == "metric_ledger" and slide.get("slide_id") == "seed_metric_ledger" for slide in slides):
            slides.append(_metric_ledger_slide(inventory))
        elif not any(slide.get("role") == "evidence_inventory" for slide in slides):
            slides.append(_evidence_inventory_slide(inventory))
        else:
            slides.append(_synthesis_slide(inventory, extra_index))
    if len(slides) >= slide_budget:
        slides = slides[: slide_budget - 1]
    slides.append(closing)

    for index, slide in enumerate(slides, start=1):
        slide["slide_index"] = index
        slide.setdefault("slide_id", f"seed_full_{index:02d}")

    return {
        "schema_version": "seed_full_deck_spec.v0",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "paper_title": inventory.get("paper", {}).get("title", ""),
        "template_id": package.get("provenance", {}).get("template_id", ""),
        "style_variant": variant["id"],
        "style": {
            "variant": variant["id"],
            "label": variant["label"],
            "visual_language": variant["visual_language"],
            "source": "ppt_master_seed_pipeline",
        },
        "purpose": "full-deck prototype for human review; not default template promotion",
        "rules": {
            "native_editability_required": True,
            "prototype_only": True,
            "default_template_promotion_allowed": False,
            "renderer_scope": "seed_pipeline_full_deck_prototype",
            "style_variant": variant["id"],
        },
        "source_contract": {
            "schema_version": contract.get("schema_version", ""),
            "deck_thesis": contract.get("deck_intent", {}).get("deck_thesis", ""),
            "slide_budget": slide_budget,
        },
        "slides": slides,
    }


def render_seed_full_deck_pptx(
    full_deck_spec: Dict[str, Any],
    package: Dict[str, Any],
    output_pptx: Path,
    *,
    promotion_gate: Optional[Dict[str, Any]] = None,
    allow_blocked_prototype: bool = False,
    style_variant: Optional[str] = None,
) -> Dict[str, Any]:
    """Render a seed full-deck prototype as native editable PPTX."""
    _assert_full_deck_render_allowed(full_deck_spec, promotion_gate, allow_blocked_prototype)
    variant = _style_variant(style_variant or full_deck_spec.get("style_variant") or DEFAULT_STYLE_VARIANT)
    spec_lock = package.get("spec_lock", {}) or {}
    brand = package.get("brand", {}) or {}
    palette = _variant_palette(_palette(spec_lock, brand), variant)
    fonts = brand.get("font_families", []) or ["Aptos", "Arial", "Calibri"]
    font_name = fonts[0]
    canvas = spec_lock.get("canvas", {}) or {}
    width = float(canvas.get("width_in", 13.333))
    height = float(canvas.get("height_in", 7.5))

    prs = Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)
    blank_layout = prs.slide_layouts[6]

    trace_slides = []
    for slide_spec in full_deck_spec.get("slides", []) or []:
        slide = prs.slides.add_slide(blank_layout)
        _set_background(slide, palette["background"])
        _render_variant_chrome(slide, palette, width, height, variant)
        trace = _render_seed_slide(slide, slide_spec, palette, font_name, width, height, variant)
        trace_slides.append(trace)

    output_pptx = Path(output_pptx)
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    return {
        "schema_version": "seed_full_deck_render_trace.v0",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "output_pptx": str(output_pptx),
        "template_id": package.get("provenance", {}).get("template_id", ""),
        "paper_title": full_deck_spec.get("paper_title", ""),
        "style_variant": variant["id"],
        "style_label": variant["label"],
        "visual_language": variant["visual_language"],
        "slide_count": len(trace_slides),
        "slides": trace_slides,
        "renderer_policy": {
            "scope": "seed_full_deck_prototype",
            "native_editable_required": True,
            "prototype_only": True,
            "default_template_promotion_allowed": False,
            "main_renderer_modified": False,
            "allow_blocked_prototype": bool(allow_blocked_prototype),
        },
    }


def write_seed_full_deck_artifacts(
    inventory_path: Path,
    contract_path: Path,
    package_dir: Path,
    output_dir: Path,
    output_pptx: Path,
    *,
    promotion_gate_path: Optional[Path] = None,
    allow_blocked_prototype: bool = False,
    slide_budget: int = 24,
    style_variant: str = DEFAULT_STYLE_VARIANT,
) -> Dict[str, str]:
    """Write spec, PPTX, and render trace for a full-deck prototype."""
    inventory = _read_json(inventory_path)
    contract = _read_json(contract_path)
    package = load_seed_template_package(package_dir)
    gate = _read_json(promotion_gate_path) if promotion_gate_path else None
    spec = build_seed_full_deck_spec(inventory, contract, package, slide_budget=slide_budget, style_variant=style_variant)

    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    spec_path = output_dir / "seed_full_deck_spec.json"
    trace_path = output_dir / "render_trace.json"
    _write_json(spec_path, spec)
    trace = render_seed_full_deck_pptx(
        spec,
        package,
        output_pptx,
        promotion_gate=gate,
        allow_blocked_prototype=allow_blocked_prototype,
        style_variant=style_variant,
    )
    _write_json(trace_path, trace)
    return {
        "seed_full_deck_spec": str(spec_path),
        "pptx": str(output_pptx),
        "render_trace": str(trace_path),
    }


def _style_variant(style_variant: Optional[str]) -> Dict[str, Any]:
    key = str(style_variant or DEFAULT_STYLE_VARIANT).strip().lower().replace("-", "_").replace(" ", "_")
    aliases = {
        "default": DEFAULT_STYLE_VARIANT,
        "editorial": DEFAULT_STYLE_VARIANT,
        "style_a": DEFAULT_STYLE_VARIANT,
        "07": DEFAULT_STYLE_VARIANT,
        "blueprint": "blueprint_system_map",
        "style_b": "blueprint_system_map",
        "console": "dark_evidence_console",
        "dark": "dark_evidence_console",
        "style_c": "dark_evidence_console",
    }
    key = aliases.get(key, key)
    if key not in STYLE_VARIANTS:
        available = ", ".join(sorted(STYLE_VARIANTS))
        raise ValueError(f"unknown style_variant {style_variant!r}; available variants: {available}")
    variant = dict(STYLE_VARIANTS[key])
    variant["id"] = key
    return variant


def _variant_palette(base_palette: Dict[str, str], variant: Dict[str, Any]) -> Dict[str, str]:
    palette = dict(base_palette)
    palette.update(variant.get("palette", {}) or {})
    palette.setdefault("background", "#F4F0E8")
    palette.setdefault("surface", "#FFFFFF")
    palette.setdefault("ink", "#171717")
    palette.setdefault("accent", "#E24A2B")
    palette.setdefault("muted", "#6B7280")
    palette.setdefault("line", "#D8D1C5")
    return palette


def _render_variant_chrome(
    slide: Any,
    palette: Dict[str, str],
    width: float,
    height: float,
    variant: Dict[str, Any],
) -> None:
    chrome = variant.get("chrome", "none")
    if chrome == "blueprint_grid":
        for x in (1.5, 3.0, 4.5, 6.0, 7.5, 9.0, 10.5, 12.0):
            _decorative_rect(slide, x, 0.0, 0.012, height, palette["line"])
        for y in (0.75, 1.5, 2.25, 3.0, 3.75, 4.5, 5.25, 6.0, 6.75):
            _decorative_rect(slide, 0.0, y, width, 0.012, palette["line"])
    elif chrome == "console_guides":
        _decorative_rect(slide, 0.38, 0.55, 0.026, height - 1.08, palette["accent"])
        _decorative_rect(slide, width - 0.42, 0.55, 0.026, height - 1.08, palette["line"])
        _decorative_rect(slide, 0.75, 0.34, width - 1.5, 0.014, palette["line"])
        _decorative_rect(slide, 0.75, height - 0.72, width - 1.5, 0.014, palette["line"])


def _decorative_rect(slide: Any, x: float, y: float, w: float, h: float, color: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()


def _assert_full_deck_render_allowed(
    spec: Dict[str, Any],
    gate: Optional[Dict[str, Any]],
    allow_blocked_prototype: bool,
) -> None:
    slides = spec.get("slides", []) or []
    if len(slides) < 10:
        raise ValueError("full-deck renderer expects at least 10 slides")
    if len(slides) > 32:
        raise ValueError("full-deck prototype refuses decks larger than 32 slides")
    if gate and not gate.get("decision", {}).get("full_deck_expansion_allowed", False) and not allow_blocked_prototype:
        raise ValueError("promotion gate blocks full-deck expansion; pass allow_blocked_prototype for a review-only draft")


def _render_seed_slide(
    slide: Any,
    spec: Dict[str, Any],
    palette: Dict[str, str],
    font_name: str,
    width: float,
    height: float,
    variant: Dict[str, Any],
) -> Dict[str, Any]:
    role = spec.get("role", "")
    if role == "cover":
        _render_full_cover(slide, spec, palette, font_name, width, height, variant)
        layout = "full_cover"
    elif role == "agenda":
        _render_full_agenda(slide, spec, palette, font_name, width, height, variant)
        layout = "full_agenda"
    elif role == "section":
        _render_full_section(slide, spec, palette, font_name, width, height, variant)
        layout = "full_section"
    elif role == "metric_ledger":
        _render_full_metric_ledger(slide, spec, palette, font_name, width, height, variant)
        layout = "full_metric_ledger"
    elif role == "evidence_inventory":
        _render_full_evidence_inventory(slide, spec, palette, font_name, width, height, variant)
        layout = "full_evidence_inventory"
    elif role == "closing":
        _render_full_closing(slide, spec, palette, font_name, width, height, variant)
        layout = "full_closing"
    else:
        _render_full_content(slide, spec, palette, font_name, width, height, variant)
        layout = spec.get("layout", "full_content")
    return {
        "slide_id": spec.get("slide_id", ""),
        "role": role,
        "title": spec.get("title", ""),
        "layout_used": layout,
        "style_variant": variant["id"],
        "native_only": False if _has_picture(spec) else True,
        "proof_type": spec.get("proof_object", {}).get("type", "none"),
    }


def _render_full_cover(
    slide: Any,
    spec: Dict[str, Any],
    palette: Dict[str, str],
    font_name: str,
    width: float,
    height: float,
    variant: Dict[str, Any],
) -> None:
    _accent_rule(slide, palette, 0.75, 0.72, 2.3)
    _text(slide, _split_title(spec.get("title", "")), 0.75, 1.05, 8.05, 1.6, 35, palette["ink"], font_name, bold=True)
    _text(slide, spec.get("claim", ""), 0.82, 3.1, 7.3, 1.05, 19, palette["ink"], font_name)
    _panel(slide, 8.85, 1.05, 3.65, 4.85, palette, fill="surface")
    _text(slide, variant["cover_label"], 9.14, 1.35, 2.9, 0.35, 12.5, palette["accent"], font_name, bold=True)
    for index, item in enumerate(spec.get("highlights", [])[:4]):
        y = 1.85 + index * 0.84
        _text(slide, f"{index + 1:02d}", 9.14, y, 0.34, 0.24, 9, palette["accent"], font_name, bold=True)
        _text(slide, _limit_words(item, 6), 9.55, y - 0.02, 2.45, 0.5, 10.3, palette["ink"], font_name)
    _footer(slide, spec, palette, font_name, width, height, variant)


def _render_full_agenda(
    slide: Any,
    spec: Dict[str, Any],
    palette: Dict[str, str],
    font_name: str,
    width: float,
    height: float,
    variant: Dict[str, Any],
) -> None:
    _slide_header(slide, spec, palette, font_name)
    sections = spec.get("sections", []) or []
    for index, section in enumerate(sections[:6]):
        y = 1.45 + index * 0.78
        _panel(slide, 0.82, y, 11.7, 0.56, palette, fill="surface")
        _accent_rule(slide, palette, 0.82, y, 0.52)
        _text(slide, f"{index + 1:02d}", 1.0, y + 0.16, 0.36, 0.18, 9, palette["accent"], font_name, bold=True)
        _text(slide, section.get("title", ""), 1.62, y + 0.12, 3.2, 0.22, 13.5, palette["ink"], font_name, bold=True)
        _text(slide, f"{section.get('slide_count', 0)} content slides", 5.2, y + 0.15, 1.6, 0.18, 9.5, palette["muted"], font_name)
        _text(slide, _limit(section.get("sample", ""), 72), 7.05, y + 0.1, 4.65, 0.28, 10, palette["muted"], font_name)
    _footer(slide, spec, palette, font_name, width, height, variant)


def _render_full_section(
    slide: Any,
    spec: Dict[str, Any],
    palette: Dict[str, str],
    font_name: str,
    width: float,
    height: float,
    variant: Dict[str, Any],
) -> None:
    _panel(slide, 0.85, 1.05, 11.55, 5.08, palette, fill="surface")
    _accent_rule(slide, palette, 1.22, 1.68, 1.35)
    _text(slide, f"{variant['section_label']} {spec.get('section_index', 0):02d}", 1.22, 2.1, 2.8, 0.24, 9, palette["accent"], font_name, bold=True)
    _text(slide, spec.get("title", ""), 1.2, 2.65, 8.5, 0.8, 34, palette["ink"], font_name, bold=True)
    _text(slide, spec.get("support", ""), 1.24, 4.15, 9.7, 0.55, 14, palette["muted"], font_name)
    _footer(slide, spec, palette, font_name, width, height, variant)


def _render_full_content(
    slide: Any,
    spec: Dict[str, Any],
    palette: Dict[str, str],
    font_name: str,
    width: float,
    height: float,
    variant: Dict[str, Any],
) -> None:
    _slide_header(slide, spec, palette, font_name)
    proof = spec.get("proof_object", {}) or {}
    layout = spec.get("layout", "claim_evidence_split")
    if layout == "table_focus":
        _text(slide, spec.get("claim", ""), 0.75, 1.48, 9.0, 0.82, 24, palette["ink"], font_name, bold=True)
        _text(slide, spec.get("support", ""), 0.78, 2.42, 5.1, 0.82, 12.5, palette["ink"], font_name)
        _render_proof(slide, proof, 0.78, 3.38, 11.72, 2.8, palette, font_name, table_compact=True)
    elif layout == "figure_focus":
        _text(slide, spec.get("claim", ""), 0.75, 1.48, 5.35, 1.04, 24, palette["ink"], font_name, bold=True)
        _text(slide, spec.get("support", ""), 0.78, 2.78, 4.92, 2.2, 12.5, palette["ink"], font_name)
        _render_proof(slide, proof, 6.25, 1.62, 6.25, 4.58, palette, font_name)
    else:
        _text(slide, spec.get("claim", ""), 0.75, 1.48, 5.75, 1.1, 24, palette["ink"], font_name, bold=True)
        _text(slide, spec.get("support", ""), 0.78, 2.88, 5.25, 2.1, 12.5, palette["ink"], font_name)
        _render_proof(slide, proof, 6.7, 1.62, 5.8, 4.52, palette, font_name)
    _footer(slide, spec, palette, font_name, width, height, variant)


def _render_full_metric_ledger(
    slide: Any,
    spec: Dict[str, Any],
    palette: Dict[str, str],
    font_name: str,
    width: float,
    height: float,
    variant: Dict[str, Any],
) -> None:
    _slide_header(slide, spec, palette, font_name)
    _text(slide, spec.get("claim", ""), 0.75, 1.28, 8.4, 0.8, 23, palette["ink"], font_name, bold=True)
    items = spec.get("proof_object", {}).get("items", []) or []
    positions = [(0.75, 2.25), (4.0, 2.25), (7.25, 2.25), (10.5, 2.25), (0.75, 4.55), (4.0, 4.55), (7.25, 4.55), (10.5, 4.55)]
    for index, (x, y) in enumerate(positions):
        item = items[index] if index < len(items) else {}
        _metric_card(slide, item, x, y, 2.45, 1.62, palette, font_name, index + 1)
    _footer(slide, spec, palette, font_name, width, height, variant)


def _render_full_evidence_inventory(
    slide: Any,
    spec: Dict[str, Any],
    palette: Dict[str, str],
    font_name: str,
    width: float,
    height: float,
    variant: Dict[str, Any],
) -> None:
    _slide_header(slide, spec, palette, font_name)
    _text(slide, spec.get("claim", ""), 0.75, 1.3, 8.5, 0.72, 22, palette["ink"], font_name, bold=True)
    items = spec.get("proof_object", {}).get("items", []) or []
    positions = [(0.75, 2.25), (4.92, 2.25), (9.08, 2.25), (0.75, 4.38), (4.92, 4.38), (9.08, 4.38)]
    for index, (x, y) in enumerate(positions):
        item = items[index] if index < len(items) else {}
        _evidence_card(slide, item, x, y, 3.36, 1.56, palette, font_name, index + 1)
    _footer(slide, spec, palette, font_name, width, height, variant)


def _render_full_closing(
    slide: Any,
    spec: Dict[str, Any],
    palette: Dict[str, str],
    font_name: str,
    width: float,
    height: float,
    variant: Dict[str, Any],
) -> None:
    _panel(slide, 0.85, 1.0, 11.65, 5.15, palette, fill="surface")
    _accent_rule(slide, palette, 1.24, 1.58, 1.45)
    _text(slide, variant["closing_label"], 1.23, 2.0, 2.8, 0.24, 9, palette["accent"], font_name, bold=True)
    _text(slide, spec.get("claim", ""), 1.2, 2.45, 9.25, 1.25, 30, palette["ink"], font_name, bold=True)
    _text(slide, spec.get("support", ""), 1.25, 4.0, 9.2, 0.8, 15, palette["muted"], font_name)
    _footer(slide, spec, palette, font_name, width, height, variant)


def _slide_header(slide: Any, spec: Dict[str, Any], palette: Dict[str, str], font_name: str) -> None:
    _accent_rule(slide, palette, 0.75, 0.62, 1.45)
    _text(slide, spec.get("title", ""), 0.75, 0.68, 8.75, 0.68, 24, palette["ink"], font_name, bold=True)
    _text(slide, spec.get("section", spec.get("role", "")).replace("_", " "), 9.9, 0.79, 2.6, 0.24, 9.5, palette["muted"], font_name, align=PP_ALIGN.RIGHT)


def _footer(
    slide: Any,
    spec: Dict[str, Any],
    palette: Dict[str, str],
    font_name: str,
    width: float,
    height: float,
    variant: Dict[str, Any],
) -> None:
    _text(slide, f"{int(spec.get('slide_index', 0)):02d}", 0.75, height - 0.48, 0.5, 0.2, 8, palette["muted"], font_name)
    _text(slide, variant["footer_label"], width - 3.15, height - 0.48, 2.4, 0.2, 8, palette["muted"], font_name, align=PP_ALIGN.RIGHT)


def _render_proof(
    slide: Any,
    proof: Dict[str, Any],
    x: float,
    y: float,
    w: float,
    h: float,
    palette: Dict[str, str],
    font_name: str,
    *,
    table_compact: bool = False,
) -> None:
    _panel(slide, x, y, w, h, palette, fill="surface")
    proof_type = proof.get("type", "text_evidence")
    _text(slide, proof_type.replace("_", " ").title(), x + 0.25, y + 0.25, 2.8, 0.28, 10.5, palette["accent"], font_name, bold=True)
    if proof_type == "figure" and proof.get("path") and Path(proof.get("path", "")).exists():
        _add_picture_fit(slide, proof.get("path", ""), x + 0.28, y + 0.75, w - 0.56, h - 1.45)
        _text(slide, _limit(proof.get("caption", ""), 120), x + 0.28, y + h - 0.55, w - 0.56, 0.3, 9.5, palette["muted"], font_name)
    elif proof_type == "table" and proof.get("rows"):
        _native_table(slide, proof.get("rows", []), x + 0.24, y + 0.76, w - 0.48, h - 1.18, palette, font_name, compact=table_compact)
        _text(slide, _limit(proof.get("caption", ""), 90), x + 0.28, y + h - 0.36, w - 0.56, 0.24, 9.5, palette["muted"], font_name)
    elif proof_type == "metric_set":
        items = proof.get("items", []) or []
        for index, item in enumerate(items[:4]):
            _metric_card(slide, item, x + 0.3 + (index % 2) * ((w - 0.75) / 2), y + 0.78 + (index // 2) * 1.62, (w - 0.95) / 2, 1.28, palette, font_name, index + 1)
    else:
        text = _proof_text(proof)
        _text(slide, text, x + 0.3, y + 0.82, w - 0.6, h - 1.2, 13, palette["ink"], font_name)


def _metric_card(slide: Any, item: Dict[str, Any], x: float, y: float, w: float, h: float, palette: Dict[str, str], font_name: str, index: int) -> None:
    _panel(slide, x, y, w, h, palette, fill="surface")
    value = str(item.get("value", "") or f"{index:02d}")
    label = str(item.get("label", "") or "Metric")
    note = str(item.get("note", "") or item.get("source_slide", "") or "source metric")
    _text(slide, _limit(value, 18), x + 0.16, y + 0.15, w - 0.32, 0.38, 21, palette["accent"], font_name, bold=True)
    _text(slide, _limit(label, 46), x + 0.17, y + 0.58, w - 0.34, 0.34, 9.8, palette["ink"], font_name, bold=True)
    _text(slide, _limit(note, 54), x + 0.17, y + 1.0, w - 0.34, 0.24, 9.2, palette["muted"], font_name)


def _evidence_card(slide: Any, item: Dict[str, Any], x: float, y: float, w: float, h: float, palette: Dict[str, str], font_name: str, index: int) -> None:
    _panel(slide, x, y, w, h, palette, fill="surface")
    kind = str(item.get("kind", "") or item.get("type", "") or "evidence").title()
    label = item.get("id", "") or item.get("label", "") or f"Evidence {index}"
    caption = item.get("caption", "") or item.get("note", "") or item.get("value", "")
    _text(slide, kind, x + 0.2, y + 0.17, w - 0.4, 0.24, 9.5, palette["accent"], font_name, bold=True)
    _text(slide, _limit(label, 44), x + 0.2, y + 0.48, w - 0.4, 0.22, 10, palette["ink"], font_name, bold=True)
    _text(slide, _limit(caption, 104), x + 0.2, y + 0.78, w - 0.4, 0.48, 8.8, palette["muted"], font_name)


def _native_table(slide: Any, rows: List[List[str]], x: float, y: float, w: float, h: float, palette: Dict[str, str], font_name: str, *, compact: bool = False) -> None:
    visible = [list(row) for row in rows[:7 if compact else 6] if row]
    if not visible:
        _text(slide, "Table rows unavailable; see source evidence.", x, y, w, h, 11, palette["ink"], font_name)
        return
    col_count = min(max(len(row) for row in visible), 4)
    visible = [(row + [""] * col_count)[:col_count] for row in visible]
    table_shape = slide.shapes.add_table(len(visible), col_count, Inches(x), Inches(y), Inches(w), Inches(max(0.8, h)))
    table = table_shape.table
    font_size = 7.2 if compact or col_count >= 4 else 8.4
    for row_idx, row in enumerate(visible):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = _limit(value, 48)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.space_after = Pt(0)
                for run in paragraph.runs:
                    run.font.name = font_name
                    run.font.size = Pt(font_size)
                    run.font.bold = row_idx == 0
                    run.font.color.rgb = _rgb(palette["ink"])


def _add_picture_fit(slide: Any, path: str, x: float, y: float, w: float, h: float) -> None:
    try:
        width_in, height_in = _fit_picture_size(path, w, h)
        left = x + max(0.0, (w - width_in) / 2)
        top = y + max(0.0, (h - height_in) / 2)
        slide.shapes.add_picture(str(Path(path)), Inches(left), Inches(top), width=Inches(width_in), height=Inches(height_in))
    except Exception:
        return


def _cover_slide(inventory: Dict[str, Any], contract: Dict[str, Any]) -> Dict[str, Any]:
    highlights = []
    for item in inventory.get("paper_highlights", [])[:4]:
        highlights.append(str(item.get("body", "") or item.get("label", "")))
    if not highlights:
        highlights = [contract.get("deck_intent", {}).get("deck_thesis", "")]
    return {
        "slide_id": "seed_full_cover",
        "role": "cover",
        "title": inventory.get("paper", {}).get("title", ""),
        "claim": contract.get("deck_intent", {}).get("deck_thesis", ""),
        "highlights": highlights,
    }


def _agenda_slide(sections: "OrderedDict[str, List[Dict[str, Any]]]") -> Dict[str, Any]:
    agenda = []
    for section, slides in sections.items():
        sample = "; ".join(_limit_words(slide.get("title", ""), 4) for slide in slides[:1])
        agenda.append({"title": section, "slide_count": len(slides), "sample": sample})
    return {"slide_id": "seed_full_agenda", "role": "agenda", "title": "Route Map", "sections": agenda}


def _section_slide(index: int, section: str, items: List[Dict[str, Any]]) -> Dict[str, Any]:
    sample = " / ".join(_limit_words(item.get("title", ""), 5) for item in items[:3])
    return {
        "slide_id": f"seed_section_{index:02d}",
        "role": "section",
        "section_index": index,
        "title": section,
        "support": sample,
    }


def _content_slide(plan: Dict[str, Any], curated: Dict[str, Any], inventory: Dict[str, Any], section: str, index: int) -> Dict[str, Any]:
    proof = _proof_object(plan, curated, inventory)
    claim = curated.get("takeaway", "") or _first_point(curated, "claim") or plan.get("title", "")
    support = " ".join(point.get("detail", "") for point in (curated.get("points", []) or [])[:3])
    support = support or plan.get("content", "")
    return {
        "slide_id": plan.get("slide_id", f"seed_content_{index:02d}"),
        "role": "content",
        "section": section,
        "title": plan.get("title", "") or curated.get("title", ""),
        "claim": _fit_claim_for_proof(claim, proof),
        "support": _limit_words(support, 20),
        "proof_object": proof,
        "layout": _layout_for_proof(proof, index),
    }


def _metric_ledger_slide(inventory: Dict[str, Any]) -> Dict[str, Any]:
    return {
        "slide_id": "seed_metric_ledger",
        "role": "metric_ledger",
        "section": "Experiments",
        "title": "Evaluation Ledger",
        "claim": "Key scores summarize the safety and reliability picture.",
        "proof_object": {"type": "metric_set", "items": (inventory.get("metrics", []) or [])[:8]},
    }


def _evidence_inventory_slide(inventory: Dict[str, Any]) -> Dict[str, Any]:
    assets = inventory.get("assets", {}) or {}
    items = (assets.get("figures", []) or [])[:3] + (assets.get("tables", []) or [])[:3]
    return {
        "slide_id": "seed_evidence_inventory",
        "role": "evidence_inventory",
        "section": "Evidence",
        "title": "Evidence Inventory",
        "claim": "The deck remains grounded in source figures, tables, and metric references.",
        "proof_object": {"type": "evidence_set", "items": items},
    }


def _synthesis_slide(inventory: Dict[str, Any], index: int) -> Dict[str, Any]:
    item = (inventory.get("summary_items", []) or [{}])[min(index, max(0, len(inventory.get("summary_items", []) or []) - 1))]
    return {
        "slide_id": f"seed_synthesis_{index:02d}",
        "role": "content",
        "section": item.get("category", "Synthesis"),
        "title": item.get("title", "Synthesis"),
        "claim": _limit_words(item.get("title", "Synthesis"), 18),
        "support": _limit_words(item.get("text", ""), 52),
        "proof_object": {"type": "text_evidence", "caption": item.get("title", ""), "note": _limit_words(item.get("text", ""), 40)},
        "layout": "claim_evidence_split",
    }


def _closing_slide(inventory: Dict[str, Any], contract: Dict[str, Any]) -> Dict[str, Any]:
    return {
        "slide_id": "seed_full_closing",
        "role": "closing",
        "title": "Takeaway",
        "claim": "GPT-5 pairs capability gains with layered safety evaluation and operational controls.",
        "support": "This is a seed full-deck prototype. Promotion still requires human preference review and a content gate before the style becomes default.",
    }


def _section_groups(plan_slides: List[Dict[str, Any]]) -> "OrderedDict[str, List[Dict[str, Any]]]":
    groups: "OrderedDict[str, List[Dict[str, Any]]]" = OrderedDict()
    for slide in plan_slides:
        section = slide.get("section") or "Content"
        groups.setdefault(section, []).append(slide)
    return groups


def _proof_object(plan: Dict[str, Any], curated: Dict[str, Any], inventory: Dict[str, Any]) -> Dict[str, Any]:
    figure = _first_figure(curated, plan)
    if figure and _prefer_figure(plan, curated):
        return {
            "type": "figure",
            "id": figure.get("title", "") or figure.get("id", ""),
            "caption": figure.get("caption", "") or figure.get("placeholder_text", ""),
            "path": figure.get("path", ""),
        }
    table = _first_table(curated, plan)
    if table:
        return {
            "type": "table",
            "id": table.get("table_id", "") or table.get("title", ""),
            "caption": table.get("caption", "") or table.get("focus", "") or table.get("title", ""),
            "rows": table.get("rows", []) or [],
        }
    if figure:
        return {
            "type": "figure",
            "id": figure.get("title", "") or figure.get("id", ""),
            "caption": figure.get("caption", "") or figure.get("placeholder_text", ""),
            "path": figure.get("path", ""),
        }
    metrics = [item for item in inventory.get("metrics", []) or [] if item.get("source_slide") == plan.get("slide_id")]
    if metrics:
        return {"type": "metric_set", "items": metrics[:4]}
    evidence = _first_point(curated, "evidence") or plan.get("title", "")
    return {"type": "text_evidence", "caption": evidence, "note": _limit_words(plan.get("content", ""), 42)}


def _fit_picture_size(path: str, max_w: float, max_h: float) -> Tuple[float, float]:
    try:
        from PIL import Image

        with Image.open(path) as image:
            px_w, px_h = image.size
    except Exception:
        return max_w, max_h
    if px_w <= 0 or px_h <= 0:
        return max_w, max_h
    image_ratio = px_w / px_h
    box_ratio = max_w / max_h if max_h else image_ratio
    if image_ratio >= box_ratio:
        width = max_w
        height = max_w / image_ratio
    else:
        height = max_h
        width = max_h * image_ratio
    return max(0.1, width), max(0.1, height)


def _prefer_figure(plan: Dict[str, Any], curated: Dict[str, Any]) -> bool:
    title = f"{plan.get('title', '')} {curated.get('title', '')}".lower()
    figure_terms = [
        "prompt injection",
        "hallucination",
        "factuality",
        "deception",
        "reward",
        "evaluation result",
        "production metrics",
    ]
    return any(term in title for term in figure_terms)


def _first_table(curated: Dict[str, Any], plan: Dict[str, Any]) -> Dict[str, Any]:
    for table in curated.get("tables", []) or []:
        if table:
            return table
    for table in plan.get("tables", []) or []:
        if table:
            return {"table_id": table.get("table_id", ""), "caption": table.get("focus", ""), "rows": []}
    return {}


def _first_figure(curated: Dict[str, Any], plan: Dict[str, Any]) -> Dict[str, Any]:
    for figure in curated.get("figures", []) or []:
        if figure:
            return figure
    for figure in plan.get("figures", []) or []:
        if figure:
            return figure
    return {}


def _layout_for_proof(proof: Dict[str, Any], index: int) -> str:
    if proof.get("type") == "table":
        return "table_focus"
    if proof.get("type") == "figure":
        return "figure_focus"
    if proof.get("type") == "metric_set":
        return "claim_evidence_split"
    return "claim_evidence_split" if index % 3 else "figure_focus"


def _proof_text(proof: Dict[str, Any]) -> str:
    parts = [proof.get("caption", ""), proof.get("note", "")]
    return _limit(" ".join(str(part) for part in parts if part), 320)


def _first_point(curated: Dict[str, Any], field: str) -> str:
    points = curated.get("points", []) or []
    if points:
        return str(points[0].get(field, ""))
    return ""


def _has_picture(spec: Dict[str, Any]) -> bool:
    proof = spec.get("proof_object", {}) or {}
    return proof.get("type") == "figure" and bool(proof.get("path"))


def _split_title(title: Any) -> str:
    value = str(title or "")
    if ": " in value:
        return value.replace(": ", ":\n", 1)
    return value


def _fit_claim(text: Any) -> str:
    claim = _limit_words(text, 9)
    if len(claim) > 78:
        return _limit_words(text, 8)
    return claim


def _fit_claim_for_proof(text: Any, proof: Dict[str, Any]) -> str:
    if proof.get("type") == "figure":
        return _limit_words(text, 7)
    return _fit_claim(text)


def _limit_words(text: Any, limit: int) -> str:
    words = str(text or "").replace("\n", " ").split()
    if len(words) <= limit:
        return " ".join(words)
    return " ".join(words[:limit]).rstrip(" ,;:") + "..."


def main(argv: List[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Render a seed-template full-deck prototype PPTX.")
    parser.add_argument("--inventory", required=True, help="Path to content_inventory.json")
    parser.add_argument("--contract", required=True, help="Path to seed_template_contract.json")
    parser.add_argument("--package-dir", required=True, help="Path to seed_template_package directory")
    parser.add_argument("--output-dir", required=True, help="Directory for spec and trace outputs")
    parser.add_argument("--output-pptx", required=True, help="Output PPTX path")
    parser.add_argument("--promotion-gate", help="Optional promotion_gate.json")
    parser.add_argument("--allow-blocked-prototype", action="store_true", help="Allow review-only rendering when promotion gate blocks full-deck expansion")
    parser.add_argument("--slide-budget", type=int, default=24, help="Target slide count")
    parser.add_argument(
        "--style-variant",
        default=DEFAULT_STYLE_VARIANT,
        help="Seed style variant: editorial_data_report, blueprint_system_map, or dark_evidence_console",
    )
    args = parser.parse_args(argv)

    paths = write_seed_full_deck_artifacts(
        Path(args.inventory),
        Path(args.contract),
        Path(args.package_dir),
        Path(args.output_dir),
        Path(args.output_pptx),
        promotion_gate_path=Path(args.promotion_gate) if args.promotion_gate else None,
        allow_blocked_prototype=args.allow_blocked_prototype,
        slide_budget=args.slide_budget,
        style_variant=args.style_variant,
    )
    print(json.dumps(paths, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
