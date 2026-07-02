"""Convert editable PPTX files into DeckIR and universal scorecards."""
from __future__ import annotations

import argparse
import json
import math
import re
from collections import Counter
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple

from paper2slides.benchmark.nonvisual_audit import EMU_PER_INCH, inspect_pptx_nonvisual

from .content_alignment import align_deck_to_checkpoints
from .deck_ir import (
    DECK_IR_SCHEMA_VERSION,
    count_metric_like_text,
    quantized_layout_signature,
    score_deck_ir_v0,
    short_text,
    top_items,
    write_schema_bundle,
)


SAFE_MARGIN_IN = 0.18


def pptx_to_deck_ir(
    pptx_path: Path,
    *,
    label: Optional[str] = None,
    generator: str = "unknown",
    native_editability_expected: bool = True,
) -> Dict[str, Any]:
    """Parse a PPTX into DeckIR v1."""
    from pptx import Presentation

    pptx_path = Path(pptx_path).resolve()
    prs = Presentation(pptx_path)
    width_in = _emu_to_inches(prs.slide_width)
    height_in = _emu_to_inches(prs.slide_height)

    slides: List[Dict[str, Any]] = []
    palette = Counter()
    fonts = Counter()
    backgrounds = Counter()

    for slide_index, slide in enumerate(prs.slides, start=1):
        objects = []
        for object_index, shape in enumerate(_iter_shapes(slide.shapes), start=1):
            record = _shape_to_object_record(shape, object_index, width_in, height_in)
            objects.append(record)
            for color in [record.get("fill_color"), record.get("line_color"), record.get("font_color")]:
                if color:
                    palette[color] += 1
            for font in record.get("font_families", []) or []:
                fonts[font] += 1
            if record.get("is_full_background") and record.get("fill_color"):
                backgrounds[record["fill_color"]] += 1

        text_model = _slide_text_model(objects)
        slide_ir = {
            "slide_index": slide_index,
            "role_guess": _guess_slide_role(slide_index, len(prs.slides), objects, text_model),
            "objects": objects,
            "text": text_model,
            "layout": _slide_layout_model(objects, width_in, height_in),
            "editability": _slide_editability_model(objects, width_in, height_in),
        }
        slide_ir["layout"]["layout_signature"] = quantized_layout_signature(slide_ir)
        slides.append(slide_ir)

    summary = _deck_summary(slides, width_in, height_in)
    return {
        "schema_version": DECK_IR_SCHEMA_VERSION,
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "source": {
            "path": str(pptx_path),
            "label": label or pptx_path.stem,
            "generator": generator,
            "artifact_kind": "pptx",
            "native_editability_expected": native_editability_expected,
        },
        "deck": {
            "slide_count": len(prs.slides),
            "width_in": round(width_in, 3),
            "height_in": round(height_in, 3),
            "theme_signals": {
                "palette": top_items(palette, 16),
                "font_families": top_items(fonts, 12),
                "dominant_backgrounds": top_items(backgrounds, 6),
            },
        },
        "slides": slides,
        "summary": summary,
    }


def write_pptx_intake_bundle(
    pptx_path: Path,
    output_dir: Path,
    *,
    label: Optional[str] = None,
    generator: str = "unknown",
    audit_path: Optional[Path] = None,
    repair_log_path: Optional[Path] = None,
    summary_checkpoint: Optional[Path] = None,
    plan_checkpoint: Optional[Path] = None,
    spec_checkpoint: Optional[Path] = None,
    run_nonvisual_audit: bool = True,
) -> Dict[str, str]:
    """Write deck_ir.json, scorecard, and supporting intake artifacts."""
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    deck_ir = pptx_to_deck_ir(pptx_path, label=label, generator=generator)

    audit = _load_json(audit_path) if audit_path else None
    if audit is None and run_nonvisual_audit:
        audit = inspect_pptx_nonvisual(Path(pptx_path))
        (output_dir / "nonvisual_audit.json").write_text(json.dumps(audit, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    elif audit is not None:
        (output_dir / "nonvisual_audit.json").write_text(json.dumps(audit, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    repair_log = _load_json(repair_log_path) if repair_log_path else None
    if repair_log is not None:
        (output_dir / "repair_log.snapshot.json").write_text(json.dumps(repair_log, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    content_alignment = align_deck_to_checkpoints(
        deck_ir,
        summary_checkpoint=summary_checkpoint,
        plan_checkpoint=plan_checkpoint,
        spec_checkpoint=spec_checkpoint,
    )
    if content_alignment is not None:
        (output_dir / "checkpoint_alignment.v0.json").write_text(
            json.dumps(content_alignment, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
        )

    scorecard = score_deck_ir_v0(deck_ir, nonvisual_audit=audit, repair_log=repair_log, content_alignment=content_alignment)

    deck_path = output_dir / "deck_ir.json"
    scorecard_path = output_dir / "universal_scorecard.v0.json"
    deck_path.write_text(json.dumps(deck_ir, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    scorecard_path.write_text(json.dumps(scorecard, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths = {
        "deck_ir": str(deck_path),
        "universal_scorecard": str(scorecard_path),
    }
    if audit is not None:
        paths["nonvisual_audit"] = str(output_dir / "nonvisual_audit.json")
    if repair_log is not None:
        paths["repair_log_snapshot"] = str(output_dir / "repair_log.snapshot.json")
    if content_alignment is not None:
        paths["checkpoint_alignment"] = str(output_dir / "checkpoint_alignment.v0.json")
    return paths


def _iter_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        child_shapes = getattr(shape, "shapes", None)
        if child_shapes is not None:
            for child in _iter_shapes(child_shapes):
                yield child
        else:
            yield shape


def _shape_to_object_record(shape: Any, object_index: int, slide_w: float, slide_h: float) -> Dict[str, Any]:
    left = _emu_to_inches(getattr(shape, "left", 0))
    top = _emu_to_inches(getattr(shape, "top", 0))
    width = _emu_to_inches(getattr(shape, "width", 0))
    height = _emu_to_inches(getattr(shape, "height", 0))
    area = max(0.0, width * height)
    text = _clean_text(getattr(shape, "text", "") if getattr(shape, "has_text_frame", False) else "")
    font_sizes, font_families, font_colors = _font_signals(shape)
    fill_color = _shape_fill_color(shape)
    line_color = _shape_line_color(shape)
    kind = _shape_kind(shape, text)
    picture = _picture_info(shape, width, height) if kind == "picture" else {}
    table = _table_info(shape) if getattr(shape, "has_table", False) else {"has_table": False}
    chart = _chart_info(shape) if getattr(shape, "has_chart", False) else {"has_chart": False}
    is_full_background = area >= slide_w * slide_h * 0.86 and not text and not table.get("has_table")
    if is_full_background:
        kind = "background"
    return {
        "object_id": f"obj_{object_index:03d}",
        "name": str(getattr(shape, "name", "") or ""),
        "kind": kind,
        "shape_type": str(getattr(shape, "shape_type", "")),
        "bbox": {
            "x": round(left, 3),
            "y": round(top, 3),
            "w": round(width, 3),
            "h": round(height, 3),
            "right": round(left + width, 3),
            "bottom": round(top + height, 3),
            "area": round(area, 3),
        },
        "text": short_text(text, 220),
        "text_chars": len(text),
        "text_words": len(text.split()),
        "font": {
            "sizes_pt": font_sizes[:8],
            "min_pt": min(font_sizes) if font_sizes else None,
            "avg_pt": round(sum(font_sizes) / len(font_sizes), 2) if font_sizes else None,
            "max_pt": max(font_sizes) if font_sizes else None,
        },
        "font_families": sorted(font_families),
        "font_color": top_items(Counter(font_colors), 1)[0] if font_colors else None,
        "fill_color": fill_color,
        "line_color": line_color,
        "has_table": bool(table.get("has_table")),
        "table": table,
        "has_chart": bool(chart.get("has_chart")),
        "chart": chart,
        "is_picture": kind == "picture",
        "picture": picture,
        "is_full_background": is_full_background,
        "native_editable": kind in {"text", "shape", "table", "chart"},
    }


def _shape_kind(shape: Any, text: str) -> str:
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "has_chart", False):
        return "chart"
    shape_type = str(getattr(shape, "shape_type", ""))
    if shape_type == "PICTURE (13)" or shape_type == "13":
        return "picture"
    if text:
        return "text"
    if _emu_to_inches(getattr(shape, "width", 0)) < 0.08 or _emu_to_inches(getattr(shape, "height", 0)) < 0.045:
        return "decorative"
    return "shape"


def _slide_text_model(objects: List[Dict[str, Any]]) -> Dict[str, Any]:
    text_objects = [obj for obj in objects if obj.get("text")]
    title_candidates = []
    body_blocks = []
    caption_candidates = []
    for obj in text_objects:
        text = obj.get("text", "")
        bbox = obj.get("bbox", {})
        avg_pt = obj.get("font", {}).get("avg_pt") or 0.0
        token = text.lower()
        if bbox.get("y", 99) < 1.25 or avg_pt >= 22.0:
            title_candidates.append(text)
        elif _looks_caption(token, avg_pt, bbox):
            caption_candidates.append(text)
        else:
            body_blocks.append(text)
    return {
        "title_candidates": title_candidates[:5],
        "body_blocks": body_blocks[:18],
        "caption_candidates": caption_candidates[:12],
    }


def _slide_layout_model(objects: List[Dict[str, Any]], width_in: float, height_in: float) -> Dict[str, Any]:
    meaningful = [obj for obj in objects if _meaningful_object(obj)]
    occupied = sum(float(obj.get("bbox", {}).get("area", 0.0) or 0.0) for obj in meaningful)
    safe_area_violations = []
    for obj in meaningful:
        bbox = obj.get("bbox", {})
        if (
            bbox.get("x", 0.0) < -0.02
            or bbox.get("y", 0.0) < -0.02
            or bbox.get("right", 0.0) > width_in + 0.02
            or bbox.get("bottom", 0.0) > height_in + 0.02
            or bbox.get("x", 0.0) < SAFE_MARGIN_IN and bbox.get("area", 0.0) > 0.08
            or bbox.get("right", 0.0) > width_in - SAFE_MARGIN_IN and bbox.get("area", 0.0) > 0.08
        ):
            safe_area_violations.append(obj.get("object_id", ""))
    return {
        "occupancy": round(occupied / max(0.01, width_in * height_in), 3),
        "alignment_groups": _alignment_groups(meaningful),
        "safe_area_violations": safe_area_violations,
    }


def _slide_editability_model(objects: List[Dict[str, Any]], width_in: float, height_in: float) -> Dict[str, Any]:
    slide_area = max(0.01, width_in * height_in)
    picture_area = sum(float(obj.get("bbox", {}).get("area", 0.0) or 0.0) for obj in objects if obj.get("kind") == "picture")
    return {
        "text_chars_native": sum(int(obj.get("text_chars", 0) or 0) for obj in objects if obj.get("kind") in {"text", "table", "chart"}),
        "raster_area_ratio": round(picture_area / slide_area, 3),
        "native_shape_count": sum(1 for obj in objects if obj.get("kind") in {"text", "shape", "decorative"}),
        "picture_count": sum(1 for obj in objects if obj.get("kind") == "picture"),
        "table_count": sum(1 for obj in objects if obj.get("kind") == "table"),
        "chart_count": sum(1 for obj in objects if obj.get("kind") == "chart"),
    }


def _deck_summary(slides: List[Dict[str, Any]], width_in: float, height_in: float) -> Dict[str, Any]:
    role_counts = Counter(slide.get("role_guess", "unknown") for slide in slides)
    object_count = sum(len(slide.get("objects", []) or []) for slide in slides)
    native_text_chars = sum(slide.get("editability", {}).get("text_chars_native", 0) for slide in slides)
    text_snippets = []
    caption_count = 0
    source_like = 0
    for slide in slides:
        text = slide.get("text", {})
        snippets = list(text.get("title_candidates", []) or []) + list(text.get("body_blocks", []) or []) + list(text.get("caption_candidates", []) or [])
        text_snippets.extend(snippets)
        caption_count += len(text.get("caption_candidates", []) or [])
        source_like += sum(1 for snippet in snippets if _looks_source_like(snippet))
    picture_area_ratio = sum(slide.get("editability", {}).get("raster_area_ratio", 0.0) for slide in slides) / max(1, len(slides))
    return {
        "object_count": object_count,
        "native_text_chars": native_text_chars,
        "native_text_words": sum(len(str(snippet).split()) for snippet in text_snippets),
        "native_shape_count": sum(slide.get("editability", {}).get("native_shape_count", 0) for slide in slides),
        "picture_count": sum(slide.get("editability", {}).get("picture_count", 0) for slide in slides),
        "table_count": sum(slide.get("editability", {}).get("table_count", 0) for slide in slides),
        "chart_count": sum(slide.get("editability", {}).get("chart_count", 0) for slide in slides),
        "raster_area_ratio": round(picture_area_ratio, 3),
        "role_counts": dict(role_counts),
        "slides_with_title": sum(1 for slide in slides if slide.get("text", {}).get("title_candidates")),
        "caption_candidate_count": caption_count,
        "source_like_text_count": source_like,
        "metric_like_text_count": count_metric_like_text(text_snippets),
        "layout_signatures": [slide.get("layout", {}).get("layout_signature", "") for slide in slides],
        "canvas_area_in": round(width_in * height_in, 3),
    }


def _guess_slide_role(slide_index: int, slide_count: int, objects: List[Dict[str, Any]], text_model: Dict[str, Any]) -> str:
    joined = " ".join(
        list(text_model.get("title_candidates", []) or [])
        + list(text_model.get("body_blocks", []) or [])
        + list(text_model.get("caption_candidates", []) or [])
    ).lower()
    if slide_index == 1:
        return "cover"
    if slide_index == slide_count and any(token in joined for token in ["takeaway", "closing", "conclusion", "next", "summary"]):
        return "closing"
    if any(token in joined for token in ["agenda", "roadmap", "deck map", "contents"]):
        return "agenda"
    if any(token in joined for token in ["section", "part ", "chapter"]):
        return "section"
    if any(obj.get("kind") == "table" for obj in objects) or any(obj.get("kind") == "picture" for obj in objects):
        return "evidence"
    if count_metric_like_text([joined]) >= 2 or any(token in joined for token in ["%", "metric", "score", "rate", "ratio"]):
        return "metric"
    if joined:
        return "content"
    return "unknown"


def _alignment_groups(objects: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    buckets: Dict[Tuple[str, float], List[str]] = {}
    for obj in objects:
        bbox = obj.get("bbox", {})
        for edge in ("x", "y", "right", "bottom"):
            key = (edge, round(float(bbox.get(edge, 0.0)) / 0.08) * 0.08)
            buckets.setdefault(key, []).append(obj.get("object_id", ""))
    groups = []
    for (edge, value), object_ids in sorted(buckets.items(), key=lambda item: (-len(item[1]), item[0][0], item[0][1])):
        if len(object_ids) >= 3:
            groups.append({"edge": edge, "value_in": round(value, 3), "object_ids": object_ids[:10], "count": len(object_ids)})
        if len(groups) >= 12:
            break
    return groups


def _font_signals(shape: Any) -> Tuple[List[float], set, List[str]]:
    sizes: List[float] = []
    families = set()
    colors: List[str] = []
    if not getattr(shape, "has_text_frame", False):
        return sizes, families, colors
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            font = run.font
            if font.size is not None:
                sizes.append(round(float(font.size.pt), 2))
            if font.name:
                families.add(str(font.name))
            color = _font_color(font)
            if color:
                colors.append(color)
    return sizes, families, colors


def _shape_fill_color(shape: Any) -> Optional[str]:
    try:
        fill = shape.fill
        if not fill:
            return None
        return _rgb_to_hex(fill.fore_color.rgb)
    except Exception:
        return None


def _shape_line_color(shape: Any) -> Optional[str]:
    try:
        return _rgb_to_hex(shape.line.color.rgb)
    except Exception:
        return None


def _font_color(font: Any) -> Optional[str]:
    try:
        return _rgb_to_hex(font.color.rgb)
    except Exception:
        return None


def _rgb_to_hex(rgb: Any) -> Optional[str]:
    if rgb is None:
        return None
    text = str(rgb)
    if not text:
        return None
    if len(text) == 6 and all(char in "0123456789ABCDEFabcdef" for char in text):
        return f"#{text.upper()}"
    return None


def _picture_info(shape: Any, width: float, height: float) -> Dict[str, Any]:
    try:
        image_w, image_h = shape.image.size
        source_aspect = image_w / image_h if image_h else None
        image_ext = shape.image.ext
    except Exception:
        image_w = image_h = None
        source_aspect = None
        image_ext = None
    box_aspect = width / height if height else None
    distortion = None
    if source_aspect and box_aspect:
        distortion = abs(math.log(max(source_aspect, box_aspect) / max(0.001, min(source_aspect, box_aspect))))
    return {
        "image_px": {"w": image_w, "h": image_h},
        "image_ext": image_ext,
        "source_aspect": round(source_aspect, 3) if source_aspect else None,
        "box_aspect": round(box_aspect, 3) if box_aspect else None,
        "aspect_distortion": round(distortion, 3) if distortion is not None else None,
    }


def _table_info(shape: Any) -> Dict[str, Any]:
    try:
        table = shape.table
        rows = []
        for row in table.rows:
            rows.append([_clean_text(cell.text) for cell in row.cells])
        return {
            "has_table": True,
            "row_count": len(table.rows),
            "column_count": len(table.columns),
            "sample_rows": rows[:4],
        }
    except Exception:
        return {"has_table": False}


def _chart_info(shape: Any) -> Dict[str, Any]:
    try:
        chart = shape.chart
        return {"has_chart": True, "chart_type": str(chart.chart_type)}
    except Exception:
        return {"has_chart": False}


def _meaningful_object(obj: Dict[str, Any]) -> bool:
    if obj.get("is_full_background"):
        return False
    if obj.get("kind") in {"decorative", "background"}:
        return False
    return bool(obj.get("text") or obj.get("has_table") or obj.get("has_chart") or obj.get("is_picture") or obj.get("kind") == "shape")


def _looks_caption(token: str, avg_pt: float, bbox: Dict[str, Any]) -> bool:
    if avg_pt and avg_pt <= 10.5:
        return True
    return any(word in token for word in ["figure", "table", "source", "caption", "note:"]) and bbox.get("h", 0.0) <= 0.55


def _looks_source_like(text: str) -> bool:
    token = str(text or "").lower()
    return any(word in token for word in ["source", "figure", "table", "appendix", "paper", "system card"])


def _load_json(path: Optional[Path]) -> Optional[Dict[str, Any]]:
    if not path:
        return None
    path = Path(path)
    if not path.exists():
        return None
    return json.loads(path.read_text(encoding="utf-8"))


def _clean_text(text: Any) -> str:
    return re.sub(r"\s+", " ", str(text or "")).strip()


def _emu_to_inches(value: Any) -> float:
    try:
        return float(value) / EMU_PER_INCH
    except Exception:
        return 0.0


def main(argv: Optional[List[str]] = None) -> int:
    parser = argparse.ArgumentParser(description="Convert PPTX into universal DeckIR and scorecard v0.")
    parser.add_argument("--pptx", required=True, help="Input PPTX path.")
    parser.add_argument("--output-dir", required=True, help="Directory to write deck_ir.json and scorecard.")
    parser.add_argument("--label", help="Human readable deck label.")
    parser.add_argument("--generator", default="unknown", help="Generator/source label.")
    parser.add_argument("--audit", help="Existing nonvisual_audit.json to reuse.")
    parser.add_argument("--repair-log", help="Existing repair_log.json to reuse.")
    parser.add_argument("--summary-checkpoint", help="Optional checkpoint_summary.json for content alignment.")
    parser.add_argument("--plan-checkpoint", help="Optional checkpoint_plan.json for content alignment.")
    parser.add_argument("--spec-checkpoint", help="Optional checkpoint_slide_spec.json for content alignment.")
    parser.add_argument("--no-audit", action="store_true", help="Do not run nonvisual audit when --audit is omitted.")
    parser.add_argument("--write-schemas", action="store_true", help="Also write DeckIR and scorecard schema JSON files.")
    args = parser.parse_args(argv)

    output_dir = Path(args.output_dir)
    paths = write_pptx_intake_bundle(
        Path(args.pptx),
        output_dir,
        label=args.label,
        generator=args.generator,
        audit_path=Path(args.audit) if args.audit else None,
        repair_log_path=Path(args.repair_log) if args.repair_log else None,
        summary_checkpoint=Path(args.summary_checkpoint) if args.summary_checkpoint else None,
        plan_checkpoint=Path(args.plan_checkpoint) if args.plan_checkpoint else None,
        spec_checkpoint=Path(args.spec_checkpoint) if args.spec_checkpoint else None,
        run_nonvisual_audit=not args.no_audit,
    )
    if args.write_schemas:
        paths.update(write_schema_bundle(output_dir))
    print(json.dumps(paths, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
