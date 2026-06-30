"""Non-visual PPTX geometry and text-capacity audit.

This module inspects PPTX metadata directly. It does not render slides, take
screenshots, or call a vision model.
"""
from __future__ import annotations

import argparse
import json
import math
import re
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple


EMU_PER_INCH = 914400

RULE_SCHEMA_VERSION = "benchmark_rule_schema.v3"

BLIND_RECTANGULAR_STYLE_SCOPES = [
    "blind_rectangular_research_board",
    "golden_baseline2_blind_rectangular_research_board",
]

DIMENSIONS = [
    "content",
    "evidence",
    "layout",
    "typography",
    "component_fit",
    "style",
    "repair_risk",
]

DEFAULT_RULE_METADATA = {
    "dimension": "layout",
    "scope": "global",
    "style_scope": [],
    "repair_mode": "suggest",
    "confidence": 0.55,
    "human_outcome": "pending_review",
}

RULE_METADATA: Dict[str, Dict[str, Any]] = {
    "low_font_size": {"dimension": "typography", "scope": "global", "repair_mode": "auto", "confidence": 0.86},
    "below_ideal_font_band": {"dimension": "typography", "scope": "global", "repair_mode": "suggest", "confidence": 0.72},
    "cover_left_typography_underpowered": {"dimension": "typography", "scope": "style_aware", "style_scope": ["golden_baseline1_from_scratch_warm_academic"], "repair_mode": "human_gated", "confidence": 0.68},
    "estimated_text_overflow": {"dimension": "layout", "scope": "global", "repair_mode": "auto", "confidence": 0.88},
    "text_exceeds_container_bounds": {"dimension": "layout", "scope": "global", "repair_mode": "auto", "confidence": 0.86},
    "near_text_capacity": {"dimension": "typography", "scope": "global", "repair_mode": "suggest", "confidence": 0.72},
    "low_text_density": {"dimension": "typography", "scope": "style_aware", "repair_mode": "suggest", "confidence": 0.58},
    "sparse_slide_risk": {"dimension": "typography", "scope": "style_aware", "repair_mode": "suggest", "confidence": 0.62},
    "card_copy_imbalance": {"dimension": "typography", "scope": "style_aware", "repair_mode": "suggest", "confidence": 0.68},
    "shape_overlap_risk": {"dimension": "layout", "scope": "global", "repair_mode": "auto", "confidence": 0.9},
    "component_boundary_inset_violation": {"dimension": "layout", "scope": "global", "repair_mode": "auto", "confidence": 0.82},
    "weak_table_grammar": {"dimension": "component_fit", "scope": "global", "repair_mode": "auto", "confidence": 0.84},
    "dense_table_readability_risk": {"dimension": "component_fit", "scope": "global", "repair_mode": "suggest", "confidence": 0.76},
    "table_exceeds_container_bounds": {"dimension": "layout", "scope": "global", "repair_mode": "auto", "confidence": 0.91},
    "table_container_height_mismatch": {"dimension": "component_fit", "scope": "global", "repair_mode": "auto_then_human_review", "confidence": 0.74},
    "table_readability_after_fit": {"dimension": "component_fit", "scope": "global", "repair_mode": "human_gated", "confidence": 0.62},
    "table_sparse_columns_rendered": {"dimension": "component_fit", "scope": "global", "repair_mode": "auto_then_human_review", "confidence": 0.78},
    "table_cell_text_wrapping_risk": {"dimension": "component_fit", "scope": "global", "repair_mode": "auto_then_human_review", "confidence": 0.78},
    "table_view_label_missing": {"dimension": "style", "scope": "style_aware", "style_scope": BLIND_RECTANGULAR_STYLE_SCOPES, "repair_mode": "auto", "confidence": 0.74},
    "table_caption_missing_or_not_centered": {"dimension": "evidence", "scope": "style_aware", "style_scope": BLIND_RECTANGULAR_STYLE_SCOPES, "repair_mode": "auto", "confidence": 0.78},
    "table_underutilized_in_evidence_panel": {"dimension": "component_fit", "scope": "style_aware", "style_scope": BLIND_RECTANGULAR_STYLE_SCOPES, "repair_mode": "auto_then_human_review", "confidence": 0.73},
    "figure_picture_aspect_distortion": {"dimension": "component_fit", "scope": "global", "repair_mode": "auto", "confidence": 0.88},
    "figure_panel_aspect_mismatch": {"dimension": "component_fit", "scope": "style_aware", "style_scope": ["golden_baseline1_from_scratch_warm_academic"], "repair_mode": "human_gated", "confidence": 0.72},
    "image_underutilized_in_wide_panel": {"dimension": "component_fit", "scope": "style_aware", "style_scope": BLIND_RECTANGULAR_STYLE_SCOPES, "repair_mode": "auto_then_human_review", "confidence": 0.7},
    "figure_caption_not_centered_in_wide_panel": {"dimension": "style", "scope": "style_aware", "style_scope": BLIND_RECTANGULAR_STYLE_SCOPES, "repair_mode": "auto", "confidence": 0.76},
    "figure_image_off_center_in_panel": {"dimension": "style", "scope": "style_aware", "style_scope": ["golden_baseline1_from_scratch_warm_academic"], "repair_mode": "human_gated", "confidence": 0.7},
    "figure_badge_identity_label_conflation": {"dimension": "style", "scope": "style_aware", "style_scope": ["golden_baseline1_from_scratch_warm_academic"], "repair_mode": "human_gated", "confidence": 0.72},
    "figure_label_anchor_drift": {"dimension": "style", "scope": "style_aware", "style_scope": ["golden_baseline1_from_scratch_warm_academic"], "repair_mode": "human_gated", "confidence": 0.72},
    "figure_label_text_alignment_off_center": {"dimension": "style", "scope": "style_aware", "style_scope": ["golden_baseline1_from_scratch_warm_academic"], "repair_mode": "human_gated", "confidence": 0.74},
    "panel_identity_label_anchor_drift": {"dimension": "style", "scope": "style_aware", "style_scope": ["golden_baseline1_from_scratch_warm_academic"], "repair_mode": "human_gated", "confidence": 0.74},
    "panel_identity_label_text_alignment_off_center": {"dimension": "style", "scope": "style_aware", "style_scope": ["golden_baseline1_from_scratch_warm_academic"], "repair_mode": "human_gated", "confidence": 0.74},
    "stacked_figure_identity_label_overcorrection": {"dimension": "repair_risk", "scope": "style_aware", "style_scope": ["golden_baseline1_from_scratch_warm_academic"], "repair_mode": "detect_only", "confidence": 0.66},
    "metric_label_gap_too_large": {"dimension": "style", "scope": "style_aware", "repair_mode": "suggest", "confidence": 0.64},
    "container_stack_off_balance": {"dimension": "style", "scope": "style_aware", "repair_mode": "suggest", "confidence": 0.64},
    "paired_label_body_gap_too_large": {"dimension": "style", "scope": "style_aware", "repair_mode": "suggest", "confidence": 0.64},
    "component_frame_overallocated_after_text_fit": {"dimension": "component_fit", "scope": "style_aware", "style_scope": ["golden_baseline1_from_scratch_warm_academic"], "repair_mode": "human_gated", "confidence": 0.66},
    "card_internal_spacing_not_scaled_to_frame": {"dimension": "style", "scope": "style_aware", "style_scope": ["golden_baseline1_from_scratch_warm_academic"], "repair_mode": "human_gated", "confidence": 0.72},
    "text_card_vertical_alignment_top_heavy": {"dimension": "style", "scope": "style_aware", "style_scope": BLIND_RECTANGULAR_STYLE_SCOPES, "repair_mode": "auto", "confidence": 0.74},
    "agenda_read_path_header_too_close": {"dimension": "style", "scope": "style_aware", "style_scope": ["golden_baseline1_from_scratch_warm_academic"], "repair_mode": "human_gated", "confidence": 0.72},
    "table_support_band_off_balance": {"dimension": "style", "scope": "style_aware", "style_scope": ["golden_baseline1_from_scratch_warm_academic"], "repair_mode": "human_gated", "confidence": 0.7},
    "academic_right_evidence_void": {"dimension": "style", "scope": "style_aware", "style_scope": ["academic"], "repair_mode": "human_gated", "confidence": 0.76},
    "academic_toc_missing_canonical_sections": {"dimension": "content", "scope": "style_aware", "style_scope": ["academic"], "repair_mode": "human_gated", "confidence": 0.78},
    "weak_fragment_point_heading": {"dimension": "content", "scope": "global", "repair_mode": "auto", "confidence": 0.78},
    "spurious_generic_metric_card": {"dimension": "content", "scope": "global", "repair_mode": "auto", "confidence": 0.76},
    "metric_improved_visual_regressed": {"dimension": "repair_risk", "scope": "human_feedback", "repair_mode": "detect_only", "confidence": 0.64, "human_outcome": "tradeoff_review"},
    "likely_overcorrection": {"dimension": "repair_risk", "scope": "human_feedback", "repair_mode": "detect_only", "confidence": 0.62, "human_outcome": "likely_overcorrection"},
    "style_scope_mismatch": {"dimension": "repair_risk", "scope": "style_aware", "repair_mode": "detect_only", "confidence": 0.68},
    "repair_introduced_new_findings": {"dimension": "repair_risk", "scope": "human_feedback", "repair_mode": "detect_only", "confidence": 0.78},
    "image_legibility_regression": {"dimension": "repair_risk", "scope": "human_feedback", "repair_mode": "detect_only", "confidence": 0.62},
    "layout_rhythm_regression": {"dimension": "repair_risk", "scope": "human_feedback", "repair_mode": "detect_only", "confidence": 0.6},
}

NONVISUAL_AUDIT_RULES = {
    "title_claim_min_pt": 20.0,
    "support_min_pt": 11.0,
    "card_text_min_pt": 9.0,
    "body_min_pt": 10.5,
    "title_claim_ideal_min_pt": 24.0,
    "support_ideal_min_pt": 12.5,
    "card_text_ideal_min_pt": 9.8,
    "body_ideal_min_pt": 12.0,
    "near_capacity_ratio": 0.82,
    "overflow_ratio": 1.03,
    "sparse_ratio": 0.16,
    "sparse_area_sq_in": 1.15,
    "overlap_ratio": 0.14,
    "container_overlap_ratio": 0.18,
    "low_occupancy_ratio": 0.18,
    "card_copy_imbalance_ratio": 2.6,
    "cover_left_min_pt": 14.0,
    "deck_typography_signal_count": 8,
    "flow_label_min_gap_in": 0.10,
    "flow_grid_max_col_gap_in": 1.75,
    "flow_grid_min_row_gap_in": 0.62,
    "flow_grid_col_to_row_ratio_max": 2.7,
    "flow_label_center_tolerance_in": 0.09,
    "agenda_read_path_header_min_gap_in": 0.22,
    "table_support_claim_gap_max_in": 0.24,
    "table_support_panel_gap_min_in": 0.28,
    "paired_text_top_delta_max_in": 0.26,
    "component_label_min_inset_in": 0.30,
    "component_frame_extra_height_max_in": 0.42,
    "card_internal_gap_max_in": 0.055,
    "card_internal_min_bottom_padding_in": 0.09,
    "card_internal_shallow_max_h_in": 1.15,
    "card_internal_narrow_max_w_in": 3.8,
    "metric_label_gap_max_in": 0.065,
    "metric_card_min_area_sq_in": 0.40,
    "balance_container_min_area_sq_in": 3.6,
    "balance_padding_tolerance_in": 0.38,
    "balance_padding_tolerance_ratio": 0.11,
    "picture_aspect_ratio_tolerance": 0.18,
    "tall_figure_aspect_max": 0.9,
    "tall_figure_panel_aspect_max": 0.95,
    "wide_figure_aspect_min": 1.9,
    "wide_figure_panel_aspect_min": 2.0,
    "figure_image_center_tolerance_in": 0.35,
    "figure_image_center_tolerance_ratio": 0.08,
    "figure_label_semantic_wide_panel_aspect_min": 1.6,
    "figure_label_semantic_tall_panel_aspect_max": 0.98,
    "figure_panel_badge_left_inset_max_in": 0.78,
    "figure_panel_badge_top_inset_min_in": 0.22,
    "figure_panel_badge_top_inset_max_in": 0.55,
    "figure_identity_label_center_tolerance_in": 0.35,
    "figure_identity_label_vertical_gap_min_in": -0.08,
    "figure_identity_label_vertical_gap_max_in": 0.18,
    "panel_identity_label_center_tolerance_in": 0.35,
    "table_container_padding_in": 0.08,
    "table_container_overflow_tolerance_in": 0.03,
    "table_container_height_fill_max": 0.96,
    "table_readable_row_height_min_in": 0.18,
    "table_readable_col_width_min_in": 0.45,
    "table_sparse_column_count_min": 8,
    "table_sparse_empty_cell_ratio": 0.34,
    "table_sparse_min_non_empty_cells": 8,
    "table_wrapping_col_width_max_in": 0.56,
    "table_wrapping_effective_units_per_in": 18.0,
    "table_underutilized_area_ratio": 0.42,
    "table_underutilized_margin_in": 0.55,
    "table_caption_center_tolerance_in": 0.25,
    "text_container_padding_in": 0.08,
    "text_container_overflow_tolerance_in": 0.04,
    "text_card_vertical_slack_ratio": 1.45,
    "text_card_vertical_extra_min_in": 0.18,
    "wide_panel_aspect_min": 2.3,
    "wide_panel_min_area_sq_in": 4.5,
    "image_underutilized_area_ratio": 0.28,
    "image_underutilized_margin_in": 0.42,
    "figure_caption_center_tolerance_in": 0.25,
}


def inspect_pptx_nonvisual(pptx_path: Path, rules: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    """Inspect a PPTX without rendering pages."""
    from pptx import Presentation

    active_rules = dict(NONVISUAL_AUDIT_RULES)
    if rules:
        active_rules.update(rules)

    prs = Presentation(pptx_path)
    slide_w = _emu_to_inches(prs.slide_width)
    slide_h = _emu_to_inches(prs.slide_height)
    findings: List[Dict[str, Any]] = []
    slide_summaries: List[Dict[str, Any]] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        records = [_shape_record(shape, idx, slide_w, slide_h) for idx, shape in enumerate(slide.shapes)]
        for record in records:
            record["role"] = _infer_role(record, slide_h)

        slide_findings: List[Dict[str, Any]] = []
        slide_findings.extend(_text_capacity_findings(slide_index, records, active_rules))
        slide_findings.extend(_cover_typography_findings(slide_index, records, active_rules))
        slide_findings.extend(_paired_text_stack_findings(slide_index, records, active_rules))
        slide_findings.extend(_copy_distribution_findings(slide_index, records, active_rules))
        slide_findings.extend(_flow_layout_findings(slide_index, records, active_rules))
        slide_findings.extend(_agenda_read_path_header_findings(slide_index, records, active_rules))
        slide_findings.extend(_table_support_band_findings(slide_index, records, active_rules))
        slide_findings.extend(_metric_stack_findings(slide_index, records, active_rules))
        slide_findings.extend(_component_boundary_findings(slide_index, records, active_rules))
        slide_findings.extend(_text_container_bounds_findings(slide_index, records, active_rules))
        slide_findings.extend(_text_card_vertical_alignment_findings(slide_index, records, active_rules))
        slide_findings.extend(_academic_right_evidence_void_findings(slide_index, records))
        slide_findings.extend(_academic_toc_canonical_sections_findings(slide_index, records))
        slide_findings.extend(_semantic_content_findings(slide_index, records))
        slide_findings.extend(_figure_label_semantics_findings(slide_index, records, active_rules))
        slide_findings.extend(_panel_identity_label_findings(slide_index, records, active_rules))
        slide_findings.extend(_component_frame_fit_findings(slide_index, records, active_rules))
        slide_findings.extend(_card_internal_spacing_findings(slide_index, records, active_rules))
        slide_findings.extend(_container_balance_findings(slide_index, records, active_rules))
        slide_findings.extend(_overlap_findings(slide_index, records, active_rules))
        slide_findings.extend(_table_findings(slide_index, records, active_rules))
        slide_findings.extend(_picture_findings(slide_index, records, active_rules))
        slide_findings.extend(_slide_density_findings(slide_index, records, slide_w, slide_h, active_rules))
        findings.extend(slide_findings)
        slide_summaries.append(_slide_summary(slide_index, records, slide_findings, slide_w, slide_h))

    return {
        "schema_version": "pptx_nonvisual_audit.v1",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "pptx_path": str(pptx_path),
        "review_mode": "non_visual_metadata_only",
        "rendering_used": False,
        "vision_model_used": False,
        "rule_schema_version": RULE_SCHEMA_VERSION,
        "slide_count": len(prs.slides),
        "slide_size": {"width_in": round(slide_w, 3), "height_in": round(slide_h, 3)},
        "rules": active_rules,
        "rule_registry": _rule_registry_snapshot(findings),
        "summary": _summarize_findings(findings, active_rules),
        "findings": findings,
        "slides": slide_summaries,
        "repair_policy": {
            "preserve_component_composition": True,
            "default_text_repair": "adjust typography, copy allocation, or line breaks before resizing visual components",
            "avoid": [
                "auto-shrinking accepted v5 components only because text density is low",
                "moving cards to a new macro-layout without accepted-reference comparison",
                "rendering every slide for visual judging",
            ],
        },
    }


def _shape_record(shape: Any, index: int, slide_w: float, slide_h: float) -> Dict[str, Any]:
    text = _clean_text(getattr(shape, "text", "") if getattr(shape, "has_text_frame", False) else "")
    table_info = _table_info(shape)
    left = _emu_to_inches(getattr(shape, "left", 0))
    top = _emu_to_inches(getattr(shape, "top", 0))
    width = _emu_to_inches(getattr(shape, "width", 0))
    height = _emu_to_inches(getattr(shape, "height", 0))
    font_sizes = _font_sizes(shape)
    area = max(0.0, width * height)
    shape_type = str(getattr(shape, "shape_type", ""))
    is_picture = shape_type == "PICTURE (13)" or shape_type == "13"
    is_line_like = width < 0.08 or height < 0.045
    is_full_background = area >= slide_w * slide_h * 0.86 and not text and not table_info.get("has_table")
    picture_info = _picture_info(shape, width, height) if is_picture else {"is_picture": False}
    return {
        "index": index,
        "name": getattr(shape, "name", ""),
        "shape_type": shape_type,
        "has_text": bool(text),
        "text": text,
        "text_chars": len(text),
        "text_words": len(text.split()),
        "has_table": bool(table_info.get("has_table")),
        "table": table_info,
        "bbox": {
            "x": round(left, 3),
            "y": round(top, 3),
            "w": round(width, 3),
            "h": round(height, 3),
            "right": round(left + width, 3),
            "bottom": round(top + height, 3),
            "area": round(area, 3),
        },
        "font": {
            "sizes_pt": font_sizes,
            "min_pt": min(font_sizes) if font_sizes else None,
            "avg_pt": round(sum(font_sizes) / len(font_sizes), 2) if font_sizes else None,
        },
        "paragraph_alignment": _paragraph_alignment(shape),
        "vertical_anchor": _vertical_anchor(shape),
        "is_picture": is_picture,
        "picture": picture_info,
        "is_line_like": is_line_like,
        "is_full_background": is_full_background,
    }


def _picture_info(shape: Any, width: float, height: float) -> Dict[str, Any]:
    try:
        image = shape.image
        image_w, image_h = image.size
        source_aspect = image_w / image_h if image_h else None
    except Exception:
        image_w = image_h = None
        source_aspect = None
    box_aspect = width / height if height else None
    if source_aspect and box_aspect:
        distortion = abs(math.log(max(source_aspect, box_aspect) / max(0.001, min(source_aspect, box_aspect))))
    else:
        distortion = None
    return {
        "is_picture": True,
        "image_px": {"w": image_w, "h": image_h},
        "source_aspect": round(source_aspect, 3) if source_aspect else None,
        "box_aspect": round(box_aspect, 3) if box_aspect else None,
        "aspect_distortion": round(distortion, 3) if distortion is not None else None,
    }


def _table_info(shape: Any) -> Dict[str, Any]:
    if not getattr(shape, "has_table", False):
        return {"has_table": False}
    table = shape.table
    rows = []
    empty_cells = 0
    non_empty_cells = 0
    max_cell_units = 0.0
    long_cell_count = 0
    for row in table.rows:
        cells = []
        for cell in row.cells:
            text = _clean_text(cell.text)
            cells.append(text)
            if text:
                non_empty_cells += 1
                units = _effective_char_units(text)
                max_cell_units = max(max_cell_units, units)
                if units >= 12.0 or any(_effective_char_units(word) >= 7.0 for word in text.split()):
                    long_cell_count += 1
            else:
                empty_cells += 1
        if any(cells):
            rows.append(cells)
    total_cells = max(1, len(table.rows) * len(table.columns))
    return {
        "has_table": True,
        "row_count": len(table.rows),
        "col_count": len(table.columns),
        "non_empty_rows": len(rows),
        "non_empty_cells": non_empty_cells,
        "empty_cells": empty_cells,
        "empty_cell_ratio": round(empty_cells / total_cells, 3),
        "max_cell_effective_units": round(max_cell_units, 1),
        "long_cell_count": long_cell_count,
        "preview_rows": rows[:4],
    }


def _font_sizes(shape: Any) -> List[float]:
    if not getattr(shape, "has_text_frame", False):
        return []
    sizes: List[float] = []
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.font.size is not None:
            sizes.append(round(paragraph.font.size.pt, 2))
        for run in paragraph.runs:
            if run.font.size is not None:
                sizes.append(round(run.font.size.pt, 2))
    return sizes


def _paragraph_alignment(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    paragraphs = getattr(shape.text_frame, "paragraphs", [])
    if not paragraphs:
        return ""
    return str(paragraphs[0].alignment or "")


def _vertical_anchor(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return str(getattr(shape.text_frame, "vertical_anchor", "") or "")


def _infer_role(record: Dict[str, Any], slide_h: float) -> str:
    text = record.get("text", "")
    bbox = record.get("bbox", {})
    font_avg = record.get("font", {}).get("avg_pt") or 0
    area = bbox.get("area", 0.0)
    x = bbox.get("x", 0.0)
    y = bbox.get("y", 0.0)
    w = bbox.get("w", 0.0)

    if record.get("is_full_background"):
        return "background"
    if record.get("has_table"):
        return "table"
    if record.get("is_picture"):
        return "picture"
    if not text:
        if record.get("is_line_like"):
            return "decorative_rule"
        if area > 0.2:
            return "container"
        return "decorative"
    if text.startswith("Sources:") or y > slide_h - 0.65 or (y > slide_h - 0.95 and font_avg <= 9):
        return "source_footer"
    if re.match(r"^\d{1,2}\s*/\s*\d{1,2}\b", text) or (y < 0.55 and w < 2.2 and font_avg <= 10):
        return "page_marker"
    normalized_label = _label_token(text)
    if normalized_label in {
        "ACADEMICPAPERREADING",
        "DECKMAP",
        "EVIDENCEMOSAIC",
        "EVIDENCENOTES",
        "KEYEVIDENCE",
        "MODULECHECKPOINT",
        "PAPERHIGHLIGHTS",
        "PROOFOBJECT",
        "SECTION",
        "SUMMARY",
    }:
        return "component_label"
    if text in {"Conclusion", "Method", "Limitations", "Reading note 1", "Reading note 2", "Reading note 3"}:
        return "card_label"
    if text.isupper() and len(text) < 70 and font_avg <= 12:
        return "component_label"
    if font_avg <= 10.5 and len(text) > 18 and y > 1.0:
        return "card_text"
    if font_avg >= 20 or (font_avg >= 16 and y < 2.8 and w > 4.0 and len(text) > 48 and bbox.get("h", 0.0) >= 0.55):
        return "title_claim"
    if len(text) > 55 and font_avg <= 15:
        return "support_body"
    if font_avg <= 10:
        return "small_text"
    return "body_text"


def _text_capacity_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    findings = []
    for record in records:
        if not record.get("has_text"):
            continue
        role = record.get("role", "")
        if role in {"source_footer", "page_marker", "component_label", "card_label"}:
            continue
        estimate = _text_capacity_estimate(record)
        record["text_capacity"] = estimate
        font_min = record.get("font", {}).get("min_pt")
        floor = _font_floor(role, rules)
        if record.get("text_chars", 0) < 28 and role in {"card_text", "small_text", "body_text"}:
            floor = 0.0
        if role == "title_claim" and (record.get("font", {}).get("avg_pt") or 0) < 18:
            floor = 0.0
        if font_min is not None and floor and font_min < floor:
            findings.append(
                _finding(
                    slide_index,
                    "low_font_size",
                    "medium" if role in {"title_claim", "support_body", "body_text"} else "low",
                    f"{role} text uses {font_min:.1f}pt below the {floor:.1f}pt metadata floor.",
                    record,
                    "Increase font size or shorten copy while preserving the component's accepted composition.",
                )
            )
        elif font_min is not None and floor:
            ideal = _font_ideal_floor(role, rules)
            if ideal and font_min < ideal and record.get("text_chars", 0) >= 28:
                findings.append(
                    _finding(
                        slide_index,
                        "below_ideal_font_band",
                        "low",
                        f"{role} text uses {font_min:.1f}pt; preferred metadata band starts near {ideal:.1f}pt.",
                        record,
                        "Prefer a small typography lift before changing component geometry.",
                    )
                )
        fill_ratio = estimate["fill_ratio"]
        if fill_ratio >= float(rules["overflow_ratio"]):
            findings.append(
                _finding(
                    slide_index,
                    "estimated_text_overflow",
                    "high",
                    f"{role} text is estimated at {fill_ratio:.2f}x box capacity.",
                    record,
                    "Shorten text, add a line break, or move secondary detail into notes; resize only if capacity remains impossible.",
                )
            )
        elif fill_ratio >= float(rules["near_capacity_ratio"]):
            findings.append(
                _finding(
                    slide_index,
                    "near_text_capacity",
                    "low",
                    f"{role} text is estimated at {fill_ratio:.2f}x box capacity.",
                    record,
                    "Keep composition fixed; slightly shorten text or raise the box's assigned font only after checking capacity.",
                )
            )
        elif (
            fill_ratio <= float(rules["sparse_ratio"])
            and record.get("bbox", {}).get("area", 0) >= float(rules["sparse_area_sq_in"])
            and role in {"support_body", "card_text", "body_text"}
            and int(record.get("text_words") or 0) > 3
        ):
            findings.append(
                _finding(
                    slide_index,
                    "low_text_density",
                    "low",
                    f"{role} text occupies only {fill_ratio:.2f}x of estimated capacity in a large box.",
                    record,
                    "Prefer stronger typography or richer evidence text; do not auto-shrink the component because low density hurt v6.",
                )
            )
    return findings


def _cover_typography_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    if slide_index != 1:
        return []
    candidates = [
        record
        for record in records
        if record.get("role") in {"body_text", "support_body"}
        and record.get("bbox", {}).get("x", 0.0) < 8.2
        and int(record.get("text_words") or 0) >= 2
    ]
    if not candidates:
        return []
    weakest = min(candidates, key=lambda item: item.get("font", {}).get("min_pt") or 99.0)
    font_min = weakest.get("font", {}).get("min_pt")
    floor = float(rules["cover_left_min_pt"])
    if font_min is None or font_min >= floor:
        return []
    return [
        _finding(
            slide_index,
            "cover_left_typography_underpowered",
            "medium",
            f"Cover left narrative text uses {font_min:.1f}pt below the {floor:.1f}pt cover comfort floor.",
            weakest,
            "Lift the cover's left-side title metadata/support typography as a group, then recheck wrapping before moving components.",
        )
    ]


def _paired_text_stack_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    if slide_index != 1:
        return []
    highlight_labels = {"Core result", "Scale", "Design edge"}
    findings = []
    for label in records:
        if label.get("text") not in highlight_labels:
            continue
        label_box = label.get("bbox", {})
        candidates = [
            record
            for record in records
            if record is not label
            and record.get("has_text")
            and int(record.get("text_words") or 0) >= 4
            and abs(record.get("bbox", {}).get("x", 0.0) - label_box.get("x", 0.0)) <= 0.08
            and 0 < record.get("bbox", {}).get("y", 0.0) - label_box.get("y", 0.0) <= 0.55
        ]
        if not candidates:
            continue
        body = min(candidates, key=lambda item: item.get("bbox", {}).get("y", 0.0))
        top_delta = body["bbox"].get("y", 0.0) - label_box.get("y", 0.0)
        if top_delta <= float(rules["paired_text_top_delta_max_in"]):
            continue
        findings.append(
            {
                "type": "paired_label_body_gap_too_large",
                "severity": "low",
                **_finding_meta("paired_label_body_gap_too_large", "low"),
                "slide_page": slide_index,
                "message": f"Cover highlight '{label.get('text')}' and its body start {top_delta:.2f}in apart, weakening the paired rail reading.",
                "evidence": {
                    "label": _shape_evidence(label),
                    "body": _shape_evidence(body),
                    "top_delta_in": round(top_delta, 3),
                },
                "repair_strategy": "Treat the highlight label and explanation as one stack; move the body closer before changing the overall rail positions.",
            }
        )
    return findings


def _copy_distribution_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    card_records = [
        record
        for record in records
        if record.get("role") == "card_text"
        and record.get("bbox", {}).get("area", 0.0) >= 0.45
        and int(record.get("text_words") or 0) > 0
    ]
    if len(card_records) < 3:
        return []
    counts = [max(1, int(record.get("text_words") or 0)) for record in card_records]
    smallest = min(counts)
    largest = max(counts)
    if largest / max(1, smallest) < float(rules["card_copy_imbalance_ratio"]) or smallest > 6:
        return []
    record = card_records[counts.index(smallest)]
    return [
        _finding(
            slide_index,
            "card_copy_imbalance",
            "low",
            f"Evidence cards have uneven copy distribution: shortest card has {smallest} words, longest has {largest}.",
            record,
            "Rebalance reading notes or enrich the sparse card copy; keep the accepted card group size unless geometry fails.",
        )
    ]


def _flow_layout_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    findings = []
    labels = [record for record in records if record.get("text") in {"Problem", "Method", "Evidence", "Takeaways"}]
    if len(labels) < 4:
        return []
    ordered = sorted(labels, key=lambda item: item.get("bbox", {}).get("x", 0.0))
    row_span = max(item["bbox"].get("y", 0.0) for item in ordered) - min(item["bbox"].get("y", 0.0) for item in ordered)
    if row_span <= 0.18:
        gaps = [
            ordered[index + 1]["bbox"].get("x", 0.0) - ordered[index]["bbox"].get("right", 0.0)
            for index in range(len(ordered) - 1)
        ]
        min_gap = min(gaps) if gaps else 0.0
        if min_gap < float(rules["flow_label_min_gap_in"]):
            findings.append(
                {
                    "type": "flow_nodes_overpacked",
                    "severity": "medium",
                    **_finding_meta("flow_nodes_overpacked", "medium"),
                    "slide_page": slide_index,
                    "message": f"Read-path labels sit in one tight row with only {min_gap:.2f}in between adjacent labels.",
                    "evidence": {
                        "labels": [_shape_evidence(item) for item in ordered],
                        "min_label_gap_in": round(min_gap, 3),
                    },
                    "repair_strategy": "Switch compact flow components to two rows when labels are crowded instead of shrinking the label typography.",
                }
            )

    nodes = [record for record in records if record.get("text") in {"P", "M", "E", "T"}]
    if len(nodes) < 4:
        return findings
    nodes_by_text = {record.get("text"): record for record in nodes}
    labels_by_text = {record.get("text"): record for record in labels}
    node_centers = {
        key: (
            item.get("bbox", {}).get("x", 0.0) + item.get("bbox", {}).get("w", 0.0) / 2,
            item.get("bbox", {}).get("y", 0.0) + item.get("bbox", {}).get("h", 0.0) / 2,
        )
        for key, item in nodes_by_text.items()
    }
    if {"P", "M", "E", "T"}.issubset(node_centers):
        col_gap = ((node_centers["M"][0] - node_centers["P"][0]) + (node_centers["T"][0] - node_centers["E"][0])) / 2
        row_gap = ((node_centers["E"][1] - node_centers["P"][1]) + (node_centers["T"][1] - node_centers["M"][1])) / 2
        ratio = col_gap / max(0.01, row_gap)
        if (
            col_gap > float(rules["flow_grid_max_col_gap_in"])
            or row_gap < float(rules["flow_grid_min_row_gap_in"])
            or ratio > float(rules["flow_grid_col_to_row_ratio_max"])
        ):
            findings.append(
                {
                    "type": "flow_grid_alignment_drift",
                    "severity": "low",
                    **_finding_meta("flow_grid_alignment_drift", "low"),
                    "slide_page": slide_index,
                    "message": f"2x2 read-path grid has column gap {col_gap:.2f}in and row gap {row_gap:.2f}in, making the component feel stretched.",
                    "evidence": {
                        "nodes": [_shape_evidence(nodes_by_text[key]) for key in ["P", "M", "E", "T"]],
                        "column_gap_in": round(col_gap, 3),
                        "row_gap_in": round(row_gap, 3),
                        "column_to_row_ratio": round(ratio, 3),
                    },
                    "repair_strategy": "Use a compact 2x2 grid: reduce column distance, give rows enough air, and keep labels centered below their nodes.",
                }
            )
    text_to_label = {"P": "Problem", "M": "Method", "E": "Evidence", "T": "Takeaways"}
    center_tolerance = float(rules["flow_label_center_tolerance_in"])
    for letter, label_text in text_to_label.items():
        node = nodes_by_text.get(letter)
        label = labels_by_text.get(label_text)
        if not node or not label:
            continue
        node_cx = node["bbox"].get("x", 0.0) + node["bbox"].get("w", 0.0) / 2
        label_cx = label["bbox"].get("x", 0.0) + label["bbox"].get("w", 0.0) / 2
        drift = abs(node_cx - label_cx)
        if drift <= center_tolerance:
            continue
        findings.append(
            {
                "type": "flow_grid_alignment_drift",
                "severity": "low",
                **_finding_meta("flow_grid_alignment_drift", "low"),
                "slide_page": slide_index,
                "message": f"Read-path label '{label_text}' is {drift:.2f}in away from the center of node {letter}.",
                "evidence": {
                    "node": _shape_evidence(node),
                    "label": _shape_evidence(label),
                    "center_drift_in": round(drift, 3),
                },
                "repair_strategy": "Center each read-path label under its circle after wrapping the flow to two rows.",
            }
        )
    return findings


def _agenda_read_path_header_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    if slide_index != 2:
        return []
    headers = [record for record in records if record.get("text") == "Read path"]
    nodes = [record for record in records if record.get("text") in {"P", "M", "E", "T"}]
    if not headers or not nodes:
        return []
    header = min(headers, key=lambda record: record.get("bbox", {}).get("y", 0.0))
    header_box = header.get("bbox", {})
    lower_nodes = [
        node
        for node in nodes
        if node.get("bbox", {}).get("y", 0.0) > header_box.get("y", 0.0)
    ]
    if not lower_nodes:
        return []
    nearest_node_top = min(node.get("bbox", {}).get("y", 0.0) for node in lower_nodes)
    gap = nearest_node_top - header_box.get("bottom", 0.0)
    min_gap = float(rules["agenda_read_path_header_min_gap_in"])
    if gap >= min_gap:
        return []
    return [
        {
            "type": "agenda_read_path_header_too_close",
            "severity": "low",
            **_finding_meta("agenda_read_path_header_too_close", "low"),
            "slide_page": slide_index,
            "message": f"Agenda Read path header sits only {gap:.2f}in above the flow nodes, making the rail feel cramped.",
            "evidence": {
                "header": _shape_evidence(header),
                "nodes": [_shape_evidence(node) for node in sorted(nodes, key=lambda item: item.get("text", ""))],
                "gap_in": round(gap, 3),
                "minimum_gap_in": round(min_gap, 3),
            },
            "repair_strategy": "Lift the Read path header slightly while preserving the flow node grid and rail composition.",
        }
    ]


def _table_support_band_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    findings = []
    tables = [record for record in records if record.get("role") == "table"]
    supports = [record for record in records if record.get("role") == "support_body"]
    claims = [record for record in records if record.get("role") == "title_claim"]
    if not tables or not supports or not claims:
        return findings

    containers = _container_records(records)
    for table in tables:
        table_box = table.get("bbox", {})
        if table_box.get("y", 0.0) < 3.45 or table_box.get("w", 0.0) < 7.0:
            continue
        panel = _smallest_containing_container(table_box, containers)
        if not panel:
            continue
        panel_box = panel.get("bbox", {})
        if panel_box.get("y", 0.0) < 3.0 or panel_box.get("w", 0.0) < 7.0:
            continue
        support_candidates = [
            support
            for support in supports
            if support.get("bbox", {}).get("bottom", 0.0) <= panel_box.get("y", 0.0) + 0.08
            and _horizontal_overlap_ratio(support.get("bbox", {}), panel_box) >= 0.45
        ]
        if not support_candidates:
            continue
        support = max(support_candidates, key=lambda item: item.get("bbox", {}).get("y", 0.0))
        support_box = support.get("bbox", {})
        claim_candidates = [
            claim
            for claim in claims
            if claim.get("bbox", {}).get("bottom", 0.0) <= support_box.get("y", 0.0) + 0.04
            and _horizontal_overlap_ratio(claim.get("bbox", {}), support_box) >= 0.35
        ]
        if not claim_candidates:
            continue
        claim = max(claim_candidates, key=lambda item: item.get("bbox", {}).get("bottom", 0.0))
        claim_box = claim.get("bbox", {})
        claim_gap = support_box.get("y", 0.0) - claim_box.get("bottom", 0.0)
        panel_gap = panel_box.get("y", 0.0) - support_box.get("bottom", 0.0)
        max_claim_gap = float(rules["table_support_claim_gap_max_in"])
        min_panel_gap = float(rules["table_support_panel_gap_min_in"])
        reasons = []
        if claim_gap > max_claim_gap:
            reasons.append(f"claim/support gap {claim_gap:.2f}in")
        if panel_gap < min_panel_gap:
            reasons.append(f"support/table-panel gap {panel_gap:.2f}in")
        if not reasons:
            continue
        findings.append(
            {
                "type": "table_support_band_off_balance",
                "severity": "low",
                **_finding_meta("table_support_band_off_balance", "low"),
                "slide_page": slide_index,
                "message": "Table-bottom support copy is optically pulled toward the table panel: " + "; ".join(reasons) + ".",
                "evidence": {
                    "claim": _shape_evidence(claim),
                    "support": _shape_evidence(support),
                    "table_panel": _shape_evidence(panel),
                    "table": _shape_evidence(table),
                    "claim_gap_in": round(claim_gap, 3),
                    "max_claim_gap_in": round(max_claim_gap, 3),
                    "panel_gap_in": round(panel_gap, 3),
                    "min_panel_gap_in": round(min_panel_gap, 3),
                },
                "repair_strategy": "Treat claim and support as one upper text stack, then lower the table panel enough to preserve a visible gutter.",
            }
        )
    return findings


def _metric_stack_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    findings = []
    for container in _container_records(records):
        bbox = container.get("bbox", {})
        area = bbox.get("area", 0.0)
        if area < float(rules["metric_card_min_area_sq_in"]) or area > 8.0:
            continue
        children = _contained_records(container, records, include_containers=False)
        values = [child for child in children if _looks_metric_value(child.get("text", ""))]
        if not values:
            continue
        value = max(values, key=lambda child: child.get("font", {}).get("avg_pt") or 0.0)
        labels = [
            child
            for child in children
            if child is not value
            and child.get("has_text")
            and child.get("bbox", {}).get("y", 0.0) > value.get("bbox", {}).get("y", 0.0)
            and child.get("role") in {"small_text", "card_text", "body_text", "support_body"}
        ]
        if not labels:
            continue
        label = min(labels, key=lambda child: child.get("bbox", {}).get("y", 0.0))
        gap = label["bbox"].get("y", 0.0) - value["bbox"].get("bottom", 0.0)
        if gap <= float(rules["metric_label_gap_max_in"]):
            continue
        findings.append(
            {
                "type": "metric_label_gap_too_large",
                "severity": "low",
                **_finding_meta("metric_label_gap_too_large", "low"),
                "slide_page": slide_index,
                "message": f"Metric value and label are separated by {gap:.2f}in inside one card.",
                "evidence": {
                    "container": _shape_evidence(container),
                    "value": _shape_evidence(value),
                    "label": _shape_evidence(label),
                    "gap_in": round(gap, 3),
                },
                "repair_strategy": "Treat value and label as one optical text stack; use a fixed stack gap instead of card-height proportional placement.",
            }
        )
    return findings


def _component_boundary_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    findings = []
    watched_labels = {"FIGURE", "TABLE"}
    containers = _container_records(records)
    min_inset = float(rules["component_label_min_inset_in"])
    for label in records:
        if label.get("role") != "component_label" or _label_token(label.get("text", "")) not in watched_labels:
            continue
        label_box = label.get("bbox", {})
        parent = _smallest_containing_container(label_box, containers)
        if not parent:
            continue
        parent_box = parent.get("bbox", {})
        insets = {
            "left": label_box.get("x", 0.0) - parent_box.get("x", 0.0),
            "top": label_box.get("y", 0.0) - parent_box.get("y", 0.0),
            "right": parent_box.get("right", 0.0) - label_box.get("right", 0.0),
            "bottom": parent_box.get("bottom", 0.0) - label_box.get("bottom", 0.0),
        }
        leading_inset = min(insets["left"], insets["top"])
        if leading_inset >= min_inset:
            continue
        findings.append(
            {
                "type": "component_boundary_inset_violation",
                "severity": "low",
                **_finding_meta("component_boundary_inset_violation", "low"),
                "slide_page": slide_index,
                "message": f"{label.get('text')} label sits only {leading_inset:.2f}in from its rounded component boundary.",
                "evidence": {
                    "label": _shape_evidence(label),
                    "container": _shape_evidence(parent),
                    "insets_in": {key: round(value, 3) for key, value in insets.items()},
                },
                "repair_strategy": "Inset component labels from rounded panel corners before fitting the figure/table content area.",
            }
        )
    return findings


def _text_container_bounds_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    findings = []
    containers = _container_records(records)
    padding = float(rules["text_container_padding_in"])
    tolerance = float(rules["text_container_overflow_tolerance_in"])
    for record in records:
        if not record.get("has_text"):
            continue
        if record.get("role") in {"component_label", "card_label", "source_footer", "page_marker"}:
            continue
        if int(record.get("text_words") or 0) < 4:
            continue
        parent = _nearest_container_for_record(record, containers)
        if not parent:
            continue
        parent_box = parent.get("bbox", {})
        record_box = record.get("bbox", {})
        if parent_box.get("area", 0.0) > 12.0:
            continue
        estimate = record.get("text_capacity") or _text_capacity_estimate(record)
        fill_ratio = float(estimate.get("fill_ratio", 0.0) or 0.0)
        expected_h = _estimated_rendered_text_height(record)
        expected_bottom = record_box.get("y", 0.0) + expected_h
        overflow_bottom = expected_bottom - (parent_box.get("bottom", 0.0) - padding)
        explicit_box_overflow = record_box.get("bottom", 0.0) - (parent_box.get("bottom", 0.0) - padding)
        if max(overflow_bottom, explicit_box_overflow) <= tolerance and fill_ratio < float(rules["overflow_ratio"]):
            continue
        severity = "high" if fill_ratio >= float(rules["overflow_ratio"]) or overflow_bottom >= 0.12 else "medium"
        findings.append(
            {
                "type": "text_exceeds_container_bounds",
                "severity": severity,
                **_finding_meta("text_exceeds_container_bounds", severity),
                "slide_page": slide_index,
                "message": f"Text is likely to render beyond its containing component by {max(overflow_bottom, explicit_box_overflow):.2f}in.",
                "evidence": {
                    "text": _shape_evidence(record),
                    "container": _shape_evidence(parent),
                    "fill_ratio": round(fill_ratio, 3),
                    "expected_text_height_in": round(expected_h, 3),
                    "expected_bottom_overflow_in": round(max(0.0, overflow_bottom), 3),
                    "textbox_bottom_overflow_in": round(max(0.0, explicit_box_overflow), 3),
                },
                "repair_strategy": "Resize the evidence card text area, split copy across cards, or shorten the card body before accepting the slide.",
            }
        )
    return findings


def _text_card_vertical_alignment_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    findings = []
    containers = _container_records(records)
    slack_ratio = float(rules["text_card_vertical_slack_ratio"])
    min_extra = float(rules["text_card_vertical_extra_min_in"])
    for record in records:
        if record.get("role") not in {"card_text", "support_body", "body_text"}:
            continue
        if int(record.get("text_words") or 0) < 5:
            continue
        parent = _nearest_container_for_record(record, containers)
        if not parent:
            continue
        parent_box = parent.get("bbox", {})
        record_box = record.get("bbox", {})
        if parent_box.get("area", 0.0) > 8.5 or parent_box.get("h", 0.0) > 1.65:
            continue
        expected_h = _estimated_rendered_text_height(record)
        extra_h = record_box.get("h", 0.0) - expected_h
        if extra_h < min_extra or record_box.get("h", 0.0) / max(0.01, expected_h) < slack_ratio:
            continue
        if _has_middle_vertical_anchor(record):
            continue
        findings.append(
            {
                "type": "text_card_vertical_alignment_top_heavy",
                "severity": "medium",
                **_finding_meta("text_card_vertical_alignment_top_heavy", "medium"),
                "slide_page": slide_index,
                "message": f"Text card has {extra_h:.2f}in vertical slack but is not middle-anchored, so copy appears top-heavy.",
                "evidence": {
                    "text": _shape_evidence(record),
                    "container": _shape_evidence(parent),
                    "expected_text_height_in": round(expected_h, 3),
                    "vertical_extra_in": round(extra_h, 3),
                    "vertical_anchor": record.get("vertical_anchor", ""),
                },
                "repair_strategy": "Use middle vertical anchoring for shallow straight-rectangle evidence cards, or reduce the card height to match text.",
            }
        )
    return findings


def _figure_label_semantics_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    findings = []
    containers = _container_records(records)
    for container in containers:
        parent_box = container.get("bbox", {})
        parent_aspect = parent_box.get("w", 0.0) / max(0.001, parent_box.get("h", 0.0))
        children = [record for record in records if record is not container and _center_inside(record.get("bbox", {}), parent_box)]
        if not any(child.get("is_picture") for child in children):
            continue
        badges = [
            child
            for child in children
            if child.get("role") == "component_label" and _label_token(child.get("text", "")) == "FIGURE"
        ]
        figure_ids = [
            child
            for child in children
            if _is_figure_identity_label(child.get("text", ""))
        ]
        figure_semantics_panel = bool(badges or figure_ids)
        figure_aspect_panel = (
            parent_aspect >= float(rules["figure_label_semantic_wide_panel_aspect_min"])
            or parent_aspect <= float(rules["figure_label_semantic_tall_panel_aspect_max"])
        )
        if not (figure_semantics_panel or figure_aspect_panel):
            continue
        for badge in badges:
            badge_box = badge.get("bbox", {})
            left_inset = badge_box.get("x", 0.0) - parent_box.get("x", 0.0)
            top_inset = badge_box.get("y", 0.0) - parent_box.get("y", 0.0)
            if (
                left_inset <= float(rules["figure_panel_badge_left_inset_max_in"])
                and float(rules["figure_panel_badge_top_inset_min_in"])
                <= top_inset
                <= float(rules["figure_panel_badge_top_inset_max_in"])
            ):
                continue
            findings.append(
                {
                    "type": "figure_badge_identity_label_conflation",
                    "severity": "low",
                    **_finding_meta("figure_badge_identity_label_conflation", "low"),
                    "slide_page": slide_index,
                    "message": "FIGURE badge is not anchored as the rounded panel's top-left type label.",
                    "evidence": {
                        "badge": _shape_evidence(badge),
                        "container": _shape_evidence(container),
                        "insets_in": {"left": round(left_inset, 3), "top": round(top_inset, 3)},
                    },
                    "repair_strategy": "Pin the green FIGURE badge to the rounded panel's inner top-left corner; reserve image-adjacent placement for the Figure N identity label.",
                }
            )
        for figure_id in figure_ids:
            label_box = figure_id.get("bbox", {})
            picture = _primary_picture_child(children)
            if not picture:
                continue
            picture_box = picture.get("bbox", {})
            if _is_stacked_figure_identity_text(figure_id.get("text", "")):
                findings.append(
                    {
                        "type": "stacked_figure_identity_label_overcorrection",
                        "severity": "low",
                        **_finding_meta("stacked_figure_identity_label_overcorrection", "low"),
                        "slide_page": slide_index,
                        "message": f"{figure_id.get('text')} is stacked vertically; reviewer preferred a horizontal Figure N label above the image.",
                        "evidence": {
                            "figure_label": _shape_evidence(figure_id),
                            "picture": _shape_evidence(picture),
                        },
                        "repair_strategy": "Use a horizontal Figure N label just above the fitted image's upper-left edge; keep the green FIGURE badge on the panel corner.",
                    }
                )
                continue
            center_delta = abs(_center_x(label_box) - _center_x(picture_box))
            vertical_gap = picture_box.get("y", 0.0) - label_box.get("bottom", 0.0)
            geometry_ok = (
                center_delta <= float(rules["figure_identity_label_center_tolerance_in"])
                and float(rules["figure_identity_label_vertical_gap_min_in"])
                <= vertical_gap
                <= float(rules["figure_identity_label_vertical_gap_max_in"])
            )
            if geometry_ok and _is_center_aligned_text(figure_id):
                continue
            if geometry_ok:
                findings.append(
                    {
                        "type": "figure_label_text_alignment_off_center",
                        "severity": "low",
                        **_finding_meta("figure_label_text_alignment_off_center", "low"),
                        "slide_page": slide_index,
                        "message": f"{figure_id.get('text')} textbox is centered above the image, but its text is not center-aligned inside the box.",
                        "evidence": {
                            "figure_label": _shape_evidence(figure_id),
                            "picture": _shape_evidence(picture),
                            "center_delta_in": round(center_delta, 3),
                            "vertical_gap_in": round(vertical_gap, 3),
                        },
                        "repair_strategy": "Center-align the Figure N paragraph inside the image-centered label textbox.",
                    }
                )
                continue
            findings.append(
                {
                    "type": "figure_label_anchor_drift",
                    "severity": "low",
                    **_finding_meta("figure_label_anchor_drift", "low"),
                    "slide_page": slide_index,
                    "message": f"{figure_id.get('text')} is not centered above the fitted image.",
                    "evidence": {
                        "figure_label": _shape_evidence(figure_id),
                        "picture": _shape_evidence(picture),
                        "center_delta_in": round(center_delta, 3),
                        "vertical_gap_in": round(vertical_gap, 3),
                    },
                    "repair_strategy": "Place Figure N horizontally above the fitted image and align it to the image centerline, outside the image itself.",
                }
            )
    return findings


def _panel_identity_label_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    findings = []
    checked_badge_tokens = {"TEXTEVIDENCE", "TABLE", "METRIC"}
    for container in _container_records(records):
        parent_box = container.get("bbox", {})
        children = _contained_records(container, records, include_containers=False)
        badges = [
            child
            for child in children
            if child.get("role") == "component_label" and _label_token(child.get("text", "")) in checked_badge_tokens
        ]
        if not badges:
            continue
        identity = _panel_identity_label_child(children, parent_box)
        if not identity:
            continue
        target = _panel_identity_target_child(children, identity)
        if not target:
            continue
        label_box = identity.get("bbox", {})
        target_box = target.get("bbox", {})
        center_delta = abs(_center_x(label_box) - _center_x(target_box))
        center_ok = center_delta <= float(rules["panel_identity_label_center_tolerance_in"])
        if center_ok and _is_center_aligned_text(identity):
            continue
        if center_ok:
            findings.append(
                {
                    "type": "panel_identity_label_text_alignment_off_center",
                    "severity": "low",
                    **_finding_meta("panel_identity_label_text_alignment_off_center", "low"),
                    "slide_page": slide_index,
                    "message": f"{identity.get('text')} label box is centered over its panel content, but its text is not center-aligned inside the box.",
                    "evidence": {
                        "identity_label": _shape_evidence(identity),
                        "target_content": _shape_evidence(target),
                        "center_delta_in": round(center_delta, 3),
                    },
                    "repair_strategy": "Center-align proof identity text inside the textbox after aligning the textbox to the underlying table, figure, metric, or explanation text.",
                }
            )
            continue
        findings.append(
            {
                "type": "panel_identity_label_anchor_drift",
                "severity": "low",
                **_finding_meta("panel_identity_label_anchor_drift", "low"),
                "slide_page": slide_index,
                "message": f"{identity.get('text')} label is not centered above the panel's main content.",
                "evidence": {
                    "identity_label": _shape_evidence(identity),
                    "target_content": _shape_evidence(target),
                    "center_delta_in": round(center_delta, 3),
                },
                "repair_strategy": "Place the proof identity label on the centerline of the underlying table, figure, metric, or explanation text; keep the green type badge as a panel-corner label.",
            }
        )
    return findings


def _component_frame_fit_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    findings = []
    for container in _container_records(records):
        bbox = container.get("bbox", {})
        if bbox.get("h", 0.0) > 2.2 or bbox.get("w", 0.0) < 1.4 or bbox.get("area", 0.0) < 1.0:
            continue
        children = [
            child
            for child in _contained_records(container, records, include_containers=False)
            if child.get("has_text") and child.get("role") not in {"source_footer", "page_marker", "component_label"}
        ]
        if len(children) < 2 or any(child.get("has_table") or child.get("is_picture") for child in children):
            continue
        if not any(child.get("role") == "card_label" for child in children):
            continue
        if not any(child.get("role") in {"card_text", "body_text", "support_body"} for child in children):
            continue
        expected_h = _expected_text_stack_height(children, bbox)
        extra_h = bbox.get("h", 0.0) - expected_h
        if extra_h <= float(rules["component_frame_extra_height_max_in"]):
            continue
        findings.append(
            {
                "type": "component_frame_overallocated_after_text_fit",
                "severity": "low",
                **_finding_meta("component_frame_overallocated_after_text_fit", "low"),
                "slide_page": slide_index,
                "message": f"Evidence card frame has about {extra_h:.2f}in more height than its fitted text stack needs.",
                "evidence": {
                    "container": _shape_evidence(container),
                    "text_children": [_shape_evidence(child) for child in sorted(children, key=lambda item: item["bbox"].get("y", 0.0))],
                    "expected_stack_height_in": round(expected_h, 3),
                    "extra_height_in": round(extra_h, 3),
                },
                "repair_strategy": "After fitting text, resize the local card frame and reflow sibling cards on the same slide instead of changing deck-wide type roles.",
            }
        )
    return findings


def _card_internal_spacing_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    findings = []
    max_gap = float(rules["card_internal_gap_max_in"])
    min_bottom = float(rules["card_internal_min_bottom_padding_in"])
    shallow_max_h = float(rules["card_internal_shallow_max_h_in"])
    narrow_max_w = float(rules["card_internal_narrow_max_w_in"])
    for container in _container_records(records):
        bbox = container.get("bbox", {})
        if (
            bbox.get("h", 0.0) > shallow_max_h
            or bbox.get("w", 0.0) > narrow_max_w
            or bbox.get("w", 0.0) < 1.4
            or bbox.get("area", 0.0) < 0.9
        ):
            continue
        children = [
            child
            for child in _contained_records(container, records, include_containers=False)
            if child.get("has_text") and child.get("role") not in {"source_footer", "page_marker", "component_label"}
        ]
        labels = [child for child in children if child.get("role") == "card_label"]
        bodies = [
            child
            for child in children
            if child.get("role") in {"card_text", "body_text", "support_body"}
            and int(child.get("text_words") or 0) >= 3
        ]
        if not labels or not bodies:
            continue
        label = min(labels, key=lambda child: child.get("bbox", {}).get("y", 0.0))
        body_candidates = [
            body
            for body in bodies
            if body.get("bbox", {}).get("y", 0.0) >= label.get("bbox", {}).get("y", 0.0)
        ]
        if not body_candidates:
            continue
        body = min(body_candidates, key=lambda child: child.get("bbox", {}).get("y", 0.0))
        label_box = label.get("bbox", {})
        body_box = body.get("bbox", {})
        gap = body_box.get("y", 0.0) - label_box.get("bottom", 0.0)
        bottom_padding = bbox.get("bottom", 0.0) - body_box.get("bottom", 0.0)
        if gap <= max_gap and bottom_padding >= min_bottom:
            continue
        reasons = []
        if gap > max_gap:
            reasons.append(f"label/body gap {gap:.2f}in")
        if bottom_padding < min_bottom:
            reasons.append(f"bottom padding {bottom_padding:.2f}in")
        findings.append(
            {
                "type": "card_internal_spacing_not_scaled_to_frame",
                "severity": "low",
                **_finding_meta("card_internal_spacing_not_scaled_to_frame", "low"),
                "slide_page": slide_index,
                "message": f"Shallow evidence card uses fixed internal spacing ({', '.join(reasons)}), making the text stack feel low in the frame.",
                "evidence": {
                    "container": _shape_evidence(container),
                    "label": _shape_evidence(label),
                    "body": _shape_evidence(body),
                    "gap_in": round(gap, 3),
                    "bottom_padding_in": round(bottom_padding, 3),
                },
                "repair_strategy": "Scale label/body gap and body box height to the local card height before changing deck-wide typography.",
            }
        )
    return findings


def _container_balance_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    findings = []
    for container in _container_records(records):
        bbox = container.get("bbox", {})
        if bbox.get("area", 0.0) < float(rules["balance_container_min_area_sq_in"]):
            continue
        if bbox.get("w", 0.0) > 8.0 or bbox.get("h", 0.0) > 6.8:
            continue
        children = _contained_records(container, records, include_containers=True)
        children = [
            child
            for child in children
            if child.get("role") not in {"decorative", "decorative_rule", "component_label", "source_footer", "page_marker"}
            and not child.get("is_full_background")
        ]
        if len(children) < 2:
            continue
        union = _bbox_union([child["bbox"] for child in children])
        top_pad = union.get("y", 0.0) - bbox.get("y", 0.0)
        bottom_pad = bbox.get("bottom", 0.0) - union.get("bottom", 0.0)
        delta = abs(top_pad - bottom_pad)
        tolerance = max(float(rules["balance_padding_tolerance_in"]), bbox.get("h", 0.0) * float(rules["balance_padding_tolerance_ratio"]))
        if delta <= tolerance:
            continue
        findings.append(
            {
                "type": "container_stack_off_balance",
                "severity": "low",
                **_finding_meta("container_stack_off_balance", "low"),
                "slide_page": slide_index,
                "message": f"Container content stack has uneven vertical padding: top {top_pad:.2f}in, bottom {bottom_pad:.2f}in.",
                "evidence": {
                    "container": _shape_evidence(container),
                    "child_count": len(children),
                    "stack_bbox": union,
                    "top_padding_in": round(top_pad, 3),
                    "bottom_padding_in": round(bottom_pad, 3),
                    "padding_delta_in": round(delta, 3),
                },
                "repair_strategy": "Rebalance the internal stack or enrich sparse metric cards before changing the whole slide macro-layout.",
            }
        )
    return findings


def _text_capacity_estimate(record: Dict[str, Any]) -> Dict[str, Any]:
    bbox = record.get("bbox", {})
    text = record.get("text", "")
    font = record.get("font", {}).get("avg_pt") or _fallback_font(record.get("role", ""))
    usable_w = max(0.1, float(bbox.get("w", 0.0)) - 0.18)
    usable_h = max(0.1, float(bbox.get("h", 0.0)) - 0.12)
    chars_per_line = max(1.0, usable_w * 72.0 / max(1.0, font * 0.50))
    line_count = max(1.0, usable_h * 72.0 / max(1.0, font * 1.18))
    capacity = max(1.0, chars_per_line * line_count * 0.92)
    effective_chars = _effective_char_units(text)
    return {
        "font_pt": round(font, 2),
        "capacity_chars": round(capacity, 1),
        "effective_chars": round(effective_chars, 1),
        "estimated_lines": round(line_count, 2),
        "fill_ratio": round(effective_chars / capacity, 3),
    }


def _estimated_rendered_text_height(record: Dict[str, Any]) -> float:
    bbox = record.get("bbox", {})
    text = record.get("text", "")
    font = record.get("font", {}).get("avg_pt") or _fallback_font(record.get("role", ""))
    usable_w = max(0.1, float(bbox.get("w", 0.0)) - 0.18)
    chars_per_line = max(1.0, usable_w * 72.0 / max(1.0, font * 0.50))
    line_count = max(1.0, math.ceil(_effective_char_units(text) / chars_per_line))
    return line_count * font * 1.18 / 72.0 + 0.08


def _overlap_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    findings = []
    for i, left in enumerate(records):
        if _ignore_overlap_record(left):
            continue
        for right in records[i + 1 :]:
            if _ignore_overlap_record(right):
                continue
            overlap = _intersection(left["bbox"], right["bbox"])
            if overlap <= 0:
                continue
            smaller = max(0.01, min(left["bbox"].get("area", 0.0), right["bbox"].get("area", 0.0)))
            ratio = overlap / smaller
            if _is_allowed_containment(left, right):
                continue
            threshold = float(rules["container_overlap_ratio"]) if "container" in {left.get("role"), right.get("role")} else float(rules["overlap_ratio"])
            if ratio >= threshold:
                severity = "high" if ratio >= 0.35 else "medium"
                findings.append(
                    {
                        "type": "shape_overlap_risk",
                        "severity": severity,
                        **_finding_meta("shape_overlap_risk", severity),
                        "slide_page": slide_index,
                        "message": f"{left.get('role')} overlaps {right.get('role')} by {ratio:.2f} of the smaller shape.",
                        "evidence": {
                            "shape_a": _shape_evidence(left),
                            "shape_b": _shape_evidence(right),
                            "overlap_area_sq_in": round(overlap, 3),
                            "overlap_ratio": round(ratio, 3),
                        },
                        "repair_strategy": "Move proof/table/card systems within their existing composition grid; do not rely on rendered screenshots to discover this.",
                    }
                )
    return findings


def _table_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    findings = []
    containers = _container_records(records)
    for record in records:
        if record.get("role") != "table":
            continue
        table = record.get("table", {})
        bbox = record.get("bbox", {})
        rows = int(table.get("row_count") or 0)
        cols = int(table.get("col_count") or 0)
        if rows <= 1 or cols <= 1:
            findings.append(
                _finding(
                    slide_index,
                    "weak_table_grammar",
                    "medium",
                    "Native table has too few rows or columns to carry table evidence.",
                    record,
                    "Check table extraction before changing layout; table proof should preserve row/column grammar.",
                )
            )
        row_height = bbox.get("h", 0) / max(1, rows)
        col_width = bbox.get("w", 0) / max(1, cols)
        if row_height < float(rules["table_readable_row_height_min_in"]) or col_width < float(rules["table_readable_col_width_min_in"]):
            findings.append(
                _finding(
                    slide_index,
                    "dense_table_readability_risk",
                    "medium",
                    f"Table cell estimate is tight: row height {row_height:.2f}in, column width {col_width:.2f}in.",
                    record,
                    "Preserve table size and simplify visible columns/rows, or move dense detail into appendix-style treatment.",
                )
            )
        empty_ratio = float(table.get("empty_cell_ratio") or 0.0)
        non_empty_cells = int(table.get("non_empty_cells") or 0)
        if (
            cols >= int(rules["table_sparse_column_count_min"])
            and empty_ratio >= float(rules["table_sparse_empty_cell_ratio"])
            and non_empty_cells >= int(rules["table_sparse_min_non_empty_cells"])
        ):
            findings.append(
                {
                    "type": "table_sparse_columns_rendered",
                    "severity": "medium",
                    **_finding_meta("table_sparse_columns_rendered", "medium"),
                    "slide_page": slide_index,
                    "message": f"Table renders {cols} columns with {empty_ratio:.2f} empty cells, producing a sparse/cramped grid.",
                    "evidence": {
                        "table": _shape_evidence(record),
                        "row_count": rows,
                        "col_count": cols,
                        "empty_cell_ratio": empty_ratio,
                        "non_empty_cells": non_empty_cells,
                    },
                    "repair_strategy": "Collapse empty columns, preserve source table spanning as an image, or split the table into focused stage blocks before accepting the slide.",
                }
            )
        wrapping_capacity_units = max(1.0, col_width * float(rules["table_wrapping_effective_units_per_in"]))
        wrapping_risk = (
            col_width <= float(rules["table_wrapping_col_width_max_in"])
            and int(table.get("long_cell_count") or 0) > 0
            and float(table.get("max_cell_effective_units") or 0.0) >= wrapping_capacity_units
        )
        if wrapping_risk:
            findings.append(
                {
                    "type": "table_cell_text_wrapping_risk",
                    "severity": "medium",
                    **_finding_meta("table_cell_text_wrapping_risk", "medium"),
                    "slide_page": slide_index,
                    "message": f"Table column width {col_width:.2f}in is too narrow for long cell text, causing wrapping/readability risk.",
                    "evidence": {
                        "table": _shape_evidence(record),
                        "col_width_in": round(col_width, 3),
                        "max_cell_effective_units": table.get("max_cell_effective_units"),
                        "long_cell_count": table.get("long_cell_count"),
                        "wrapping_capacity_units": round(wrapping_capacity_units, 1),
                    },
                    "repair_strategy": "Increase effective column width, reduce visible columns, or render the source table/diagram as a fitted image when native cells become unreadable.",
                }
            )
        parent = _nearest_container_for_record(record, containers)
        if not parent:
            continue
        parent_box = parent.get("bbox", {})
        view_label = _nearest_table_view_label(record, parent, records)
        if not view_label:
            findings.append(
                {
                    "type": "table_view_label_missing",
                    "severity": "medium",
                    **_finding_meta("table_view_label_missing", "medium"),
                    "slide_page": slide_index,
                    "message": "Table evidence lacks a Focused table view label above the table.",
                    "evidence": {
                        "table": _shape_evidence(record),
                        "container": _shape_evidence(parent),
                    },
                    "repair_strategy": "Add a compact Focused table view label above every table evidence object in the straight-rectangle style.",
                }
            )
        caption = _nearest_table_caption(record, parent, records)
        caption_center_dx = abs(_center_x(caption.get("bbox", {})) - _center_x(parent_box)) if caption else 999.0
        if not caption or not _is_center_aligned_text(caption) or caption_center_dx > float(rules["table_caption_center_tolerance_in"]):
            findings.append(
                {
                    "type": "table_caption_missing_or_not_centered",
                    "severity": "medium",
                    **_finding_meta("table_caption_missing_or_not_centered", "medium"),
                    "slide_page": slide_index,
                    "message": "Table evidence lacks a centered explanatory note below the table.",
                    "evidence": {
                        "table": _shape_evidence(record),
                        "container": _shape_evidence(parent),
                        "caption": _shape_evidence(caption) if caption else None,
                        "caption_center_delta_in": round(caption_center_dx, 3) if caption else None,
                    },
                    "repair_strategy": "Place a short paper/table caption or proof focus note below the table and center-align it.",
                }
            )
        table_area_ratio = bbox.get("area", 0.0) / max(0.01, parent_box.get("area", 0.0))
        table_margins = _inner_margins(bbox, parent_box)
        if (
            parent_box.get("area", 0.0) >= float(rules["wide_panel_min_area_sq_in"])
            and bbox.get("area", 0.0) >= 0.8
            and table_area_ratio <= float(rules["table_underutilized_area_ratio"])
            and max(table_margins.values()) >= float(rules["table_underutilized_margin_in"])
        ):
            findings.append(
                {
                    "type": "table_underutilized_in_evidence_panel",
                    "severity": "medium",
                    **_finding_meta("table_underutilized_in_evidence_panel", "medium"),
                    "slide_page": slide_index,
                    "message": f"Table uses only {table_area_ratio:.2f} of its evidence panel while leaving large internal margins.",
                    "evidence": {
                        "table": _shape_evidence(record),
                        "container": _shape_evidence(parent),
                        "table_area_ratio": round(table_area_ratio, 3),
                        "inner_margins_in": table_margins,
                    },
                    "repair_strategy": "Use a larger focused table or route dense/long tables into the bottom wide evidence rectangle.",
                }
            )
        padding = float(rules["table_container_padding_in"])
        tolerance = float(rules["table_container_overflow_tolerance_in"])
        overflow = {
            "left": round(max(0.0, parent_box.get("x", 0.0) + padding - bbox.get("x", 0.0)), 3),
            "top": round(max(0.0, parent_box.get("y", 0.0) + padding - bbox.get("y", 0.0)), 3),
            "right": round(max(0.0, bbox.get("right", 0.0) - (parent_box.get("right", 0.0) - padding)), 3),
            "bottom": round(max(0.0, bbox.get("bottom", 0.0) - (parent_box.get("bottom", 0.0) - padding)), 3),
        }
        max_overflow = max(overflow.values())
        if max_overflow > tolerance:
            severity = "high" if max_overflow >= 0.12 else "medium"
            findings.append(
                {
                    "type": "table_exceeds_container_bounds",
                    "severity": severity,
                    **_finding_meta("table_exceeds_container_bounds", severity),
                    "slide_page": slide_index,
                    "message": f"Native table exceeds its nearest container padding by up to {max_overflow:.2f}in.",
                    "evidence": {
                        "table": _shape_evidence(record),
                        "container": _shape_evidence(parent),
                        "padding_in": padding,
                        "overflow_in": overflow,
                    },
                    "repair_strategy": "Treat this as a global correctness issue: resize or reflow the table inside the container before applying style polish.",
                }
            )
        table_to_container_ratio = bbox.get("h", 0.0) / max(0.01, parent_box.get("h", 0.0))
        if (
            max_overflow <= tolerance
            and table_to_container_ratio > float(rules["table_container_height_fill_max"])
            and row_height < float(rules["table_readable_row_height_min_in"]) * 1.12
        ):
            findings.append(
                {
                    "type": "table_container_height_mismatch",
                    "severity": "medium",
                    **_finding_meta("table_container_height_mismatch", "medium"),
                    "slide_page": slide_index,
                    "message": f"Table uses {table_to_container_ratio:.2f} of container height while estimated row height is only {row_height:.2f}in.",
                    "evidence": {
                        "table": _shape_evidence(record),
                        "container": _shape_evidence(parent),
                        "height_ratio": round(table_to_container_ratio, 3),
                        "row_height_in": round(row_height, 3),
                    },
                    "repair_strategy": "Compare a taller table container, fewer visible rows, or appendix-style treatment before accepting the fitted result.",
                }
            )
    return findings


def _picture_findings(slide_index: int, records: List[Dict[str, Any]], rules: Dict[str, Any]) -> List[Dict[str, Any]]:
    findings = []
    tolerance = float(rules["picture_aspect_ratio_tolerance"])
    containers = _container_records(records)
    for record in records:
        if not record.get("is_picture"):
            continue
        picture = record.get("picture", {}) or {}
        source_aspect = picture.get("source_aspect")
        parent = _smallest_containing_container(record.get("bbox", {}), containers)
        if source_aspect and parent:
            parent_box = parent.get("bbox", {})
            parent_aspect = parent_box.get("w", 0.0) / max(0.001, parent_box.get("h", 0.0))
            center_dx = abs(_center_x(record.get("bbox", {})) - _center_x(parent_box))
            center_tolerance = max(
                float(rules["figure_image_center_tolerance_in"]),
                parent_box.get("w", 0.0) * float(rules["figure_image_center_tolerance_ratio"]),
            )
            if center_dx > center_tolerance and record.get("bbox", {}).get("area", 0.0) >= 0.8:
                findings.append(
                    _finding(
                        slide_index,
                        "figure_image_off_center_in_panel",
                        "low",
                        f"Figure image center is {center_dx:.2f}in away from the proof panel center, suggesting label space is steering the image off-center.",
                        record,
                        "Treat figure labels as compact annotations, not reserved layout columns; center the image and caption in the full proof panel.",
                    )
                )
            image_area_ratio = record.get("bbox", {}).get("area", 0.0) / max(0.01, parent_box.get("area", 0.0))
            margins = _inner_margins(record.get("bbox", {}), parent_box)
            if (
                parent_aspect >= float(rules["wide_panel_aspect_min"])
                and parent_box.get("area", 0.0) >= float(rules["wide_panel_min_area_sq_in"])
                and record.get("bbox", {}).get("area", 0.0) >= 0.45
                and image_area_ratio <= float(rules["image_underutilized_area_ratio"])
                and max(margins.values()) >= float(rules["image_underutilized_margin_in"])
            ):
                findings.append(
                    {
                        "type": "image_underutilized_in_wide_panel",
                        "severity": "medium",
                        **_finding_meta("image_underutilized_in_wide_panel", "medium"),
                        "slide_page": slide_index,
                        "message": f"Wide panel aspect {parent_aspect:.2f} has a figure using only {image_area_ratio:.2f} of panel area.",
                        "evidence": {
                            "picture": _shape_evidence(record),
                            "container": _shape_evidence(parent),
                            "source_aspect": source_aspect,
                            "parent_aspect": round(parent_aspect, 3),
                            "image_area_ratio": round(image_area_ratio, 3),
                            "inner_margins_in": margins,
                        },
                        "repair_strategy": "Do not blindly keep the bottom-band repair. Compare readable area against a right-panel-large layout and send this slide to human visual review.",
                    }
                )
            caption = _nearest_caption_for_picture(record, parent, records)
            if caption and parent_aspect >= float(rules["wide_panel_aspect_min"]) and parent_box.get("area", 0.0) >= float(rules["wide_panel_min_area_sq_in"]):
                caption_center_dx = abs(_center_x(caption.get("bbox", {})) - _center_x(parent_box))
                if not _is_center_aligned_text(caption) or caption_center_dx > float(rules["figure_caption_center_tolerance_in"]):
                    findings.append(
                        {
                            "type": "figure_caption_not_centered_in_wide_panel",
                            "severity": "medium",
                            **_finding_meta("figure_caption_not_centered_in_wide_panel", "medium"),
                            "slide_page": slide_index,
                            "message": "Wide-panel figure caption is not center-aligned under the image.",
                            "evidence": {
                                "picture": _shape_evidence(record),
                                "caption": _shape_evidence(caption),
                                "container": _shape_evidence(parent),
                                "caption_center_delta_in": round(caption_center_dx, 3),
                            },
                            "repair_strategy": "Center the caption across the proof panel under the fitted image.",
                        }
                    )
            if (
                source_aspect <= float(rules["tall_figure_aspect_max"])
                and parent_aspect > float(rules["tall_figure_panel_aspect_max"])
            ):
                findings.append(
                    _finding(
                        slide_index,
                        "figure_panel_aspect_mismatch",
                        "medium",
                        f"Tall figure source aspect {source_aspect:.2f} is placed in a near-square/wide panel aspect {parent_aspect:.2f}.",
                        record,
                        "Route tall figures to a vertical proof panel and keep figure labels as compact annotations so they do not steal image center.",
                    )
                )
            elif (
                source_aspect >= float(rules["wide_figure_aspect_min"])
                and parent_aspect < float(rules["wide_figure_panel_aspect_min"])
            ):
                findings.append(
                    _finding(
                        slide_index,
                        "figure_panel_aspect_mismatch",
                        "medium",
                        f"Wide figure source aspect {source_aspect:.2f} is placed in a narrow/near-square panel aspect {parent_aspect:.2f}.",
                        record,
                        "Route clearly-wide figures to a horizontal proof panel and keep labels as compact annotations rather than a reserved side column.",
                    )
                )
        distortion = picture.get("aspect_distortion")
        if distortion is None or distortion <= tolerance:
            continue
        findings.append(
            _finding(
                slide_index,
                "figure_picture_aspect_distortion",
                "medium",
                f"Picture box aspect ratio {picture.get('box_aspect')} differs from source image aspect ratio {picture.get('source_aspect')}.",
                record,
                "Fit the figure inside the component box while preserving the source image aspect ratio; route non-square figures to panels that match their source shape.",
            )
        )
    return findings


def _academic_right_evidence_void_findings(slide_index: int, records: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Detect academic bottom-table pages that leave the upper-right evidence area empty."""
    has_key_message = any(_clean_text(record.get("text", "")).lower() == "key message" for record in records)
    if not has_key_message:
        return []
    bottom_tables = [
        record
        for record in records
        if record.get("has_table")
        and record.get("bbox", {}).get("x", 0.0) <= 1.05
        and record.get("bbox", {}).get("w", 0.0) >= 9.8
        and record.get("bbox", {}).get("y", 0.0) >= 4.6
    ]
    if not bottom_tables:
        return []
    right_region = {"x": 6.25, "y": 1.05, "right": 12.5, "bottom": 4.75}
    right_content = []
    for record in records:
        if not _meaningful_for_occupancy(record) or record.get("has_table"):
            continue
        if record.get("role") in {"source_footer", "page_marker", "background", "decorative", "decorative_rule"}:
            continue
        box = record.get("bbox", {})
        if box.get("area", 0.0) < 0.12:
            continue
        if _intersection(box, right_region) >= 0.18:
            right_content.append(record)
    if right_content:
        return []
    return [
        {
            "type": "academic_right_evidence_void",
            "severity": "medium",
            **_finding_meta("academic_right_evidence_void", "medium"),
            "slide_page": slide_index,
            "message": "Academic bottom-table slide leaves the upper-right evidence region empty.",
            "evidence": {
                "bottom_table": _shape_evidence(bottom_tables[0]),
                "right_region": {"x_min": 6.25, "y_min": 1.05, "y_max": 4.75},
            },
            "repair_strategy": "Render metric cards, evidence notes, or a compact proof summary in the right evidence column instead of leaving the region blank.",
        }
    ]


def _academic_toc_canonical_sections_findings(slide_index: int, records: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Detect original academic Contents slides that lost the six-module golden-baseline route."""
    texts = [_clean_text(record.get("text", "")) for record in records if record.get("has_text")]
    text_set = {text for text in texts if text}
    has_academic_contents = "Contents" in text_set and any(
        "A sectioned route through the paper" in text for text in texts
    )
    if not has_academic_contents:
        return []

    expected = ["Motivation", "Method", "Analysis", "Ablations", "Results", "Conclusion"]
    missing = [section for section in expected if section not in text_set]
    if not missing:
        return []
    title_record = next((record for record in records if _clean_text(record.get("text", "")) == "Contents"), records[0])
    return [
        {
            "type": "academic_toc_missing_canonical_sections",
            "severity": "medium",
            **_finding_meta("academic_toc_missing_canonical_sections", "medium"),
            "slide_page": slide_index,
            "message": "Original academic Contents slide does not preserve the six-module golden-baseline route.",
            "evidence": {
                "missing_sections": missing,
                "expected_sections": expected,
                "contents_title": _shape_evidence(title_record),
            },
            "repair_strategy": "Render the original academic table of contents with Motivation, Method, Analysis, Ablations, Results, and Conclusion.",
        }
    ]


def _semantic_content_findings(slide_index: int, records: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Catch content-quality failures that are visible in PPTX text metadata."""
    findings: List[Dict[str, Any]] = []
    for record in records:
        text = _clean_text(record.get("text", ""))
        if _looks_like_fragment_heading(text, record):
            findings.append(
                _finding(
                    slide_index,
                    "weak_fragment_point_heading",
                    "medium",
                    f"Numbered point/title text looks like an unfinished clause: '{text}'.",
                    record,
                    "Regenerate the point claim from the supporting sentence; use a complete idea phrase, not a boilerplate sentence prefix.",
                )
            )
    findings.extend(_spurious_metric_card_findings(slide_index, records))
    return findings


def _looks_like_fragment_heading(text: str, record: Dict[str, Any]) -> bool:
    text = _clean_text(text).strip(" .;:-")
    if not text:
        return False
    role = record.get("role", "")
    bbox = record.get("bbox", {})
    if role == "source_footer" or bbox.get("y", 0.0) > 6.7:
        return False
    words = text.lower().split()
    if len(words) > 10:
        return False
    weak_exact = {
        "in short",
        "taken together",
        "the paper addresses the problem of",
        "its goal is to make it",
        "this meant practitioners could not reliably",
        "even with improved initialization and batch",
        "they can help but are not",
    }
    lower = text.lower()
    if re.fullmatch(r"(?:slide|page)\s*\d+", lower):
        return True
    if lower in weak_exact:
        return True
    if lower.startswith(("the paper addresses", "this paper addresses", "its goal is", "in short", "taken together")):
        return True
    if lower.startswith(("this meant", "even with", "they can help")):
        return True
    weak_endings = {
        "a", "an", "the", "of", "to", "in", "on", "for", "with", "by", "and", "or", "that", "which",
        "it", "its", "not", "but", "while",
    }
    if words and words[-1].strip(" ,;:-()[]").lower() in weak_endings:
        return True
    if text.count("(") != text.count(")"):
        return True
    return False


def _spurious_metric_card_findings(slide_index: int, records: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    findings: List[Dict[str, Any]] = []
    value_records = [
        record
        for record in records
        if record.get("has_text")
        and record.get("role") != "source_footer"
        and record.get("bbox", {}).get("y", 0.0) < 6.7
        and re.fullmatch(r"\d+(?:\.\d+)?", _clean_text(record.get("text", "")))
    ]
    generic_labels = {"accuracy", "rating", "score", "metric", "key metric", "key number"}
    for value_record in value_records:
        value = _clean_text(value_record.get("text", ""))
        try:
            number = float(value)
        except ValueError:
            continue
        if number > 100:
            continue
        value_box = value_record.get("bbox", {})
        for label_record in records:
            if label_record is value_record or not label_record.get("has_text"):
                continue
            label = _clean_text(label_record.get("text", "")).lower()
            if label not in generic_labels:
                continue
            label_box = label_record.get("bbox", {})
            vertical_gap = label_box.get("y", 0.0) - value_box.get("bottom", 0.0)
            x_overlap = _axis_overlap(
                value_box.get("x", 0.0),
                value_box.get("right", 0.0),
                label_box.get("x", 0.0),
                label_box.get("right", 0.0),
            )
            min_width = max(0.01, min(value_box.get("w", 0.0), label_box.get("w", 0.0)))
            if -0.05 <= vertical_gap <= 0.55 and x_overlap / min_width >= 0.35:
                findings.append(
                    _finding(
                        slide_index,
                        "spurious_generic_metric_card",
                        "medium",
                        f"Metric card shows bare value '{value}' with generic label '{label_record.get('text')}', which often comes from dataset/model names rather than a real paper result.",
                        value_record,
                        "Drop the metric card or replace it with a sourced value that has a unit, percentage, layer/count noun, or explicit table/figure provenance.",
                    )
                )
                break
    return findings


def _slide_density_findings(
    slide_index: int,
    records: List[Dict[str, Any]],
    slide_w: float,
    slide_h: float,
    rules: Dict[str, Any],
) -> List[Dict[str, Any]]:
    meaningful = [r for r in records if _meaningful_for_occupancy(r)]
    occupied = sum(r["bbox"].get("area", 0.0) for r in meaningful)
    occupancy = occupied / max(0.01, slide_w * slide_h)
    words = sum(int(r.get("text_words") or 0) for r in meaningful)
    if occupancy < float(rules["low_occupancy_ratio"]) and words < 28 and len(meaningful) <= 5:
        return [
            {
                "type": "sparse_slide_risk",
                "severity": "low",
                **_finding_meta("sparse_slide_risk", "low"),
                "slide_page": slide_index,
                "message": f"Slide has low metadata occupancy ({occupancy:.2f}) and only {words} visible words.",
                "evidence": {"occupancy_ratio": round(occupancy, 3), "visible_words": words, "meaningful_shape_count": len(meaningful)},
                "repair_strategy": "First enrich or rebalance text hierarchy; avoid shrinking well-composed components solely to fill whitespace.",
            }
        ]
    return []


def _slide_summary(
    slide_index: int,
    records: List[Dict[str, Any]],
    findings: List[Dict[str, Any]],
    slide_w: float,
    slide_h: float,
) -> Dict[str, Any]:
    meaningful = [r for r in records if _meaningful_for_occupancy(r)]
    occupied = sum(r["bbox"].get("area", 0.0) for r in meaningful)
    role_counts: Dict[str, int] = {}
    for record in records:
        role_counts[record.get("role", "unknown")] = role_counts.get(record.get("role", "unknown"), 0) + 1
    return {
        "page": slide_index,
        "shape_count": len(records),
        "meaningful_shape_count": len(meaningful),
        "visible_words": sum(int(r.get("text_words") or 0) for r in meaningful),
        "occupancy_ratio": round(occupied / max(0.01, slide_w * slide_h), 3),
        "role_counts": role_counts,
        "finding_count": len(findings),
        "finding_types": sorted({finding["type"] for finding in findings}),
    }


def _summarize_findings(findings: List[Dict[str, Any]], rules: Dict[str, Any]) -> Dict[str, Any]:
    by_type: Dict[str, int] = {}
    by_severity: Dict[str, int] = {}
    by_problem_type: Dict[str, int] = {}
    by_dimension: Dict[str, int] = {}
    dimension_penalty: Dict[str, float] = {}
    severity_weight = {"high": 12.0, "medium": 6.0, "low": 2.0}
    for finding in findings:
        by_type[finding["type"]] = by_type.get(finding["type"], 0) + 1
        by_severity[finding["severity"]] = by_severity.get(finding["severity"], 0) + 1
        problem_type = finding.get("problem_type", _problem_type_for_finding(finding["type"]))
        by_problem_type[problem_type] = by_problem_type.get(problem_type, 0) + 1
        dimension = finding.get("dimension", _rule_metadata(finding["type"]).get("dimension", "layout"))
        by_dimension[dimension] = by_dimension.get(dimension, 0) + 1
        dimension_penalty[dimension] = dimension_penalty.get(dimension, 0.0) + severity_weight.get(finding.get("severity", "low"), 2.0)
    dimension_scores = {
        dimension: max(0.0, round(100.0 - dimension_penalty.get(dimension, 0.0), 1))
        for dimension in DIMENSIONS
    }
    return {
        "finding_count": len(findings),
        "by_type": by_type,
        "by_severity": by_severity,
        "by_problem_type": by_problem_type,
        "by_dimension": by_dimension,
        "dimension_scores": dimension_scores,
        "dimension_score_model": "100 minus severity penalties: high=12, medium=6, low=2",
        "high_risk_pages": sorted({finding["slide_page"] for finding in findings if finding["severity"] == "high"}),
        "medium_risk_pages": sorted({finding["slide_page"] for finding in findings if finding["severity"] == "medium"}),
        "typography_pages": sorted({finding["slide_page"] for finding in findings if finding.get("problem_type") == "typography"}),
        "copy_fitting_pages": sorted({finding["slide_page"] for finding in findings if finding.get("problem_type") == "copy_fitting"}),
        "geometry_pages": sorted({finding["slide_page"] for finding in findings if finding.get("problem_type") == "geometry"}),
        "optical_balance_pages": sorted({finding["slide_page"] for finding in findings if finding.get("problem_type") == "optical_balance"}),
        "deck_flags": _deck_flags(by_type, rules),
    }


def _finding(
    slide_index: int,
    kind: str,
    severity: str,
    message: str,
    record: Dict[str, Any],
    repair_strategy: str,
) -> Dict[str, Any]:
    return {
        "type": kind,
        "severity": severity,
        **_finding_meta(kind, severity),
        "slide_page": slide_index,
        "message": message,
        "evidence": {"shape": _shape_evidence(record)},
        "repair_strategy": repair_strategy,
    }


def _shape_evidence(record: Dict[str, Any]) -> Dict[str, Any]:
    text = record.get("text", "")
    return {
        "index": record.get("index"),
        "role": record.get("role"),
        "bbox": record.get("bbox"),
        "font": record.get("font"),
        "paragraph_alignment": record.get("paragraph_alignment", ""),
        "vertical_anchor": record.get("vertical_anchor", ""),
        "text_preview": text[:90],
        "text_words": record.get("text_words", 0),
        "text_capacity": record.get("text_capacity", {}),
    }


def _font_floor(role: str, rules: Dict[str, Any]) -> float:
    if role == "title_claim":
        return float(rules["title_claim_min_pt"])
    if role == "support_body":
        return float(rules["support_min_pt"])
    if role == "card_text":
        return float(rules["card_text_min_pt"])
    if role == "body_text":
        return float(rules["body_min_pt"])
    return 0.0


def _font_ideal_floor(role: str, rules: Dict[str, Any]) -> float:
    if role == "title_claim":
        return float(rules["title_claim_ideal_min_pt"])
    if role == "support_body":
        return float(rules["support_ideal_min_pt"])
    if role == "card_text":
        return float(rules["card_text_ideal_min_pt"])
    if role == "body_text":
        return float(rules["body_ideal_min_pt"])
    return 0.0


def _finding_meta(kind: str, severity: str) -> Dict[str, Any]:
    problem_type = _problem_type_for_finding(kind)
    rule_meta = _rule_metadata(kind)
    return {
        "problem_type": problem_type,
        "dimension": rule_meta["dimension"],
        "scope": rule_meta["scope"],
        "style_scope": rule_meta["style_scope"],
        "repair_mode": rule_meta["repair_mode"],
        "confidence": rule_meta["confidence"],
        "human_outcome": rule_meta["human_outcome"],
        "repair_priority": _repair_priority(kind, severity),
        "geometry_repair_allowed": problem_type == "geometry" or kind == "component_frame_overallocated_after_text_fit",
        "preferred_repair_scope": _preferred_repair_scope(problem_type),
    }


def _problem_type_for_finding(kind: str) -> str:
    if kind in {"low_font_size", "below_ideal_font_band", "cover_left_typography_underpowered"}:
        return "typography"
    if kind in {"estimated_text_overflow", "near_text_capacity", "low_text_density", "sparse_slide_risk", "card_copy_imbalance"}:
        return "copy_fitting"
    if kind in {
        "shape_overlap_risk",
        "text_exceeds_container_bounds",
        "weak_table_grammar",
        "dense_table_readability_risk",
        "table_exceeds_container_bounds",
        "table_container_height_mismatch",
        "table_readability_after_fit",
        "flow_nodes_overpacked",
        "flow_grid_alignment_drift",
        "component_boundary_inset_violation",
        "figure_picture_aspect_distortion",
    }:
        return "geometry"
    if kind in {
        "figure_panel_aspect_mismatch",
        "image_underutilized_in_wide_panel",
        "figure_caption_not_centered_in_wide_panel",
        "table_view_label_missing",
        "table_caption_missing_or_not_centered",
        "table_underutilized_in_evidence_panel",
        "text_card_vertical_alignment_top_heavy",
        "table_sparse_columns_rendered",
        "table_cell_text_wrapping_risk",
        "metric_label_gap_too_large",
        "container_stack_off_balance",
        "paired_label_body_gap_too_large",
        "component_frame_overallocated_after_text_fit",
        "card_internal_spacing_not_scaled_to_frame",
        "agenda_read_path_header_too_close",
        "table_support_band_off_balance",
        "figure_image_off_center_in_panel",
        "figure_badge_identity_label_conflation",
        "figure_label_anchor_drift",
        "figure_label_text_alignment_off_center",
        "panel_identity_label_anchor_drift",
        "panel_identity_label_text_alignment_off_center",
        "stacked_figure_identity_label_overcorrection",
        "academic_right_evidence_void",
        "academic_toc_missing_canonical_sections",
    }:
        return "optical_balance"
    if kind in {
        "metric_improved_visual_regressed",
        "likely_overcorrection",
        "style_scope_mismatch",
        "repair_introduced_new_findings",
        "image_legibility_regression",
        "layout_rhythm_regression",
    }:
        return "repair_risk"
    if kind in {"weak_fragment_point_heading", "spurious_generic_metric_card"}:
        return "content_semantics"
    return "metadata"


def _repair_priority(kind: str, severity: str) -> str:
    if severity == "high":
        return "P1"
    if kind in {
        "shape_overlap_risk",
        "text_exceeds_container_bounds",
        "weak_table_grammar",
        "dense_table_readability_risk",
        "table_exceeds_container_bounds",
        "table_container_height_mismatch",
        "table_view_label_missing",
        "table_caption_missing_or_not_centered",
        "table_sparse_columns_rendered",
        "table_cell_text_wrapping_risk",
        "table_underutilized_in_evidence_panel",
        "component_boundary_inset_violation",
        "figure_picture_aspect_distortion",
    }:
        return "P2"
    if kind in {
        "figure_panel_aspect_mismatch",
        "image_underutilized_in_wide_panel",
        "figure_caption_not_centered_in_wide_panel",
        "table_view_label_missing",
        "table_caption_missing_or_not_centered",
        "text_card_vertical_alignment_top_heavy",
        "metric_improved_visual_regressed",
    }:
        return "P2"
    if kind in {"low_font_size", "estimated_text_overflow", "near_text_capacity"}:
        return "P2" if severity == "medium" else "P3"
    if kind in {"weak_fragment_point_heading", "spurious_generic_metric_card"}:
        return "P2"
    return "P3"


def _preferred_repair_scope(problem_type: str) -> str:
    if problem_type == "typography":
        return "font hierarchy, weight, and line-height before layout changes"
    if problem_type == "copy_fitting":
        return "copy allocation, line breaks, richer notes, or notes splitting before resizing"
    if problem_type == "geometry":
        return "minimal component position/size repair only because a structural signal failed"
    if problem_type == "optical_balance":
        return "internal text-stack spacing and padding before macro-layout changes"
    if problem_type == "content_semantics":
        return "content curation and sourced claim/metric repair before layout changes"
    if problem_type == "repair_risk":
        return "human review and style-scope validation before accepting the repair"
    return "metadata review"


def _rule_metadata(kind: str) -> Dict[str, Any]:
    meta = dict(DEFAULT_RULE_METADATA)
    meta.update(RULE_METADATA.get(kind, {}))
    meta["confidence"] = round(float(meta.get("confidence", 0.55)), 2)
    meta["style_scope"] = list(meta.get("style_scope") or [])
    return meta


def _rule_registry_snapshot(findings: List[Dict[str, Any]]) -> Dict[str, Any]:
    observed = {finding.get("type", "") for finding in findings if finding.get("type")}
    rules = observed | set(RULE_METADATA)
    return {kind: _rule_metadata(kind) for kind in sorted(rules)}


def _deck_flags(by_type: Dict[str, int], rules: Dict[str, Any]) -> List[str]:
    flags = []
    typography_signals = (
        by_type.get("below_ideal_font_band", 0)
        + by_type.get("low_font_size", 0)
        + by_type.get("cover_left_typography_underpowered", 0)
    )
    if typography_signals >= int(rules["deck_typography_signal_count"]):
        flags.append("deck_type_scale_under_comfort_band")
    return flags


def _fallback_font(role: str) -> float:
    return {
        "title_claim": 24.0,
        "support_body": 12.0,
        "card_text": 9.0,
        "body_text": 11.0,
        "small_text": 8.0,
    }.get(role, 10.0)


def _effective_char_units(text: str) -> float:
    total = 0.0
    for char in text:
        if char.isspace():
            total += 0.35
        elif ord(char) > 127:
            total += 1.75
        elif char in ".,;:!|":
            total += 0.35
        else:
            total += 1.0
    return total


def _container_records(records: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    return [
        record
        for record in records
        if record.get("role") == "container"
        and not record.get("is_full_background")
        and record.get("bbox", {}).get("area", 0.0) > 0.01
    ]


def _contained_records(container: Dict[str, Any], records: List[Dict[str, Any]], include_containers: bool) -> List[Dict[str, Any]]:
    result = []
    outer = container.get("bbox", {})
    for record in records:
        if record is container:
            continue
        if record.get("role") == "container" and not include_containers:
            continue
        if record.get("is_full_background"):
            continue
        bbox = record.get("bbox", {})
        if bbox.get("area", 0.0) <= 0.0:
            continue
        if _center_inside(bbox, outer):
            result.append(record)
    return result


def _nearest_container_for_record(record: Dict[str, Any], containers: List[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
    bbox = record.get("bbox", {})
    containing = [container for container in containers if _center_inside(bbox, container.get("bbox", {}))]
    if containing:
        return min(containing, key=lambda item: item.get("bbox", {}).get("area", 0.0))
    overlapping = []
    for container in containers:
        container_box = container.get("bbox", {})
        overlap = _intersection(bbox, container_box)
        if overlap <= 0:
            continue
        overlapping.append((overlap / max(0.01, bbox.get("area", 0.0)), container))
    if overlapping:
        return max(overlapping, key=lambda item: item[0])[1]
    return None


def _smallest_containing_container(label_box: Dict[str, float], containers: List[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
    containing = [container for container in containers if _center_inside(label_box, container.get("bbox", {}))]
    if not containing:
        return None
    return min(containing, key=lambda item: item.get("bbox", {}).get("area", 0.0))


def _expected_text_stack_height(children: List[Dict[str, Any]], container_box: Dict[str, float]) -> float:
    sorted_children = sorted(children, key=lambda item: item.get("bbox", {}).get("y", 0.0))
    usable_w = max(0.1, container_box.get("w", 0.0) - 0.36)
    text_height = 0.0
    for child in sorted_children:
        font = child.get("font", {}).get("avg_pt") or _fallback_font(child.get("role", ""))
        char_units = _effective_char_units(child.get("text", ""))
        chars_per_line = max(1.0, usable_w * 72.0 / max(1.0, font * 0.50))
        lines_needed = max(1.0, math.ceil(char_units / chars_per_line))
        text_height += lines_needed * font * 1.18 / 72.0
    gaps = max(0, len(sorted_children) - 1) * 0.12
    vertical_margins = 0.34
    return text_height + gaps + vertical_margins


def _bbox_union(boxes: List[Dict[str, float]]) -> Dict[str, float]:
    if not boxes:
        return {"x": 0.0, "y": 0.0, "w": 0.0, "h": 0.0, "right": 0.0, "bottom": 0.0, "area": 0.0}
    left = min(float(box.get("x", 0.0)) for box in boxes)
    top = min(float(box.get("y", 0.0)) for box in boxes)
    right = max(float(box.get("right", 0.0)) for box in boxes)
    bottom = max(float(box.get("bottom", 0.0)) for box in boxes)
    return {
        "x": round(left, 3),
        "y": round(top, 3),
        "w": round(max(0.0, right - left), 3),
        "h": round(max(0.0, bottom - top), 3),
        "right": round(right, 3),
        "bottom": round(bottom, 3),
        "area": round(max(0.0, right - left) * max(0.0, bottom - top), 3),
    }


def _looks_metric_value(text: str) -> bool:
    return bool(re.fullmatch(r"\d+(?:\.\d+)?(?:%|[TBMK])?", _clean_text(text)))


def _ignore_overlap_record(record: Dict[str, Any]) -> bool:
    return record.get("role") in {"background", "decorative", "decorative_rule", "source_footer", "page_marker", "component_label"} or record.get("bbox", {}).get("area", 0) <= 0.01


def _meaningful_for_occupancy(record: Dict[str, Any]) -> bool:
    if record.get("role") in {"background", "decorative", "decorative_rule", "source_footer", "page_marker", "component_label"}:
        return False
    return bool(record.get("has_text") or record.get("has_table") or record.get("is_picture") or record.get("role") == "container")


def _is_allowed_containment(a: Dict[str, Any], b: Dict[str, Any]) -> bool:
    roles = {a.get("role"), b.get("role")}
    if roles == {"container"}:
        return True
    if "container" not in roles:
        return False
    container = a if a.get("role") == "container" else b
    contained = b if a.get("role") == "container" else a
    if contained.get("role") in {"picture", "table"}:
        return _center_inside(contained["bbox"], container["bbox"])
    if contained.get("role") == "title_claim":
        return _center_inside(contained["bbox"], container["bbox"])
    if contained.get("role") in {"card_label", "card_text", "small_text", "body_text", "support_body"}:
        return _center_inside(contained["bbox"], container["bbox"])
    return False


def _center_inside(inner: Dict[str, float], outer: Dict[str, float]) -> bool:
    cx = inner.get("x", 0.0) + inner.get("w", 0.0) / 2
    cy = inner.get("y", 0.0) + inner.get("h", 0.0) / 2
    return outer.get("x", 0.0) <= cx <= outer.get("right", 0.0) and outer.get("y", 0.0) <= cy <= outer.get("bottom", 0.0)


def _center_x(box: Dict[str, float]) -> float:
    return box.get("x", 0.0) + box.get("w", 0.0) / 2


def _inner_margins(inner: Dict[str, float], outer: Dict[str, float]) -> Dict[str, float]:
    return {
        "left": round(max(0.0, inner.get("x", 0.0) - outer.get("x", 0.0)), 3),
        "top": round(max(0.0, inner.get("y", 0.0) - outer.get("y", 0.0)), 3),
        "right": round(max(0.0, outer.get("right", 0.0) - inner.get("right", 0.0)), 3),
        "bottom": round(max(0.0, outer.get("bottom", 0.0) - inner.get("bottom", 0.0)), 3),
    }


def _nearest_caption_for_picture(
    picture: Dict[str, Any],
    parent: Dict[str, Any],
    records: List[Dict[str, Any]],
) -> Optional[Dict[str, Any]]:
    picture_box = picture.get("bbox", {})
    parent_box = parent.get("bbox", {})
    candidates = []
    for record in records:
        if record is picture or not record.get("has_text"):
            continue
        if record.get("role") in {"source_footer", "page_marker", "component_label"}:
            continue
        bbox = record.get("bbox", {})
        if not _center_inside(bbox, parent_box):
            continue
        if bbox.get("y", 0.0) < picture_box.get("bottom", 0.0) - 0.08:
            continue
        if bbox.get("h", 0.0) > 0.48 or (record.get("font", {}).get("avg_pt") or 99) > 10.5:
            continue
        candidates.append(record)
    if not candidates:
        return None
    return min(candidates, key=lambda item: abs(item.get("bbox", {}).get("y", 0.0) - picture_box.get("bottom", 0.0)))


def _nearest_table_view_label(
    table: Dict[str, Any],
    parent: Dict[str, Any],
    records: List[Dict[str, Any]],
) -> Optional[Dict[str, Any]]:
    table_box = table.get("bbox", {})
    parent_box = parent.get("bbox", {})
    candidates = []
    for record in records:
        if record is table or not record.get("has_text"):
            continue
        text_token = _label_token(record.get("text", ""))
        if "TABLEVIEW" not in text_token:
            continue
        bbox = record.get("bbox", {})
        if not _center_inside(bbox, parent_box):
            continue
        if bbox.get("bottom", 0.0) > table_box.get("y", 0.0) + 0.12:
            continue
        candidates.append(record)
    if not candidates:
        return None
    return min(candidates, key=lambda item: abs(item.get("bbox", {}).get("bottom", 0.0) - table_box.get("y", 0.0)))


def _nearest_table_caption(
    table: Dict[str, Any],
    parent: Dict[str, Any],
    records: List[Dict[str, Any]],
) -> Optional[Dict[str, Any]]:
    table_box = table.get("bbox", {})
    parent_box = parent.get("bbox", {})
    candidates = []
    for record in records:
        if record is table or not record.get("has_text"):
            continue
        if record.get("role") in {"source_footer", "page_marker", "component_label"}:
            continue
        bbox = record.get("bbox", {})
        if not _center_inside(bbox, parent_box):
            continue
        if bbox.get("y", 0.0) < table_box.get("bottom", 0.0) - 0.08:
            continue
        if bbox.get("h", 0.0) > 0.55 or (record.get("font", {}).get("avg_pt") or 99) > 10.8:
            continue
        candidates.append(record)
    if not candidates:
        return None
    return min(candidates, key=lambda item: abs(item.get("bbox", {}).get("y", 0.0) - table_box.get("bottom", 0.0)))


def _has_middle_vertical_anchor(record: Dict[str, Any]) -> bool:
    return "MIDDLE" in str(record.get("vertical_anchor", "")).upper()


def _intersection(a: Dict[str, float], b: Dict[str, float]) -> float:
    x1 = max(float(a.get("x", 0.0)), float(b.get("x", 0.0)))
    y1 = max(float(a.get("y", 0.0)), float(b.get("y", 0.0)))
    x2 = min(float(a.get("right", 0.0)), float(b.get("right", 0.0)))
    y2 = min(float(a.get("bottom", 0.0)), float(b.get("bottom", 0.0)))
    if x2 <= x1 or y2 <= y1:
        return 0.0
    return (x2 - x1) * (y2 - y1)


def _horizontal_overlap_ratio(a: Dict[str, float], b: Dict[str, float]) -> float:
    x1 = max(float(a.get("x", 0.0)), float(b.get("x", 0.0)))
    x2 = min(float(a.get("right", 0.0)), float(b.get("right", 0.0)))
    if x2 <= x1:
        return 0.0
    return (x2 - x1) / max(0.01, min(float(a.get("w", 0.0)), float(b.get("w", 0.0))))


def _axis_overlap(start_a: float, end_a: float, start_b: float, end_b: float) -> float:
    left = max(float(start_a), float(start_b))
    right = min(float(end_a), float(end_b))
    if right <= left:
        return 0.0
    return right - left


def _emu_to_inches(value: Any) -> float:
    try:
        return float(value) / EMU_PER_INCH
    except Exception:
        return 0.0


def _clean_text(text: Any) -> str:
    return re.sub(r"\s+", " ", str(text or "")).strip()


def _label_token(text: str) -> str:
    return re.sub(r"[^A-Z0-9]+", "", str(text or "").upper())


def _is_figure_identity_label(text: str) -> bool:
    return bool(re.match(r"^Figure\s+\d+\b", str(text or ""), flags=re.IGNORECASE)) or bool(re.match(r"^FIGURE\d+$", _label_token(text)))


def _is_stacked_figure_identity_text(text: str) -> bool:
    return bool(re.match(r"^F\s+I\s+G\s+U\s+R\s+E\s+\d+\b", _clean_text(text), flags=re.IGNORECASE))


def _is_center_aligned_text(record: Dict[str, Any]) -> bool:
    return "CENTER" in str(record.get("paragraph_alignment", "")).upper()


def _panel_identity_label_child(children: List[Dict[str, Any]], parent_box: Dict[str, float]) -> Optional[Dict[str, Any]]:
    header_floor = parent_box.get("y", 0.0) + 1.18
    candidates = []
    for child in children:
        text = _clean_text(child.get("text", ""))
        if not text or child.get("role") == "component_label":
            continue
        if child.get("role") in {"page_marker", "source_footer"}:
            continue
        if child.get("bbox", {}).get("y", 0.0) > header_floor:
            continue
        if len(text.split()) > 6 or len(text) > 48:
            continue
        candidates.append(child)
    if not candidates:
        return None
    return max(candidates, key=lambda item: (item.get("font", {}).get("avg_pt") or 0.0, -item.get("bbox", {}).get("y", 0.0)))


def _panel_identity_target_child(children: List[Dict[str, Any]], identity: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    identity_bottom = identity.get("bbox", {}).get("bottom", 0.0)
    structured = [
        child
        for child in children
        if child is not identity
        and (child.get("has_table") or child.get("is_picture"))
        and child.get("bbox", {}).get("y", 0.0) >= identity_bottom - 0.08
    ]
    if structured:
        return max(structured, key=lambda item: item.get("bbox", {}).get("area", 0.0))
    text_targets = [
        child
        for child in children
        if child is not identity
        and child.get("has_text")
        and child.get("role") != "component_label"
        and child.get("bbox", {}).get("y", 0.0) >= identity_bottom - 0.04
        and len(_clean_text(child.get("text", "")).split()) > 3
    ]
    if not text_targets:
        return None
    return max(text_targets, key=lambda item: item.get("bbox", {}).get("area", 0.0))


def _primary_picture_child(records: List[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
    pictures = [record for record in records if record.get("is_picture")]
    if not pictures:
        return None
    return max(pictures, key=lambda record: record.get("bbox", {}).get("area", 0.0))


def main(argv: Optional[List[str]] = None) -> int:
    parser = argparse.ArgumentParser(description="Inspect PPTX geometry and text capacity without rendering slides.")
    parser.add_argument("--pptx", required=True, help="Path to PPTX file.")
    parser.add_argument("--output", help="Optional path to write nonvisual audit JSON.")
    args = parser.parse_args(argv)

    audit = inspect_pptx_nonvisual(Path(args.pptx))
    if args.output:
        output = Path(args.output)
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_text(json.dumps(audit, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(audit["summary"], ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
