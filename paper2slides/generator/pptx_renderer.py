"""
Deterministic PPTX renderer for structured slide specs.

The renderer intentionally stays native-PPTX and editable, but it behaves more
like a small layout engine than a fixed template: it applies a restrained deck
theme, reserves safe regions, adapts type size to text length, and uses
different compositions for metric, visual, table, and cover slides.
"""
from __future__ import annotations

import re
from dataclasses import dataclass
from dataclasses import replace
from datetime import datetime
from pathlib import Path
from typing import Iterable, List, Sequence

from .slide_schema import MetricBlock, PresentationSpec, SlideSpec, TextBlock
from .style_presets import get_style_preset


ACADEMIC_BASELINE_SECTIONS = [
    "Motivation",
    "Method",
    "Analysis",
    "Ablations",
    "Results",
    "Conclusion",
]


@dataclass(frozen=True)
class DeckTheme:
    background: tuple[int, int, int] = (247, 248, 250)
    surface: tuple[int, int, int] = (255, 255, 255)
    surface_alt: tuple[int, int, int] = (242, 246, 248)
    ink: tuple[int, int, int] = (24, 32, 45)
    muted: tuple[int, int, int] = (84, 96, 112)
    rule: tuple[int, int, int] = (205, 214, 224)
    primary: tuple[int, int, int] = (10, 115, 112)
    secondary: tuple[int, int, int] = (45, 76, 138)
    accent: tuple[int, int, int] = (128, 92, 45)
    pale_primary: tuple[int, int, int] = (229, 244, 242)
    pale_secondary: tuple[int, int, int] = (232, 238, 248)
    pale_neutral: tuple[int, int, int] = (241, 243, 246)
    title_font: str = "Aptos Display"
    body_font: str = "Aptos"
    layout_family: str = "academic"
    display_name: str = "Academic Golden Baseline"


class PptxRenderer:
    """Render structured slide specs into an editable PPTX file."""

    def __init__(self, theme: DeckTheme | None = None, style: str = "academic"):
        self.style = (style or "academic").strip().lower()
        self.theme = theme or self._theme_from_style(self.style)

    def _theme_from_style(self, style: str) -> DeckTheme:
        preset = get_style_preset(style)
        theme = DeckTheme()
        if preset.theme:
            theme = replace(theme, **preset.theme)
        return replace(theme, layout_family=preset.layout_family, display_name=preset.display_name)

    def render(self, spec: PresentationSpec, output_path: Path) -> Path:
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.util import Inches, Pt
        except ImportError as exc:
            raise RuntimeError(
                "python-pptx is required for PPTX rendering. Install dependencies from requirements.txt."
            ) from exc

        self.RGBColor = RGBColor
        self.MSO_AUTO_SHAPE_TYPE = MSO_AUTO_SHAPE_TYPE
        self.PP_ALIGN = PP_ALIGN
        self.MSO_ANCHOR = MSO_ANCHOR
        self.Inches = Inches
        self.Pt = Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]
        render_index = 1
        sections = self._section_sequence(spec)

        title_slide = prs.slides.add_slide(blank_layout)
        self._paint_background(title_slide, prs, title=True)
        self._render_title_page(title_slide, spec, sections, render_index)
        render_index += 1

        toc_slide = prs.slides.add_slide(blank_layout)
        self._paint_background(toc_slide, prs)
        self._render_toc(toc_slide, sections, render_index)
        render_index += 1

        content_slides = list(spec.slides)
        if content_slides and content_slides[0].section_type == "opening":
            content_slides = content_slides[1:]

        last_section = ""
        for slide_spec in content_slides:
            section = (slide_spec.section_label or "").strip()
            if section and section != last_section and slide_spec.section_type != "opening":
                divider = prs.slides.add_slide(blank_layout)
                self._paint_background(divider, prs)
                self._render_section_divider(divider, section, render_index)
                render_index += 1
                last_section = section
            elif section:
                last_section = section

            slide = prs.slides.add_slide(blank_layout)
            self._paint_background(slide, prs)

            layout = self._normalized_layout(slide_spec)
            if layout == "cover":
                self._render_cover(slide, slide_spec, render_index)
            elif layout in {"statement", "metric_focus", "closing"} and not slide_spec.image_blocks and not slide_spec.table_blocks:
                self._render_statement(slide, slide_spec, render_index, closing=(layout == "closing"))
            elif layout == "table_focus":
                self._render_table_focus(slide, slide_spec, render_index)
            else:
                self._render_visual_or_mixed(slide, slide_spec, render_index, visual_left=(layout == "visual_left"))
            render_index += 1

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        return output_path

    def _normalized_layout(self, slide_spec: SlideSpec) -> str:
        layout = (slide_spec.layout or "auto").lower()
        allowed = {"cover", "statement", "metric_focus", "visual_right", "visual_left", "table_focus", "quote", "closing"}
        if slide_spec.section_type == "opening":
            return "cover"
        if slide_spec.section_type == "ending":
            return "closing"
        if layout not in allowed and layout not in {"section", "auto", ""}:
            if slide_spec.table_blocks:
                return "table_focus"
            if slide_spec.image_blocks:
                return "visual_right"
            if slide_spec.metric_blocks:
                return "metric_focus"
            return "statement"
        if layout == "auto":
            if slide_spec.table_blocks:
                return "table_focus"
            if slide_spec.image_blocks:
                return "visual_right"
            if slide_spec.metric_blocks:
                return "metric_focus"
            return "statement"
        if layout in {"visual_left", "visual_right"} and not slide_spec.image_blocks:
            if slide_spec.table_blocks:
                return "table_focus"
            if slide_spec.metric_blocks:
                return "metric_focus"
            return "statement"
        if layout == "table_focus" and not slide_spec.table_blocks:
            if slide_spec.image_blocks:
                return "visual_right"
            if slide_spec.metric_blocks:
                return "metric_focus"
            return "statement"
        if layout == "quote" and not slide_spec.image_blocks and not slide_spec.table_blocks:
            return "metric_focus" if slide_spec.metric_blocks else "statement"
        return layout

    def _paint_background(self, slide, prs, title: bool = False) -> None:
        t = self.theme
        family = t.layout_family
        bg = slide.shapes.add_shape(self.MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._rgb(t.background)
        bg.line.fill.background()

        if family == "editorial":
            self._rounded_rect(slide, self.Inches(0.72), self.Inches(0.42), self.Inches(11.9), self.Inches(0.015), t.rule, t.rule)
            self._rounded_rect(slide, self.Inches(0.72), self.Inches(6.86), self.Inches(11.9), self.Inches(0.015), t.rule, t.rule)
            if title:
                self._rounded_rect(slide, self.Inches(0.72), self.Inches(0.88), self.Inches(0.1), self.Inches(5.35), t.primary, t.primary)
            return

        rail = slide.shapes.add_shape(self.MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, self.Inches(0.06))
        rail.fill.solid()
        rail.fill.fore_color.rgb = self._rgb(t.primary)
        rail.line.fill.background()

        short_rail = slide.shapes.add_shape(self.MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, self.Inches(1.75), self.Inches(0.06))
        short_rail.fill.solid()
        short_rail.fill.fore_color.rgb = self._rgb(t.secondary)
        short_rail.line.fill.background()

        if title:
            band = slide.shapes.add_shape(
                self.MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                self.Inches(0),
                self.Inches(0.06),
                self.Inches(3.05),
                int(prs.slide_height - self.Inches(0.06)),
            )
            band.fill.solid()
            band.fill.fore_color.rgb = self._rgb(t.pale_secondary)
            band.line.fill.background()
            accent = slide.shapes.add_shape(
                self.MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                self.Inches(2.86),
                self.Inches(0.06),
                self.Inches(0.16),
                int(prs.slide_height - self.Inches(0.06)),
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = self._rgb(t.secondary)
            accent.line.fill.background()

        if family == "systems":
            for i in range(1, 12):
                x = self.Inches(0.72 + i * 1.0)
                self._rounded_rect(slide, x, self.Inches(0.2), self.Inches(0.006), self.Inches(6.55), t.rule, t.rule)
            for i in range(1, 7):
                y = self.Inches(0.35 + i * 0.95)
                self._rounded_rect(slide, self.Inches(0.56), y, self.Inches(12.05), self.Inches(0.006), t.rule, t.rule)

    def _section_sequence(self, spec: PresentationSpec) -> List[str]:
        sections: List[str] = []
        for slide in spec.slides:
            label = (slide.section_label or "").strip()
            if not label or slide.section_type == "opening":
                continue
            if label not in sections:
                sections.append(label)
        if self._uses_mature_academic_toc(sections):
            return list(ACADEMIC_BASELINE_SECTIONS)
        if not sections:
            sections = ["Motivation", "Core Ideas", "Method", "Results", "Conclusion"]
        return sections[:7]

    def _uses_mature_academic_toc(self, sections: Sequence[str]) -> bool:
        if self.style != "academic" or self.theme.layout_family != "academic":
            return False
        labels = set(sections)
        return (
            "Motivation" in labels
            and "Method" in labels
            and ("Results" in labels or "Conclusion" in labels)
        )

    def _deck_identity(self, spec: PresentationSpec) -> tuple[str, str, str, str]:
        first = spec.slides[0] if spec.slides else None
        source = " ".join(
            part
            for part in [
                first.title if first else "",
                first.takeaway if first else "",
                " ".join(block.text for block in first.text_blocks[:2]) if first else "",
            ]
            if part
        )

        title = (spec.title or "").strip()
        if not title or title.lower() in {"paper2slides presentation", "presentation"}:
            match = re.search(r"Title:\s*(.*?)(?:\s+Authors?:|$)", source, flags=re.IGNORECASE)
            title = match.group(1).strip() if match else (first.title if first else "Paper Presentation")
        title = re.sub(r"^\s*Title:\s*", "", title, flags=re.IGNORECASE)
        title = re.split(r"\s+Title:\s*", title, maxsplit=1, flags=re.IGNORECASE)[0].strip()

        author_match = re.search(r"Authors?:\s*(.*?)(?:\s+Affiliations?:|$)", source, flags=re.IGNORECASE)
        authors = author_match.group(1).strip() if author_match else ""
        authors = re.split(r"\s+Title:\s*", authors, maxsplit=1, flags=re.IGNORECASE)[0].strip()
        affiliation_match = re.search(r"Affiliations?:\s*(.*)$", source, flags=re.IGNORECASE)
        affiliation = affiliation_match.group(1).strip() if affiliation_match else ""

        if not authors:
            authors = "Authors from source paper"
        if not affiliation:
            affiliation = "Academic paper presentation"
        date_text = datetime.now().strftime("%Y-%m-%d")
        return self._truncate(title, 118), self._truncate(authors, 76), self._truncate(affiliation, 76), date_text

    def _render_title_page(self, slide, spec: PresentationSpec, sections: Sequence[str], slide_index: int) -> None:
        t = self.theme
        title, authors, affiliation, date_text = self._deck_identity(spec)
        if t.layout_family == "editorial":
            self._render_editorial_title_page(slide, spec, sections, slide_index, title, authors, affiliation, date_text)
            return
        if t.layout_family == "systems":
            self._render_systems_title_page(slide, spec, sections, slide_index, title, authors, affiliation, date_text)
            return
        if t.layout_family in {"conference", "report", "explainer"}:
            self._render_presentation_style_title_page(slide, spec, sections, slide_index, title, authors, affiliation, date_text)
            return

        self._add_text(
            slide,
            "Research Paper Presentation",
            self.Inches(0.72),
            self.Inches(0.9),
            self.Inches(1.9),
            self.Inches(0.62),
            size=10.5,
            bold=True,
            color=t.secondary,
            font=t.body_font,
            max_lines=2,
            alignment=self.PP_ALIGN.CENTER,
        )
        self._add_text(
            slide,
            title,
            self.Inches(3.5),
            self.Inches(1.08),
            self.Inches(8.4),
            self.Inches(1.55),
            size=self._fit_title_size(title, base=32, min_size=24),
            bold=True,
            color=t.ink,
            font=t.title_font,
            max_lines=2,
        )
        self._rounded_rect(slide, self.Inches(3.52), self.Inches(2.86), self.Inches(1.25), self.Inches(0.07), t.primary, t.primary)
        self._rounded_rect(slide, self.Inches(4.88), self.Inches(2.86), self.Inches(2.75), self.Inches(0.07), t.secondary, t.secondary)

        self._add_text(slide, authors, self.Inches(3.55), self.Inches(3.35), self.Inches(7.5), self.Inches(0.38), 15, True, t.primary, t.body_font, 1)
        self._add_text(slide, affiliation, self.Inches(3.55), self.Inches(3.82), self.Inches(7.5), self.Inches(0.32), 11, False, t.muted, t.body_font, 1)
        self._add_text(slide, date_text, self.Inches(3.55), self.Inches(4.26), self.Inches(2.4), self.Inches(0.28), 10, False, t.muted, t.body_font, 1)

        if sections:
            preview = " / ".join(sections[:4])
            self._add_text(
                slide,
                preview,
                self.Inches(3.55),
                self.Inches(5.34),
                self.Inches(7.9),
                self.Inches(0.34),
                size=10.5,
                bold=True,
                color=t.secondary,
                font=t.body_font,
                max_lines=1,
            )
        self._render_title_summary_tiles(slide, spec, sections)
        self._footer(slide, slide_index)

    def _render_editorial_title_page(
        self,
        slide,
        spec: PresentationSpec,
        sections: Sequence[str],
        slide_index: int,
        title: str,
        authors: str,
        affiliation: str,
        date_text: str,
    ) -> None:
        t = self.theme
        self._add_text(slide, "TECHNICAL REPORT", self.Inches(1.08), self.Inches(0.86), self.Inches(3.4), self.Inches(0.24), 9.5, True, t.primary, t.body_font, 1)
        self._add_text(
            slide,
            title,
            self.Inches(1.05),
            self.Inches(1.34),
            self.Inches(9.9),
            self.Inches(1.68),
            size=self._fit_title_size(title, base=36, min_size=25),
            bold=True,
            color=t.ink,
            font=t.title_font,
            max_lines=2,
        )
        self._rounded_rect(slide, self.Inches(1.08), self.Inches(3.18), self.Inches(1.25), self.Inches(0.035), t.primary, t.primary)
        self._add_text(slide, authors, self.Inches(1.08), self.Inches(3.58), self.Inches(8.1), self.Inches(0.35), 14, True, t.primary, t.body_font, 1)
        self._add_text(slide, affiliation, self.Inches(1.08), self.Inches(4.0), self.Inches(8.1), self.Inches(0.3), 10.5, False, t.muted, t.body_font, 1)
        self._add_text(slide, date_text, self.Inches(1.08), self.Inches(4.38), self.Inches(2.3), self.Inches(0.25), 9.5, False, t.muted, t.body_font, 1)
        if sections:
            self._add_text(slide, " / ".join(sections[:5]), self.Inches(1.08), self.Inches(5.52), self.Inches(9.7), self.Inches(0.32), 10.5, True, t.secondary, t.body_font, 1)
        self._render_title_summary_tiles(slide, spec, sections)
        self._footer(slide, slide_index)

    def _render_systems_title_page(
        self,
        slide,
        spec: PresentationSpec,
        sections: Sequence[str],
        slide_index: int,
        title: str,
        authors: str,
        affiliation: str,
        date_text: str,
    ) -> None:
        t = self.theme
        self._rounded_rect(slide, self.Inches(0.78), self.Inches(0.8), self.Inches(11.75), self.Inches(5.75), t.surface, t.rule)
        self._add_text(slide, "SYSTEMS / AGENTIC WORKFLOW", self.Inches(1.08), self.Inches(1.08), self.Inches(4.2), self.Inches(0.24), 9.0, True, t.secondary, t.body_font, 1)
        self._add_text(
            slide,
            title,
            self.Inches(1.06),
            self.Inches(1.58),
            self.Inches(7.7),
            self.Inches(1.45),
            size=self._fit_title_size(title, base=32, min_size=24),
            bold=True,
            color=t.ink,
            font=t.title_font,
            max_lines=2,
        )
        self._add_text(slide, authors, self.Inches(1.08), self.Inches(3.28), self.Inches(6.9), self.Inches(0.33), 13.5, True, t.primary, t.body_font, 1)
        self._add_text(slide, affiliation, self.Inches(1.08), self.Inches(3.72), self.Inches(6.9), self.Inches(0.3), 10.2, False, t.muted, t.body_font, 1)
        labels = sections[:4] or ["Context", "Data", "Training", "Evaluation"]
        for index, label in enumerate(labels, start=1):
            x = self.Inches(8.9)
            y = self.Inches(1.22 + (index - 1) * 0.9)
            self._rounded_rect(slide, x, y, self.Inches(2.35), self.Inches(0.46), t.pale_primary if index % 2 else t.pale_secondary, t.rule)
            self._add_text(slide, f"{index:02d}", int(x + self.Inches(0.16)), int(y + self.Inches(0.12)), self.Inches(0.34), self.Inches(0.18), 8, True, t.primary, t.title_font, 1)
            self._add_text(slide, self._truncate(label, 24), int(x + self.Inches(0.62)), int(y + self.Inches(0.1)), self.Inches(1.55), self.Inches(0.22), 8.3, True, t.ink, t.body_font, 1)
        self._add_text(slide, date_text, self.Inches(1.08), self.Inches(5.95), self.Inches(2.3), self.Inches(0.25), 9, False, t.muted, t.body_font, 1)
        self._footer(slide, slide_index)

    def _render_presentation_style_title_page(
        self,
        slide,
        spec: PresentationSpec,
        sections: Sequence[str],
        slide_index: int,
        title: str,
        authors: str,
        affiliation: str,
        date_text: str,
    ) -> None:
        t = self.theme
        label = {
            "conference": "CONFERENCE ORAL",
            "report": "BENCHMARK REPORT",
            "explainer": "VISUAL EXPLAINER",
        }.get(t.layout_family, "PRESENTATION")
        self._add_text(slide, label, self.Inches(0.86), self.Inches(0.86), self.Inches(3.4), self.Inches(0.24), 9.0, True, t.primary, t.body_font, 1)
        self._add_text(slide, title, self.Inches(0.84), self.Inches(1.28), self.Inches(8.4), self.Inches(1.36), self._fit_title_size(title, 34, 24), True, t.ink, t.title_font, 2)
        self._render_title_summary_tiles(slide, spec, sections)
        self._add_text(slide, authors, self.Inches(0.88), self.Inches(3.24), self.Inches(7.2), self.Inches(0.32), 13.5, True, t.primary, t.body_font, 1)
        self._add_text(slide, affiliation, self.Inches(0.88), self.Inches(3.66), self.Inches(7.2), self.Inches(0.26), 10, False, t.muted, t.body_font, 1)
        if sections:
            self._render_stage_cards(slide, [TextBlock(text=s, claim=s, detail="Section route") for s in sections[:4]], self.Inches(0.9), self.Inches(4.55), self.Inches(8.5), self.Inches(0.92), horizontal=True)
        self._add_text(slide, date_text, self.Inches(0.88), self.Inches(6.26), self.Inches(2.1), self.Inches(0.2), 8.5, False, t.muted, t.body_font, 1)
        self._footer(slide, slide_index)

    def _render_title_summary_tiles(self, slide, spec: PresentationSpec, sections: Sequence[str]) -> None:
        t = self.theme
        content_slides = [slide_spec for slide_spec in spec.slides if slide_spec.section_type != "opening"]
        figure_count = sum(len(slide_spec.image_blocks) for slide_spec in spec.slides)
        items = [
            (str(len(sections)), "Sections", t.pale_primary, t.primary),
            (str(len(content_slides)), "Content slides", t.pale_secondary, t.secondary),
            (str(figure_count), "Source figures", t.pale_neutral, t.accent),
        ]
        positions = [
            (self.Inches(9.05), self.Inches(4.58), self.Inches(2.55), self.Inches(0.52)),
            (self.Inches(9.05), self.Inches(5.18), self.Inches(1.85), self.Inches(0.52)),
            (self.Inches(9.05), self.Inches(5.78), self.Inches(1.25), self.Inches(0.52)),
        ]
        for (value, label, fill, accent), (left, top, width, height) in zip(items, positions):
            self._rounded_rect(slide, left, top, width, height, fill, self.theme.rule)
            self._add_text(slide, value, int(left + self.Inches(0.16)), int(top + self.Inches(0.1)), self.Inches(0.55), self.Inches(0.24), 13, True, accent, t.title_font, 1)
            self._add_text(slide, label, int(left + self.Inches(0.76)), int(top + self.Inches(0.16)), int(width - self.Inches(0.86)), self.Inches(0.2), 7.2, True, t.muted, t.body_font, 1)

    def _render_toc(self, slide, sections: Sequence[str], slide_index: int) -> None:
        t = self.theme
        title_text = "Contents" if t.layout_family != "systems" else "System Route"
        self._add_text(slide, title_text, self.Inches(0.78), self.Inches(0.62), self.Inches(4.5), self.Inches(0.68), 30, True, t.ink, t.title_font, 1)
        self._add_text(
            slide,
            "A sectioned route through the paper: why it matters, what is new, how it works, and what it proves."
            if t.layout_family != "editorial"
            else "A report-style reading path: argument first, evidence second, implications last.",
            self.Inches(0.82),
            self.Inches(1.35),
            self.Inches(9.4),
            self.Inches(0.42),
            12.5,
            True,
            t.primary,
            t.body_font,
            1,
        )

        palette = [t.primary, t.secondary, t.accent, t.ink, t.primary, t.secondary, t.accent]
        section_count = min(7, len(sections))
        start_top = self.Inches(2.05)
        row_h = self.Inches(0.56 if section_count >= 7 else 0.68)
        gap = self.Inches(0.11 if section_count >= 7 else 0.18)
        for index, section in enumerate(sections[:7], start=1):
            y = int(start_top + (index - 1) * (row_h + gap))
            accent = palette[(index - 1) % len(palette)]
            marker_w = self.Inches(0.62 if t.layout_family != "editorial" else 0.42)
            self._rounded_rect(slide, self.Inches(0.86), y, marker_w, row_h, accent, accent)
            self._add_text(
                slide,
                f"{index:02d}",
                self.Inches(0.97),
                y + self.Inches(0.14),
                self.Inches(0.42),
                self.Inches(0.22),
                10.5,
                True,
                (255, 255, 255),
                t.title_font,
                1,
                alignment=self.PP_ALIGN.CENTER,
            )
            self._add_text(slide, section, self.Inches(1.72), y + self.Inches(0.08), self.Inches(5.4), self.Inches(0.32), 14.5 if section_count >= 7 else 15.5, True, t.ink, t.title_font, 1)
            line_left = self.Inches(7.12)
            line_width = self.Inches(3.95)
            self._rounded_rect(slide, line_left, y + self.Inches(0.28), line_width, self.Inches(0.025), t.rule, t.rule)
            self._rounded_rect(slide, line_left, y + self.Inches(0.28), int(line_width * index / max(1, section_count)), self.Inches(0.025), accent, accent)

        self._footer(slide, slide_index)

    def _render_cover(self, slide, slide_spec: SlideSpec, slide_index: int) -> None:
        t = self.theme
        if t.layout_family in {"editorial", "systems", "conference", "report", "explainer"}:
            self._render_style_cover(slide, slide_spec, slide_index)
            return
        title = self._clean_title(slide_spec.title)
        self._add_text(
            slide,
            title,
            self.Inches(0.78),
            self.Inches(0.82),
            self.Inches(6.75),
            self.Inches(1.35),
            size=self._fit_title_size(title, base=34, min_size=25),
            bold=True,
            color=t.ink,
            font=t.title_font,
            max_lines=2,
        )

        if slide_spec.takeaway:
            self._add_text(
                slide,
                slide_spec.takeaway,
                self.Inches(0.82),
                self.Inches(2.42),
                self.Inches(6.25),
                self.Inches(0.72),
                size=17,
                bold=True,
                color=t.primary,
                font=t.body_font,
                max_lines=2,
            )

        bullets_top = self.Inches(3.42)
        self._add_bullet_list(
            slide,
            slide_spec.text_blocks[:3],
            self.Inches(0.9),
            bullets_top,
            self.Inches(5.8),
            self.Inches(1.45),
            size=13,
            max_items=3,
        )

        if slide_spec.metric_blocks:
            self._render_metric_band(
                slide,
                slide_spec.metric_blocks[:3],
                self.Inches(0.88),
                self.Inches(5.38),
                self.Inches(5.9),
                self.Inches(0.72),
            )

        if slide_spec.image_blocks:
            self._render_images(
                slide,
                slide_spec,
                self.Inches(7.35),
                self.Inches(0.88),
                self.Inches(5.15),
                self.Inches(5.55),
                caption=True,
            )

        self._footer(slide, slide_index)

    def _render_style_cover(self, slide, slide_spec: SlideSpec, slide_index: int) -> None:
        t = self.theme
        title = self._clean_title(slide_spec.title)
        claim = slide_spec.takeaway or (slide_spec.text_blocks[0].text if slide_spec.text_blocks else "")
        if t.layout_family == "systems":
            self._render_header(slide, slide_spec)
            self._render_key_message(slide, claim, self.Inches(0.9), self.Inches(1.35), self.Inches(5.2), self.Inches(0.88))
            self._render_stage_cards(slide, slide_spec.text_blocks[:4], self.Inches(0.9), self.Inches(2.55), self.Inches(7.2), self.Inches(2.75), horizontal=True)
            if slide_spec.metric_blocks:
                self._render_metric_column(slide, slide_spec.metric_blocks[:4], self.Inches(8.65), self.Inches(1.35), self.Inches(3.2), self.Inches(4.4))
            elif slide_spec.image_blocks:
                self._render_images(slide, slide_spec, self.Inches(8.25), self.Inches(1.45), self.Inches(3.8), self.Inches(3.8), caption=True)
            self._footer(slide, slide_index)
            return

        if t.layout_family == "conference":
            self._add_text(slide, "OPENING CLAIM", self.Inches(0.82), self.Inches(0.75), self.Inches(2.4), self.Inches(0.22), 8.5, True, t.primary, t.body_font, 1)
            self._add_text(slide, title, self.Inches(0.78), self.Inches(1.08), self.Inches(10.9), self.Inches(1.1), self._fit_title_size(title, 34, 24), True, t.ink, t.title_font, 2)
            self._render_key_message(slide, claim, self.Inches(0.82), self.Inches(2.48), self.Inches(6.8), self.Inches(0.88))
            if slide_spec.image_blocks:
                self._render_images(slide, slide_spec, self.Inches(7.95), self.Inches(2.34), self.Inches(4.35), self.Inches(3.0), caption=True)
            self._render_metric_band(slide, slide_spec.metric_blocks[:3], self.Inches(0.86), self.Inches(5.45), self.Inches(10.9), self.Inches(0.72))
            self._footer(slide, slide_index)
            return

        if t.layout_family == "report":
            self._render_header(slide, slide_spec)
            self._add_text(slide, "SCORECARD", self.Inches(0.9), self.Inches(1.28), self.Inches(2.1), self.Inches(0.2), 8.5, True, t.secondary, t.body_font, 1)
            self._render_metric_band(slide, slide_spec.metric_blocks[:4], self.Inches(0.9), self.Inches(1.68), self.Inches(11.1), self.Inches(0.96))
            self._render_key_message(slide, claim, self.Inches(0.9), self.Inches(3.05), self.Inches(6.2), self.Inches(0.86))
            self._render_stage_cards(slide, slide_spec.text_blocks[:3], self.Inches(0.9), self.Inches(4.22), self.Inches(10.9), self.Inches(1.45), horizontal=True)
            self._footer(slide, slide_index)
            return

        if t.layout_family == "explainer":
            if slide_spec.image_blocks:
                self._render_images(slide, slide_spec, self.Inches(6.4), self.Inches(0.95), self.Inches(5.75), self.Inches(4.85), caption=True)
            self._add_text(slide, title, self.Inches(0.78), self.Inches(0.98), self.Inches(5.15), self.Inches(1.2), self._fit_title_size(title, 30, 22), True, t.ink, t.title_font, 2)
            self._render_key_message(slide, claim, self.Inches(0.82), self.Inches(2.48), self.Inches(4.85), self.Inches(0.94))
            self._render_stage_cards(slide, slide_spec.text_blocks[:3], self.Inches(0.86), self.Inches(3.78), self.Inches(4.9), self.Inches(1.75), horizontal=False)
            self._footer(slide, slide_index)
            return

        self._add_text(slide, "TECHNICAL WHITE PAPER", self.Inches(0.9), self.Inches(0.82), self.Inches(3.2), self.Inches(0.22), 8.5, True, t.primary, t.body_font, 1)
        self._add_text(slide, title, self.Inches(0.88), self.Inches(1.18), self.Inches(9.5), self.Inches(1.36), self._fit_title_size(title, 36, 25), True, t.ink, t.title_font, 2)
        self._render_key_message(slide, claim, self.Inches(0.92), self.Inches(2.98), self.Inches(7.1), self.Inches(0.78))
        self._render_stage_cards(slide, slide_spec.text_blocks[:3], self.Inches(0.95), self.Inches(4.18), self.Inches(7.25), self.Inches(1.34), horizontal=True)
        if slide_spec.image_blocks:
            self._render_images(slide, slide_spec, self.Inches(8.55), self.Inches(1.36), self.Inches(3.7), self.Inches(3.8), caption=True)
        elif slide_spec.metric_blocks:
            self._render_metric_column(slide, slide_spec.metric_blocks[:3], self.Inches(8.65), self.Inches(1.45), self.Inches(3.0), self.Inches(3.4))
        self._footer(slide, slide_index)

    def _render_section_divider(self, slide, section: str, slide_index: int) -> None:
        t = self.theme
        if t.layout_family == "editorial":
            self._add_text(slide, "CHAPTER", self.Inches(0.98), self.Inches(1.52), self.Inches(1.6), self.Inches(0.24), 9, True, t.primary, t.body_font, 1)
            self._add_text(
                slide,
                section,
                self.Inches(0.95),
                self.Inches(2.1),
                self.Inches(10.6),
                self.Inches(0.96),
                size=self._fit_title_size(section, base=38, min_size=27),
                bold=True,
                color=t.ink,
                font=t.title_font,
                max_lines=1,
            )
            self._rounded_rect(slide, self.Inches(0.98), self.Inches(3.42), self.Inches(7.4), self.Inches(0.035), t.primary, t.primary)
            self._footer(slide, slide_index)
            return
        if t.layout_family == "systems":
            self._rounded_rect(slide, self.Inches(0.86), self.Inches(1.58), self.Inches(10.9), self.Inches(2.0), t.surface, t.rule)
            self._add_text(slide, "STAGE TRANSITION", self.Inches(1.12), self.Inches(1.9), self.Inches(2.6), self.Inches(0.24), 9, True, t.secondary, t.body_font, 1)
            self._add_text(
                slide,
                section,
                self.Inches(1.1),
                self.Inches(2.32),
                self.Inches(9.7),
                self.Inches(0.68),
                size=self._fit_title_size(section, base=32, min_size=24),
                bold=True,
                color=t.ink,
                font=t.title_font,
                max_lines=1,
            )
            self._rounded_rect(slide, self.Inches(1.1), self.Inches(3.22), self.Inches(2.4), self.Inches(0.07), t.primary, t.primary)
            self._rounded_rect(slide, self.Inches(3.65), self.Inches(3.22), self.Inches(4.8), self.Inches(0.07), t.secondary, t.secondary)
            self._footer(slide, slide_index)
            return
        self._add_text(
            slide,
            "Next section",
            self.Inches(0.94),
            self.Inches(1.72),
            self.Inches(3.0),
            self.Inches(0.28),
            size=12,
            bold=True,
            color=t.primary,
            font=t.body_font,
            max_lines=1,
        )
        self._add_text(
            slide,
            section,
            self.Inches(0.9),
            self.Inches(2.15),
            self.Inches(10.8),
            self.Inches(0.9),
            size=self._fit_title_size(section, base=34, min_size=26),
            bold=True,
            color=t.ink,
            font=t.title_font,
            max_lines=1,
        )
        self._rounded_rect(slide, self.Inches(0.92), self.Inches(3.25), self.Inches(5.8), self.Inches(0.09), t.primary, t.primary)
        self._rounded_rect(slide, self.Inches(6.85), self.Inches(3.25), self.Inches(4.0), self.Inches(0.09), t.secondary, t.secondary)
        self._footer(slide, slide_index)

    def _render_statement(self, slide, slide_spec: SlideSpec, slide_index: int, closing: bool = False) -> None:
        t = self.theme
        if t.layout_family in {"editorial", "systems", "conference", "report", "explainer"}:
            self._render_style_statement(slide, slide_spec, slide_index, closing=closing)
            return
        self._render_header(slide, slide_spec)

        claim = slide_spec.takeaway or (slide_spec.text_blocks[0].text if slide_spec.text_blocks else "")
        has_metrics = bool(slide_spec.metric_blocks)
        claim_width = self.Inches(7.15 if has_metrics else 10.95)
        bullet_width = self.Inches(7.2 if has_metrics else 11.0)

        self._render_key_message(
            slide,
            claim,
            self.Inches(0.86),
            self.Inches(1.28),
            claim_width,
            self.Inches(0.92 if not closing else 1.08),
            accent=t.secondary if closing else t.primary,
        )

        self._add_bullet_list(
            slide,
            slide_spec.text_blocks[:5],
            self.Inches(0.96),
            self.Inches(2.42),
            bullet_width,
            self.Inches(3.95),
            size=13.0,
            max_items=5,
        )

        if has_metrics:
            self._render_metric_column(
                slide,
                slide_spec.metric_blocks[:4],
                self.Inches(8.62),
                self.Inches(1.34),
                self.Inches(3.28),
                self.Inches(4.88),
            )

        self._footer(slide, slide_index)

    def _render_style_statement(self, slide, slide_spec: SlideSpec, slide_index: int, closing: bool = False) -> None:
        t = self.theme
        self._render_header(slide, slide_spec)
        claim = slide_spec.takeaway or (slide_spec.text_blocks[0].text if slide_spec.text_blocks else "")

        if t.layout_family == "editorial":
            self._add_text(slide, "Argument", self.Inches(0.9), self.Inches(1.38), self.Inches(1.4), self.Inches(0.2), 8, True, t.primary, t.body_font, 1)
            self._render_key_message(slide, claim, self.Inches(0.9), self.Inches(1.72), self.Inches(7.6), self.Inches(0.82 if not closing else 1.0))
            self._render_stage_cards(slide, slide_spec.text_blocks[:3], self.Inches(0.95), self.Inches(3.0), self.Inches(7.6), self.Inches(2.6), horizontal=False)
            if slide_spec.metric_blocks:
                self._render_metric_column(slide, slide_spec.metric_blocks[:4], self.Inches(9.25), self.Inches(1.5), self.Inches(2.7), self.Inches(4.55))
            self._footer(slide, slide_index)
            return

        if t.layout_family == "systems":
            self._render_key_message(slide, claim, self.Inches(0.86), self.Inches(1.24), self.Inches(5.6), self.Inches(0.78))
            self._render_stage_cards(slide, slide_spec.text_blocks[:4], self.Inches(0.92), self.Inches(2.32), self.Inches(7.65), self.Inches(3.42), horizontal=True)
            if slide_spec.metric_blocks:
                self._render_metric_column(slide, slide_spec.metric_blocks[:4], self.Inches(9.0), self.Inches(1.28), self.Inches(2.9), self.Inches(4.7))
            else:
                self._add_system_lane(slide, self.Inches(9.0), self.Inches(1.42), self.Inches(2.9), self.Inches(4.28), slide_spec)
            self._footer(slide, slide_index)
            return

        if t.layout_family == "conference":
            self._add_text(slide, self._body_sentence(claim, 190), self.Inches(0.9), self.Inches(1.3), self.Inches(10.4), self.Inches(0.98), 21, True, t.ink, t.title_font, 2)
            self._rounded_rect(slide, self.Inches(0.9), self.Inches(2.55), self.Inches(3.2), self.Inches(0.07), t.primary, t.primary)
            self._render_stage_cards(slide, slide_spec.text_blocks[:3], self.Inches(0.9), self.Inches(3.1), self.Inches(10.4), self.Inches(2.1), horizontal=True)
            if slide_spec.metric_blocks:
                self._render_metric_band(slide, slide_spec.metric_blocks[:4], self.Inches(0.9), self.Inches(5.82), self.Inches(10.8), self.Inches(0.62))
            self._footer(slide, slide_index)
            return

        if t.layout_family == "report":
            if slide_spec.metric_blocks:
                self._render_metric_band(slide, slide_spec.metric_blocks[:4], self.Inches(0.86), self.Inches(1.18), self.Inches(11.15), self.Inches(0.98))
                claim_top = self.Inches(2.55)
            else:
                claim_top = self.Inches(1.28)
            self._render_key_message(slide, claim, self.Inches(0.86), claim_top, self.Inches(6.7), self.Inches(0.82))
            self._render_stage_cards(slide, slide_spec.text_blocks[:3], self.Inches(0.92), int(claim_top + self.Inches(1.16)), self.Inches(10.85), self.Inches(2.8), horizontal=False)
            self._footer(slide, slide_index)
            return

        if slide_spec.image_blocks:
            self._render_images(slide, slide_spec, self.Inches(0.82), self.Inches(1.28), self.Inches(6.35), self.Inches(4.65), caption=True)
            text_left = self.Inches(7.55)
            text_width = self.Inches(4.65)
        else:
            text_left = self.Inches(0.92)
            text_width = self.Inches(7.8)
        self._render_key_message(slide, claim, text_left, self.Inches(1.28), text_width, self.Inches(0.86))
        self._render_stage_cards(slide, slide_spec.text_blocks[:3], text_left, self.Inches(2.55), text_width, self.Inches(3.05), horizontal=False)
        if slide_spec.metric_blocks and not slide_spec.image_blocks:
            self._render_metric_column(slide, slide_spec.metric_blocks[:3], self.Inches(9.3), self.Inches(1.3), self.Inches(2.75), self.Inches(3.6))
        self._footer(slide, slide_index)

    def _render_visual_or_mixed(self, slide, slide_spec: SlideSpec, slide_index: int, visual_left: bool = False) -> None:
        t = self.theme
        if t.layout_family in {"conference", "report", "explainer"}:
            self._render_style_visual(slide, slide_spec, slide_index, visual_left=visual_left)
            return
        self._render_header(slide, slide_spec)

        if t.layout_family == "editorial":
            image_left = self.Inches(6.0 if not visual_left else 0.78)
            text_left = self.Inches(0.86 if not visual_left else 7.1)
            visual_width = self.Inches(6.4 if not visual_left else 5.95)
            text_width = self.Inches(4.75 if not visual_left else 4.55)
        elif t.layout_family == "systems":
            image_left = self.Inches(0.82 if visual_left else 6.35)
            text_left = self.Inches(6.9 if visual_left else 0.82)
            visual_width = self.Inches(5.75 if visual_left else 5.9)
            text_width = self.Inches(5.2 if visual_left else 5.15)
        else:
            image_left = self.Inches(0.75 if visual_left else 6.2)
            text_left = self.Inches(6.85 if visual_left else 0.82)
            visual_width = self.Inches(5.65 if visual_left else 6.25)
            text_width = self.Inches(5.6 if visual_left else 5.0)

        has_table = bool(slide_spec.table_blocks)
        self._render_images(
            slide,
            slide_spec,
            image_left,
            self.Inches(1.58 if t.layout_family == "editorial" else 1.72),
            visual_width,
            self.Inches(3.05 if has_table else (4.55 if t.layout_family == "editorial" else 4.12)),
            caption=True,
        )

        if slide_spec.takeaway:
            self._render_key_message(
                slide,
                slide_spec.takeaway,
                text_left,
                self.Inches(1.28),
                text_width,
                self.Inches(0.82),
            )
            points_top = self.Inches(2.34)
            points_height = self.Inches(2.9 if has_table else 3.08)
        else:
            points_top = self.Inches(1.55)
            points_height = self.Inches(3.35 if has_table else 4.15)

        self._add_bullet_list(
            slide,
            slide_spec.text_blocks[:5],
            text_left,
            points_top,
            text_width,
            points_height,
            size=11.5,
            max_items=4,
        )

        if has_table:
            if slide_spec.metric_blocks and not slide_spec.image_blocks:
                self._render_metric_column(
                    slide,
                    slide_spec.metric_blocks[:3],
                    image_left,
                    self.Inches(1.72),
                    visual_width,
                    self.Inches(3.05),
                )
            elif not slide_spec.image_blocks:
                self._render_evidence_column(
                    slide,
                    slide_spec.text_blocks[2:5] or slide_spec.text_blocks[:2],
                    image_left,
                    self.Inches(1.72),
                    visual_width,
                    self.Inches(3.05),
                )
            self._render_table(
                slide,
                slide_spec,
                self.Inches(0.85),
                self.Inches(5.35),
                self.Inches(11.6),
                self.Inches(1.25),
            )
        elif slide_spec.metric_blocks:
            self._render_metric_band(
                slide,
                slide_spec.metric_blocks[:3],
                text_left,
                self.Inches(5.72),
                text_width,
                self.Inches(0.66),
            )

        self._footer(slide, slide_index)

    def _render_evidence_column(self, slide, blocks: Sequence[TextBlock], left, top, width, height) -> None:
        blocks = list(blocks)[:3]
        if not blocks:
            return
        self._rounded_rect(slide, left, top, width, height, self.theme.pale_neutral, self.theme.rule)
        self._add_text(
            slide,
            "Evidence notes",
            int(left + self.Inches(0.28)),
            int(top + self.Inches(0.22)),
            int(width - self.Inches(0.56)),
            self.Inches(0.18),
            7.6,
            True,
            self.theme.primary,
            self.theme.body_font,
            1,
        )
        self._render_stage_cards(
            slide,
            blocks,
            int(left + self.Inches(0.22)),
            int(top + self.Inches(0.58)),
            int(width - self.Inches(0.44)),
            int(height - self.Inches(0.78)),
            horizontal=False,
        )

    def _render_style_visual(self, slide, slide_spec: SlideSpec, slide_index: int, visual_left: bool = False) -> None:
        t = self.theme
        self._render_header(slide, slide_spec)
        claim = slide_spec.takeaway or (slide_spec.text_blocks[0].text if slide_spec.text_blocks else "")
        if t.layout_family == "explainer":
            self._render_images(slide, slide_spec, self.Inches(0.82), self.Inches(1.28), self.Inches(7.05), self.Inches(4.85), caption=True)
            self._render_key_message(slide, claim, self.Inches(8.25), self.Inches(1.28), self.Inches(3.75), self.Inches(0.86))
            self._render_stage_cards(slide, slide_spec.text_blocks[:4], self.Inches(8.25), self.Inches(2.52), self.Inches(3.75), self.Inches(3.05), horizontal=False)
        elif t.layout_family == "conference":
            self._add_text(slide, self._body_sentence(claim, 180), self.Inches(0.9), self.Inches(1.24), self.Inches(10.3), self.Inches(0.72), 18.5, True, t.ink, t.title_font, 2)
            self._render_images(slide, slide_spec, self.Inches(1.0), self.Inches(2.18), self.Inches(6.4), self.Inches(3.65), caption=True)
            self._render_stage_cards(slide, slide_spec.text_blocks[:3], self.Inches(7.8), self.Inches(2.3), self.Inches(3.95), self.Inches(2.85), horizontal=False)
            if slide_spec.metric_blocks:
                self._render_metric_band(slide, slide_spec.metric_blocks[:3], self.Inches(0.95), self.Inches(6.08), self.Inches(10.3), self.Inches(0.52))
        else:
            if slide_spec.metric_blocks:
                self._render_metric_band(slide, slide_spec.metric_blocks[:4], self.Inches(0.9), self.Inches(1.12), self.Inches(11.0), self.Inches(0.74))
                visual_top = self.Inches(2.12)
            else:
                visual_top = self.Inches(1.34)
            self._render_images(slide, slide_spec, self.Inches(0.92), visual_top, self.Inches(5.8), self.Inches(3.6), caption=True)
            self._render_key_message(slide, claim, self.Inches(7.15), visual_top, self.Inches(4.55), self.Inches(0.72))
            self._render_stage_cards(slide, slide_spec.text_blocks[:3], self.Inches(7.15), int(visual_top + self.Inches(1.03)), self.Inches(4.55), self.Inches(2.35), horizontal=False)
            if slide_spec.table_blocks:
                self._render_table(slide, slide_spec, self.Inches(0.95), self.Inches(5.72), self.Inches(10.9), self.Inches(0.78))
        self._footer(slide, slide_index)

    def _render_table_focus(self, slide, slide_spec: SlideSpec, slide_index: int) -> None:
        if self.theme.layout_family in {"editorial", "systems", "conference", "report", "explainer"}:
            self._render_style_table_focus(slide, slide_spec, slide_index)
            return
        self._render_header(slide, slide_spec)
        has_image = bool(slide_spec.image_blocks)
        if has_image:
            self._render_images(
                slide,
                slide_spec,
                self.Inches(0.85),
                self.Inches(1.55),
                self.Inches(5.15),
                self.Inches(2.25),
                caption=False,
            )
            bullet_left = self.Inches(6.35)
            bullet_width = self.Inches(5.9)
        else:
            bullet_left = self.Inches(0.9)
            bullet_width = self.Inches(11.45)

        if slide_spec.takeaway:
            self._render_key_message(
                slide,
                slide_spec.takeaway,
                bullet_left,
                self.Inches(1.25),
                bullet_width,
                self.Inches(0.72),
            )
            bullet_top = self.Inches(2.12)
            bullet_height = self.Inches(1.2 if has_image else 1.05)
        else:
            bullet_top = self.Inches(1.58)
            bullet_height = self.Inches(1.78 if has_image else 1.6)

        self._add_bullet_list(
            slide,
            slide_spec.text_blocks[:4],
            bullet_left,
            bullet_top,
            bullet_width,
            bullet_height,
            size=10.6,
            max_items=1 if slide_spec.metric_blocks else (2 if slide_spec.takeaway else 3),
        )

        if slide_spec.metric_blocks and not has_image:
            self._render_metric_band(
                slide,
                slide_spec.metric_blocks[:4],
                self.Inches(0.9),
                self.Inches(3.02),
                self.Inches(11.35),
                self.Inches(0.58),
            )

        table_top = self.Inches(4.08 if has_image else 3.82)
        table_height = self.Inches(2.18)
        self._render_table(slide, slide_spec, self.Inches(0.9), table_top, self.Inches(11.45), table_height)
        self._footer(slide, slide_index)

    def _render_style_table_focus(self, slide, slide_spec: SlideSpec, slide_index: int) -> None:
        t = self.theme
        self._render_header(slide, slide_spec)
        claim = slide_spec.takeaway or (slide_spec.text_blocks[0].text if slide_spec.text_blocks else "")
        if t.layout_family == "report":
            if slide_spec.metric_blocks:
                self._render_metric_band(slide, slide_spec.metric_blocks[:4], self.Inches(0.88), self.Inches(1.12), self.Inches(11.0), self.Inches(0.72))
                top = self.Inches(2.12)
            else:
                top = self.Inches(1.3)
            self._render_key_message(slide, claim, self.Inches(0.9), top, self.Inches(10.7), self.Inches(0.62))
            self._render_table(slide, slide_spec, self.Inches(0.9), int(top + self.Inches(0.95)), self.Inches(11.2), self.Inches(2.8))
        elif t.layout_family == "conference":
            self._render_key_message(slide, claim, self.Inches(0.9), self.Inches(1.24), self.Inches(10.6), self.Inches(0.72))
            self._render_table(slide, slide_spec, self.Inches(0.9), self.Inches(2.25), self.Inches(11.0), self.Inches(3.45))
            self._render_stage_cards(slide, slide_spec.text_blocks[:2], self.Inches(0.9), self.Inches(6.0), self.Inches(8.0), self.Inches(0.55), horizontal=True)
        elif t.layout_family == "systems":
            self._render_key_message(slide, claim, self.Inches(0.9), self.Inches(1.22), self.Inches(5.6), self.Inches(0.74))
            self._render_stage_cards(slide, slide_spec.text_blocks[:3], self.Inches(0.9), self.Inches(2.25), self.Inches(5.6), self.Inches(2.8), horizontal=False)
            self._render_table(slide, slide_spec, self.Inches(6.85), self.Inches(1.32), self.Inches(5.25), self.Inches(4.35))
        else:
            self._render_key_message(slide, claim, self.Inches(0.9), self.Inches(1.25), self.Inches(6.2), self.Inches(0.72))
            if slide_spec.image_blocks and t.layout_family == "explainer":
                self._render_images(slide, slide_spec, self.Inches(7.35), self.Inches(1.3), self.Inches(4.65), self.Inches(2.4), caption=True)
            self._render_table(slide, slide_spec, self.Inches(0.9), self.Inches(2.45), self.Inches(11.2), self.Inches(2.9))
            self._render_stage_cards(slide, slide_spec.text_blocks[:2], self.Inches(0.9), self.Inches(5.75), self.Inches(9.0), self.Inches(0.75), horizontal=True)
        self._footer(slide, slide_index)

    def _render_header(self, slide, slide_spec: SlideSpec) -> None:
        t = self.theme
        title = self._clean_title(slide_spec.title)
        if t.layout_family in {"report", "conference"}:
            title = self._truncate(title, 66)
        if t.layout_family == "editorial":
            section = self._truncate(slide_spec.section_label or "Research note", 36)
            self._add_text(slide, section.upper(), self.Inches(0.78), self.Inches(0.34), self.Inches(3.1), self.Inches(0.2), 7.4, True, t.primary, t.body_font, 1)
            self._add_text(
                slide,
                title,
                self.Inches(0.78),
                self.Inches(0.58),
                self.Inches(10.7),
                self.Inches(0.44),
                size=self._fit_title_size(title, base=22, min_size=16),
                bold=True,
                color=t.ink,
                font=t.title_font,
                max_lines=1,
            )
            self._rounded_rect(slide, self.Inches(0.78), self.Inches(1.08), self.Inches(10.95), self.Inches(0.02), t.rule, t.rule)
            return
        if t.layout_family == "systems":
            self._rounded_rect(slide, self.Inches(0.68), self.Inches(0.28), self.Inches(11.88), self.Inches(0.68), t.surface, t.rule)
            self._rounded_rect(slide, self.Inches(0.82), self.Inches(0.44), self.Inches(0.52), self.Inches(0.22), t.primary, t.primary)
            self._add_text(slide, "SYS", self.Inches(0.9), self.Inches(0.48), self.Inches(0.35), self.Inches(0.12), 6.5, True, (255, 255, 255), t.body_font, 1, alignment=self.PP_ALIGN.CENTER)
            self._add_text(
                slide,
                title,
                self.Inches(1.55),
                self.Inches(0.4),
                self.Inches(10.5),
                self.Inches(0.52),
                size=self._fit_title_size(title, base=21, min_size=16),
                bold=True,
                color=t.ink,
                font=t.title_font,
                max_lines=1,
            )
            return
        self._rounded_rect(slide, self.Inches(0.68), self.Inches(0.28), self.Inches(11.88), self.Inches(0.64), t.surface, t.rule)
        self._rounded_rect(slide, self.Inches(0.68), self.Inches(0.28), self.Inches(0.13), self.Inches(0.64), t.primary, t.primary)
        self._add_text(
            slide,
            title,
            self.Inches(0.92),
            self.Inches(0.38),
            self.Inches(11.25),
            self.Inches(0.52),
            size=self._fit_title_size(title, base=22, min_size=17),
            bold=True,
            color=t.ink,
            font=t.title_font,
            max_lines=1,
        )
        # The takeaway is rendered as a larger key-message block in each layout.

    def _render_images(self, slide, slide_spec: SlideSpec, image_left, image_top, image_width, max_height, caption: bool) -> None:
        count = max(1, min(2, len(slide_spec.image_blocks)))
        slot_height = int(max_height / count)

        for index, block in enumerate(slide_spec.image_blocks[:2]):
            top = int(image_top + slot_height * index)
            picture_height = int(slot_height - (self.Inches(0.32) if caption else self.Inches(0.08)))
            image_path = Path(block.path) if block.path else None

            if image_path and image_path.exists():
                left, fitted_top, width, height = self._fit_picture(
                    image_path,
                    image_left,
                    top,
                    image_width,
                    picture_height,
                )
                slide.shapes.add_picture(str(image_path), left, fitted_top, width=width, height=height)
            else:
                self._placeholder(slide, image_left, top, image_width, picture_height, block.placeholder_text or block.title)

            caption_text = self._truncate(f"{block.title}: {block.caption}" if block.title and block.caption else block.title or block.caption, 96)
            if caption and caption_text:
                self._add_text(
                    slide,
                    caption_text,
                    image_left,
                    int(top + slot_height - self.Inches(0.32)),
                    image_width,
                    self.Inches(0.3),
                    size=8.0,
                    color=self.theme.muted,
                    font=self.theme.body_font,
                    max_lines=1,
                    alignment=self.PP_ALIGN.CENTER,
                )

    def _render_metric_band(self, slide, metrics: Sequence[MetricBlock], left, top, width, height) -> None:
        metrics = list(metrics)[:4]
        if not metrics:
            return
        gap = int(width * 0.018)
        slot_width = int((int(width) - gap * (len(metrics) - 1)) / len(metrics))
        fills = [self.theme.pale_primary, self.theme.pale_secondary, self.theme.pale_neutral, (238, 242, 244)]
        accents = [self.theme.primary, self.theme.secondary, self.theme.ink, self.theme.accent]
        for index, metric in enumerate(metrics):
            x = int(left + index * (slot_width + gap))
            if self.theme.layout_family == "editorial":
                self._rounded_rect(slide, x, top, slot_width, self.Inches(0.03), accents[index % len(accents)], accents[index % len(accents)])
                self._add_text(slide, self._metric_value(metric.value), x, int(top + height * 0.1), int(slot_width * 0.86), int(height * 0.38), 16, True, accents[index % len(accents)], self.theme.title_font, 1)
                self._add_text(slide, self._metric_label(metric), x, int(top + height * 0.55), int(slot_width * 0.92), int(height * 0.34), 7.8, False, self.theme.muted, self.theme.body_font, 1)
            elif self.theme.layout_family == "report":
                self._rounded_rect(slide, x, top, slot_width, height, fills[index % len(fills)], self.theme.rule)
                self._add_text(slide, self._metric_value(metric.value), x + int(slot_width * 0.07), int(top + height * 0.06), int(slot_width * 0.86), int(height * 0.46), 16.8, True, accents[index % len(accents)], self.theme.title_font, 1)
                self._add_text(slide, self._truncate(self._metric_label(metric), 22), x + int(slot_width * 0.07), int(top + height * 0.57), int(slot_width * 0.86), int(height * 0.3), 7.2, True, self.theme.ink, self.theme.body_font, 1)
            else:
                self._rounded_rect(slide, x, top, slot_width, height, fills[index % len(fills)], self.theme.rule)
                self._add_text(slide, self._metric_value(metric.value), x + int(slot_width * 0.07), int(top + height * 0.08), int(slot_width * 0.86), int(height * 0.42), 14, True, accents[index % len(accents)], self.theme.title_font, 1)
                self._add_text(slide, self._metric_label(metric), x + int(slot_width * 0.07), int(top + height * 0.53), int(slot_width * 0.86), int(height * 0.32), 7.8, False, self.theme.muted, self.theme.body_font, 1)

    def _render_metric_column(self, slide, metrics: Sequence[MetricBlock], left, top, width, height) -> None:
        metrics = list(metrics)[:4]
        if not metrics:
            return
        gap = self.Inches(0.18)
        slot_height = min(self.Inches(1.0), int((int(height) - int(gap) * (len(metrics) - 1)) / len(metrics)))
        for index, metric in enumerate(metrics):
            y = int(top + index * (slot_height + int(gap)))
            fill = self.theme.pale_primary if index % 2 == 0 else self.theme.pale_secondary
            self._rounded_rect(slide, left, y, width, slot_height, fill, self.theme.rule)
            value_size = 20 if self.theme.layout_family in {"report", "systems"} else 18
            self._add_text(slide, self._metric_value(metric.value), int(left + width * 0.08), y + int(slot_height * 0.16), int(width * 0.84), int(slot_height * 0.36), value_size, True, self.theme.primary, self.theme.title_font, 1)
            self._add_text(slide, self._metric_label(metric), int(left + width * 0.08), y + int(slot_height * 0.58), int(width * 0.84), int(slot_height * 0.28), 8.2, True, self.theme.muted, self.theme.body_font, 1)

    def _render_table(self, slide, slide_spec: SlideSpec, table_left, table_top, table_width, table_height) -> None:
        if not slide_spec.table_blocks:
            return
        table_block = slide_spec.table_blocks[0]
        rows = (table_block.rows or [["No table data"]])[:6]
        columns = min(max(len(row) for row in rows), 4)
        rows = [(row + [""] * columns)[:columns] for row in rows]

        shape = slide.shapes.add_table(len(rows), columns, table_left, table_top, table_width, table_height)
        table = shape.table
        for row_index, row in enumerate(rows):
            for col_index, value in enumerate(row):
                cell = table.cell(row_index, col_index)
                cell.text = self._truncate(str(value), 44)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = self.Pt(9.5 if row_index else 10.5)
                paragraph.font.bold = row_index == 0
                paragraph.font.name = self.theme.body_font
                paragraph.font.color.rgb = self._rgb(self.theme.ink)
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._rgb(self.theme.pale_primary if row_index == 0 else self.theme.surface if row_index % 2 else self.theme.surface_alt)

        if table_block.caption:
            self._add_text(slide, self._truncate(table_block.caption, 120), table_left, int(table_top + table_height + self.Pt(5)), table_width, self.Pt(32), 8.5, False, self.theme.muted, self.theme.body_font, 1)

    def _add_bullet_list(self, slide, blocks: Sequence[TextBlock], left, top, width, height, size: float, max_items: int) -> None:
        blocks = list(blocks)[:max_items]
        if not blocks:
            return
        self._render_numbered_points(slide, blocks, left, top, width, height, size=size, max_items=max_items)

    def _render_stage_cards(self, slide, blocks: Sequence[TextBlock], left, top, width, height, horizontal: bool) -> None:
        blocks = list(blocks)
        if not blocks:
            return
        t = self.theme
        max_count = 4 if t.layout_family == "systems" else 3
        count = min(max_count, len(blocks))
        gap = self.Inches(0.12)
        palette = [t.primary, t.secondary, t.accent, t.ink]
        if horizontal:
            slot_w = int((int(width) - int(gap) * (count - 1)) / count)
            slot_h = int(height)
            for idx, block in enumerate(blocks[:count], start=1):
                x = int(left + (idx - 1) * (slot_w + int(gap)))
                self._render_one_stage_card(slide, block, idx, x, top, slot_w, slot_h, palette[(idx - 1) % len(palette)])
        else:
            slot_h = int((int(height) - int(gap) * (count - 1)) / count)
            slot_h = max(slot_h, self.Inches(0.48))
            for idx, block in enumerate(blocks[:count], start=1):
                y = int(top + (idx - 1) * (slot_h + int(gap)))
                self._render_one_stage_card(slide, block, idx, left, y, width, slot_h, palette[(idx - 1) % len(palette)])

    def _render_one_stage_card(self, slide, block: TextBlock, index: int, left, top, width, height, accent) -> None:
        t = self.theme
        lead, detail = self._point_parts(block)
        fill = t.surface if t.layout_family in {"editorial", "conference", "explainer"} else t.pale_neutral
        if t.layout_family == "systems":
            fill = t.pale_primary if index % 2 else t.pale_secondary
        if t.layout_family == "report":
            fill = t.surface_alt if index % 2 else t.surface
        self._rounded_rect(slide, left, top, width, height, fill, t.rule)
        marker_w = min(self.Inches(0.44), int(width * 0.18))
        self._rounded_rect(slide, left, top, marker_w, height, accent, accent)
        self._add_text(slide, f"{index:02d}", int(left + marker_w * 0.08), int(top + height * 0.18), int(marker_w * 0.82), self.Inches(0.18), 7.2, True, (255, 255, 255), t.title_font, 1, alignment=self.PP_ALIGN.CENTER)
        text_left = int(left + marker_w + self.Inches(0.12))
        text_width = int(width - marker_w - self.Inches(0.22))
        compact = int(height) < self.Inches(0.82) or int(width) < self.Inches(2.25)
        lead_size = 8.6 if compact else 10.2
        detail_size = 7.2 if compact else 8.3
        lead_text = self._headline_text(lead or block.text, 34 if compact else 46)
        self._add_text(slide, lead_text, text_left, int(top + height * 0.16), text_width, int(height * (0.56 if compact else 0.3)), lead_size, True, accent, t.body_font, 1)
        if not compact and int(height) >= self.Inches(0.74):
            self._add_text(slide, self._body_sentence(detail or block.text, 92), text_left, int(top + height * 0.48), text_width, int(height * 0.38), detail_size, False, t.ink, t.body_font, 2)

    def _add_system_lane(self, slide, left, top, width, height, slide_spec: SlideSpec) -> None:
        t = self.theme
        self._rounded_rect(slide, left, top, width, height, t.surface, t.rule)
        self._add_text(slide, "Workflow Reading", int(left + self.Inches(0.18)), int(top + self.Inches(0.18)), int(width - self.Inches(0.36)), self.Inches(0.2), 8.2, True, t.secondary, t.body_font, 1)
        labels = [slide_spec.section_label or "Input", "Evidence", "Decision", "Output"]
        for idx, label in enumerate(labels[:4], start=1):
            y = int(top + self.Inches(0.72) + (idx - 1) * self.Inches(0.74))
            self._rounded_rect(slide, int(left + self.Inches(0.2)), y, int(width - self.Inches(0.4)), self.Inches(0.38), t.pale_primary if idx % 2 else t.pale_secondary, t.rule)
            self._add_text(slide, self._truncate(str(label), 24), int(left + self.Inches(0.38)), int(y + self.Inches(0.1)), int(width - self.Inches(0.76)), self.Inches(0.16), 7.6, True, t.ink, t.body_font, 1)

    def _render_key_message(self, slide, text: str, left, top, width, height, accent=None) -> None:
        if not text:
            return
        t = self.theme
        accent = accent or t.primary
        if t.layout_family == "editorial":
            self._rounded_rect(slide, left, top, self.Inches(0.05), height, accent, accent)
            self._add_text(
                slide,
                self._body_sentence(text, 240),
                int(left + self.Inches(0.24)),
                int(top + self.Inches(0.04)),
                int(width - self.Inches(0.28)),
                int(height - self.Inches(0.06)),
                self._fit_title_size(text, base=14.2, min_size=10.5),
                True,
                t.ink,
                t.title_font,
                2,
            )
            return
        if t.layout_family == "systems":
            self._rounded_rect(slide, left, top, width, height, t.surface, t.rule)
            self._rounded_rect(slide, left, top, self.Inches(0.18), height, accent, accent)
            self._add_text(slide, "STATE", int(left + self.Inches(0.34)), int(top + self.Inches(0.1)), int(width - self.Inches(0.46)), self.Inches(0.16), 6.8, True, accent, t.body_font, 1)
            self._add_text(
                slide,
                self._body_sentence(text, 220),
                int(left + self.Inches(0.34)),
                int(top + self.Inches(0.3)),
                int(width - self.Inches(0.48)),
                int(height - self.Inches(0.34)),
                self._fit_title_size(text, base=12.6, min_size=9.8),
                True,
                t.ink,
                t.body_font,
                2,
            )
            return
        self._rounded_rect(slide, left, top, width, height, t.pale_primary, t.rule)
        self._rounded_rect(slide, left, top, self.Inches(0.12), height, accent, accent)
        self._add_text(
            slide,
            "Key message",
            int(left + self.Inches(0.28)),
            int(top + self.Inches(0.12)),
            int(width - self.Inches(0.42)),
            self.Inches(0.18),
            7.5,
            True,
            accent,
            t.body_font,
            1,
        )
        self._add_text(
            slide,
            self._body_sentence(text, 230),
            int(left + self.Inches(0.28)),
            int(top + self.Inches(0.33)),
            int(width - self.Inches(0.46)),
            int(height - self.Inches(0.38)),
            self._fit_title_size(text, base=13.0, min_size=10),
            True,
            t.ink,
            t.title_font,
            2,
        )

    def _render_numbered_points(self, slide, blocks: Sequence[TextBlock], left, top, width, height, size: float, max_items: int) -> None:
        t = self.theme
        blocks = list(blocks)[:max_items]
        if not blocks:
            return
        count = len(blocks)
        gap = self.Inches(0.12)
        slot_height = int((int(height) - int(gap) * (count - 1)) / count)
        min_slot = int(self.Inches(0.58))
        slot_height = max(min_slot, slot_height)
        palette = [t.primary, t.secondary, t.accent, t.ink, t.primary]

        for index, block in enumerate(blocks, start=1):
            y = int(top + (index - 1) * (slot_height + int(gap)))
            accent = palette[(index - 1) % len(palette)]
            lead, detail = self._point_parts(block)
            text_top = y + int(slot_height * 0.15)
            number_size = min(self.Inches(0.36), int(slot_height * 0.62))

            self._rounded_rect(slide, left, y + int(slot_height * 0.1), self.Inches(0.46), number_size, accent, accent)
            self._add_text(
                slide,
                str(index),
                int(left + self.Inches(0.06)),
                y + int(slot_height * 0.18),
                self.Inches(0.34),
                self.Inches(0.18),
                8.5,
                True,
                (255, 255, 255),
                t.title_font,
                1,
                alignment=self.PP_ALIGN.CENTER,
            )
            text_left = int(left + self.Inches(0.72))
            text_width = int(width - self.Inches(0.76))
            if detail and slot_height >= self.Inches(0.62):
                self._add_text(slide, lead, text_left, text_top, text_width, max(self.Inches(0.26), int(slot_height * 0.28)), min(size, 11.4), True, accent, t.body_font, 1)
                self._add_text(
                    slide,
                    self._body_sentence(detail, 220),
                    text_left,
                    y + int(slot_height * 0.44),
                    text_width,
                    max(self.Inches(0.28), int(slot_height * 0.42)),
                    min(size - 0.8, 10.8),
                    False,
                    t.ink,
                    t.body_font,
                    2,
                )
            else:
                self._add_text(
                    slide,
                    self._body_sentence(block.text, 220),
                    text_left,
                    text_top,
                    text_width,
                    max(self.Inches(0.32), int(slot_height * 0.68)),
                    min(size, 11.8),
                    False,
                    t.ink,
                    t.body_font,
                    2,
                )

    def _point_parts(self, block: TextBlock) -> tuple[str, str]:
        claim = self._headline_text(getattr(block, "claim", "") or "", 54)
        detail = self._body_sentence(getattr(block, "detail", "") or "", 170)
        if claim and detail:
            return claim, detail
        return self._split_point(block.text)

    def _add_text(self, slide, text: str, left, top, width, height, size: float, bold: bool = False, color=None, font: str | None = None, max_lines: int = 2, alignment=None):
        text = self._truncate_lines(text or "", max_lines=max_lines)
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = self.Inches(0.02)
        frame.margin_right = self.Inches(0.02)
        frame.margin_top = self.Inches(0.01)
        frame.margin_bottom = self.Inches(0.01)
        frame.vertical_anchor = self.MSO_ANCHOR.TOP
        frame.clear()
        lines = text.split("\n") or [""]
        for index, line in enumerate(lines):
            p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            p.text = line
            p.font.size = self.Pt(size)
            p.font.bold = bold
            p.font.name = font or self.theme.body_font
            p.font.color.rgb = self._rgb(color or self.theme.ink)
            if alignment is not None:
                p.alignment = alignment
            p.space_after = self.Pt(5 if len(lines) > 1 else 0)
        return box

    def _rounded_rect(self, slide, left, top, width, height, fill, line) -> None:
        shape = slide.shapes.add_shape(self.MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(fill)
        shape.line.color.rgb = self._rgb(line)

    def _placeholder(self, slide, left, top, width, height, text: str) -> None:
        self._rounded_rect(slide, left, top, width, height, self.theme.surface_alt, self.theme.rule)
        self._add_text(slide, text or "Original figure", left, int(top + height * 0.43), width, self.Inches(0.28), 10, False, self.theme.muted, self.theme.body_font, 1)

    def _footer(self, slide, slide_index: int) -> None:
        self._add_text(slide, str(slide_index), self.Inches(11.95), self.Inches(7.03), self.Inches(0.45), self.Inches(0.18), 8, False, self.theme.muted, self.theme.body_font, 1)

    def _fit_picture(self, image_path: Path, box_left, box_top, box_width, box_height):
        try:
            from PIL import Image

            with Image.open(image_path) as img:
                source_width, source_height = img.size
            scale = min(int(box_width) / max(1, source_width), int(box_height) / max(1, source_height))
            width = max(1, int(source_width * scale))
            height = max(1, int(source_height * scale))
            left = int(box_left + (int(box_width) - width) / 2)
            top = int(box_top + (int(box_height) - height) / 2)
            return left, top, width, height
        except Exception:
            return box_left, box_top, box_width, box_height

    def _fit_title_size(self, text: str, base: int, min_size: int) -> int:
        length = len(text or "")
        if length <= 42:
            return base
        if length <= 62:
            return max(min_size, base - 3)
        if length <= 82:
            return max(min_size, base - 6)
        return min_size

    def _clean_title(self, text: str) -> str:
        cleaned = (text or "").replace(" - Cover", "")
        cleaned = cleaned.replace("**", "")
        cleaned = re.sub(r"^\s*Title:\s*", "", cleaned, flags=re.IGNORECASE)
        return self._truncate(cleaned, 82)

    def _split_point(self, text: str) -> tuple[str, str]:
        text = self._body_sentence((text or "").replace("**", ""), 170)
        if ":" not in text:
            lead = self._auto_point_lead(text)
            detail = text
            if lead and detail and lead.lower().rstrip(".") != detail.lower().rstrip("."):
                return lead, detail
            return text.strip(" -"), ""
        lead, detail = text.split(":", 1)
        lead = lead.strip(" -")
        detail = detail.strip(" -")
        if not lead or len(lead.split()) > 12:
            auto_lead = self._auto_point_lead(text)
            return (auto_lead, text) if auto_lead else (text.strip(" -"), "")
        return self._headline_text(lead, 54), self._body_sentence(detail, 132)

    def _auto_point_lead(self, text: str) -> str:
        text = self._strip_ellipsis(text)
        before_colon = text.split(":", 1)[0].strip(" -")
        if before_colon and 2 <= len(before_colon.split()) <= 12:
            return self._headline_text(before_colon, 54)
        first_clause = re.split(r",|;|\sdue to\s|\sbecause\s|\bthat\s", text, maxsplit=1, flags=re.IGNORECASE)[0].strip(" -")
        if 2 <= len(first_clause.split()) <= 8:
            return self._headline_text(first_clause, 54)
        words = text.split()
        if len(words) < 8:
            return ""
        stop_words = {
            "the", "a", "an", "of", "to", "in", "and", "or", "for", "with", "that", "this", "these", "those",
            "is", "are", "was", "were", "be", "been", "being", "has", "have", "had", "can", "could",
            "such",
        }
        lead_words = []
        for word in words:
            clean = word.strip(" ,.;:()[]").lower()
            lead_words.append(word.strip(" ,.;:"))
            if len(lead_words) >= 4 and clean not in stop_words:
                break
            if len(lead_words) >= 7:
                break
        return self._headline_text(" ".join(lead_words), 54)

    def _body_sentence(self, text: str, max_len: int) -> str:
        text = self._strip_ellipsis(text)
        if len(text) <= max_len:
            return self._ensure_sentence(text)
        cut = text[:max_len].rstrip()
        sentence_end = max(cut.rfind("."), cut.rfind(";"), cut.rfind("!"), cut.rfind("?"))
        if sentence_end >= max_len * 0.45:
            cut = cut[: sentence_end + 1]
        else:
            cut = cut.rsplit(" ", 1)[0].rstrip(" ,;:-")
            weak_endings = {
                "a", "an", "the", "of", "to", "in", "on", "for", "with", "by", "and", "or", "that", "which",
                "due", "associated", "including", "such", "as", "from", "through", "into", "across", "between",
                "during",
            }
            words = cut.split()
            while len(words) > 6 and words[-1].strip(" ,;:-").lower() in weak_endings:
                words.pop()
            cut = " ".join(words)
        return self._ensure_sentence(cut)

    def _headline_text(self, text: str, max_len: int) -> str:
        text = self._strip_ellipsis(text).strip(" .;:-")
        if len(text) > max_len:
            text = text[:max_len].rsplit(" ", 1)[0].strip(" .;:-")
        words = text.split()
        weak_endings = {"a", "an", "the", "of", "to", "in", "on", "for", "with", "by", "and", "or", "that", "which"}
        while len(words) > 2 and words[-1].strip(" ,;:-").lower() in weak_endings:
            words.pop()
        return " ".join(words)

    def _ensure_sentence(self, text: str) -> str:
        text = self._strip_ellipsis(text).strip()
        if not text:
            return ""
        if text[-1] in ".!?":
            return text
        return text + "."

    def _strip_ellipsis(self, text: str) -> str:
        return re.sub(r"\.{2,}\s*$", "", " ".join((text or "").split())).strip()

    def _truncate_lines(self, text: str, max_lines: int) -> str:
        lines = [line.strip() for line in str(text).split("\n") if line.strip()]
        if len(lines) > max_lines:
            lines = lines[:max_lines]
            lines[-1] = self._truncate(lines[-1], max(12, len(lines[-1]) - 3))
        return "\n".join(lines)

    def _truncate(self, text: str, max_len: int) -> str:
        text = " ".join((text or "").split())
        if len(text) <= max_len:
            return text
        return text[: max_len - 1].rstrip(" ,;:-") + "..."

    def _metric_value(self, value: str) -> str:
        return self._truncate(value or "", 18)

    def _metric_label(self, metric: MetricBlock) -> str:
        label = metric.label or metric.note
        if not label and metric.note:
            label = metric.note
        return self._headline_text(label or "", 34)

    def _rgb(self, value):
        return self.RGBColor(*value)
