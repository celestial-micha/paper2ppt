"""
Lightweight PPTX layout QA for generated native decks.

This is not a full PowerPoint renderer, but it catches common defects before
delivery: out-of-bounds shapes, very long title text, empty slides, and likely
text overflow based on box size and word count.
"""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import re
from typing import Dict, List, Optional

from .slide_schema import PresentationSpec


@dataclass
class PptxQaResult:
    passed: bool
    warnings: List[str]
    slide_count: int

    def to_dict(self) -> Dict:
        return {
            "passed": self.passed,
            "warnings": self.warnings,
            "slide_count": self.slide_count,
        }


def inspect_pptx_layout(pptx_path: Path) -> PptxQaResult:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for PPTX QA.") from exc

    prs = Presentation(str(pptx_path))
    warnings: List[str] = []
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)

    for slide_index, slide in enumerate(prs.slides, start=1):
        visible_text_shapes = 0
        visible_pictures = 0
        visible_tables = 0

        for shape in slide.shapes:
            left = int(shape.left)
            top = int(shape.top)
            right = left + int(shape.width)
            bottom = top + int(shape.height)

            if right < 0 or bottom < 0 or left > slide_width or top > slide_height:
                warnings.append(f"slide {slide_index}: shape is fully outside slide bounds")
                continue

            if left < -1000 or top < -1000 or right > slide_width + 1000 or bottom > slide_height + 1000:
                warnings.append(f"slide {slide_index}: shape partly exceeds slide bounds")

            text = ""
            if getattr(shape, "has_text_frame", False):
                text = (shape.text or "").strip()
                if text:
                    visible_text_shapes += 1
                    _check_text_box(slide_index, shape, text, warnings)

            if getattr(shape, "has_table", False):
                visible_tables += 1
            if shape.shape_type == 13:
                visible_pictures += 1

        if visible_text_shapes == 0 and visible_pictures == 0 and visible_tables == 0:
            warnings.append(f"slide {slide_index}: slide appears empty")

    severe = [w for w in warnings if "outside" in w or "exceeds" in w or "empty" in w]
    return PptxQaResult(passed=not severe, warnings=warnings, slide_count=len(prs.slides))


def evaluate_presentation_spec(
    spec: PresentationSpec,
    layout_result: Optional[PptxQaResult] = None,
) -> Dict:
    """Evaluate semantic/spec quality plus optional PPTX layout warnings."""
    warnings: List[str] = []
    failed_slides: set[int] = set()

    if not spec.slides:
        warnings.append("deck: slide spec is empty")

    for slide_index, slide in enumerate(spec.slides, start=1):
        slide_warnings = _evaluate_slide_spec(slide_index, slide)
        if slide_warnings:
            if any(_is_severe_warning(warning) for warning in slide_warnings):
                failed_slides.add(slide_index)
            warnings.extend(slide_warnings)

    layout_warnings = list(layout_result.warnings if layout_result else [])
    for warning in layout_warnings:
        match = re.search(r"slide\s+(\d+)", warning, flags=re.IGNORECASE)
        if match and _is_severe_warning(warning):
            failed_slides.add(int(match.group(1)))
        warnings.append(f"layout: {warning}")

    severe = [warning for warning in warnings if _is_severe_warning(warning)]
    return {
        "passed": not severe,
        "warnings": warnings,
        "failed_slides": sorted(failed_slides),
        "slide_count": len(spec.slides),
        "layout": layout_result.to_dict() if layout_result else None,
        "checks": [
            "empty components",
            "meaningless decoration/content placeholders",
            "truncated ellipsis",
            "numbered point claim/detail/evidence",
            "metric label/value quality",
            "layout QA",
        ],
    }


def _evaluate_slide_spec(slide_index: int, slide) -> List[str]:
    warnings: List[str] = []
    title = (slide.title or "").strip()
    if not title:
        warnings.append(f"slide {slide_index}: empty title")

    has_content = any(
        [
            (slide.takeaway or "").strip(),
            slide.text_blocks,
            slide.image_blocks,
            slide.table_blocks,
            slide.metric_blocks,
        ]
    )
    if not has_content:
        warnings.append(f"slide {slide_index}: slide spec appears empty")

    for point_index, point in enumerate(slide.text_blocks, start=1):
        claim = (getattr(point, "claim", "") or "").strip()
        detail = (getattr(point, "detail", "") or "").strip()
        evidence = (getattr(point, "evidence", "") or "").strip()
        text = (point.text or "").strip()
        if not claim:
            warnings.append(f"slide {slide_index}: point {point_index} missing claim")
        if not detail:
            warnings.append(f"slide {slide_index}: point {point_index} missing detail")
        if not evidence:
            warnings.append(f"slide {slide_index}: point {point_index} missing evidence")
        if _has_truncated_ellipsis(" ".join([text, claim, detail, evidence])):
            warnings.append(f"slide {slide_index}: point {point_index} contains truncated ellipsis")
        if claim and detail and claim.lower().rstrip(".") == detail.lower().rstrip("."):
            warnings.append(f"slide {slide_index}: point {point_index} claim repeats detail")

    for metric_index, metric in enumerate(slide.metric_blocks, start=1):
        label = (metric.label or "").strip()
        value = (metric.value or "").strip()
        note = (metric.note or "").strip()
        if not label or label.lower() in {"metric", "key metric", "key number", "number"}:
            warnings.append(f"slide {slide_index}: metric {metric_index} has meaningless label")
        if not value:
            warnings.append(f"slide {slide_index}: metric {metric_index} missing value")
        elif not re.search(r"\d|%|=|x|k|m|b", value, flags=re.IGNORECASE):
            warnings.append(f"slide {slide_index}: metric {metric_index} value looks non-quantitative")
        if _has_truncated_ellipsis(" ".join([label, value, note])):
            warnings.append(f"slide {slide_index}: metric {metric_index} contains truncated ellipsis")

    for image_index, image in enumerate(slide.image_blocks, start=1):
        if not any([(image.path or "").strip(), (image.title or "").strip(), (image.caption or "").strip(), (image.placeholder_text or "").strip()]):
            warnings.append(f"slide {slide_index}: image {image_index} is an empty component")
        placeholder = (image.placeholder_text or "").strip().lower()
        if placeholder in {"original figure", "figure", "image", "placeholder"} and not image.caption:
            warnings.append(f"slide {slide_index}: image {image_index} has meaningless decoration placeholder")

    for table_index, table in enumerate(slide.table_blocks, start=1):
        if not table.rows:
            warnings.append(f"slide {slide_index}: table {table_index} is empty")

    if _has_truncated_ellipsis(" ".join([slide.title or "", slide.takeaway or ""])):
        warnings.append(f"slide {slide_index}: title/takeaway contains truncated ellipsis")

    return warnings


def _has_truncated_ellipsis(text: str) -> bool:
    return bool(re.search(r"(\.\.\.|…)\s*$|(\.\.\.|…)\s+[A-Z0-9]", text or ""))


def _is_severe_warning(warning: str) -> bool:
    return any(
        marker in warning.lower()
        for marker in (
            "missing claim",
            "missing detail",
            "empty",
            "meaningless",
            "truncated",
            "outside",
            "exceeds",
            "appears empty",
        )
    )


def _check_text_box(slide_index: int, shape, text: str, warnings: List[str]) -> None:
    width_inches = int(shape.width) / 914400
    height_inches = int(shape.height) / 914400
    words = len(text.split())
    lines = [line for line in text.splitlines() if line.strip()]

    if width_inches > 8 and height_inches < 0.35 and len(text) > 90:
        warnings.append(f"slide {slide_index}: long title/subtitle may wrap or clip")

    if height_inches < 0.25 and len(text) > 32:
        warnings.append(f"slide {slide_index}: very small text box contains long text")

    if height_inches < 0.75 and words > 22:
        warnings.append(f"slide {slide_index}: text box may overflow vertically")

    if len(lines) >= 5 and height_inches < 2.0:
        warnings.append(f"slide {slide_index}: dense bullet list may overflow")
