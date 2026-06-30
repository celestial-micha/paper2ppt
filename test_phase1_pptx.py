import unittest
import json
import uuid
from pathlib import Path
from unittest.mock import patch

from pptx import Presentation

from paper2slides.generator.detailed_tex import generate_detailed_tex_deck
from paper2slides.generator.config import GenerationConfig, GenerationInput, OutputType, SlidesLength, StyleType
from paper2slides.generator.content_planner import ContentPlanner
from paper2slides.generator.content_planner import ContentPlan, FigureRef, Section, TableRef
from paper2slides.generator.pptx_qa import _has_truncated_ellipsis, evaluate_presentation_spec, inspect_pptx_layout
from paper2slides.generator.pptx_renderer import PptxRenderer
from paper2slides.generator.text_pptx_workflow import _build_speaker_script, _compact_metric_blocks, _ensure_structured_points, _extract_metrics_from_slide, _get_figure_analysis_model, _get_pptx_llm_model, _normalize_slide_layout, _qa_repair_node, _repair_slide_title
from paper2slides.generator.spec_builder import build_presentation_spec
from paper2slides.generator.slide_schema import MetricBlock, PresentationSpec, SlideSpec, TableBlock, TextBlock
from paper2slides.benchmark.qa_summary import QaRunResult, classify_warning, summarize_layout_qa
from paper2slides.benchmark.papers import expand_paper_set, validate_paper_files
from paper2slides.benchmark.runner import _build_command, _preflight_environment, summarize_run_results
from paper2slides.benchmark.from_scratch import build_content_inventory, build_rough_draft_spec, render_rough_draft_pptx, write_from_scratch_artifacts
from paper2slides.benchmark.fourway import _audit_speaker_script, _build_blind_style_contract, build_style_drift_report_payload, render_blind_blueprint_pptx
from paper2slides.benchmark.human_feedback import badcase_ids, load_human_feedback_benchmark, summarize_human_feedback_benchmark
from paper2slides.benchmark.nonvisual_audit import _academic_right_evidence_void_findings, _academic_toc_canonical_sections_findings, _semantic_content_findings, inspect_pptx_nonvisual
from paper2slides.core.stages.rag_stage import _run_fast_queries_by_category
from paper2slides.summary import FigureInfo, GeneralContent, OriginalElements, TableInfo


class Phase1PptxSmokeTest(unittest.TestCase):
    def test_builds_slide_spec_and_pptx(self):
        plan = ContentPlan(
            output_type="slides",
            sections=[
                Section(
                    id="slide_01",
                    title="Phase 1 Smoke Test",
                    section_type="opening",
                    content="This is the first editable slide. It should render as body text in a PPTX.",
                    tables=[TableRef(table_id="Table 1", focus="Key metrics")],
                    figures=[FigureRef(figure_id="Figure 1", focus="Keep as reference")],
                )
            ],
            tables_index={
                "Table 1": TableInfo(
                    table_id="Table 1",
                    caption="Example table",
                    html_content="<table><tr><th>Metric</th><th>Value</th></tr><tr><td>Slides</td><td>1</td></tr></table>",
                )
            },
            figures_index={
                "Figure 1": FigureInfo(
                    figure_id="Figure 1",
                    caption="Missing image is allowed",
                    image_path="missing.png",
                )
            },
            metadata={"page_range": [1, 1]},
        )

        spec = build_presentation_spec(plan, title="Smoke Test")
        self.assertEqual(len(spec.slides), 1)
        self.assertTrue(spec.slides[0].table_blocks)
        self.assertTrue(spec.slides[0].image_blocks)

        temp_root = Path(__file__).parent / "outputs" / "tmp"
        temp_root.mkdir(parents=True, exist_ok=True)
        output_path = temp_root / f"slides_{uuid.uuid4().hex}.pptx"
        renderer = PptxRenderer()
        renderer.render(spec, output_path)
        self.assertTrue(output_path.exists())

    def test_toc_with_seven_sections_stays_inside_slide_bounds(self):
        slides = []
        for index in range(7):
            slides.append(
                SlideSpec(
                    slide_id=f"slide_{index + 1:02d}",
                    title=f"Section {index + 1} Slide",
                    section_label=f"Section {index + 1}",
                    layout="statement",
                    takeaway="This slide checks table of contents spacing.",
                    text_blocks=[
                        TextBlock(
                            text="Claim: The generated deck keeps section labels readable.",
                            claim="The generated deck keeps section labels readable.",
                            detail="Seven section entries should fit inside the table of contents slide.",
                            evidence="Renderer layout QA",
                        )
                    ],
                )
            )
        spec = PresentationSpec(title="TOC Fit", slides=slides)

        temp_root = Path(__file__).parent / "outputs" / "tmp"
        temp_root.mkdir(parents=True, exist_ok=True)
        output_path = temp_root / f"toc_{uuid.uuid4().hex}.pptx"
        PptxRenderer().render(spec, output_path)

        qa = inspect_pptx_layout(output_path)
        self.assertTrue(qa.passed, qa.warnings)

    def test_academic_toc_preserves_mhc_six_module_route(self):
        spec = PresentationSpec(
            title="Academic Route",
            slides=[
                SlideSpec(slide_id="slide_01", title="Cover", section_type="opening", section_label="Overview"),
                SlideSpec(slide_id="slide_02", title="Problem", section_label="Motivation"),
                SlideSpec(slide_id="slide_03", title="Approach", section_label="Method"),
                SlideSpec(slide_id="slide_04", title="Findings", section_label="Results"),
                SlideSpec(slide_id="slide_05", title="Wrap-up", section_type="ending", section_label="Conclusion"),
            ],
        )

        sections = PptxRenderer(style="academic")._section_sequence(spec)

        self.assertEqual(
            sections,
            ["Motivation", "Method", "Analysis", "Ablations", "Results", "Conclusion"],
        )

    def test_academic_toc_missing_canonical_sections_is_detected(self):
        def record(shape_id, text, x, y, w, h):
            return {
                "shape_id": shape_id,
                "role": "body_text",
                "has_text": True,
                "text": text,
                "has_table": False,
                "is_picture": False,
                "bbox": {"x": x, "y": y, "w": w, "h": h, "right": x + w, "bottom": y + h, "area": w * h},
                "font": {"avg_pt": 12.0},
            }

        records = [
            record(1, "Contents", 0.78, 0.62, 4.5, 0.68),
            record(2, "A sectioned route through the paper: why it matters, what is new, how it works, and what it proves.", 0.82, 1.35, 9.4, 0.42),
            record(3, "Motivation", 1.72, 2.27, 5.4, 0.34),
            record(4, "Method", 1.72, 3.13, 5.4, 0.34),
            record(5, "Results", 1.72, 3.99, 5.4, 0.34),
            record(6, "Conclusion", 1.72, 4.85, 5.4, 0.34),
        ]

        findings = _academic_toc_canonical_sections_findings(2, records)

        self.assertEqual(findings[0]["type"], "academic_toc_missing_canonical_sections")
        self.assertEqual(findings[0]["evidence"]["missing_sections"], ["Analysis", "Ablations"])

    def test_slide_planning_uses_text_only_figure_manifest(self):
        planner = ContentPlanner(api_key="test-key", base_url="http://localhost", model="test-model")

        captured = {}

        def fake_text_llm(prompt):
            captured["prompt"] = prompt
            return """
```json
{"slides":[{"id":"slide_01","title":"Text-only","content":"Uses an original figure.","tables":[],"figures":[{"figure_id":"Figure 1","focus":"Pipeline"}]}]}
```
"""

        def fail_multimodal(*args, **kwargs):
            raise AssertionError("Slides planning must not call the multimodal/image path")

        planner._call_text_llm = fake_text_llm
        planner._call_multimodal_llm = fail_multimodal
        planner._load_figure_images = fail_multimodal

        gen_input = GenerationInput(
            config=GenerationConfig(
                output_type=OutputType.SLIDES,
                slides_length=SlidesLength.SHORT,
                style=StyleType.ACADEMIC,
            ),
            content=GeneralContent(content="A paper summary with a method pipeline."),
            origin=OriginalElements(
                figures=[
                    FigureInfo(
                        figure_id="Figure 1",
                        caption="Pipeline overview",
                        image_path="images/pipeline.png",
                    )
                ]
            ),
        )

        plan = planner.plan(gen_input)
        self.assertEqual(len(plan.sections), 1)
        self.assertIn("Source image: images/pipeline.png", captured["prompt"])
        self.assertNotIn("[FIGURE_IMAGES]", captured["prompt"])

    def test_qa_repair_compacts_risky_slide(self):
        long_text = " ".join(["overflow"] * 30)
        spec = PresentationSpec(
            title="QA Repair",
            slides=[
                SlideSpec(
                    slide_id="slide_01",
                    title="This title is intentionally much too long for a normal presentation header",
                    layout="statement",
                    takeaway=long_text,
                    text_blocks=[TextBlock(text=long_text) for _ in range(5)],
                    metric_blocks=[MetricBlock(label="Very long label", value="1234567890") for _ in range(4)],
                )
            ],
        )

        repaired = _qa_repair_node(
            {
                "spec": spec,
                "qa_warnings": ["slide 1: long title/subtitle may wrap or clip"],
                "validation_warnings": [],
            }
        )
        slide = repaired["spec"].slides[0]
        self.assertEqual(repaired["qa_attempt"], 1)
        self.assertLessEqual(len(slide.text_blocks), 3)
        self.assertLessEqual(len(slide.title.split()), 8)
        self.assertLessEqual(len(slide.takeaway.split()), 12)

    def test_builds_speaker_script_from_final_spec(self):
        spec = PresentationSpec(
            title="Script Test",
            slides=[
                SlideSpec(
                    slide_id="slide_01",
                    title="Opening",
                    takeaway="The method fails visually.",
                    text_blocks=[TextBlock(text="Screenshots drive most errors.")],
                    metric_blocks=[MetricBlock(label="Success rate", value="5.36%")],
                    notes=["Referenced figures: Figure 1"],
                )
            ],
        )

        script = _build_speaker_script(spec, ["attempt 1: slide 1 compressed"])
        self.assertIn("# Script Test", script)
        self.assertIn("## Slide 1: Opening", script)
        self.assertIn("Suggested narration", script)
        self.assertIn("Success rate: 5.36%", script)

    def test_blind_blueprint_renderer_avoids_protected_baseline_tokens(self):
        inventory = {
            "paper": {"title": "AGI Is Coming"},
            "paper_highlights": [
                {"label": "Core claim", "text": "Wordle-style environments expose whether agents can reason under uncertainty."}
            ],
            "assets": {"figures": [], "tables": []},
            "metrics": [{"label": "Tasks", "value": "12"}],
            "curated_slides": [],
        }
        rough = {
            "slides": [
                {
                    "title": "Wordle stresses reasoning traces",
                    "slide_role": "thesis",
                    "claim": "The benchmark emphasizes deliberate hypothesis updates.",
                    "support": "Agents must revise beliefs from sparse feedback instead of only matching surface patterns.",
                    "proof_object": {
                        "type": "text_evidence",
                        "id": "Motivation",
                        "focus": "Wordle offers compact but adversarial feedback loops for reasoning evaluation.",
                    },
                    "source_evidence": [{"source": "summary", "id": "motivation"}],
                }
            ]
        }
        temp_root = Path(__file__).parent / "outputs" / "tmp"
        temp_root.mkdir(parents=True, exist_ok=True)
        output_path = temp_root / f"blind_blueprint_{uuid.uuid4().hex}.pptx"

        render_blind_blueprint_pptx(inventory, rough, output_path, repair_mode=True)
        self.assertTrue(output_path.exists())
        prs = Presentation(output_path)
        all_text = "\n".join(getattr(shape, "text", "") for slide in prs.slides for shape in slide.shapes).upper()
        self.assertNotIn("PROOF OBJECT", all_text)
        self.assertNotIn("DECK MAP", all_text)
        self.assertNotIn("ACADEMIC PAPER READING", all_text)

        report = build_style_drift_report_payload(
            pptx_path=output_path,
            audit={"findings": []},
            style_id="blind_experimental_blueprint",
            repair_profile="experimental_from_scratch_loop",
            style_scope="experimental",
        )
        self.assertEqual(report["drift_risk"], "low")
        self.assertEqual(report["baseline_similarity_signals"], {})

    def test_blind_style_contract_is_run_scoped(self):
        inventory = {"paper": {"title": "Deep Residual Learning for Image Recognition"}}
        contract_a = _build_blind_style_contract(inventory, Path("benchmark_runs/deep_residual_fourway_20260616_0003/routes/04_blind_experimental_loop"))
        contract_b = _build_blind_style_contract(inventory, Path("benchmark_runs/deep_residual_fourway_20260616_0004/routes/04_blind_experimental_loop"))

        self.assertNotEqual(contract_a["style_id"], contract_b["style_id"])
        self.assertIn("previous_blind_experimental_candidates", contract_a["forgets_styles"])

    def test_speaker_script_audit_flags_placeholder_copy(self):
        temp_root = Path(__file__).parent / "outputs" / "tmp"
        temp_root.mkdir(parents=True, exist_ok=True)
        output_path = temp_root / f"script_audit_{uuid.uuid4().hex}.pptx"
        spec = PresentationSpec(
            title="Script Audit",
            slides=[
                SlideSpec(
                    slide_id="slide_01",
                    title="Very deep training challenge",
                    takeaway="Training very deep networks is difficult without shortcut paths.",
                    text_blocks=[
                        TextBlock(
                            text="Shortcut paths preserve optimization stability.",
                            claim="Shortcut paths preserve optimization stability",
                            detail="Shortcut paths preserve optimization stability.",
                            evidence="Motivation section",
                        )
                    ],
                )
            ],
        )
        PptxRenderer(style="academic").render(spec, output_path)
        audit = _audit_speaker_script(
            "# Script Audit\n\n## Slide 1: Very deep training challenge\n\nSuggested narration: The paper addresses the problem of\n",
            output_path,
            "academic",
        )

        self.assertTrue(any(f["type"] == "script_generic_placeholder_copy" for f in audit["findings"]))

    def test_style_drift_report_keeps_golden_baseline1_rules_scoped(self):
        temp_root = Path(__file__).parent / "outputs" / "tmp"
        temp_root.mkdir(parents=True, exist_ok=True)
        output_path = temp_root / f"scope_{uuid.uuid4().hex}.pptx"
        render_blind_blueprint_pptx(
            {"paper": {"title": "Scope Test"}, "paper_highlights": [], "assets": {"figures": [], "tables": []}, "metrics": [], "curated_slides": []},
            {
                "slides": [
                    {
                        "title": "Scoped rule",
                        "slide_role": "evidence",
                        "claim": "Style-specific polish should not mutate academic decks by default.",
                        "support": "The report can detect rounded-panel findings while keeping auto repair disabled.",
                        "proof_object": {"type": "text_evidence", "id": "scope", "focus": "detect only"},
                    }
                ]
            },
            output_path,
            repair_mode=True,
        )
        report = build_style_drift_report_payload(
            pptx_path=output_path,
            audit={"findings": [{"type": "panel_identity_label_anchor_drift", "severity": "low"}]},
            style_id="academic",
            repair_profile="global_correctness_repair",
            style_scope="academic",
        )
        self.assertEqual(report["golden_baseline1_scoped_findings"]["panel_identity_label_anchor_drift"], 1)
        self.assertIn("panel_identity_label_anchor_drift", report["forbidden_auto_repairs_detected"])
        self.assertEqual(report["forbidden_auto_repairs_applied"], [])

    def test_academic_bottom_table_void_is_detected(self):
        records = [
            {
                "role": "component_label",
                "has_text": True,
                "text": "Key message",
                "has_table": False,
                "is_picture": False,
                "bbox": {"x": 1.1, "y": 1.4, "w": 4.5, "h": 0.2, "right": 5.6, "bottom": 1.6, "area": 0.9},
                "font": {},
            },
            {
                "role": "table",
                "has_text": False,
                "text": "",
                "has_table": True,
                "is_picture": False,
                "bbox": {"x": 0.85, "y": 5.35, "w": 11.6, "h": 1.25, "right": 12.45, "bottom": 6.6, "area": 14.5},
                "font": {},
            },
            {
                "role": "page_marker",
                "has_text": True,
                "text": "9",
                "has_table": False,
                "is_picture": False,
                "bbox": {"x": 11.8, "y": 7.0, "w": 0.2, "h": 0.1, "right": 12.0, "bottom": 7.1, "area": 0.02},
                "font": {},
            },
        ]

        findings = _academic_right_evidence_void_findings(9, records)
        self.assertEqual(findings[0]["type"], "academic_right_evidence_void")

    def test_academic_bottom_table_with_right_picture_is_not_void(self):
        records = [
            {
                "role": "component_label",
                "has_text": True,
                "text": "Key message",
                "has_table": False,
                "is_picture": False,
                "bbox": {"x": 1.1, "y": 1.4, "w": 4.5, "h": 0.2, "right": 5.6, "bottom": 1.6, "area": 0.9},
                "font": {},
            },
            {
                "role": "picture",
                "has_text": False,
                "text": "",
                "has_table": False,
                "is_picture": True,
                "bbox": {"x": 6.2, "y": 1.72, "w": 5.8, "h": 3.0, "right": 12.0, "bottom": 4.72, "area": 17.4},
                "font": {},
            },
            {
                "role": "table",
                "has_text": False,
                "text": "",
                "has_table": True,
                "is_picture": False,
                "bbox": {"x": 0.85, "y": 5.35, "w": 11.6, "h": 1.25, "right": 12.45, "bottom": 6.6, "area": 14.5},
                "font": {},
            },
        ]

        self.assertEqual(_academic_right_evidence_void_findings(12, records), [])

    def test_academic_table_metric_page_fills_right_evidence_region(self):
        spec = PresentationSpec(
            title="Academic Table Metrics",
            slides=[
                SlideSpec(
                    slide_id="slide_01",
                    title="Overall Performance",
                    layout="metric_focus",
                    takeaway="CUA solves only 5.36% of Wordle games; average 3.25 guesses when solved.",
                    text_blocks=[
                        TextBlock(
                            text="Only about 1 in 19 games solved by the agent.",
                            claim="Extremely low success rate",
                            detail="Only about 1 in 19 games solved by the agent.",
                            evidence="Wordle task table",
                        ),
                        TextBlock(
                            text="When solved, averages 3.25 guesses vs. human around 4.",
                            claim="Guesses below human average",
                            detail="When solved, averages 3.25 guesses vs. human around 4.",
                            evidence="Wordle task table",
                        ),
                    ],
                    table_blocks=[
                        TableBlock(
                            title="Overall metrics",
                            rows=[
                                ["Metric", "CUA Agent"],
                                ["Avg. guesses per solved puzzle", "3.25"],
                                ["Success rate", "5.36%"],
                            ],
                            caption="Overall metrics on Wordle task.",
                        )
                    ],
                    metric_blocks=[
                        MetricBlock(label="Success rate", value="5.36%", note="Wordle task"),
                        MetricBlock(label="Avg. solved guesses", value="3.25", note="Solved games"),
                    ],
                )
            ],
        )
        temp_root = Path(__file__).parent / "outputs" / "tmp"
        temp_root.mkdir(parents=True, exist_ok=True)
        output_path = temp_root / f"academic_table_metric_{uuid.uuid4().hex}.pptx"

        PptxRenderer(style="academic").render(spec, output_path)
        audit = inspect_pptx_nonvisual(output_path)

        self.assertNotIn("academic_right_evidence_void", audit["summary"]["by_type"])

    def test_drops_metrics_without_visible_values(self):
        metrics = [
            MetricBlock(label="Not really a metric", value="", note="qualitative point"),
            MetricBlock(label="Success rate", value="5.36%", note="overall"),
        ]

        compact = _compact_metric_blocks(metrics)
        self.assertEqual(len(compact), 1)
        self.assertEqual(compact[0].value, "5.36%")

    def test_drops_phrase_metric_values_that_are_not_numbers(self):
        metrics = [
            MetricBlock(label="56-layer training error", value="higher than 20-layer", note="on CIFAR-10"),
            MetricBlock(label="ResNet-34 top-1 err", value="25.03% top-1", note="10-crop"),
        ]

        compact = _compact_metric_blocks(metrics)
        self.assertEqual(len(compact), 1)
        self.assertEqual(compact[0].value, "25.03%")

    def test_drops_generic_bare_integer_metric_cards(self):
        metrics = [
            MetricBlock(label="Accuracy", value="50", note="model depth"),
            MetricBlock(label="Top-1 error", value="25.03%", note="ImageNet"),
        ]

        compact = _compact_metric_blocks(metrics)
        self.assertEqual(len(compact), 1)
        self.assertEqual(compact[0].label, "Top-1 error")

    def test_renderer_headline_truncation_avoids_weak_endings(self):
        renderer = PptxRenderer(style="academic")

        self.assertEqual(renderer._headline_text("Ensemble top-5 error on ImageNet", 24), "Ensemble top-5 error")

    def test_metric_extraction_skips_dataset_number_tokens(self):
        slide = SlideSpec(
            slide_id="slide_01",
            title="Training challenge",
            takeaway="The paper studies very deep networks on ImageNet and CIFAR-10, but this slide states no numeric result.",
            text_blocks=[
                TextBlock(
                    text="Training depth is hard even when experiments include CIFAR-10.",
                    claim="Very deep training challenge",
                    detail="Training depth is hard even when experiments include CIFAR-10.",
                    evidence="Motivation section",
                )
            ],
        )

        self.assertEqual(_extract_metrics_from_slide(slide), [])

    def test_repairs_fragmented_point_claims_into_complete_ideas(self):
        slide = SlideSpec(
            slide_id="slide_01",
            title="The paper addresses the problem of",
            takeaway="The paper addresses the problem of training very deep convolutional neural networks for visual recognition.",
            text_blocks=[
                TextBlock(
                    text="The paper addresses the problem of training very deep convolutional neural networks for visual recognition.",
                    claim="The paper addresses the problem of",
                    detail="The paper addresses the problem of training very deep convolutional neural networks for visual recognition.",
                    evidence="Motivation section",
                )
            ],
            section_label="Motivation",
        )

        repaired = _ensure_structured_points(slide.text_blocks, slide)
        self.assertEqual(len(repaired), 1)
        self.assertNotEqual(repaired[0].claim, "The paper addresses the problem of")
        self.assertIn("challenge", repaired[0].claim.lower())

    def test_semantic_audit_detects_fragment_headings_and_spurious_metric_cards(self):
        records = [
            {
                "index": 1,
                "role": "title_claim",
                "has_text": True,
                "text": "The paper addresses the problem of",
                "text_words": 6,
                "bbox": {"x": 1.0, "y": 0.6, "w": 4.8, "h": 0.3, "right": 5.8, "bottom": 0.9, "area": 1.44},
                "font": {"avg_pt": 22.0},
            },
            {
                "index": 2,
                "role": "card_text",
                "has_text": True,
                "text": "50",
                "text_words": 1,
                "bbox": {"x": 8.9, "y": 1.5, "w": 0.6, "h": 0.3, "right": 9.5, "bottom": 1.8, "area": 0.18},
                "font": {"avg_pt": 18.0},
            },
            {
                "index": 3,
                "role": "small_text",
                "has_text": True,
                "text": "Accuracy",
                "text_words": 1,
                "bbox": {"x": 8.85, "y": 1.9, "w": 1.2, "h": 0.2, "right": 10.05, "bottom": 2.1, "area": 0.24},
                "font": {"avg_pt": 9.0},
            },
        ]

        findings = _semantic_content_findings(4, records)
        kinds = {finding["type"] for finding in findings}
        self.assertIn("weak_fragment_point_heading", kinds)
        self.assertIn("spurious_generic_metric_card", kinds)

    def test_placeholder_slide_number_titles_are_repaired_and_audited(self):
        slide = SlideSpec(
            slide_id="slide_02",
            title="Slide 2",
            takeaway="Deeper plain networks suffer from higher training error than shallower networks.",
            text_blocks=[
                TextBlock(
                    text="Plain networks with more layers exhibit higher training error than their shallower counterparts.",
                    claim="Deeper networks degrade accuracy",
                    detail="Plain networks with more layers exhibit higher training error than their shallower counterparts.",
                    evidence="Training error comparison",
                )
            ],
            section_label="Motivation",
        )

        self.assertEqual(_repair_slide_title(slide), "Deeper networks degrade accuracy")

        findings = _semantic_content_findings(
            4,
            [
                {
                    "index": 1,
                    "role": "title",
                    "has_text": True,
                    "text": "Slide 2",
                    "text_words": 2,
                    "bbox": {"x": 1.0, "y": 0.5, "w": 2.0, "h": 0.4, "right": 3.0, "bottom": 0.9, "area": 0.8},
                    "font": {"avg_pt": 24.0},
                }
            ],
        )

        self.assertTrue(any(finding["type"] == "weak_fragment_point_heading" for finding in findings))

    def test_structures_numbered_points_for_evaluator(self):
        slide = SlideSpec(
            slide_id="slide_01",
            title="Result",
            takeaway="The new model improves inference quality.",
            text_blocks=[TextBlock(text="Accuracy improves because the model uses stronger supervision.")],
            section_label="Results",
        )
        slide.text_blocks = _ensure_structured_points(slide.text_blocks, slide)

        point = slide.text_blocks[0]
        self.assertTrue(point.claim)
        self.assertTrue(point.detail)
        self.assertTrue(point.evidence)
        report = evaluate_presentation_spec(PresentationSpec(title="Eval", slides=[slide]))
        self.assertTrue(report["passed"])

    def test_evaluator_flags_missing_point_fields_and_bad_metric(self):
        spec = PresentationSpec(
            title="Eval",
            slides=[
                SlideSpec(
                    slide_id="slide_01",
                    title="Weak point",
                    text_blocks=[TextBlock(text="Incomplete...")],
                    metric_blocks=[MetricBlock(label="Key number", value="large")],
                )
            ],
        )
        report = evaluate_presentation_spec(spec)
        self.assertFalse(report["passed"])
        self.assertIn(1, report["failed_slides"])
        self.assertTrue(any("missing claim" in warning for warning in report["warnings"]))
        self.assertTrue(any("meaningless label" in warning for warning in report["warnings"]))

    def test_unicode_ellipsis_inside_real_title_is_not_truncation(self):
        self.assertFalse(_has_truncated_ellipsis("AGI Is Coming… Right After AI Learns to Play Wordle"))
        self.assertTrue(_has_truncated_ellipsis("This title was truncated... Before the noun"))
        self.assertTrue(_has_truncated_ellipsis("This title was truncated…"))

    def test_visual_layout_without_image_is_rejected_and_normalized(self):
        slide = SlideSpec(
            slide_id="slide_01",
            title="No Visual",
            layout="visual_left",
            takeaway="This slide has only text.",
            text_blocks=[
                TextBlock(
                    text="Claim: A complete detail sentence.",
                    claim="Claim",
                    detail="A complete detail sentence.",
                    evidence="source section",
                )
            ],
        )
        report = evaluate_presentation_spec(PresentationSpec(title="Eval", slides=[slide]))
        self.assertFalse(report["passed"])
        self.assertIn(1, report["failed_slides"])
        self.assertEqual(_normalize_slide_layout(slide), "statement")

    def test_unknown_content_layout_without_visual_normalizes_to_statement(self):
        slide = SlideSpec(
            slide_id="slide_01",
            title="Content Layout",
            layout="content",
            takeaway="This slide should use the full text canvas.",
            text_blocks=[
                TextBlock(
                    text="Claim: A complete detail sentence.",
                    claim="Claim",
                    detail="A complete detail sentence.",
                    evidence="source section",
                )
            ],
        )
        report = evaluate_presentation_spec(PresentationSpec(title="Eval", slides=[slide]))
        self.assertFalse(report["passed"])
        self.assertIn(1, report["failed_slides"])
        self.assertEqual(_normalize_slide_layout(slide), "statement")

    def test_generates_detailed_tex_sidecar(self):
        plan = ContentPlan(
            output_type="slides",
            sections=[
                Section(
                    id="slide_01",
                    title="Detailed Companion",
                    section_type="opening",
                    content="This slide keeps more context for a longer academic explanation. It cites the table evidence.",
                    tables=[TableRef(table_id="Table 1", focus="Evidence")],
                )
            ],
            tables_index={
                "Table 1": TableInfo(
                    table_id="Table 1",
                    caption="Example detailed table",
                    html_content="<table><tr><th>Metric</th><th>Value</th></tr><tr><td>Accuracy</td><td>92%</td></tr></table>",
                )
            },
        )
        spec = PresentationSpec(
            title="Detailed Test",
            slides=[
                SlideSpec(
                    slide_id="slide_01",
                    title="Detailed Companion",
                    takeaway="Detailed sidecar preserves extra context.",
                    text_blocks=[TextBlock(text="More evidence is shown than in the compact deck.")],
                )
            ],
        )
        temp_root = Path(__file__).parent / "outputs" / "tmp" / f"tex_{uuid.uuid4().hex}"
        result = generate_detailed_tex_deck(plan, spec, temp_root, compile_pdf=False)
        tex_path = Path(result["detailed_tex_path"])
        self.assertTrue(tex_path.exists())
        tex = tex_path.read_text(encoding="utf-8")
        self.assertIn("\\documentclass", tex)
        self.assertIn("Detailed sidecar preserves extra context.", tex)
        self.assertIn("Accuracy", tex)

    def test_pptx_model_can_follow_main_llm_model(self):
        with patch.dict(
            "os.environ",
            {
                "LLM_MODEL": "deepseek-v4-flash",
            },
            clear=True,
        ):
            self.assertEqual(_get_pptx_llm_model(), "deepseek-v4-flash")

    def test_pptx_model_override_and_figure_model_priority(self):
        with patch.dict(
            "os.environ",
            {
                "LLM_MODEL": "deepseek-v4-flash",
                "PPTX_LLM_MODEL": "deepseek-v4-pro",
                "PPTX_VISION_MODEL": "vision-compatible-model",
            },
            clear=True,
        ):
            self.assertEqual(_get_pptx_llm_model(), "deepseek-v4-pro")
            self.assertEqual(_get_figure_analysis_model(), "vision-compatible-model")

    def test_benchmark_warning_classifier(self):
        self.assertEqual(classify_warning("slide 2: text box may overflow vertically"), "text_overflow")
        self.assertEqual(classify_warning("slide 3: visual layout has no image"), "layout_payload_mismatch")
        self.assertEqual(classify_warning("slide 4: point 1 missing claim"), "structured_point")
        self.assertEqual(classify_warning("slide 5: metric 1 has meaningless label"), "metric_quality")

    def test_benchmark_summary_aggregates_qa_runs(self):
        summary = summarize_layout_qa(
            [
                QaRunResult(
                    path="outputs/a/layout_qa.json",
                    project="a",
                    passed=True,
                    slide_count=10,
                    warning_count=1,
                    failed_slides=[],
                    categories={"text_overflow": 1},
                    warnings=["slide 1: text box may overflow vertically"],
                ),
                QaRunResult(
                    path="outputs/b/layout_qa.json",
                    project="b",
                    passed=False,
                    slide_count=5,
                    warning_count=2,
                    failed_slides=[3],
                    categories={"empty_content": 1, "metric_quality": 1},
                    warnings=["slide 3: slide appears empty", "slide 3: metric 1 missing value"],
                ),
            ]
        )
        self.assertEqual(summary["run_count"], 2)
        self.assertEqual(summary["passed_runs"], 1)
        self.assertEqual(summary["failed_runs"], 1)
        self.assertEqual(summary["total_slides"], 15)
        self.assertEqual(summary["category_counts"]["metric_quality"], 1)

    def test_benchmark_manifest_expands_included_sets(self):
        manifest = {
            "sets": {
                "local_1": {"papers": [{"id": "a", "path": "a.pdf"}]},
                "ai2": {
                    "includes": ["local_1"],
                    "additional_papers": [{"id": "b", "path": "b.pdf"}],
                },
            }
        }
        papers = expand_paper_set(manifest, "ai2")
        self.assertEqual([paper["id"] for paper in papers], ["a", "b"])

    def test_benchmark_manifest_validation_reports_missing_files(self):
        root = Path(__file__).parent / "outputs" / "tmp" / f"manifest_{uuid.uuid4().hex}"
        root.mkdir(parents=True, exist_ok=True)
        present = root / "present.pdf"
        present.write_bytes(b"%PDF-1.4\n")
        manifest = root / "papers.json"
        manifest.write_text(
            """
{
  "sets": {
    "sample": {
      "papers": [
        {"id": "present", "path": "__PRESENT__"},
        {"id": "missing", "path": "__MISSING__"}
      ]
    }
  }
}
""".replace("__PRESENT__", str(present).replace("\\", "\\\\")).replace(
                "__MISSING__", str(root / "missing.pdf").replace("\\", "\\\\")
            ),
            encoding="utf-8",
        )
        result = validate_paper_files(manifest, "sample")
        self.assertEqual(result["total"], 2)
        self.assertEqual(len(result["present"]), 1)
        self.assertEqual(len(result["missing"]), 1)

    def test_human_feedback_benchmark_registry_loads(self):
        data = load_human_feedback_benchmark()
        summary = summarize_human_feedback_benchmark(data)
        self.assertEqual(summary["accepted_reference"], "rough_draft_v5")
        self.assertEqual(summary["golden_baseline1_style_id"], "golden_baseline1_from_scratch_warm_academic")
        self.assertEqual(summary["golden_baseline1_status"], "promoted_after_human_acceptance")
        self.assertEqual(summary["candidate_style_reference"], "rough_draft_v10_component_reflow")
        self.assertEqual(summary["candidate_style_status"], "promoted_to_golden_baseline1_after_cross_paper_validation")
        self.assertGreaterEqual(summary["candidate_validation_case_count"], 2)
        self.assertIn("rough_draft_v6", summary["avoid_versions"])
        self.assertGreaterEqual(summary["badcase_count"], 8)
        self.assertGreaterEqual(summary["autonomous_workflow_stage_count"], 6)
        self.assertGreaterEqual(summary["non_visual_detectable_count"], 8)
        self.assertGreaterEqual(summary["repair_priority_count"], 6)
        self.assertIn("component_overlap", badcase_ids(data))
        self.assertIn("overoptimized_density_regression", badcase_ids(data))
        self.assertIn("card_font_too_small", badcase_ids(data))
        self.assertIn("sparse_card_copy", badcase_ids(data))
        self.assertIn("geometry_changed_without_structural_need", badcase_ids(data))
        self.assertIn("cover_left_typography_underpowered", badcase_ids(data))
        self.assertIn("flow_nodes_overpacked", badcase_ids(data))
        self.assertIn("text_stack_off_optical_center", badcase_ids(data))
        self.assertIn("metric_label_gap_too_large", badcase_ids(data))
        self.assertIn("deck_type_scale_under_comfort_band", badcase_ids(data))
        self.assertIn("paired_label_body_gap_too_large", badcase_ids(data))
        self.assertIn("flow_grid_alignment_drift", badcase_ids(data))
        self.assertIn("component_frame_overallocated_after_text_fit", badcase_ids(data))
        self.assertIn("component_boundary_inset_violation", badcase_ids(data))
        self.assertIn("cross_paper_text_compression_overflow", badcase_ids(data))
        self.assertIn("wide_figure_forced_into_side_panel", badcase_ids(data))
        self.assertIn("figure_panel_aspect_mismatch", badcase_ids(data))
        self.assertIn("figure_image_off_center_in_panel", badcase_ids(data))
        self.assertIn("figure_label_anchor_drift", badcase_ids(data))
        self.assertIn("figure_badge_identity_label_conflation", badcase_ids(data))
        self.assertIn("stacked_figure_identity_label_overcorrection", badcase_ids(data))
        self.assertIn("figure_label_text_alignment_off_center", badcase_ids(data))
        self.assertIn("panel_identity_label_anchor_drift", badcase_ids(data))
        self.assertIn("panel_identity_label_text_alignment_off_center", badcase_ids(data))
        self.assertIn("figure_picture_aspect_distortion", badcase_ids(data))
        self.assertIn("inline_table_payload_not_indexed", badcase_ids(data))
        self.assertIn("card_internal_spacing_not_scaled_to_frame", badcase_ids(data))
        self.assertIn("agenda_read_path_header_too_close", badcase_ids(data))
        self.assertIn("table_support_band_off_balance", badcase_ids(data))
        self.assertIn("proof_caption_overflow_after_cross_paper_transfer", badcase_ids(data))

    def test_from_scratch_routes_wide_figures_and_inline_tables(self):
        from PIL import Image

        from paper2slides.benchmark.from_scratch import _figure_path_index, _layout_family, _limit_text_for_box, _table_index

        root = Path(__file__).parent / "outputs" / "tmp" / f"wide_figure_{uuid.uuid4().hex}"
        root.mkdir(parents=True, exist_ok=True)
        figure_path = root / "wide.png"
        medium_wide_path = root / "medium_wide.png"
        mildly_wide_path = root / "mildly_wide.png"
        tall_path = root / "tall.png"
        Image.new("RGB", (800, 100), color=(255, 255, 255)).save(figure_path)
        Image.new("RGB", (620, 300), color=(255, 255, 255)).save(medium_wide_path)
        Image.new("RGB", (510, 300), color=(255, 255, 255)).save(mildly_wide_path)
        Image.new("RGB", (420, 680), color=(255, 255, 255)).save(tall_path)
        inventory = {
            "paper": {"title": "Wide Figure Paper"},
            "assets": {
                "figures": [
                    {"id": "Figure 4", "caption": "Wide schedule", "path": str(figure_path)},
                    {"id": "Figure 5", "caption": "Medium-wide timeline", "path": str(medium_wide_path)},
                    {"id": "Figure 6", "caption": "Mildly-wide chart", "path": str(mildly_wide_path)},
                    {"id": "Figure 7", "caption": "Tall reasoning trace", "path": str(tall_path)},
                ],
                "tables": [],
            },
            "curated_slides": [
                {
                    "tables": [
                        {
                            "title": "Inline Table",
                            "caption": "Rows from slide spec",
                            "rows": [["Metric", "Value"], ["Loss", "-0.027"]],
                        }
                    ]
                }
            ],
        }
        slide_data = {"proof_object": {"type": "figure", "id": "Figure 4", "focus": "wide schedule"}}
        medium_wide_slide = {"proof_object": {"type": "figure", "id": "Figure 5", "focus": "medium-wide timeline"}}
        mildly_wide_slide = {"proof_object": {"type": "figure", "id": "Figure 6", "focus": "mildly-wide chart"}}
        tall_slide = {"proof_object": {"type": "figure", "id": "Figure 7", "focus": "tall reasoning trace"}}

        self.assertEqual(_layout_family(slide_data, 1, _figure_path_index(inventory)), "figure_bottom_wide")
        self.assertEqual(_layout_family(medium_wide_slide, 1, _figure_path_index(inventory)), "figure_bottom_wide")
        self.assertEqual(_layout_family(mildly_wide_slide, 1, _figure_path_index(inventory)), "visual_right")
        self.assertEqual(_layout_family(tall_slide, 1, _figure_path_index(inventory)), "figure_tall_right")
        self.assertEqual(_layout_family(tall_slide, 2, _figure_path_index(inventory)), "figure_tall_left")
        inline_table = _table_index(inventory)["Inline Table"]
        self.assertEqual(inline_table["row_count"], 2)
        self.assertEqual(inline_table["column_count"], 2)
        long_caption = (
            "| Illustration of a very long figure caption that describes implementation details, "
            "baseline comparisons, and additional caveats that should not overflow a short proof caption box."
        )
        fitted_caption = _limit_text_for_box(long_caption, 4.8, 0.42, 10.0, fill_ratio=0.78)
        self.assertLess(len(fitted_caption), len(long_caption))
        self.assertTrue(fitted_caption.endswith("."))

    def test_benchmark_runner_builds_expected_command(self):
        command = _build_command(
            python_executable="python",
            paper_path=Path("test_papers/example.pdf"),
            output_dir=Path("outputs"),
            style="academic",
            length="medium",
            slides=24,
            fast=True,
            from_stage="generate",
        )
        self.assertEqual(command[:3], ["python", "-m", "paper2slides"])
        self.assertIn("--fast", command)
        self.assertIn("--from-stage", command)
        self.assertIn("generate", command)
        self.assertIn("--slides", command)
        self.assertIn("24", command)

    def test_benchmark_runner_summary_counts_failures_and_warnings(self):
        summary = summarize_run_results(
            [
                {
                    "style": "academic",
                    "returncode": 0,
                    "qa_passed": True,
                    "elapsed_seconds": 10,
                    "warning_count": 1,
                    "slide_count": 5,
                    "warning_categories": {"text_overflow": 1},
                },
                {
                    "style": "academic",
                    "returncode": 0,
                    "qa_passed": False,
                    "elapsed_seconds": 20,
                    "warning_count": 2,
                    "slide_count": 4,
                    "warning_categories": {"metric_quality": 2},
                },
                {
                    "style": "visual",
                    "returncode": 1,
                    "qa_passed": None,
                    "elapsed_seconds": 5,
                    "warning_count": 0,
                    "slide_count": 0,
                    "warning_categories": {},
                },
            ]
        )
        self.assertEqual(summary["total_runs"], 3)
        self.assertEqual(summary["succeeded"], 1)
        self.assertEqual(summary["qa_failed"], 1)
        self.assertEqual(summary["command_failed"], 1)
        self.assertEqual(summary["total_warnings"], 3)
        self.assertEqual(summary["warning_categories"]["metric_quality"], 2)

    def test_fast_rag_defaults_to_text_only_without_image_payloads(self):
        import asyncio

        class FakeCompletions:
            def __init__(self):
                self.messages = None

            def create(self, **kwargs):
                self.messages = kwargs["messages"]

                class Choice:
                    message = type("Message", (), {"content": "answer"})

                return type("Response", (), {"choices": [Choice()]})

        class FakeClient:
            def __init__(self):
                self.completions = FakeCompletions()
                self.chat = type("Chat", (), {"completions": self.completions})

        temp_root = Path(__file__).parent / "outputs" / "tmp" / f"rag_{uuid.uuid4().hex}"
        temp_root.mkdir(parents=True, exist_ok=True)
        md_path = temp_root / "paper.md"
        md_path.write_text("Text before image. ![Figure](figure.png) Caption text.", encoding="utf-8")
        (temp_root / "figure.png").write_bytes(b"not a real png, but should not be read")

        client = FakeClient()
        result = asyncio.run(
            _run_fast_queries_by_category(
                client=client,
                markdown_content="",
                markdown_paths=[str(md_path)],
                queries_by_category={"summary": ["What is the method?"]},
                model="deepseek-v4-flash",
                max_concurrency=1,
                include_images=False,
            )
        )
        self.assertEqual(result["summary"][0]["mode"], "fast_direct_text")
        user_content = client.completions.messages[1]["content"]
        self.assertFalse(any(part.get("type") == "image_url" for part in user_content))

    def test_fast_rag_uses_vision_model_when_images_are_enabled(self):
        import asyncio

        class FakeCompletions:
            def __init__(self):
                self.model = None
                self.messages = None

            def create(self, **kwargs):
                self.model = kwargs["model"]
                self.messages = kwargs["messages"]

                class Choice:
                    message = type("Message", (), {"content": "answer"})

                return type("Response", (), {"choices": [Choice()]})

        class FakeClient:
            def __init__(self):
                self.completions = FakeCompletions()
                self.chat = type("Chat", (), {"completions": self.completions})

        temp_root = Path(__file__).parent / "outputs" / "tmp" / f"rag_vision_{uuid.uuid4().hex}"
        temp_root.mkdir(parents=True, exist_ok=True)
        md_path = temp_root / "paper.md"
        md_path.write_text("![Figure](figure.png) Caption text.", encoding="utf-8")
        # A minimal PNG signature is enough for the base64 payload test.
        (temp_root / "figure.png").write_bytes(b"\x89PNG\r\n\x1a\n")

        client = FakeClient()
        result = asyncio.run(
            _run_fast_queries_by_category(
                client=client,
                markdown_content="",
                markdown_paths=[str(md_path)],
                queries_by_category={"summary": ["What is the figure?"]},
                model="deepseek-v4-flash",
                max_concurrency=1,
                include_images=True,
                vision_model="gpt-5-mini",
            )
        )
        self.assertEqual(client.completions.model, "gpt-5-mini")
        self.assertEqual(result["summary"][0]["mode"], "fast_direct_with_vision")
        self.assertEqual(result["summary"][0]["model"], "gpt-5-mini")
        user_content = client.completions.messages[1]["content"]
        self.assertTrue(any(part.get("type") == "image_url" for part in user_content))

    def test_benchmark_preflight_requires_vision_key_when_images_enabled(self):
        with patch.dict("os.environ", {}, clear=True):
            with self.assertRaises(RuntimeError):
                _preflight_environment(
                    {
                        "RAG_FAST_INCLUDE_IMAGES": "1",
                        "PPTX_ENABLE_FIGURE_ANALYSIS": "auto",
                    }
                )

    def test_benchmark_preflight_allows_split_text_and_vision_keys(self):
        with patch.dict("os.environ", {}, clear=True):
            _preflight_environment(
                {
                    "RAG_FAST_INCLUDE_IMAGES": "1",
                    "PPTX_ENABLE_FIGURE_ANALYSIS": "auto",
                    "RAG_VISION_API_KEY": "vision-key",
                }
            )

    def test_from_scratch_inventory_reuses_existing_checkpoints(self):
        root = Path(__file__).parent / "outputs" / "tmp" / f"inventory_{uuid.uuid4().hex}"
        summary_path = root / "checkpoint_summary.json"
        plan_path = root / "checkpoint_plan.json"
        spec_path = root / "checkpoint_slide_spec.json"
        summary_path.parent.mkdir(parents=True, exist_ok=True)
        summary_path.write_text(
            """
{
  "content_type": "paper",
  "content": {
    "paper_info": "**Title**: Example Agent Paper  **Authors**: Team",
    "motivation": "## RESEARCH PROBLEM\\nAgents need reliable tool-use planning.",
    "solution": "## FRAMEWORK OVERVIEW\\nThe system combines planning, retrieval, and execution.",
    "results": "## MAIN RESULTS\\nThe model reaches 87% success on a benchmark.",
    "contributions": "## MAIN CONTRIBUTIONS\\nOpen checkpoints and training recipe."
  },
  "origin": {
    "base_path": "paper_assets",
    "figures": [{"id": "Figure 1", "caption": "Pipeline overview", "path": "fig1.png"}],
    "tables": [{"id": "Table 1", "caption": "Benchmark scores", "html": "<table><tr><td>87%</td></tr></table>"}]
  }
}
""",
            encoding="utf-8",
        )
        plan_path.write_text(
            """
{
  "plan": {
    "output_type": "slides",
    "sections": [
      {
        "id": "slide_01",
        "title": "Example Agent Paper",
        "type": "opening",
        "section": "Overview",
        "content": "An overview of the agent system.",
        "figures": [{"figure_id": "Figure 1", "focus": "Pipeline"}],
        "tables": []
      },
      {
        "id": "slide_02",
        "title": "Benchmark Result",
        "type": "content",
        "section": "Results",
        "content": "The model reaches 87% success.",
        "figures": [],
        "tables": [{"table_id": "Table 1", "focus": "Success rate"}]
      }
    ]
  }
}
""",
            encoding="utf-8",
        )
        spec_path.write_text(
            """
{
  "title": "Example Agent Paper",
  "slides": [
    {
      "slide_id": "slide_01",
      "title": "Example Agent Paper",
      "layout": "cover",
      "takeaway": "The agent system combines planning and execution.",
      "section_type": "opening",
      "section_label": "Overview",
      "text_blocks": [
        {"text": "Unified agent loop: Planning and execution are joined.", "claim": "Unified agent loop", "detail": "Planning and execution are joined.", "evidence": "Figure 1"}
      ],
      "metric_blocks": []
    },
    {
      "slide_id": "slide_02",
      "title": "Benchmark Result",
      "layout": "table_focus",
      "takeaway": "The model reaches 87% success.",
      "section_type": "content",
      "section_label": "Results",
      "text_blocks": [
        {"text": "Strong benchmark result: The model reaches 87% success.", "claim": "Strong benchmark result", "detail": "The model reaches 87% success.", "evidence": "Table 1"}
      ],
      "metric_blocks": [{"label": "Success rate", "value": "87%", "note": "benchmark"}]
    }
  ]
}
""",
            encoding="utf-8",
        )

        inventory = build_content_inventory(summary_path, plan_path, spec_path)
        self.assertEqual(inventory["paper"]["title"], "Example Agent Paper")
        self.assertEqual(inventory["coverage"]["plan_slide_count"], 2)
        self.assertEqual(inventory["coverage"]["figure_count"], 1)
        self.assertEqual(inventory["coverage"]["metric_count"], 1)
        self.assertTrue(inventory["coverage"]["has_core_sections"]["method"])
        self.assertGreaterEqual(len(inventory["paper_highlights"]), 2)
        self.assertTrue(any("87%" in item["body"] or "planning" in item["body"].lower() for item in inventory["paper_highlights"]))
        self.assertEqual(inventory["assets"]["tables"][0]["row_count"], 1)
        self.assertEqual(inventory["assets"]["tables"][0]["rows"][0][0], "87%")

        rough = build_rough_draft_spec(inventory)
        self.assertEqual(len(rough["slides"]), 2)
        self.assertEqual(rough["slides"][0]["proof_object"]["type"], "figure")
        self.assertEqual(rough["slides"][1]["proof_object"]["type"], "table")
        self.assertTrue(rough["rules"]["baseline_skeleton_forbidden"])

    def test_from_scratch_artifact_writer_outputs_inventory_and_rough_draft(self):
        root = Path(__file__).parent / "outputs" / "tmp" / f"inventory_write_{uuid.uuid4().hex}"
        summary_path = root / "checkpoint_summary.json"
        output_dir = root / "from_scratch"
        pptx_path = output_dir / "rough_draft.pptx"
        summary_path.parent.mkdir(parents=True, exist_ok=True)
        summary_path.write_text(
            """
{
  "content_type": "paper",
  "content": {
    "paper_info": "Title: Minimal Paper",
    "motivation": "A short motivation paragraph with enough content to keep.",
    "solution": "A short method paragraph with enough content to keep."
  },
  "origin": {"figures": [], "tables": []}
}
""",
            encoding="utf-8",
        )

        paths = write_from_scratch_artifacts(summary_path, None, None, output_dir, pptx_output=pptx_path)
        self.assertTrue(Path(paths["content_inventory"]).exists())
        self.assertTrue(Path(paths["rough_draft_spec"]).exists())
        self.assertTrue(Path(paths["rough_draft_pptx"]).exists())
        self.assertTrue(Path(paths["nonvisual_audit"]).exists())
        self.assertTrue(Path(paths["visual_audit"]).exists())
        self.assertGreater(Path(paths["rough_draft_pptx"]).stat().st_size, 1000)
        audit = json.loads(Path(paths["visual_audit"]).read_text(encoding="utf-8"))
        self.assertFalse(audit["non_visual_review_manifest"]["requires_rendered_screenshots"])
        self.assertFalse(audit["non_visual_review_manifest"]["requires_vision_model"])
        self.assertFalse(audit["visual_review_manifest"]["requires_rendered_screenshots"])
        self.assertGreaterEqual(len(audit["visual_review_manifest"]["render_requests"]), 2)
        nonvisual = json.loads(Path(paths["nonvisual_audit"]).read_text(encoding="utf-8"))
        self.assertEqual(nonvisual["review_mode"], "non_visual_metadata_only")
        self.assertFalse(nonvisual["rendering_used"])

    def test_from_scratch_renderer_outputs_plain_pptx(self):
        root = Path(__file__).parent / "outputs" / "tmp" / f"rough_pptx_{uuid.uuid4().hex}"
        pptx_path = root / "rough_draft.pptx"
        inventory = {
            "paper": {"title": "Plain Rough Draft"},
            "assets": {
                "figures": [],
                "tables": [
                    {
                        "id": "Table 1",
                        "caption": "Result table",
                        "rows": [["Metric", "Value"], ["Success", "87%"]],
                    }
                ],
            },
        }
        rough = {
            "slides": [
                {
                    "slide_id": "slide_01",
                    "title": "Result",
                    "slide_role": "table_interpretation",
                    "claim": "The result is strong enough to carry the slide.",
                    "support": "This support text is intentionally plain and content-first.",
                    "proof_object": {"type": "table", "id": "Table 1", "focus": "Result table"},
                    "source_evidence": [{"source": "summary", "id": "results_01"}],
                }
            ]
        }
        render_rough_draft_pptx(inventory, rough, pptx_path)
        self.assertTrue(pptx_path.exists())
        self.assertGreater(pptx_path.stat().st_size, 1000)
        prs = Presentation(pptx_path)
        self.assertEqual(len(prs.slides), 5)
        self.assertGreaterEqual(sum(1 for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_table", False)), 1)
        nonvisual = inspect_pptx_nonvisual(pptx_path)
        self.assertEqual(nonvisual["slide_count"], 5)
        self.assertFalse(nonvisual["rendering_used"])
        self.assertTrue(any("table" in slide["role_counts"] for slide in nonvisual["slides"]))
        self.assertIn("finding_count", nonvisual["summary"])
        self.assertIn("by_problem_type", nonvisual["summary"])
        self.assertIn("typography_pages", nonvisual["summary"])
        self.assertIn("optical_balance_pages", nonvisual["summary"])
        self.assertIn("deck_flags", nonvisual["summary"])
        self.assertTrue(all("problem_type" in finding for finding in nonvisual["findings"]))

    def test_nonvisual_audit_detects_optical_balance_and_flow_risks(self):
        from paper2slides.benchmark.nonvisual_audit import (
            NONVISUAL_AUDIT_RULES,
            _agenda_read_path_header_findings,
            _card_internal_spacing_findings,
            _component_boundary_findings,
            _component_frame_fit_findings,
            _container_balance_findings,
            _figure_label_semantics_findings,
            _flow_layout_findings,
            _metric_stack_findings,
            _panel_identity_label_findings,
            _paired_text_stack_findings,
            _picture_findings,
            _table_support_band_findings,
        )

        def record(
            index,
            role,
            text,
            x,
            y,
            w,
            h,
            font=10.0,
            is_picture=False,
            picture=None,
            paragraph_alignment="",
            has_table=False,
        ):
            return {
                "index": index,
                "role": role,
                "has_text": bool(text),
                "text": text,
                "text_words": len(text.split()),
                "font": {"sizes_pt": [font] if text else [], "min_pt": font if text else None, "avg_pt": font if text else None},
                "bbox": {"x": x, "y": y, "w": w, "h": h, "right": x + w, "bottom": y + h, "area": w * h},
                "is_picture": is_picture,
                "picture": picture or {},
                "has_table": has_table,
                "paragraph_alignment": paragraph_alignment,
                "is_full_background": False,
            }

        flow_records = [
            record(1, "small_text", "Problem", 9.1, 4.8, 0.72, 0.2, 7.2),
            record(2, "small_text", "Method", 9.87, 4.8, 0.72, 0.2, 7.2),
            record(3, "small_text", "Evidence", 10.64, 4.8, 0.72, 0.2, 7.2),
            record(4, "small_text", "Takeaways", 11.41, 4.8, 0.72, 0.2, 7.2),
        ]
        self.assertEqual(_flow_layout_findings(2, flow_records, NONVISUAL_AUDIT_RULES)[0]["type"], "flow_nodes_overpacked")

        wide_grid_records = [
            record(20, "small_text", "P", 9.25, 4.42, 0.34, 0.34, 8.5),
            record(21, "small_text", "M", 11.56, 4.42, 0.34, 0.34, 8.5),
            record(22, "small_text", "E", 9.25, 4.97, 0.34, 0.34, 8.5),
            record(23, "small_text", "T", 11.56, 4.97, 0.34, 0.34, 8.5),
            record(24, "small_text", "Problem", 9.03, 4.79, 0.84, 0.2, 7.5),
            record(25, "small_text", "Method", 11.34, 4.79, 0.84, 0.2, 7.5),
            record(26, "small_text", "Evidence", 9.03, 5.34, 0.84, 0.2, 7.5),
            record(27, "small_text", "Takeaways", 11.34, 5.34, 0.84, 0.2, 7.5),
        ]
        self.assertEqual(_flow_layout_findings(2, wide_grid_records, NONVISUAL_AUDIT_RULES)[0]["type"], "flow_grid_alignment_drift")

        agenda_header_records = [
            record(30, "body_text", "Read path", 9.3, 4.05, 2.2, 0.25, 11),
            record(31, "small_text", "P", 9.7, 4.42, 0.34, 0.34, 8.5),
            record(32, "small_text", "M", 11.11, 4.42, 0.34, 0.34, 8.5),
            record(33, "small_text", "E", 9.7, 5.1, 0.34, 0.34, 8.5),
            record(34, "small_text", "T", 11.11, 5.1, 0.34, 0.34, 8.5),
        ]
        self.assertEqual(
            _agenda_read_path_header_findings(2, agenda_header_records, NONVISUAL_AUDIT_RULES)[0]["type"],
            "agenda_read_path_header_too_close",
        )

        paired_records = [
            record(28, "body_text", "Core result", 9.73, 2.10, 1.6, 0.25, 11),
            record(29, "card_text", "Kimi K2 achieves open-source SOTA on agentic tasks.", 9.73, 2.40, 1.9, 0.45, 10),
        ]
        self.assertEqual(_paired_text_stack_findings(1, paired_records, NONVISUAL_AUDIT_RULES)[0]["type"], "paired_label_body_gap_too_large")

        metric_records = [
            record(5, "container", "", 7.3, 2.1, 4.3, 1.2),
            record(6, "title_claim", "70.6", 7.5, 2.25, 3.9, 0.35, 22),
            record(7, "small_text", "Tau2 retail", 7.5, 2.72, 3.9, 0.2, 10),
        ]
        self.assertEqual(_metric_stack_findings(22, metric_records, NONVISUAL_AUDIT_RULES)[0]["type"], "metric_label_gap_too_large")

        balance_records = [
            record(8, "container", "", 7.0, 1.3, 5.0, 4.4),
            record(9, "container", "", 7.3, 2.0, 4.3, 1.1),
            record(10, "container", "", 7.3, 3.4, 4.3, 0.6),
        ]
        self.assertEqual(_container_balance_findings(29, balance_records, NONVISUAL_AUDIT_RULES)[0]["type"], "container_stack_off_balance")

        boundary_records = [
            record(11, "container", "", 7.0, 1.18, 5.45, 5.15),
            record(12, "component_label", "FIGURE", 7.20, 1.34, 1.2, 0.24, 8),
        ]
        self.assertEqual(_component_boundary_findings(13, boundary_records, NONVISUAL_AUDIT_RULES)[0]["type"], "component_boundary_inset_violation")

        panel_identity_drift_records = [
            record(101, "container", "", 0.75, 5.10, 11.7, 1.08),
            record(102, "component_label", "T E X T _ E V I D E N C E", 0.97, 5.28, 1.8, 0.24, 8),
            record(103, "body_text", "Motivation", 2.85, 5.25, 2.4, 0.30, 13.0),
            record(104, "support_body", "The proof explanation spans the full panel below.", 0.97, 5.58, 11.25, 0.78, 12.5),
        ]
        self.assertEqual(
            _panel_identity_label_findings(5, panel_identity_drift_records, NONVISUAL_AUDIT_RULES)[0]["type"],
            "panel_identity_label_anchor_drift",
        )

        panel_identity_text_alignment_records = [
            record(105, "container", "", 0.65, 1.15, 6.35, 5.25),
            record(106, "component_label", "T A B L E", 1.01, 1.47, 5.63, 0.24, 8),
            record(107, "body_text", "Doc Table 1", 1.01, 1.84, 5.63, 0.38, 18.0),
            record(108, "table", "", 0.9, 2.17, 5.85, 3.58, has_table=True),
        ]
        self.assertEqual(
            _panel_identity_label_findings(25, panel_identity_text_alignment_records, NONVISUAL_AUDIT_RULES)[0][
                "type"
            ],
            "panel_identity_label_text_alignment_off_center",
        )

        panel_identity_centered_records = [
            record(109, "container", "", 0.65, 1.15, 6.35, 5.25),
            record(110, "component_label", "T A B L E", 1.01, 1.47, 5.63, 0.24, 8),
            record(111, "body_text", "Doc Table 1", 0.9, 1.84, 5.85, 0.38, 18.0, paragraph_alignment="CENTER (2)"),
            record(112, "table", "", 0.9, 2.17, 5.85, 3.58, has_table=True),
        ]
        self.assertFalse(_panel_identity_label_findings(25, panel_identity_centered_records, NONVISUAL_AUDIT_RULES))

        frame_records = [
            record(13, "container", "", 0.75, 4.55, 3.6, 1.48),
            record(14, "card_label", "Method", 0.98, 4.67, 3.1, 0.22, 11),
            record(15, "card_text", "Before RL, the model undergoes supervised fine-tuning on high-quality instruction data.", 0.98, 4.99, 3.1, 0.9, 10),
        ]
        self.assertEqual(_component_frame_fit_findings(15, frame_records, NONVISUAL_AUDIT_RULES)[0]["type"], "component_frame_overallocated_after_text_fit")

        shallow_card_records = [
            record(16, "container", "", 9.0, 1.9, 3.35, 0.96),
            record(17, "card_label", "Conclusion", 9.18, 2.02, 2.99, 0.24, 10),
            record(18, "card_text", "This concluding slide summarizes the main contributions, key findings.", 9.18, 2.34, 2.99, 0.44, 10),
        ]
        self.assertEqual(
            _card_internal_spacing_findings(30, shallow_card_records, NONVISUAL_AUDIT_RULES)[0]["type"],
            "card_internal_spacing_not_scaled_to_frame",
        )

        table_support_records = [
            record(19, "title_claim", "mHC consistently improves over the baseline and typically outperforms HC.", 0.65, 1.55, 11.8, 0.87, 24),
            record(20, "support_body", "mHC outperforms Baseline and generally surpasses HC on key tasks like BBH and DROP.", 0.67, 2.72, 10.9, 0.78, 13.5),
            record(21, "container", "", 0.75, 3.50, 11.85, 2.95),
            record(22, "table", "", 1.0, 4.52, 11.35, 1.28),
        ]
        self.assertEqual(
            _table_support_band_findings(26, table_support_records, NONVISUAL_AUDIT_RULES)[0]["type"],
            "table_support_band_off_balance",
        )
        table_support_fixed_records = [
            record(23, "title_claim", "mHC consistently improves over the baseline and typically outperforms HC.", 0.65, 1.55, 11.8, 0.87, 24),
            record(24, "support_body", "mHC outperforms Baseline and generally surpasses HC on key tasks like BBH and DROP.", 0.67, 2.58, 10.9, 0.78, 13.5),
            record(25, "container", "", 0.75, 3.66, 11.85, 2.79),
            record(26, "table", "", 1.0, 4.68, 11.35, 1.14),
        ]
        self.assertFalse(_table_support_band_findings(26, table_support_fixed_records, NONVISUAL_AUDIT_RULES))

        tall_figure_bad_records = [
            record(27, "container", "", 7.0, 1.18, 5.45, 5.15),
            record(
                28,
                "picture",
                "",
                8.52,
                2.2,
                1.72,
                2.77,
                is_picture=True,
                picture={"source_aspect": 0.62, "box_aspect": 0.62, "aspect_distortion": 0.0},
            ),
        ]
        self.assertEqual(
            _picture_findings(4, tall_figure_bad_records, NONVISUAL_AUDIT_RULES)[0]["type"],
            "figure_panel_aspect_mismatch",
        )
        off_center_figure_records = [
            record(29, "container", "", 0.75, 3.43, 11.85, 3.07),
            record(
                30,
                "picture",
                "",
                7.6,
                4.0,
                4.5,
                1.8,
                is_picture=True,
                picture={"source_aspect": 2.5, "box_aspect": 2.5, "aspect_distortion": 0.0},
            ),
        ]
        self.assertEqual(
            _picture_findings(15, off_center_figure_records, NONVISUAL_AUDIT_RULES)[0]["type"],
            "figure_image_off_center_in_panel",
        )
        conflated_figure_label_records = [
            record(33, "container", "", 0.75, 3.45, 11.85, 3.12),
            record(34, "component_label", "F I G U R E", 3.96, 3.71, 0.82, 0.22, 8),
            record(35, "body_text", "Figure 5", 4.82, 3.66, 1.6, 0.34, 13.5),
            record(
                36,
                "picture",
                "",
                3.96,
                4.03,
                5.41,
                2.12,
                is_picture=True,
                picture={"source_aspect": 2.55, "box_aspect": 2.55, "aspect_distortion": 0.0},
            ),
        ]
        self.assertEqual(
            _figure_label_semantics_findings(13, conflated_figure_label_records, NONVISUAL_AUDIT_RULES)[0]["type"],
            "figure_badge_identity_label_conflation",
        )
        stacked_figure_label_records = [
            record(37, "container", "", 0.75, 3.45, 11.85, 3.12),
            record(38, "component_label", "F I G U R E", 1.11, 3.75, 0.82, 0.22, 8),
            record(39, "body_text", "F I G U R E 5", 3.70, 4.05, 0.24, 1.58, 7.8),
            record(
                40,
                "picture",
                "",
                3.96,
                4.03,
                5.41,
                2.12,
                is_picture=True,
                picture={"source_aspect": 2.55, "box_aspect": 2.55, "aspect_distortion": 0.0},
            ),
        ]
        self.assertEqual(
            _figure_label_semantics_findings(13, stacked_figure_label_records, NONVISUAL_AUDIT_RULES)[0]["type"],
            "stacked_figure_identity_label_overcorrection",
        )
        off_center_text_figure_label_records = [
            record(41, "container", "", 0.75, 3.45, 11.85, 3.12),
            record(42, "component_label", "F I G U R E", 1.11, 3.75, 0.82, 0.22, 8),
            record(43, "body_text", "Figure 5", 5.94, 3.73, 1.45, 0.30, 13.0, paragraph_alignment="LEFT (1)"),
            record(
                44,
                "picture",
                "",
                3.96,
                4.03,
                5.41,
                2.12,
                is_picture=True,
                picture={"source_aspect": 2.55, "box_aspect": 2.55, "aspect_distortion": 0.0},
            ),
        ]
        self.assertEqual(
            _figure_label_semantics_findings(13, off_center_text_figure_label_records, NONVISUAL_AUDIT_RULES)[0][
                "type"
            ],
            "figure_label_text_alignment_off_center",
        )
        standard_figure_panel_label_records = [
            record(45, "container", "", 0.65, 1.18, 5.45, 5.15),
            record(46, "component_label", "F I G U R E", 1.01, 1.52, 0.82, 0.22, 8),
            record(47, "body_text", "Figure 4", 1.01, 2.17, 4.73, 0.30, 13.0),
            record(
                48,
                "picture",
                "",
                0.97,
                2.49,
                4.81,
                2.93,
                is_picture=True,
                picture={"source_aspect": 1.64, "box_aspect": 1.64, "aspect_distortion": 0.0},
            ),
        ]
        self.assertEqual(
            _figure_label_semantics_findings(12, standard_figure_panel_label_records, NONVISUAL_AUDIT_RULES)[0][
                "type"
            ],
            "figure_label_text_alignment_off_center",
        )
        separated_figure_label_records = [
            record(41, "container", "", 0.75, 3.45, 11.85, 3.12),
            record(42, "component_label", "F I G U R E", 1.11, 3.75, 0.82, 0.22, 8),
            record(43, "body_text", "Figure 5", 5.94, 3.73, 1.45, 0.30, 13.0, paragraph_alignment="CENTER (2)"),
            record(
                44,
                "picture",
                "",
                3.96,
                4.03,
                5.41,
                2.12,
                is_picture=True,
                picture={"source_aspect": 2.55, "box_aspect": 2.55, "aspect_distortion": 0.0},
            ),
        ]
        self.assertFalse(_figure_label_semantics_findings(13, separated_figure_label_records, NONVISUAL_AUDIT_RULES))
        wide_figure_fixed_records = [
            record(45, "container", "", 0.75, 3.43, 11.85, 3.07),
            record(
                46,
                "picture",
                "",
                3.0,
                4.0,
                6.4,
                2.1,
                is_picture=True,
                picture={"source_aspect": 3.1, "box_aspect": 3.05, "aspect_distortion": 0.02},
            ),
        ]
        self.assertFalse(_picture_findings(15, wide_figure_fixed_records, NONVISUAL_AUDIT_RULES))


if __name__ == "__main__":
    unittest.main()
